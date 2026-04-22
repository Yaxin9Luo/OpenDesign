"""No-API smoke test: imports + schemas + fonts + a real composite call.

Run with:
    python -m longcat_design.smoke

Generates `out/smoke/` containing poster.psd, poster.svg, preview.png produced
from a fake (solid-color) background + 2 real text layers rendered via Pillow.
This proves the whole pipeline below the LLM/Gemini layer works without keys.
"""

from __future__ import annotations

import json
import sys
from datetime import datetime
from pathlib import Path

from PIL import Image
from pydantic import ValidationError

from .schema import (
    AgentTraceStep, CompositionArtifacts, CritiqueResult, DesignSpec,
    DistillTrajectory, LayerNode, SafeZone, TextEffect, ThinkingBlockRecord,
    ToolResultRecord, TrainingMetadata,
)


def _fail(msg: str) -> None:
    print(f"  FAIL  {msg}")
    sys.exit(1)


def _ok(msg: str) -> None:
    print(f"  ok    {msg}")


def check_imports() -> None:
    print("[1/20] imports")
    from . import chat, cli, config, critic, planner, runner, schema, session  # noqa
    from .tools import (
        TOOL_HANDLERS, TOOL_SCHEMAS, ToolContext,
        composite, critique_tool, edit_layer, fetch_brand_asset, finalize,
        generate_background, propose_design_spec, render_text_layer,
        switch_artifact_type,
    )  # noqa
    from .util import ids, io, logging  # noqa
    _ok("all modules import (incl. chat + session + edit_layer)")


def check_tool_registry() -> None:
    print("[2/20] tool registry")
    from .tools import TOOL_HANDLERS, TOOL_SCHEMAS

    expected = {"switch_artifact_type", "propose_design_spec",
                "generate_background", "generate_image",
                "render_text_layer", "edit_layer",
                "fetch_brand_asset", "composite", "critique", "finalize",
                "ingest_document"}
    schema_names = {s["name"] for s in TOOL_SCHEMAS}
    handler_names = set(TOOL_HANDLERS.keys())
    missing = expected - (schema_names & handler_names)
    extra = (schema_names | handler_names) - expected
    if missing:
        _fail(f"missing tools: {missing}")
    if extra:
        _fail(f"unexpected tools: {extra}")
    for s in TOOL_SCHEMAS:
        for k in ("name", "description", "input_schema"):
            if k not in s:
                _fail(f"tool '{s.get('name','?')}' missing key '{k}'")
        if s["input_schema"].get("type") != "object":
            _fail(f"tool '{s['name']}' input_schema.type != 'object'")
    # switch_artifact_type should be first in TOOL_SCHEMAS (pedagogical ordering)
    if TOOL_SCHEMAS[0]["name"] != "switch_artifact_type":
        _fail(f"switch_artifact_type should be first in TOOL_SCHEMAS; got {TOOL_SCHEMAS[0]['name']}")
    _ok(f"{len(expected)} tools wired (schemas + handlers): {sorted(schema_names)}")


def check_pydantic_roundtrip() -> None:
    """v2 trajectory schema roundtrip — covers the new DistillTrajectory
    plus all step types (input / reasoning / tool_call / tool_result /
    finalize), ToolResultRecord (success + error variants), and
    ThinkingBlockRecord (plain + redacted)."""
    print("[3/20] pydantic schema round-trip (v2)")
    plain_thinking = ThinkingBlockRecord(
        thinking="I should declare poster type then propose a 3:4 spec.",
        signature="sig_opaque_anthropic",
        is_redacted=False,
    )
    redacted_thinking = ThinkingBlockRecord(
        thinking="",
        signature="enc_payload_bytes",
        is_redacted=True,
    )
    ok_result = ToolResultRecord(
        status="ok",
        payload={"layer_id": "L1", "sha256": "deadbeef", "width": 1536, "height": 2048},
    )
    err_result = ToolResultRecord(
        status="error",
        error_message="DesignSpec validation failed: missing canvas.w_px",
        error_category="validation",
        payload={"hint": None},
    )
    crit = CritiqueResult(iteration=1, verdict="pass", score=0.86,
                          issues=[], rationale="reads as a coherent poster")
    crit_payload = ToolResultRecord(status="ok", payload=crit.model_dump(mode="json"))

    trace = [
        AgentTraceStep(step_idx=1, actor="user", type="input",
                       text="design a 3:4 poster"),
        AgentTraceStep(step_idx=2, actor="planner", type="reasoning",
                       thinking_blocks=[plain_thinking, redacted_thinking],
                       model="claude-opus-4-7", stop_reason="tool_use",
                       usage={"input": 1234, "output": 200, "cache_read": 1024,
                              "cache_create": 0}),
        AgentTraceStep(step_idx=3, actor="planner", type="tool_call",
                       tool_use_id="toolu_x", tool_name="render_text_layer",
                       tool_args={"layer_id": "L1", "text": "国宝回家"},
                       model="claude-opus-4-7"),
        AgentTraceStep(step_idx=4, actor="tool", type="tool_result",
                       tool_use_id="toolu_x", tool_name="render_text_layer",
                       tool_result=ok_result),
        AgentTraceStep(step_idx=5, actor="planner", type="tool_call",
                       tool_use_id="toolu_y", tool_name="propose_design_spec",
                       tool_args={"design_spec": "stringified-by-mistake"}),
        AgentTraceStep(step_idx=6, actor="tool", type="tool_result",
                       tool_use_id="toolu_y", tool_name="propose_design_spec",
                       tool_result=err_result),
        AgentTraceStep(step_idx=7, actor="planner", type="tool_call",
                       tool_use_id="toolu_z", tool_name="critique"),
        AgentTraceStep(step_idx=8, actor="tool", type="tool_result",
                       tool_use_id="toolu_z", tool_name="critique",
                       tool_result=crit_payload),
        AgentTraceStep(step_idx=9, actor="critic", type="reasoning",
                       thinking_blocks=[plain_thinking],
                       model="claude-opus-4-7"),
        AgentTraceStep(step_idx=10, actor="planner", type="finalize",
                       text="all done"),
    ]

    traj = DistillTrajectory(
        run_id="smoke",
        brief="design a 3:4 poster",
        agent_trace=trace,
        final_reward=0.86,
        terminal_status="pass",
        metadata=TrainingMetadata(
            schema_version="v2",
            planner_model="claude-opus-4-7",
            critic_model="claude-opus-4-7",
            image_model="gemini-3-pro-image-preview",
            planner_thinking_budget=10000,
            critic_thinking_budget=10000,
            interleaved_thinking=True,
            total_input_tokens=1234,
            total_output_tokens=200,
            total_cache_read_tokens=1024,
            total_cache_creation_tokens=0,
            estimated_cost_usd=2.34,
            wall_time_s=187.5,
            source="agent_run",
        ),
    )
    dumped = traj.model_dump(mode="json")
    serialized = json.dumps(dumped, ensure_ascii=False)
    try:
        reloaded = DistillTrajectory.model_validate(json.loads(serialized))
    except ValidationError as e:
        _fail(f"DistillTrajectory round-trip: {e.errors()[:3]}")

    if reloaded.metadata.schema_version != "v2":
        _fail(f"schema_version != v2: {reloaded.metadata.schema_version}")
    if reloaded.final_reward != 0.86:
        _fail(f"final_reward not preserved: {reloaded.final_reward}")
    if reloaded.terminal_status != "pass":
        _fail(f"terminal_status not preserved: {reloaded.terminal_status}")

    reasoning_steps = [s for s in reloaded.agent_trace if s.type == "reasoning"]
    if len(reasoning_steps) != 2:
        _fail(f"expected 2 reasoning steps, got {len(reasoning_steps)}")
    planner_reasoning = next(s for s in reasoning_steps if s.actor == "planner")
    if not planner_reasoning.thinking_blocks or len(planner_reasoning.thinking_blocks) != 2:
        _fail("thinking_blocks lost in roundtrip")
    if planner_reasoning.thinking_blocks[0].thinking != plain_thinking.thinking:
        _fail("thinking text corrupted")
    if not planner_reasoning.thinking_blocks[1].is_redacted:
        _fail("redacted thinking flag lost")
    if planner_reasoning.usage["cache_read"] != 1024:
        _fail("usage.cache_read lost")

    err_step = next(s for s in reloaded.agent_trace
                    if s.tool_result and s.tool_result.status == "error")
    if not err_step.tool_result.error_message:
        _fail("error_message stripped from tool_result")
    if err_step.tool_result.error_category != "validation":
        _fail(f"error_category lost: {err_step.tool_result.error_category}")

    crit_step = next(s for s in reloaded.agent_trace
                     if s.tool_name == "critique" and s.type == "tool_result")
    if crit_step.tool_result.payload.get("score") != 0.86:
        _fail("critique payload (verdict/score/issues) lost")

    _ok(f"DistillTrajectory v2 + 10 trace steps roundtrip ({len(serialized)} bytes)")


def check_fonts() -> None:
    print("[4/20] fonts")
    from PIL import ImageFont
    from .config import REPO_ROOT
    for fname in ("NotoSansSC-Bold.otf", "NotoSerifSC-Bold.otf"):
        path = REPO_ROOT / "assets" / "fonts" / fname
        if not path.exists():
            _fail(f"missing font: {path}")
        try:
            ImageFont.truetype(str(path), size=80)
        except Exception as e:
            _fail(f"font {fname} load failed: {e}")
        _ok(f"loaded {fname} ({path.stat().st_size // 1024} KB)")


def check_composite_no_api() -> None:
    """Build a fake background + 2 real text layers, run composite end-to-end.

    Also exercises switch_artifact_type → propose_design_spec plumbing
    (artifact_type fallback from ctx.state when spec omits it).
    """
    print("[5/20] composite (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.render_text_layer import render_text_layer
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub",
        anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7",
        critic_model="claude-opus-4-7",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke")

    # First: switch_artifact_type — verifies ctx.state update + valid type
    obs = switch_artifact_type({"type": "poster"}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"switch_artifact_type: {(obs.error_message or str(obs.payload))}")
    if ctx.state.get("artifact_type") != "poster":
        _fail(f"ctx.state.artifact_type not set; got {ctx.state.get('artifact_type')}")

    # Reject invalid type (paranoia — catches the enum drift)
    bad = switch_artifact_type({"type": "billboard"}, ctx=ctx)
    if bad.status != "error":
        _fail(f"switch_artifact_type should reject 'billboard'; got status={bad.status}")

    # Spec omits artifact_type on purpose — should fall back to ctx.state value
    spec_args = {"design_spec": {
        "brief": "smoke test poster",
        "canvas": {"w_px": 768, "h_px": 1024, "dpi": 150,
                   "aspect_ratio": "3:4", "color_mode": "RGB"},
        "palette": ["#1a1a1a", "#ffffff", "#a02018"],
        "typography": {"title_font": "NotoSerifSC-Bold",
                       "subtitle_font": "NotoSansSC-Bold"},
        "mood": ["test"],
        "composition_notes": "smoke",
        "layer_graph": [],
    }}
    obs = propose_design_spec(spec_args, ctx=ctx)
    if obs.status != "ok":
        _fail(f"propose_design_spec: {(obs.error_message or str(obs.payload))}")

    # Fallback check: spec omitted artifact_type → should inherit ctx.state value
    stored_spec = ctx.state["design_spec"]
    if stored_spec.artifact_type.value != "poster":
        _fail(f"artifact_type fallback failed; got {stored_spec.artifact_type.value!r}")

    bg_path = layers_dir / "bg_smoke.png"
    Image.new("RGB", (768, 1024), (28, 14, 10)).save(bg_path)
    ctx.state["rendered_layers"]["L0_bg"] = {
        "layer_id": "L0_bg", "name": "background", "kind": "background", "z_index": 0,
        "bbox": {"x": 0, "y": 0, "w": 768, "h": 1024},
        "src_path": str(bg_path), "prompt": "(stub)",
        "aspect_ratio": "3:4", "image_size": "1K",
        "safe_zones": [], "sha256": "stub",
    }

    obs = render_text_layer({
        "layer_id": "L1_title", "name": "title", "text": "国宝回家",
        "font_family": "NotoSerifSC-Bold", "font_size_px": 110, "fill": "#fafafa",
        "bbox": {"x": 48, "y": 80, "w": 672, "h": 180}, "align": "center",
        "z_index": 1,
        "effects": {"shadow": {"color": "#00000080", "dx": 0, "dy": 4, "blur": 12}},
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"render_text_layer title: {(obs.error_message or str(obs.payload))}")

    obs = render_text_layer({
        "layer_id": "L2_subtitle", "name": "subtitle",
        "text": "National Treasures Return Home",
        "font_family": "NotoSansSC-Bold", "font_size_px": 36, "fill": "#c9a45a",
        "bbox": {"x": 64, "y": 280, "w": 640, "h": 60}, "align": "center",
        "z_index": 2,
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"render_text_layer subtitle: {(obs.error_message or str(obs.payload))}")

    obs = composite({}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"composite: {(obs.error_message or str(obs.payload))}")

    comp = ctx.state["composition"]
    for label, p in [("PSD", comp.psd_path), ("SVG", comp.svg_path),
                     ("HTML", comp.html_path), ("preview", comp.preview_path)]:
        if p is None:
            _fail(f"{label} path missing on CompositionArtifacts")
        path = Path(p)
        if not path.exists() or path.stat().st_size == 0:
            _fail(f"{label} not written: {p}")
        _ok(f"{label} {path.name} ({path.stat().st_size // 1024} KB)")


def check_svg_text_is_vector() -> None:
    print("[6/20] SVG + HTML content (vector text, contenteditable, inline fonts)")
    from .config import REPO_ROOT
    out_dir = REPO_ROOT / "out" / "smoke"

    # --- SVG -------------------------------------------------------------
    svg_path = out_dir / "poster.svg"
    if not svg_path.exists():
        _fail("smoke SVG not found — composite step likely failed")
    svg = svg_path.read_text(encoding="utf-8")
    if "<text" not in svg:
        _fail("SVG has no <text> element — text was rasterized!")
    if "国宝回家" not in svg:
        _fail("Chinese title not present as vector text in SVG")
    if "@font-face" not in svg:
        print("  warn  SVG missing @font-face (font subsetting may have failed; "
              "will rely on system fonts)")
    else:
        _ok("SVG: @font-face block embedded with subsetted WOFF2")
    _ok("SVG contains <text>国宝回家</text> as real vector")

    # --- HTML ------------------------------------------------------------
    html_path = out_dir / "poster.html"
    if not html_path.exists():
        _fail("smoke HTML not found — html_renderer step likely failed")
    html_text = html_path.read_text(encoding="utf-8")

    required_markers = {
        "canvas container":     '<div class="canvas"',
        "contenteditable text": 'contenteditable="true"',
        "inline fonts":         "@font-face",
        "WOFF2 data URI":       "data:font/woff2;base64,",
        "data-layer-id attr":   "data-layer-id=",
        "data-kind attr":       "data-kind=",
        "Chinese title text":   "国宝回家",
        "bg data URI":          "data:image/png;base64,",
        "generator meta":       '<meta name="generator" content="LongcatDesign"',
        # Edit toolbar markers (v1.0 #6 edit UX)
        "bbox data attrs":      'data-bbox-x=',
        "font-size data attr":  'data-font-size-px=',
        "fill data attr":       'data-fill=',
        "font-family data attr": 'data-font-family=',
        "drag handle":          'class="ld-drag-handle"',
        "toolbar container":    'class="ld-toolbar"',
        "font select":          'id="ld-family"',
        "size input":           'id="ld-size"',
        "color input":          'id="ld-color"',
        "save button":          'id="ld-save"',
        "save modal":           'id="ld-modal-backdrop"',
        "copy button":          'id="ld-copy"',
        "download button":      'id="ld-download"',
        "apply-edits hint":     "longcat-design apply-edits",
    }
    for label, needle in required_markers.items():
        if needle not in html_text:
            _fail(f"HTML missing expected marker — {label}: {needle!r}")
    _ok(f"HTML poster.html ({html_path.stat().st_size // 1024} KB) — "
        f"all {len(required_markers)} markers present "
        "(canvas / contenteditable / inline fonts+images / data-* attrs)")

    # Structure sanity: exactly one <canvas> div and at least 2 text layers
    if html_text.count('class="canvas"') != 1:
        _fail(f"HTML should contain exactly 1 .canvas div; got "
              f"{html_text.count('class=\"canvas\"')}")
    if html_text.count('class="layer text"') < 2:
        _fail(f"HTML should contain ≥2 text layers (smoke has title + subtitle); "
              f"got {html_text.count('class=\"layer text\"')}")
    _ok("HTML structure: 1 .canvas + ≥2 text layers, all attributed")


def check_chat_session_roundtrip() -> None:
    """ChatSession pydantic + save/load cycle — no API calls."""
    print("[7/20] chat session save/load")
    from .config import REPO_ROOT
    from .session import (
        ChatMessage, ChatSession, TrajectoryRef,
        load_session, new_session_id, save_session, session_path, list_sessions,
    )
    from .schema import ArtifactType

    tmp_dir = REPO_ROOT / "out" / "smoke_sessions"
    tmp_dir.mkdir(parents=True, exist_ok=True)

    sid = new_session_id()
    if not sid.startswith("session_"):
        _fail(f"new_session_id() should start with 'session_'; got {sid!r}")

    # Build a realistic session with 1 user + 1 assistant msg + 1 trajectory ref
    session = ChatSession(session_id=sid)
    session.append_user("design a 3:4 poster for 国宝回家")
    ref = TrajectoryRef(
        run_id="20260418-smoke-test",
        artifact_type=ArtifactType.POSTER,
        created_at=session.created_at,
        trajectory_path="/tmp/fake-trajectory.json",
        preview_path="/tmp/fake-preview.png",
        psd_path="/tmp/fake.psd",
        svg_path="/tmp/fake.svg",
        n_layers=5,
        verdict="pass",
        score=0.86,
        cost_usd=1.41,
        wall_s=100.0,
    )
    session.trajectories.append(ref)
    session.append_assistant("produced poster · 5 layers · pass(0.86)",
                              trajectory_id=ref.run_id)

    path = save_session(session, tmp_dir)
    if not path.exists():
        _fail(f"save_session did not write: {path}")
    _ok(f"saved {path.name} ({path.stat().st_size} bytes)")

    loaded = load_session(tmp_dir, sid)
    if loaded.session_id != sid:
        _fail(f"round-trip: session_id mismatch; {loaded.session_id} != {sid}")
    if len(loaded.message_history) != 2:
        _fail(f"round-trip: message count wrong; {len(loaded.message_history)} != 2")
    if len(loaded.trajectories) != 1 or loaded.trajectories[0].verdict != "pass":
        _fail("round-trip: trajectory ref did not survive")
    _ok(f"round-trip: {len(loaded.message_history)} msgs, "
        f"{len(loaded.trajectories)} trajectory ref(s), cost ${loaded.total_cost_usd()}")

    listing = list_sessions(tmp_dir)
    if sid not in [s[0] for s in listing]:
        _fail("list_sessions did not include the newly-saved session")
    _ok(f"list_sessions finds {len(listing)} session(s) incl. this one")

    # Clean up smoke session file
    path.unlink(missing_ok=True)


def check_edit_layer_no_api() -> None:
    """edit_layer semantics — subset-merge, delegates re-render, refuses non-text."""
    print("[8/20] edit_layer (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.edit_layer import edit_layer
    from .tools.propose_design_spec import propose_design_spec
    from .tools.render_text_layer import render_text_layer
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_edit"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub",
        anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7",
        critic_model="claude-opus-4-7",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke_edit")

    # Seed: switch type + spec + one text layer + one fake bg layer
    if switch_artifact_type({"type": "poster"}, ctx=ctx).status != "ok":
        _fail("seed: switch_artifact_type")

    spec_args = {"design_spec": {
        "brief": "edit_layer smoke",
        "canvas": {"w_px": 768, "h_px": 1024, "dpi": 150,
                   "aspect_ratio": "3:4", "color_mode": "RGB"},
        "palette": ["#1a1a1a", "#ffffff"],
        "typography": {"title_font": "NotoSerifSC-Bold",
                       "subtitle_font": "NotoSansSC-Bold"},
        "mood": ["test"], "composition_notes": "", "layer_graph": [],
    }}
    if propose_design_spec(spec_args, ctx=ctx).status != "ok":
        _fail("seed: propose_design_spec")

    # Fake bg layer (non-text — should be refused by edit_layer)
    ctx.state["rendered_layers"]["L0_bg"] = {
        "layer_id": "L0_bg", "name": "background", "kind": "background",
        "z_index": 0, "bbox": {"x": 0, "y": 0, "w": 768, "h": 1024},
        "src_path": "/tmp/fake.png", "sha256": "stub",
    }

    # Real text layer via render_text_layer
    obs = render_text_layer({
        "layer_id": "L1_title", "name": "title", "text": "原标题",
        "font_family": "NotoSerifSC-Bold", "font_size_px": 100, "fill": "#000000",
        "bbox": {"x": 48, "y": 80, "w": 672, "h": 180}, "align": "center",
        "z_index": 1,
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"seed: render_text_layer: {(obs.error_message or str(obs.payload))}")

    before = ctx.state["rendered_layers"]["L1_title"]
    before_sha = before["sha256"]
    before_path = Path(before["src_path"])
    if not before_path.exists():
        _fail("seed: initial PNG missing before edit")

    # --- Happy path: multi-field diff -------------------------------------
    obs = edit_layer({
        "layer_id": "L1_title",
        "diff": {"text": "新标题！", "font_size_px": 140, "fill": "#ff0000"},
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"edit_layer happy path: status={obs.status} summary={(obs.error_message or str(obs.payload))}")

    after = ctx.state["rendered_layers"]["L1_title"]
    if after["font_size_px"] != 140:
        _fail(f"font_size_px not applied: got {after['font_size_px']}")
    if after["fill"] != "#ff0000":
        _fail(f"fill not applied: got {after['fill']}")
    if after["text"] != "新标题！":
        _fail(f"text not applied: got {after['text']!r}")
    # Unchanged fields preserved
    if after["name"] != "title":
        _fail(f"name should be preserved: got {after['name']!r}")
    if after["align"] != "center":
        _fail(f"align should be preserved: got {after['align']!r}")
    if after["bbox"]["w"] != 672:
        _fail(f"bbox should be preserved: got {after['bbox']}")
    # PNG should have been rewritten (different content → different sha)
    if after["sha256"] == before_sha:
        _fail("PNG sha256 unchanged after edit — render_text_layer didn't re-run")
    _ok("happy path: text+font_size_px+fill applied, other fields preserved, PNG rewritten")

    # --- Partial bbox merge ------------------------------------------------
    obs = edit_layer({
        "layer_id": "L1_title",
        "diff": {"bbox": {"y": 200}},  # only y changes
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"bbox partial merge: {(obs.error_message or str(obs.payload))}")
    bbox = ctx.state["rendered_layers"]["L1_title"]["bbox"]
    if not (bbox["x"] == 48 and bbox["y"] == 200 and bbox["w"] == 672 and bbox["h"] == 180):
        _fail(f"bbox partial merge broken: {bbox}")
    _ok("partial bbox merge: y updated, x/w/h preserved")

    # --- Missing layer_id --------------------------------------------------
    obs = edit_layer({"layer_id": "nope", "diff": {"text": "x"}}, ctx=ctx)
    if obs.status != "error" or obs.error_category != "not_found":
        _fail(f"unknown layer should return error/not_found, got status={obs.status} cat={obs.error_category}")
    _ok("unknown layer_id → error/not_found")

    # --- Non-text layer rejected ------------------------------------------
    obs = edit_layer({"layer_id": "L0_bg", "diff": {"text": "x"}}, ctx=ctx)
    if obs.status != "error":
        _fail(f"bg layer edit should error, got {obs.status}")
    if obs.error_category != "validation":
        _fail(f"bg-error category should be validation; got: {obs.error_category}")
    _ok("non-text layer → error/validation")

    # --- Empty / unknown diff fields --------------------------------------
    if edit_layer({"layer_id": "L1_title", "diff": {}}, ctx=ctx).status != "error":
        _fail("empty diff should error")
    if edit_layer({"layer_id": "L1_title",
                   "diff": {"color": "#ff0000"}}, ctx=ctx).status != "error":
        _fail("unknown diff field should error (caught 'color' instead of 'fill')")
    _ok("empty diff + unknown field both rejected")


def check_apply_edits_roundtrip() -> None:
    """HTML → apply-edits → new PSD/SVG/HTML/preview with same semantic content."""
    print("[9/20] apply-edits round-trip (no API)")
    from .apply_edits import apply_edits
    from .config import REPO_ROOT, Settings

    src_html = REPO_ROOT / "out" / "smoke" / "poster.html"
    if not src_html.exists():
        _fail("smoke HTML missing — [5/16] composite step must run first")

    # Simulate a user edit: bump font size + change color on the title.
    # Re-emit with a changed data-font-size-px/data-fill so round-trip should
    # reflect the new values in the regenerated layer graph.
    raw = src_html.read_text(encoding="utf-8")
    edited = (raw
              .replace('data-font-size-px="110"', 'data-font-size-px="140"')
              .replace('data-fill="#fafafa"', 'data-fill="#ff3366"'))
    if edited == raw:
        _fail("could not seed edits into smoke HTML — assumed markers missing")
    edited_html_path = REPO_ROOT / "out" / "smoke_apply" / "edited.html"
    edited_html_path.parent.mkdir(parents=True, exist_ok=True)
    edited_html_path.write_text(edited, encoding="utf-8")

    settings = Settings(
        anthropic_api_key="sk-stub",
        anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7",
        critic_model="claude-opus-4-7",
    )
    out_dir = REPO_ROOT / "out" / "smoke_apply" / "restored"

    traj, traj_path, run_dir, restored_ids, skipped = apply_edits(
        edited_html_path, settings=settings, out_dir=out_dir,
    )

    # --- v2 trajectory assertions ----------------------------------------
    if traj.metadata.source != "apply_edits":
        _fail(f"metadata.source != 'apply_edits'; got {traj.metadata.source!r}")
    if traj.agent_trace:
        _fail(f"apply-edits trajectory should have empty agent_trace; got {len(traj.agent_trace)}")
    if traj.terminal_status != "abort":
        _fail(f"apply-edits should set terminal_status=abort; got {traj.terminal_status}")
    _ok(f"trajectory: source=apply_edits, terminal=abort, "
        f"{len(restored_ids)} layers restored")

    # --- artifact files exist --------------------------------------------
    for fname in ("poster.psd", "poster.svg", "poster.html", "preview.png"):
        p = run_dir / fname
        if not p.exists() or p.stat().st_size == 0:
            _fail(f"{fname} not written: {p}")
    _ok("PSD+SVG+HTML+preview all regenerated in new run_dir")

    # --- edits landed in re-rendered run dir -----------------------------
    # v2 trajectory has no layer_graph; verify the rendered HTML on disk
    # contains the edits instead.
    rendered_html = (run_dir / "poster.html").read_text(encoding="utf-8")
    if "140" not in rendered_html:
        _fail("expected font_size 140 to appear in re-rendered HTML")
    if "#ff3366" not in rendered_html.lower():
        _fail("expected fill #ff3366 to appear in re-rendered HTML")
    _ok("edits preserved in re-rendered HTML: font_size 140, fill #ff3366")

    # --- bg was decoded from data: URI ----------------------------------
    layers_dir = run_dir / "layers"
    bgs = list(layers_dir.glob("bg_*.png")) if layers_dir.exists() else []
    if bgs:
        _ok(f"bg decoded from data URI → {bgs[0].name} ({bgs[0].stat().st_size} B)")


def check_landing_mode() -> None:
    """Landing end-to-end: section-tree spec → HTML + preview → apply-edits roundtrip.

    v1.3: expanded fixture to 4 sections (hero + features + pricing + cta)
    with a `kind="cta"` child, plus a footer-variant section to exercise
    the `<footer>` auto-upgrade. 4 sections triggers auto-nav, and the
    round-trip must preserve CTA nodes with href + variant.
    """
    print("[10/20] landing mode (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.switch_artifact_type import switch_artifact_type
    from .apply_edits import apply_edits
    from .schema import ArtifactType

    out_dir = REPO_ROOT / "out" / "smoke_landing"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7", critic_model="claude-opus-4-7",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-landing")

    # --- switch + propose landing spec ----------------------------------
    if switch_artifact_type({"type": "landing"}, ctx=ctx).status != "ok":
        _fail("switch_artifact_type(landing)")

    spec_args = {"design_spec": {
        "brief": "LongcatDesign v1.0 landing",
        "artifact_type": "landing",
        "canvas": {"w_px": 1200, "h_px": 2400, "dpi": 96,
                   "aspect_ratio": "1:2", "color_mode": "RGB"},
        "palette": ["#0f172a", "#f8fafc"],
        "typography": {"title_font": "NotoSerifSC-Bold", "body_font": "NotoSansSC-Bold"},
        "mood": ["minimal"], "composition_notes": "dark hero, light features, dark cta",
        "layer_graph": [
            {"layer_id": "S1", "name": "hero", "kind": "section", "z_index": 1,
             "children": [
                 {"layer_id": "H1", "name": "hero_headline", "kind": "text", "z_index": 1,
                  "text": "LongcatDesign", "font_family": "NotoSerifSC-Bold",
                  "font_size_px": 96, "align": "center",
                  "effects": {"fill": "#f8fafc"}},
                 {"layer_id": "H2", "name": "hero_subhead", "kind": "text", "z_index": 2,
                  "text": "Open source conversational design agent.",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 28,
                  "align": "center", "effects": {"fill": "#94a3b8"}},
             ]},
            {"layer_id": "S2", "name": "features", "kind": "section", "z_index": 2,
             "children": [
                 {"layer_id": "F1", "name": "features_title", "kind": "text", "z_index": 1,
                  "text": "Three outputs, one conversation",
                  "font_family": "NotoSerifSC-Bold", "font_size_px": 48,
                  "effects": {"fill": "#0f172a"}},
                 {"layer_id": "F2", "name": "feature_1", "kind": "text", "z_index": 2,
                  "text": "Poster · PSD + SVG + HTML, fully layered and editable.",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 20,
                  "effects": {"fill": "#334155"}},
             ]},
            {"layer_id": "S3", "name": "pricing", "kind": "section", "z_index": 3,
             "children": [
                 {"layer_id": "P1", "name": "pricing_title", "kind": "text", "z_index": 1,
                  # v2.3 — embed inline + display math so the KaTeX path gets
                  # exercised in smoke. Free-form strings, no rendering happens
                  # at write time; KaTeX is client-side JS.
                  "text": "Pricing model: cost is $O(n)$ per token; total budget $$B = n \\cdot c$$",
                  "font_family": "NotoSerifSC-Bold", "font_size_px": 40,
                  "align": "center", "effects": {"fill": "#0f172a"}},
             ]},
            {"layer_id": "S4", "name": "cta", "kind": "section", "z_index": 4,
             "children": [
                 {"layer_id": "C1", "name": "cta_text", "kind": "text", "z_index": 1,
                  "text": "pip install longcat-design",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 36,
                  "align": "center", "effects": {"fill": "#f8fafc"}},
                 {"layer_id": "C2", "name": "cta_button", "kind": "cta", "z_index": 2,
                  "text": "Get started", "href": "#sec-features",
                  "variant": "primary"},
             ]},
            {"layer_id": "S5", "name": "footer", "kind": "section", "z_index": 5,
             "children": [
                 {"layer_id": "FT1", "name": "copyright", "kind": "text", "z_index": 1,
                  "text": "© 2026 LongcatDesign · MIT",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 14,
                  "align": "center", "effects": {"fill": "#94a3b8"}},
             ]},
        ],
    }}
    if propose_design_spec(spec_args, ctx=ctx).status != "ok":
        _fail("propose_design_spec(landing)")
    if ctx.state["design_spec"].artifact_type != ArtifactType.LANDING:
        _fail("design_spec.artifact_type not LANDING after propose")

    # --- composite landing ----------------------------------------------
    obs = composite({}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"landing composite: {(obs.error_message or str(obs.payload))}")
    comp = ctx.state["composition"]
    if comp.psd_path is not None or comp.svg_path is not None:
        _fail("landing should NOT produce PSD/SVG — got non-None paths")
    for label, p in [("HTML", comp.html_path), ("preview", comp.preview_path)]:
        if not p or not Path(p).exists() or Path(p).stat().st_size == 0:
            _fail(f"landing {label} missing: {p}")
    _ok(f"landing composite: HTML + preview, NO PSD/SVG (correct)")

    # --- HTML structure --------------------------------------------------
    html_text = Path(comp.html_path).read_text(encoding="utf-8")
    required_markers = {
        "landing container":     '<main class="ld-landing"',
        "landing mode meta":     '<meta name="ld-artifact-type" content="landing"',
        "hero section":          'data-section-variant="hero"',
        "features section":      'data-section-variant="features"',
        "cta section":           'data-section-variant="cta"',
        "contenteditable text":  'contenteditable="true"',
        "data-font-size attr":   'data-font-size-px=',
        "toolbar container":     'class="ld-toolbar"',
        "save modal":            'id="ld-modal-backdrop"',
        "landing hides drag":    "/* Drag handle hidden in landing",
        # v1.3 interactive markers
        "auto-nav header":       '<header class="ld-header"',
        "nav anchor":            'data-nav-target="sec-',
        "section id":            'id="sec-',
        "reveal attr":           'data-reveal="true"',
        "cta anchor":             '<a class="ld-cta ld-cta--primary"',
        "cta href":               'href="#sec-features"',
        "cta data-kind":          'data-kind="cta"',
        "semantic footer":        '<footer class="ld-section"',
        "interactive JS":         'IntersectionObserver',
        # v2.3 — KaTeX injected when math delimiters present in any text layer
        "katex stylesheet":       '<style id="ld-katex-css">',
        "katex init":             'renderMathInElement',
        "katex font data URI":    'data:font/woff2;base64',
        "inline math preserved":  'cost is $O(n)$',
        "display math preserved": 'B = n \\cdot c',
    }
    for label, needle in required_markers.items():
        if needle not in html_text:
            _fail(f"landing HTML missing marker — {label}: {needle!r}")
    _ok(f"landing HTML ({Path(comp.html_path).stat().st_size // 1024} KB) — "
        f"all {len(required_markers)} markers present")

    # --- apply-edits round-trip on a seeded-edit landing ---------------
    edited_html = html_text.replace(
        'data-font-size-px="96"', 'data-font-size-px="128"'
    ).replace(
        'data-fill="#f8fafc"', 'data-fill="#38bdf8"', 1  # only first occurrence
    )
    if edited_html == html_text:
        _fail("could not seed landing edits — markers missing")
    edited_path = out_dir / "edited.html"
    edited_path.write_text(edited_html, encoding="utf-8")

    traj, traj_path, run_dir, restored_ids, skipped = apply_edits(
        edited_path, settings=settings,
        out_dir=out_dir / "restored",
    )
    if traj.metadata.source != "apply_edits":
        _fail(f"landing round-trip lost source label: got {traj.metadata.source!r}")
    # v2: trajectory has no design_spec / layer_graph; verify the edits
    # landed in the regenerated HTML on disk instead.
    rendered_html = (run_dir / "index.html").read_text(encoding="utf-8")
    if "128" not in rendered_html:
        _fail("landing edit lost: font_size 128 missing from rendered HTML")
    if "#38bdf8" not in rendered_html.lower():
        _fail("landing edit lost: fill #38bdf8 missing from rendered HTML")
    # CTA must survive (data-* attrs encode href + variant in HTML).
    if 'data-href="#sec-features"' not in rendered_html:
        _fail("CTA href lost on round-trip")
    if 'data-variant="primary"' not in rendered_html:
        _fail("CTA variant lost on round-trip")
    _ok(f"landing round-trip: edits applied (hero_headline 96px → 128px, "
        f"fill→#38bdf8), CTA href+variant intact in regenerated HTML")


def check_design_system_styles() -> None:
    """Render a landing in each of the 6 bundled styles, verify the matching
    CSS got inlined and the style-specific signature tokens are present."""
    print("[11/20] design-system styles (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.switch_artifact_type import switch_artifact_type

    # A tiny per-style signature: a CSS token string that MUST appear in the
    # rendered HTML only when that style's CSS is loaded.
    style_signatures = {
        "minimalist":     "--ld-shadow-soft",
        "editorial":      "--ld-rule",
        "neubrutalism":   "--ld-shadow-hard",
        "glassmorphism":  "prefers-reduced-transparency",
        "claymorphism":   "--ld-shadow-clay",
        "liquid-glass":   "cubic-bezier(0.34, 1.56",
    }

    base_spec = {
        "brief": "design-system smoke",
        "artifact_type": "landing",
        "canvas": {"w_px": 1200, "h_px": 2400, "dpi": 96,
                   "aspect_ratio": "1:2", "color_mode": "RGB"},
        "palette": ["#0f172a", "#f8fafc"],
        "typography": {"title_font": "NotoSerifSC-Bold", "body_font": "NotoSansSC-Bold"},
        "mood": ["test"], "composition_notes": "",
        "layer_graph": [
            {"layer_id": "S1", "name": "hero", "kind": "section", "z_index": 1,
             "children": [
                 {"layer_id": "H1", "name": "hero_headline", "kind": "text",
                  "z_index": 1, "text": "Test", "font_family": "NotoSansSC-Bold",
                  "font_size_px": 80, "align": "center",
                  "effects": {"fill": "#f8fafc"}},
             ]},
            {"layer_id": "S2", "name": "features", "kind": "section", "z_index": 2,
             "children": [
                 {"layer_id": "F1", "name": "features_title", "kind": "text",
                  "z_index": 1, "text": "Feature",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 36,
                  "effects": {"fill": "#0f172a"}},
             ]},
        ],
    }

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7", critic_model="claude-opus-4-7",
    )

    for style, signature in style_signatures.items():
        out_dir = REPO_ROOT / "out" / "smoke_styles" / style
        layers_dir = out_dir / "layers"
        layers_dir.mkdir(parents=True, exist_ok=True)
        ctx = ToolContext(settings=settings, run_dir=out_dir,
                          layers_dir=layers_dir, run_id=f"smoke-style-{style}")

        if switch_artifact_type({"type": "landing"}, ctx=ctx).status != "ok":
            _fail(f"[{style}] switch_artifact_type")
        spec_args = {"design_spec": {**base_spec,
                                     "design_system": {"style": style}}}
        if propose_design_spec(spec_args, ctx=ctx).status != "ok":
            _fail(f"[{style}] propose_design_spec")
        if composite({}, ctx=ctx).status != "ok":
            _fail(f"[{style}] composite")

        comp = ctx.state["composition"]
        html = Path(comp.html_path).read_text(encoding="utf-8")

        if f'<meta name="ld-design-system" content="{style}">' not in html:
            _fail(f"[{style}] missing design-system meta tag")
        if f'data-ld-style="{style}"' not in html:
            _fail(f"[{style}] missing data-ld-style attr")
        if signature not in html:
            _fail(f"[{style}] CSS signature {signature!r} not inlined — "
                  "assets/design-systems/<style>.css probably not loaded")
        _ok(f"  {style:<14} ({Path(comp.html_path).stat().st_size // 1024} KB) — "
            f"meta + data-attr + CSS signature all present")

    # Final: accent_color override landed in the CSS
    out_dir = REPO_ROOT / "out" / "smoke_styles" / "accent_override"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-accent")
    switch_artifact_type({"type": "landing"}, ctx=ctx)
    spec_args = {"design_spec": {**base_spec,
                                 "design_system": {"style": "minimalist",
                                                   "accent_color": "#fe11ba"}}}
    propose_design_spec(spec_args, ctx=ctx)
    composite({}, ctx=ctx)
    html = Path(ctx.state["composition"].html_path).read_text(encoding="utf-8")
    if "--ld-accent: #fe11ba" not in html:
        _fail("accent_color override did not reach CSS (--ld-accent: #fe11ba missing)")
    _ok("accent_color override propagated to --ld-accent token")


def check_landing_with_images() -> None:
    """Landing mode with image children in sections. No NBP call —
    pre-stages a stub PNG in rendered_layers and asserts the renderer
    inlines it + apply-edits round-trips the image layer."""
    print("[12/20] landing with images (no API)")
    from .apply_edits import apply_edits
    from .config import REPO_ROOT, Settings
    from .schema import ArtifactType
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_landing_img"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Stub a hero image PNG so generate_image doesn't need to be called.
    hero_img_path = layers_dir / "img_H0_stub.png"
    Image.new("RGB", (600, 800), (244, 200, 180)).save(hero_img_path)
    feat_img_path = layers_dir / "img_F_stub.png"
    Image.new("RGB", (256, 256), (220, 214, 247)).save(feat_img_path)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7", critic_model="claude-opus-4-7",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-landing-img")

    # Seed rendered_layers as if generate_image was called
    ctx.state["rendered_layers"]["H0_img"] = {
        "layer_id": "H0_img", "name": "hero_image", "kind": "image",
        "z_index": 1, "bbox": None, "src_path": str(hero_img_path),
        "prompt": "(stub)", "aspect_ratio": "3:4",
        "image_size": "2K", "sha256": "stub",
    }
    ctx.state["rendered_layers"]["F1_img"] = {
        "layer_id": "F1_img", "name": "feature_1_icon", "kind": "image",
        "z_index": 1, "bbox": None, "src_path": str(feat_img_path),
        "prompt": "(stub)", "aspect_ratio": "1:1",
        "image_size": "1K", "sha256": "stub",
    }

    if switch_artifact_type({"type": "landing"}, ctx=ctx).status != "ok":
        _fail("switch_artifact_type")

    spec_args = {"design_spec": {
        "brief": "landing with images test",
        "artifact_type": "landing",
        "design_system": {"style": "claymorphism"},
        "canvas": {"w_px": 1200, "h_px": 2400, "dpi": 96,
                   "aspect_ratio": "1:2", "color_mode": "RGB"},
        "palette": ["#f4f0ea", "#3a2f4a"],
        "typography": {}, "mood": ["test"], "composition_notes": "",
        "layer_graph": [
            {"layer_id": "S1", "name": "hero", "kind": "section", "z_index": 1,
             "children": [
                 {"layer_id": "H0_img", "name": "hero_image", "kind": "image",
                  "z_index": 1, "aspect_ratio": "3:4"},
                 {"layer_id": "H1", "name": "hero_headline", "kind": "text",
                  "z_index": 2, "text": "Test Hero",
                  "font_family": "NotoSerifSC-Bold", "font_size_px": 88,
                  "align": "center", "effects": {"fill": "#3a2f4a"}},
             ]},
            {"layer_id": "S2", "name": "features", "kind": "section", "z_index": 2,
             "children": [
                 {"layer_id": "F1_img", "name": "feature_1_icon", "kind": "image",
                  "z_index": 1, "aspect_ratio": "1:1"},
                 {"layer_id": "F1", "name": "feature_1", "kind": "text",
                  "z_index": 2, "text": "First feature copy.",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 20,
                  "effects": {"fill": "#3a2f4a"}},
             ]},
        ],
    }}
    if propose_design_spec(spec_args, ctx=ctx).status != "ok":
        _fail("propose_design_spec")
    if composite({}, ctx=ctx).status != "ok":
        _fail("composite (with images)")

    comp = ctx.state["composition"]
    html_text = Path(comp.html_path).read_text(encoding="utf-8")

    required = {
        "figure.layer.image":     '<figure class="layer image"',
        "<img tag":               "<img src=\"data:image/",
        "data-has-image":         'data-has-image="true"',
        "hero image layer_id":    'data-layer-id="H0_img"',
        "feature image layer_id": 'data-layer-id="F1_img"',
        "aspect-ratio data":      'data-aspect-ratio=',
    }
    for label, needle in required.items():
        if needle not in html_text:
            _fail(f"landing+images HTML missing — {label}: {needle!r}")
    _ok(f"landing HTML ({Path(comp.html_path).stat().st_size // 1024} KB) — "
        f"2 <figure> image layers inlined with data URIs")

    # Apply-edits round-trip: reparse HTML, rebuild section tree with images
    traj, _, restored_dir, restored_ids, skipped = apply_edits(
        Path(comp.html_path), settings=settings,
        out_dir=out_dir / "restored",
    )
    # v2: layer info no longer in trajectory; verify the regenerated run_dir
    # has the image PNGs decoded from data: URIs onto disk.
    img_files = sorted((restored_dir / "layers").glob("img_*.png")) \
        if (restored_dir / "layers").exists() else []
    if len(img_files) < 2:
        _fail(f"round-trip lost images: expected ≥2 PNG files, got {len(img_files)}")
    for img in img_files:
        if img.stat().st_size == 0:
            _fail(f"round-trip image is empty: {img}")
    _ok(f"round-trip: {len(img_files)} image layers decoded from data: URI "
        f"into {restored_dir}/layers/")


def check_deck_mode() -> None:
    """Deck end-to-end: slide-tree spec → PPTX + per-slide PNGs + preview grid.
    No API — python-pptx writes a real .pptx that we reopen + verify."""
    print("[13/20] deck mode (no API)")
    from pptx import Presentation as _Reopen

    from .config import REPO_ROOT, Settings
    from .schema import ArtifactType
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_deck"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Stub a hero image PNG that one slide embeds, so we exercise the image path.
    hero_png = layers_dir / "img_slide02_hero.png"
    Image.new("RGB", (960, 540), (80, 140, 220)).save(hero_png)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-deck")

    if switch_artifact_type({"type": "deck"}, ctx=ctx).status != "ok":
        _fail("switch_artifact_type(deck)")

    # Pre-stage the image in rendered_layers so _hydrate_deck_image_srcs can
    # pick it up (mirrors the real two-step generate_image → composite flow).
    ctx.state["rendered_layers"]["img_slide02_hero"] = {
        "layer_id": "img_slide02_hero",
        "name": "hero_image",
        "kind": "image",
        "z_index": 5,
        "bbox": None,
        "src_path": str(hero_png),
        "prompt": "",
        "aspect_ratio": "16:9",
        "image_size": "960x540",
        "sha256": "stub",
    }

    spec = {
        "brief": "Pitch deck smoke — 3 slides: cover, problem with hero image, outro.",
        "artifact_type": "deck",
        "canvas": {"w_px": 1920, "h_px": 1080, "dpi": 96,
                   "aspect_ratio": "16:9", "color_mode": "RGB"},
        "palette": ["#0f172a", "#f8fafc", "#38bdf8"],
        "typography": {"title_font": "NotoSerifSC-Bold",
                       "body_font": "NotoSansSC-Bold"},
        "mood": ["clean", "investor-ready"],
        "composition_notes": "16:9 pitch deck. Title top, body or image below.",
        "layer_graph": [
            {
                "layer_id": "slide_01", "name": "cover", "kind": "slide",
                "z_index": 1,
                "children": [
                    {"layer_id": "slide_01_title", "name": "title",
                     "kind": "text", "z_index": 10,
                     "bbox": {"x": 120, "y": 420, "w": 1680, "h": 160},
                     "text": "MilkCloud",
                     "font_family": "NotoSerifSC-Bold",
                     "font_size_px": 96, "align": "left",
                     "effects": {"fill": "#0f172a"}},
                    {"layer_id": "slide_01_tagline", "name": "tagline",
                     "kind": "text", "z_index": 10,
                     "bbox": {"x": 120, "y": 620, "w": 1680, "h": 80},
                     "text": "the calmest bubble tea brand",
                     "font_family": "NotoSansSC-Bold",
                     "font_size_px": 40, "align": "left",
                     "effects": {"fill": "#64748b"}},
                ],
            },
            {
                "layer_id": "slide_02", "name": "problem_with_image",
                "kind": "slide", "z_index": 2,
                "children": [
                    {"layer_id": "slide_02_title", "name": "title",
                     "kind": "text", "z_index": 10,
                     "bbox": {"x": 120, "y": 80, "w": 1680, "h": 140},
                     "text": "The problem",
                     "font_family": "NotoSerifSC-Bold",
                     "font_size_px": 72, "align": "left",
                     "effects": {"fill": "#0f172a"}},
                    {"layer_id": "img_slide02_hero", "name": "hero_image",
                     "kind": "image", "z_index": 5,
                     "bbox": {"x": 120, "y": 260, "w": 1680, "h": 720},
                     "aspect_ratio": "16:9"},
                ],
                # v2.3 speaker notes — populate notes_slide.notes_text_frame
                "speaker_notes": "Open with a show-of-hands: who had a stressful bubble-tea run this week? Describe the three overlapping signage problems (logos / menu density / color war). Pause 5s on the pain statement. Transition: that's why we built MilkCloud.",
            },
            {
                "layer_id": "slide_03", "name": "thank_you", "kind": "slide",
                "z_index": 3,
                "children": [
                    {"layer_id": "slide_03_title", "name": "title",
                     "kind": "text", "z_index": 10,
                     "bbox": {"x": 120, "y": 440, "w": 1680, "h": 200},
                     "text": "Thank you",
                     "font_family": "NotoSerifSC-Bold",
                     "font_size_px": 128, "align": "center",
                     "effects": {"fill": "#0f172a"}},
                ],
            },
        ],
    }

    obs = propose_design_spec({"design_spec": spec}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"propose_design_spec(deck): {(obs.error_message or str(obs.payload))}")

    obs = composite({}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"composite(deck): {(obs.error_message or str(obs.payload))}")

    comp = ctx.state["composition"]
    if comp.psd_path is not None or comp.svg_path is not None or comp.html_path is not None:
        _fail(f"deck should produce ONLY pptx + preview — got "
              f"psd={comp.psd_path} svg={comp.svg_path} html={comp.html_path}")
    if not comp.pptx_path or not Path(comp.pptx_path).exists():
        _fail(f"deck PPTX missing: {comp.pptx_path}")
    if not comp.preview_path or not Path(comp.preview_path).exists():
        _fail(f"deck preview missing: {comp.preview_path}")

    # Per-slide PNGs exist.
    slides_dir = Path(comp.pptx_path).parent / "slides"
    slide_pngs = sorted(slides_dir.glob("slide_*.png"))
    if len(slide_pngs) != 3:
        _fail(f"expected 3 slide PNGs in {slides_dir}, got {len(slide_pngs)}")
    for p in slide_pngs:
        if p.stat().st_size == 0:
            _fail(f"empty slide PNG: {p}")

    # Reopen with python-pptx and check structure.
    prs = _Reopen(comp.pptx_path)
    if len(prs.slides) != 3:
        _fail(f"pptx reopen: expected 3 slides, got {len(prs.slides)}")

    # Slide 1 should contain both "MilkCloud" and the tagline text as native runs.
    s1_text = " ".join(
        run.text for shape in prs.slides[0].shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs for run in para.runs
    )
    if "MilkCloud" not in s1_text:
        _fail(f"slide 1 text missing 'MilkCloud' — got: {s1_text!r}")
    if "calmest bubble tea brand" not in s1_text:
        _fail(f"slide 1 tagline missing — got: {s1_text!r}")

    # Slide 2 should contain one picture shape (the hero image).
    pic_count = sum(1 for shape in prs.slides[1].shapes if shape.shape_type == 13)  # MSO_SHAPE_TYPE.PICTURE
    if pic_count != 1:
        _fail(f"slide 2 expected 1 picture shape, got {pic_count}")

    # v2.3 — speaker notes on slide 2 populated notes_slide.notes_text_frame
    s2_notes = prs.slides[1].notes_slide.notes_text_frame.text
    if "show-of-hands" not in s2_notes:
        _fail(f"slide 2 speaker notes missing or wrong — got: {s2_notes[:120]!r}")
    # Slides 1 + 3 (no speaker_notes in fixture) should have empty notes
    s1_notes = prs.slides[0].notes_slide.notes_text_frame.text.strip()
    s3_notes = prs.slides[2].notes_slide.notes_text_frame.text.strip()
    if s1_notes or s3_notes:
        _fail(f"slide 1/3 should have empty notes — got s1={s1_notes!r} s3={s3_notes!r}")

    pptx_size_kb = Path(comp.pptx_path).stat().st_size // 1024
    preview_size_kb = Path(comp.preview_path).stat().st_size // 1024
    _ok(f"deck composite: 3 slides → {pptx_size_kb} KB .pptx + "
        f"3 slide PNGs + {preview_size_kb} KB preview grid")
    _ok("pptx reopen: slide count + native text runs + picture shape + speaker notes all OK")


def check_reasoning_step_roundtrip() -> None:
    """v2 training-data: derive artifact_type / design_spec / critique
    verdict / layer count from agent_trace alone (no top-level fields).

    Asserts the chat.py helpers (_last_artifact_type / _last_design_spec /
    _last_critique_payload / _count_unique_layers) correctly recover state
    from a synthetic v2 trajectory shape.
    """
    print("[14/20] v2 trajectory: derive metadata from agent_trace only")
    from .chat import (
        _last_artifact_type, _last_design_spec, _last_critique_payload,
        _count_unique_layers,
    )

    crit_payload = {
        "iteration": 1, "verdict": "pass", "score": 0.86,
        "issues": [], "rationale": "ok",
    }
    spec_input = {
        "brief": "test", "artifact_type": "landing",
        "canvas": {"w_px": 1200, "h_px": 2400, "dpi": 96,
                   "aspect_ratio": "1:2", "color_mode": "RGB"},
        "layer_graph": [],
    }

    trace = [
        AgentTraceStep(step_idx=1, actor="user", type="input", text="brief"),
        AgentTraceStep(step_idx=2, actor="planner", type="tool_call",
                       tool_use_id="t1", tool_name="switch_artifact_type",
                       tool_args={"type": "landing"}),
        AgentTraceStep(step_idx=3, actor="tool", type="tool_result",
                       tool_use_id="t1", tool_name="switch_artifact_type",
                       tool_result=ToolResultRecord(status="ok",
                                                   payload={"type": "landing"})),
        AgentTraceStep(step_idx=4, actor="planner", type="tool_call",
                       tool_use_id="t2", tool_name="propose_design_spec",
                       tool_args={"design_spec": spec_input}),
        AgentTraceStep(step_idx=5, actor="planner", type="tool_call",
                       tool_use_id="t3", tool_name="generate_image",
                       tool_args={"layer_id": "L1", "prompt": "hero"}),
        AgentTraceStep(step_idx=6, actor="planner", type="tool_call",
                       tool_use_id="t4", tool_name="render_text_layer",
                       tool_args={"layer_id": "L2", "text": "Hello"}),
        AgentTraceStep(step_idx=7, actor="planner", type="tool_call",
                       tool_use_id="t5", tool_name="critique"),
        AgentTraceStep(step_idx=8, actor="tool", type="tool_result",
                       tool_use_id="t5", tool_name="critique",
                       tool_result=ToolResultRecord(status="ok",
                                                   payload=crit_payload)),
    ]
    traj = DistillTrajectory(
        run_id="smoke-v2-helpers",
        brief="brief",
        agent_trace=trace,
        final_reward=0.86,
        terminal_status="pass",
        metadata=TrainingMetadata(
            schema_version="v2",
            planner_model="m", critic_model="m", image_model="m",
            planner_thinking_budget=0, critic_thinking_budget=0,
            interleaved_thinking=False,
            total_input_tokens=0, total_output_tokens=0,
            estimated_cost_usd=0.0, wall_time_s=0.0,
            source="agent_run",
        ),
    )

    if _last_artifact_type(traj) != "landing":
        _fail(f"_last_artifact_type: {_last_artifact_type(traj)}")
    spec = _last_design_spec(traj)
    if not spec or spec.get("artifact_type") != "landing":
        _fail(f"_last_design_spec lost: {spec}")
    crit = _last_critique_payload(traj)
    if not crit or crit.get("score") != 0.86:
        _fail(f"_last_critique_payload lost: {crit}")
    n = _count_unique_layers(traj)
    if n != 2:
        _fail(f"_count_unique_layers: expected 2 (L1, L2), got {n}")

    _ok("artifact_type, design_spec, critique payload, layer count "
        "all recoverable from agent_trace alone")


def check_ingest_document_markdown() -> None:
    """Markdown ingestion: seed a stub .md with a relative image ref, verify
    ingest_document registers the image in rendered_layers + returns the raw
    text. No API — markdown path doesn't call Anthropic."""
    print("[15/20] ingest_document markdown (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import ingest_document

    out_dir = REPO_ROOT / "out" / "smoke_ingest_md"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Stage a stub MD with an embedded image reference (relative path).
    src_dir = out_dir / "src"
    src_dir.mkdir(parents=True, exist_ok=True)
    img_path = src_dir / "diagram.png"
    Image.new("RGB", (640, 360), (90, 130, 200)).save(img_path)
    md_path = src_dir / "notes.md"
    md_path.write_text(
        "# My project notes\n\n"
        "## Overview\nA few thoughts below.\n\n"
        "![System diagram](diagram.png)\n\n"
        "## Plans\n- Ship v1\n- Iterate\n"
        "\n![Missing](./nowhere.png)\n",
        encoding="utf-8",
    )

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-ingest-md")

    obs = ingest_document({"file_paths": [str(md_path)]}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"ingest_document(markdown): {(obs.error_message or str(obs.payload))}")

    ingested = ctx.state.get("ingested") or []
    if len(ingested) != 1 or ingested[0]["type"] != "markdown":
        _fail(f"ingested manifest missing markdown entry: {ingested}")
    registered = ingested[0]["registered_layer_ids"]
    skipped = ingested[0]["skipped_images"]
    if len(registered) != 1:
        _fail(f"expected 1 image layer registered from MD, got {registered}")
    if len(skipped) != 1:
        _fail(f"expected 1 skipped bad ref, got {skipped}")

    layer_id = registered[0]
    rec = ctx.state["rendered_layers"].get(layer_id)
    if not rec or not rec.get("src_path") or not Path(rec["src_path"]).exists():
        _fail(f"markdown-registered image not hydrated: {rec}")
    _ok(f"markdown ingested: {ingested[0]['n_chars']} chars, "
        f"1 image layer ({layer_id}), 1 bad ref skipped")


def check_ingest_document_image() -> None:
    """Standalone image ingestion: seed a PNG, verify ingest_document copies
    into layers_dir + registers a passthrough layer with correct shape."""
    print("[16/20] ingest_document image (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import ingest_document

    out_dir = REPO_ROOT / "out" / "smoke_ingest_image"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    src_dir = out_dir / "src"
    src_dir.mkdir(parents=True, exist_ok=True)
    img_path = src_dir / "logo.png"
    Image.new("RGB", (512, 512), (240, 80, 60)).save(img_path)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-ingest-image")

    obs = ingest_document({"file_paths": [str(img_path)]}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"ingest_document(image): {(obs.error_message or str(obs.payload))}")

    ingested = ctx.state.get("ingested") or []
    if len(ingested) != 1 or ingested[0]["type"] != "image":
        _fail(f"ingested manifest missing image entry: {ingested}")
    if (ingested[0]["width"], ingested[0]["height"]) != (512, 512):
        _fail(f"image dims not captured: {ingested[0]}")

    registered = ingested[0]["registered_layer_ids"]
    if len(registered) != 1:
        _fail(f"expected 1 image layer, got {registered}")
    layer_id = registered[0]
    rec = ctx.state["rendered_layers"].get(layer_id)
    if rec["kind"] != "image":
        _fail(f"image layer kind != image: {rec['kind']}")
    if rec.get("source") != "ingested":
        _fail(f"image source tag wrong: {rec.get('source')}")
    if not Path(rec["src_path"]).exists():
        _fail(f"copied image path missing: {rec['src_path']}")
    _ok(f"image ingested: {layer_id}, 512×512, source=ingested, "
        f"sha256[:8]={rec['sha256'][:8]}")


def check_ingest_document_docx() -> None:
    """Docx ingestion (v1.2.5): build a minimal Word doc with headings +
    an inline image, verify ingest_document extracts sections + figures
    without any VLM call."""
    print("[17/20] ingest_document docx (no API)")
    from docx import Document
    from docx.shared import Inches
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import ingest_document

    out_dir = REPO_ROOT / "out" / "smoke_ingest_docx"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    src_dir = out_dir / "src"
    src_dir.mkdir(parents=True, exist_ok=True)
    img_path = src_dir / "icon.png"
    Image.new("RGB", (64, 64), (200, 80, 80)).save(img_path)

    docx_path = src_dir / "notes.docx"
    doc = Document()
    doc.add_heading("Quick Tour", level=0)
    doc.add_heading("Motivation", level=1)
    doc.add_paragraph("First thought about the project. Second thought.")
    doc.add_paragraph("- Short bullet.")
    doc.add_heading("Architecture", level=1)
    doc.add_paragraph("Planner + critic + composer.")
    doc.add_picture(str(img_path), width=Inches(1.0))
    doc.save(str(docx_path))

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-ingest-docx")

    obs = ingest_document({"file_paths": [str(docx_path)]}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"ingest_document(docx): {(obs.error_message or str(obs.payload))}")

    ingested = ctx.state.get("ingested") or []
    if len(ingested) != 1 or ingested[0]["type"] != "docx":
        _fail(f"ingested manifest missing docx entry: {ingested}")
    m = ingested[0]["manifest"]
    if m["title"] != "Quick Tour":
        _fail(f"docx title wrong: {m['title']!r}")
    headings = [s["heading"] for s in m["sections"]]
    if headings != ["Motivation", "Architecture"]:
        _fail(f"docx sections wrong: {headings}")
    figure_ids = ingested[0]["registered_figure_ids"]
    if len(figure_ids) != 1 or not figure_ids[0].startswith("ingest_fig_"):
        _fail(f"docx figure registration wrong: {figure_ids}")
    rec = ctx.state["rendered_layers"][figure_ids[0]]
    if rec.get("source") != "ingested_docx" or not Path(rec["src_path"]).exists():
        _fail(f"docx figure record wrong: {rec}")
    _ok(f"docx ingested: title='{m['title']}', "
        f"{len(m['sections'])} section(s), 1 figure ({figure_ids[0]})")


def check_ingest_document_pptx() -> None:
    """Pptx ingestion (v1.2.5): build a 2-slide PowerPoint with a title,
    body bullets, and an embedded picture; verify slides become sections
    and the picture becomes an ingest_fig_NN layer."""
    print("[18/20] ingest_document pptx (no API)")
    from pptx import Presentation
    from pptx.util import Inches
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import ingest_document

    out_dir = REPO_ROOT / "out" / "smoke_ingest_pptx"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    src_dir = out_dir / "src"
    src_dir.mkdir(parents=True, exist_ok=True)
    img_path = src_dir / "chart.png"
    Image.new("RGB", (128, 128), (50, 160, 90)).save(img_path)

    pptx_path = src_dir / "deck.pptx"
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "LongcatDesign Pitch"
    s1.placeholders[1].text = "Open-source conversational design agent"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Why now?"
    s2.placeholders[1].text = "Claude Design shipped. We go terminal-first."
    s2.shapes.add_picture(str(img_path), Inches(4), Inches(2), width=Inches(2))
    prs.save(str(pptx_path))

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-ingest-pptx")

    obs = ingest_document({"file_paths": [str(pptx_path)]}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"ingest_document(pptx): {(obs.error_message or str(obs.payload))}")

    ingested = ctx.state.get("ingested") or []
    if len(ingested) != 1 or ingested[0]["type"] != "pptx":
        _fail(f"ingested manifest missing pptx entry: {ingested}")
    m = ingested[0]["manifest"]
    if m["title"] != "LongcatDesign Pitch":
        _fail(f"pptx title wrong: {m['title']!r}")
    if len(m["sections"]) != 2:
        _fail(f"pptx slides→sections wrong: {[s['heading'] for s in m['sections']]}")
    figure_ids = ingested[0]["registered_figure_ids"]
    if len(figure_ids) != 1 or not figure_ids[0].startswith("ingest_fig_"):
        _fail(f"pptx figure registration wrong: {figure_ids}")
    rec = ctx.state["rendered_layers"][figure_ids[0]]
    if rec.get("source") != "ingested_pptx" or rec.get("source_ref") != "slide=2":
        _fail(f"pptx figure record wrong: {rec}")
    _ok(f"pptx ingested: title='{m['title']}', "
        f"2 slide(s), 1 figure ({figure_ids[0]})")


def check_sub_figure_registration() -> None:
    """v2.3 sub-figure extraction: VLM detects multi-panel composite figures
    and we Pillow-crop each panel into its own `ingest_fig_NN_<label>` layer.
    Tests the registration helper directly (no API) with a synthetic
    2-panel composite PNG + a stub VLM response.

    Asserts:
    - 2 panel entries → 2 new rendered_layers records (parent stays)
    - Each panel's crop file exists on disk with correct dims
    - parent_layer_id breadcrumb set on children
    - Layer_id naming convention `ingest_fig_NN_<label>` holds
    """
    print("[19/20] sub-figure extraction (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import _register_sub_panels

    out_dir = REPO_ROOT / "out" / "smoke_sub_figs"
    layers_dir = out_dir / "layers"
    if out_dir.exists():
        import shutil
        shutil.rmtree(out_dir)
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Synthesize a composite parent: 600×300 image split into two
    # horizontal halves (panel a on left, panel b on right).
    parent_path = layers_dir / "img_ingest_fig_01.png"
    composite_img = Image.new("RGB", (600, 300), (255, 255, 255))
    # Left half red, right half blue — easy to eyeball crop correctness
    left = Image.new("RGB", (300, 300), (200, 60, 60))
    right = Image.new("RGB", (300, 300), (60, 100, 200))
    composite_img.paste(left, (0, 0))
    composite_img.paste(right, (300, 0))
    composite_img.save(parent_path, "PNG")

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-subfig")
    # Pre-stage the parent so sub-panel registration can check collisions
    ctx.state["rendered_layers"]["ingest_fig_01"] = {
        "layer_id": "ingest_fig_01", "kind": "image", "name": "figure_1",
        "src_path": str(parent_path),
    }

    panels = [
        {"label": "a", "bbox": [0, 0, 300, 300],
         "caption": "Left panel (text input)", "short_caption": "Text"},
        {"label": "b", "bbox": [300, 0, 600, 300],
         "caption": "Right panel (image input)", "short_caption": "Image"},
    ]
    created = _register_sub_panels(
        parent_layer_id="ingest_fig_01",
        parent_path=parent_path,
        parent_caption="Composite — text + image inputs",
        panels=panels,
        ctx=ctx,
        pdf_path=Path("/fake/paper.pdf"),
        source_page=3,
    )

    if created != ["ingest_fig_01_a", "ingest_fig_01_b"]:
        _fail(f"sub-panel layer_ids wrong: {created}")

    for sub_id, expected_color, expected_caption in [
        ("ingest_fig_01_a", (200, 60, 60), "Text"),
        ("ingest_fig_01_b", (60, 100, 200), "Image"),
    ]:
        rec = ctx.state["rendered_layers"].get(sub_id)
        if rec is None:
            _fail(f"sub-panel {sub_id} not registered")
        if rec.get("parent_layer_id") != "ingest_fig_01":
            _fail(f"{sub_id} missing parent_layer_id breadcrumb")
        if rec.get("extract_strategy") != "sub_panel":
            _fail(f"{sub_id} extract_strategy wrong: {rec.get('extract_strategy')!r}")
        if rec.get("caption_short") != expected_caption:
            _fail(f"{sub_id} caption_short wrong: {rec.get('caption_short')!r}")
        crop_path = Path(rec["src_path"])
        if not crop_path.exists() or crop_path.stat().st_size == 0:
            _fail(f"{sub_id} crop file missing or empty: {crop_path}")
        with Image.open(crop_path) as im:
            w, h = im.size
            if (w, h) != (300, 300):
                _fail(f"{sub_id} crop dims wrong: {w}×{h}, expected 300×300")
            # Center pixel should match the panel color (parent split in half)
            px = im.getpixel((150, 150))
            # Drop alpha if present
            if isinstance(px, tuple) and len(px) >= 3:
                px = px[:3]
            if px != expected_color:
                _fail(f"{sub_id} center pixel {px} != expected {expected_color}")

    _ok(f"sub-panel registration: 2 panels cropped + registered "
        f"({len(created)} layers, parent breadcrumb intact, crops correct)")


def check_versioning_no_api() -> None:
    """v2.2 versioning: revise loops + edit_layer must NOT clobber prior
    intermediate state. Asserts:
      - render_text_layer writes layers/<id>.v<N>.png (not <id>.png)
      - edit_layer (re-render) bumps to v<N+1>; v<N> still on disk
      - 2 composite() calls produce composites/iter_01/ AND composites/iter_02/
      - final/ symlinks point at iter_02 (the latest)
      - tool_result.payload exposes relative_path / version / supersedes_*
    """
    print("[20/20] versioning + revise-loop preservation (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.render_text_layer import render_text_layer
    from .tools.edit_layer import edit_layer
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_versioning"
    layers_dir = out_dir / "layers"
    if out_dir.exists():
        # Clean prior smoke run so version counters start at 0.
        import shutil
        shutil.rmtree(out_dir)
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-vers")

    # Stub a background so composite has 2 layers
    bg_path = layers_dir / "bg_seed.png"
    Image.new("RGB", (768, 1024), (10, 10, 30)).save(bg_path)
    ctx.state["rendered_layers"]["L0_bg"] = {
        "layer_id": "L0_bg", "name": "bg", "kind": "background", "z_index": 0,
        "bbox": {"x": 0, "y": 0, "w": 768, "h": 1024},
        "src_path": str(bg_path), "prompt": "(stub)",
        "aspect_ratio": "3:4", "image_size": "1K",
        "safe_zones": [], "sha256": "seed",
    }

    switch_artifact_type({"type": "poster"}, ctx=ctx)
    propose_design_spec({"design_spec": {
        "brief": "v",
        "canvas": {"w_px": 768, "h_px": 1024, "dpi": 96,
                   "aspect_ratio": "3:4", "color_mode": "RGB"},
        "palette": ["#000", "#fff"],
        "typography": {}, "mood": [], "composition_notes": "",
        "layer_graph": [],
    }}, ctx=ctx)

    # ── v1 layer write ─────────────────────────────────────────────────
    r1 = render_text_layer({
        "layer_id": "L1", "name": "title", "text": "v1 text",
        "font_family": "NotoSansSC-Bold", "font_size_px": 60, "fill": "#fff",
        "bbox": {"x": 50, "y": 50, "w": 600, "h": 120},
        "align": "center", "z_index": 1,
    }, ctx=ctx)
    if r1.payload.get("version") != 1:
        _fail(f"first render should be v1; got version={r1.payload.get('version')}")
    if r1.payload.get("relative_path") != "layers/text_L1.v1.png":
        _fail(f"v1 relative_path wrong: {r1.payload.get('relative_path')}")
    if r1.payload.get("supersedes_sha256") is not None:
        _fail("first render should not have supersedes_sha256")
    v1_file = layers_dir / "text_L1.v1.png"
    if not v1_file.exists():
        _fail(f"v1 file missing: {v1_file}")
    v1_sha = r1.payload["sha256"]

    # ── v2 layer write (via edit_layer → re-render) ────────────────────
    r2 = edit_layer({"layer_id": "L1", "diff": {"text": "v2 text"}}, ctx=ctx)
    if r2.payload.get("version") != 2:
        _fail(f"edit_layer should bump to v2; got {r2.payload.get('version')}")
    if r2.payload.get("supersedes_sha256") != v1_sha:
        _fail(f"v2 supersedes_sha256 should equal v1 sha; got {r2.payload.get('supersedes_sha256')!r}")
    v2_file = layers_dir / "text_L1.v2.png"
    if not v2_file.exists():
        _fail(f"v2 file missing: {v2_file}")
    if not v1_file.exists():
        _fail(f"⚠ v1 file was clobbered (this is the bug we're guarding against): {v1_file}")
    _ok(f"layer versioning: v1 + v2 both on disk; v2.supersedes = v1.sha256[:8] {v1_sha[:8]}")

    # ── Composite iter 1 ──────────────────────────────────────────────
    c1 = composite({}, ctx=ctx)
    if c1.status != "ok":
        _fail(f"composite iter 1: {(c1.error_message or c1.payload)}")
    if c1.payload.get("iteration") != 1:
        _fail(f"composite iter 1 should report iteration=1; got {c1.payload.get('iteration')}")
    iter1_dir = out_dir / "composites" / "iter_01"
    for f in ("poster.html", "poster.psd", "poster.svg", "preview.png"):
        if not (iter1_dir / f).exists():
            _fail(f"iter_01/{f} missing")
    iter1_preview_sha = c1.payload["preview_sha256"]

    # ── Edit a layer + composite iter 2 ───────────────────────────────
    edit_layer({"layer_id": "L1", "diff": {"fill": "#ff00ff"}}, ctx=ctx)
    c2 = composite({}, ctx=ctx)
    if c2.status != "ok":
        _fail(f"composite iter 2: {(c2.error_message or c2.payload)}")
    if c2.payload.get("iteration") != 2:
        _fail(f"composite iter 2 should report iteration=2; got {c2.payload.get('iteration')}")
    iter2_dir = out_dir / "composites" / "iter_02"
    if not (iter2_dir / "preview.png").exists():
        _fail("iter_02/preview.png missing")
    if not (iter1_dir / "preview.png").exists():
        _fail("⚠ iter_01/preview.png was clobbered by iter_02 — versioning broken")
    if c2.payload.get("supersedes_preview_sha256") != iter1_preview_sha:
        _fail(f"iter_02 supersedes_preview_sha256 should be iter_01's preview sha")
    _ok(f"composite versioning: iter_01 + iter_02 both intact; "
        f"iter_02.supersedes_preview = iter_01.preview_sha[:8] {iter1_preview_sha[:8]}")

    # ── final/ symlinks point at the latest iter ──────────────────────
    final_dir = out_dir / "final"
    if not (final_dir / "preview.png").is_symlink():
        _fail("final/preview.png should be a symlink")
    if not (final_dir / "poster.html").is_symlink():
        _fail("final/poster.html should be a symlink")
    final_preview = (final_dir / "preview.png").resolve()
    if final_preview.parent.name != "iter_02":
        _fail(f"final/preview.png should resolve to iter_02; got {final_preview.parent.name}")
    _ok(f"final/ symlinks resolve to iter_02 (the latest)")


def main() -> int:
    check_imports()
    check_tool_registry()
    check_pydantic_roundtrip()
    check_fonts()
    check_composite_no_api()
    check_svg_text_is_vector()
    check_chat_session_roundtrip()
    check_edit_layer_no_api()
    check_apply_edits_roundtrip()
    check_landing_mode()
    check_design_system_styles()
    check_landing_with_images()
    check_deck_mode()
    check_reasoning_step_roundtrip()
    check_ingest_document_markdown()
    check_ingest_document_image()
    check_ingest_document_docx()
    check_ingest_document_pptx()
    check_sub_figure_registration()
    check_versioning_no_api()
    print("\n  smoke test passed.")
    print("  artifacts in: out/smoke/, out/smoke_edit/, out/smoke_apply/, "
          "out/smoke_landing/, out/smoke_styles/, out/smoke_landing_img/, "
          "out/smoke_deck/, out/smoke_ingest_md/, out/smoke_ingest_image/, "
          "out/smoke_ingest_docx/, out/smoke_ingest_pptx/, out/smoke_sub_figs/")
    return 0


if __name__ == "__main__":
    sys.exit(main())
