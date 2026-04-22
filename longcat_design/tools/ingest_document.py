"""ingest_document — v1.2 paper2any entry point.

v1.1 asked a VLM to locate figure bboxes on rasterized PDF pages. That
was unreliable: the model returned half-page screenshots, clipped
diagrams, and hallucinated "figures" on text-only pages.

v1.2 separates two concerns:

1. **Figure localization** is now done by **pymupdf directly** —
   `doc.extract_image(xref)` pulls embedded raster images at native
   resolution (e.g. 1890×1211 PNGs the paper author uploaded), and
   `page.get_drawings()` + proximity clustering + 300 dpi rendering
   catches vector-drawn architecture diagrams. No VLM guessing.

2. **Reading / matching** is still VLM work — but the default model
   is now Qwen-VL-Max via OpenRouter (~5× cheaper and faster than
   Claude Sonnet for this non-reasoning workload). Two calls:

     a. **Structure extraction**: render pages at 144 dpi, send as
        multi-image request, receive title / authors / abstract /
        sections / figures / tables JSON. `figures` now carries only
        `{caption, page, description}` — no bbox or idx (we already
        have bboxes from pymupdf).

     b. **Caption matching**: for each pymupdf candidate, ask the VLM
        which caption matches and whether it's a real figure. Fake
        candidates (logos, page headers, equation renders) are
        filtered at this step.

The `rendered_layers` record shape is unchanged (see the downstream
contract in `docs/DECISIONS.md` and `longcat_design/tools/composite.py`
hydration helpers): callers reference ingested figures by the stable
`layer_id` we register (e.g. `ingest_fig_01`), and the renderer
hydrates `src_path` / `bbox` / `caption` from the layer registry.

Markdown and image branches are untouched.
"""

from __future__ import annotations

import re
import shutil
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from typing import Any

import fitz  # pymupdf

from ._contract import ToolContext, obs_error, obs_ok
from ..schema import ToolResultRecord
from ..util.io import sha256_file
from ..util.logging import log
from ..util.pdf import (
    PdfFigureCandidate,
    PdfTableCandidate,
    ScannedPdfError,
    dedup_raster_vector,
    dedup_tables_against_figures,
    detect_scanned_pdf,
    extract_embedded_rasters,
    extract_page_text,
    extract_table_candidates,
    extract_vector_clusters,
    page_count,
    render_page_png,
)
from ..util.table_png import render_table_png
from ..util.vlm import VlmImage, vlm_call_json


# Max PDF bytes we accept in one call (belt-and-suspenders — pymupdf
# itself can open almost anything, but ingest touches every page and
# we want to fail fast on pathological inputs rather than spin).
_MAX_PDF_BYTES = 24 * 1024 * 1024  # 24 MB
_MAX_PDF_PAGES = 80

_MD_IMG_RE = re.compile(r"!\[([^\]]*)\]\(([^)\s]+)(?:\s+\"[^\"]*\")?\)")

# Max number of figures listed in the ingest tool_result summary shown
# to the planner. Papers with 40+ figures would otherwise blow the
# planner's remaining output budget on turn 2 just for a catalog it
# mostly doesn't need; the planner can still reference ANY `ingest_fig_NN`
# layer_id — the rendered_layers dict has them all. 20 is enough for
# the planner to pick a diverse 4-8 for a poster.
_PLANNER_FIG_CATALOG_CAP = 20

# Max parallel caption-matching calls. VLM calls are HTTP-bound; 4–6 in
# flight keeps wall time tight without tripping OpenRouter rate limits.
_CAPTION_MATCH_PARALLELISM = 6
# Confidence floor for accepting a VLM caption match — below this we
# still keep the figure but flag it; `is_real_figure=false` drops it
# regardless of confidence.
_CAPTION_MATCH_MIN_CONFIDENCE = 0.35
# DPI for the cover-page image handed to the VLM during structure
# extraction. We only send ONE image (the front page) — the rest of
# the paper reaches the VLM as pymupdf-extracted text, which is both
# cheaper and dodges Qwen-VL-Max's per-request image-count cap
# (DashScope rejects >10 images with HTTP 400 "invalid_parameter").
_STRUCTURE_PAGE_DPI = 144
# Total text budget across all pages sent to the structure extractor.
# ~60k chars fits comfortably in the ~16k-token context window Qwen
# uses for this call while leaving headroom for the JSON response.
_STRUCTURE_TOTAL_TEXT_CAP = 60_000
# Scanned-PDF OCR fallback (v1.2.5). When `detect_scanned_pdf` returns
# True we render each page at this DPI and hand the PNG to the VLM for
# text extraction — 200 dpi is the sweet spot: dense enough for small
# body text, not so heavy that a 40-page doc burns minutes. We still
# cap at `_MAX_PDF_PAGES` so runaway OCR cost is impossible.
_OCR_PAGE_DPI = 200
_OCR_PAGE_PARALLELISM = 6
_OCR_PER_PAGE_TIMEOUT_S = 120.0


_OCR_PROMPT = """\
You are an OCR engine. The image is ONE page from a PDF. Extract every
readable word exactly as it appears, preserving reading order, newlines
between paragraphs, and bullets / numbered lists. Ignore page numbers,
running headers, and watermarks unless they carry content.

Output **a single fenced JSON code block, nothing else**:

```json
{"text": "<extracted text with \\n newlines>"}
```

Rules:
- Use `\\n\\n` between paragraphs, `\\n` between lines inside a paragraph.
- Do NOT translate. Preserve the original language.
- Do NOT invent words the image doesn't show.
- If the page is blank or unreadable, return `{"text": ""}`.
"""


_INGEST_STRUCTURE_PROMPT = """\
You are a document-structure extractor for LongcatDesign. You will be
given ONE cover-page image (page 1) for visual grounding, followed by
the full extracted text of the paper (all pages, marked with
[PAGE N] headers). Return a STRICT JSON manifest that downstream tools
consume verbatim.

Output **a single fenced JSON code block, nothing else**:

```json
{
  "title": "<paper/doc title>",
  "authors": ["<Author Name>", ...],
  "venue": "<conference / publication / null>",
  "abstract": "<2-4 sentence abstract>",
  "sections": [
    {"idx": 1, "heading": "Introduction",
     "summary": "<2-3 sentences of what this section argues>",
     "key_points": ["<one-liner>", "<one-liner>", ...]},
    ...
  ],
  "figures": [
    {"caption": "<figure's full caption text, as it appears in the doc>",
     "page": <1-indexed page where the figure is anchored>,
     "description": "<1 sentence describing what's in the figure; used later to match it to a crop>"},
    ...
  ],
  "tables": [
    {"caption": "<table caption>", "page": <int>},
    ...
  ],
  "key_quotes": ["<memorable line from the doc>", ...]
}
```

Rules:
- Titles: use the human-facing title, not the first line.
- Sections: include each top-level heading (or the logical equivalent
  if the doc doesn't have explicit headings). Max ~10 sections; if the
  doc has more, collapse aggressively (group related subsections).
- Figures: only include REAL visual figures (diagrams, charts,
  screenshots, photos). Ignore logos, page headers, decorative borders,
  page numbers, watermarks, and inline equation renders. Captions as
  they literally appear in the doc. Include sub-panels only if they
  have their own caption line.
- Pages: 1-indexed.
- Empty lists are fine. Don't guess.
- No extra prose outside the fenced JSON block.
"""


_TABLE_PARSE_PROMPT = """\
You are a table parser. You will see ONE image cropped from a PDF,
plus the raw cell text that pymupdf's table finder guessed for the
same region (the cell splits may be WRONG — trust the image). You
will also see a short list of table-caption candidates pulled from
the same paper so you can match the table to its caption.

Your job:

1. Decide whether the image is actually a data table — a grid of
   rows/columns with comparable values. Diagrams, figure panels,
   math-equation arrays, OCR example screenshots, text paragraphs,
   and decorative layout artifacts are NOT data tables.

2. If it is a table, output clean structured data: pick the header
   row, expand merged cells so every (row, col) has a value, and
   preserve numeric values as-is (no rounding, no re-formatting).
   Short dash / em-dash entries (—, -) stay as the literal string "—".

3. Match it to one of the caption candidates if any fits; otherwise
   return `matched_idx=null` and set a short `title` you extracted
   from the image.

Output **a single fenced JSON code block, nothing else**:

```json
{
  "is_table": <true | false>,
  "matched_idx": <int index into the caption candidate list, or null>,
  "title": "<short title or caption; empty string when unknown>",
  "headers": ["<col1>", "<col2>", ...],
  "rows": [
    ["<r1c1>", "<r1c2>", ...],
    ["<r2c1>", "<r2c2>", ...],
    ...
  ],
  "col_highlight_rule": ["", "max", "max", "min", ...],
  "reason": "<short explanation>"
}
```

Rules:
- Every row in `rows` must have the same length as `headers` (pad
  with "—" if necessary). If you are not confident about the header
  row, leave `headers: []` and put everything in `rows` (first row
  will be treated as header downstream).
- **Two-row headers**: if the table has a parent header spanning
  sub-columns (e.g. "Understanding" over MMMU / MathVista / etc.),
  flatten into a single row using "Parent / Child" format (e.g.
  "Understanding / MMMU"). Do NOT emit a two-row header.
- Do NOT invent rows/columns not visible in the image.
- If `is_table=false`, set `headers` and `rows` to empty lists.
- `col_highlight_rule`: same length as `headers`. For each column,
  emit `"max"` when higher values are better (accuracy, F1, win
  rate), `"min"` when lower is better (loss, error rate, latency),
  or `""` for label / non-numeric / ambiguous columns. The
  downstream renderer bolds the winning row per column — so emitting
  this honestly for benchmark tables is high leverage.
"""


_CAPTION_MATCH_PROMPT = """\
You are a figure↔caption matcher. You will see ONE image cropped from
a PDF plus a short list of figure-caption candidates pulled from the
same paper. Pick the caption that belongs to the image, or report that
the image is not a real figure.

Output **a single fenced JSON code block, nothing else**:

```json
{
  "matched_idx": <int index into the candidate list, or null>,
  "confidence": <float 0.0–1.0>,
  "is_real_figure": <true | false>,
  "reason": "<short explanation>"
}
```

Rules:
- `matched_idx` indexes the candidate list the user gives you (0-based).
- If none of the captions match OR the image is a logo, page header,
  publisher mark, decorative border, watermark, or equation render
  (not a real figure), set `matched_idx=null` and
  `is_real_figure=false`.
- If the image is a real figure but its caption isn't in the
  candidate list (the list was truncated), set `matched_idx=null` and
  `is_real_figure=true`. Downstream will keep it with an empty caption.
- Prefer confidence ≥ 0.7 when you're sure; otherwise be honest and
  use a lower value.
"""


def ingest_document(args: dict[str, Any], *, ctx: ToolContext) -> ToolResultRecord:
    raw = args.get("file_paths")
    if not raw or not isinstance(raw, list):
        return obs_error(
            "ingest_document needs 'file_paths': list[str]",
            category="validation",
        )

    summaries: list[dict[str, Any]] = []

    for fp_str in raw:
        fp = Path(str(fp_str)).expanduser()
        if not fp.is_absolute():
            fp = fp.resolve()
        if not fp.exists():
            return obs_error(f"file not found: {fp}", category="not_found")
        if not fp.is_file():
            return obs_error(f"not a regular file: {fp}", category="validation")

        ext = fp.suffix.lower()
        log("ingest.start", file=str(fp), ext=ext, bytes=fp.stat().st_size)

        try:
            if ext == ".pdf":
                s = _ingest_pdf(fp, ctx)
            elif ext == ".docx":
                s = _ingest_docx(fp, ctx)
            elif ext == ".pptx":
                s = _ingest_pptx(fp, ctx)
            elif ext in (".md", ".markdown", ".txt"):
                s = _ingest_markdown(fp, ctx)
            elif ext in (".png", ".jpg", ".jpeg", ".webp"):
                s = _ingest_image(fp, ctx)
            else:
                return obs_error(
                    f"unsupported file type {ext!r}; supported: "
                    ".pdf, .docx, .pptx, .md/.markdown/.txt, "
                    ".png/.jpg/.jpeg/.webp",
                    category="unsupported_format",
                )
        except ScannedPdfError as e:
            return obs_error(f"ingest failed on {fp.name}: {e}", category="parse_error")
        except RuntimeError as e:
            return obs_error(f"ingest failed on {fp.name}: {e}", category="parse_error")

        summaries.append(s)

    ctx.state.setdefault("ingested", []).extend(summaries)
    log("ingest.done", files=len(summaries),
        total_figures=sum(len(s.get("registered_layer_ids", [])) for s in summaries))

    # Build a structured payload with per-figure metadata. The policy needs
    # this to pick figures meaningfully (caption + dims + source page = the
    # actual environment state). NOT prose; the policy can iterate it.
    rendered = ctx.state.get("rendered_layers") or {}
    files_payload: list[dict[str, Any]] = []
    figures_payload: list[dict[str, Any]] = []
    tables_payload: list[dict[str, Any]] = []

    for s in summaries:
        f = Path(s["file"]).name
        t = s["type"]
        figure_ids = s.get("registered_figure_ids") or s.get("registered_layer_ids") or []
        table_ids = s.get("registered_table_ids") or []

        file_entry: dict[str, Any] = {
            "name": f,
            "type": t,
            "n_figures": len(figure_ids),
            "n_tables": len(table_ids),
        }
        if t in ("pdf", "docx", "pptx"):
            m = s.get("manifest") or {}
            file_entry["title"] = m.get("title")
            file_entry["n_sections"] = len(m.get("sections") or [])
            if t == "pdf":
                file_entry["authors"] = list(m.get("authors") or [])
        elif t == "markdown":
            file_entry["n_chars"] = s.get("n_chars")
        elif t == "image":
            file_entry["width"] = s.get("width")
            file_entry["height"] = s.get("height")
        files_payload.append(file_entry)

        for fid in figure_ids:
            rec = rendered.get(fid) or {}
            figures_payload.append({
                "layer_id": fid,
                "source_file": f,
                "source_page": rec.get("source_page"),
                "source_ref": rec.get("source_ref"),
                "image_size": rec.get("image_size"),
                "caption": rec.get("caption"),
                "sha256": rec.get("sha256"),
            })
        for tid in table_ids:
            rec = rendered.get(tid) or {}
            tables_payload.append({
                "layer_id": tid,
                "source_file": f,
                "source_page": rec.get("source_page"),
                "n_rows": len(rec.get("rows") or []),
                "n_cols": len(rec.get("headers") or []) or (
                    len((rec.get("rows") or [[]])[0])
                ),
                "caption": rec.get("caption") or rec.get("title"),
            })

    return obs_ok({
        "n_files": len(summaries),
        "n_figures": len(figures_payload),
        "n_tables": len(tables_payload),
        "files": files_payload,
        "figures": figures_payload,
        "tables": tables_payload,
    })


# ───────────────────────────── PDF branch ──────────────────────────────

def _ingest_pdf(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    size = fp.stat().st_size
    if size > _MAX_PDF_BYTES:
        raise RuntimeError(
            f"PDF too large ({size / 1_048_576:.1f} MB > "
            f"{_MAX_PDF_BYTES / 1_048_576:.0f} MB). Split the document."
        )
    pages = page_count(fp)
    if pages > _MAX_PDF_PAGES:
        raise RuntimeError(
            f"PDF has {pages} pages (cap {_MAX_PDF_PAGES}). Trim the document."
        )

    doc = fitz.open(fp)
    is_scanned = False
    try:
        is_scanned = detect_scanned_pdf(doc)
        if is_scanned:
            # v1.2.5 — OCR fallback via Qwen-VL-Max. Scanned PDFs have
            # no embedded rasters (the pages ARE the images) so figure
            # extraction is skipped; we render each page and OCR it so
            # structure extraction still has something to chew on.
            page_texts = _ocr_scanned_pdf(fp, doc, ctx)
            candidates: list[PdfFigureCandidate] = []
            table_candidates: list[PdfTableCandidate] = []
        else:
            # 1. pymupdf figure candidates (no LLM).
            candidates = []
            candidates.extend(extract_embedded_rasters(doc, ctx.layers_dir))
            candidates.extend(extract_vector_clusters(doc, ctx.layers_dir))
            candidates = dedup_raster_vector(candidates)
            log("ingest.pdf.candidates", file=fp.name,
                raster=sum(1 for c in candidates if c.strategy == "raster"),
                vector=sum(1 for c in candidates if c.strategy == "vector"))

            # 1b. pymupdf table candidates (localization only; VLM parses
            # cells). Produced alongside figures — planner picks either a
            # `kind="image"` figure or a `kind="table"` structured layer.
            # Dedup against vector figures on same page: if a figure bbox
            # covers ≥70% of a "table" bbox, it's almost certainly the
            # figure proper (find_tables occasionally trips on composite
            # diagrams). Prefer the figure path.
            table_candidates = extract_table_candidates(
                doc, ctx.layers_dir,
            )
            pre_dedup = len(table_candidates)
            table_candidates = dedup_tables_against_figures(
                table_candidates, candidates,
            )
            log("ingest.pdf.table_candidates",
                file=fp.name, n=len(table_candidates),
                dropped_by_figure_overlap=pre_dedup - len(table_candidates))

            # 2. Page text (lossless, free). Feeds the VLM as text instead
            # of shipping 43 page images — Qwen-VL-Max on DashScope rejects
            # multi-image payloads over ~10 images, and text is denser
            # information per token than rasterized pages anyway.
            page_texts = extract_page_text(doc)
    finally:
        doc.close()

    # Cover-page image: one rasterized page just for visual title/logo
    # grounding. The rest of the paper reaches the VLM via text.
    cover_png = ctx.layers_dir / "ingest_page_001.png"
    try:
        render_page_png(fp, 1, cover_png, dpi=_STRUCTURE_PAGE_DPI)
    except Exception as e:
        log("ingest.pdf.render_fail", page=1, error=str(e))
        cover_png = None  # type: ignore[assignment]

    manifest = _extract_structure(page_texts, cover_png, ctx, fp)
    _normalize_manifest_lists(manifest)

    # 3. Caption matching — per figure candidate, parallelized.
    registered_layer_ids: list[str] = []
    if candidates:
        matches = _match_captions_parallel(candidates, manifest, ctx)
        registered_layer_ids = _register_candidates(
            candidates=candidates, matches=matches, ctx=ctx, pdf_path=fp,
        )
    else:
        log("ingest.pdf.no_candidates", file=fp.name,
            note="pymupdf returned 0 figures; VLM manifest may still be useful")

    # 4. Table parsing — per table candidate, parallelized. VLM reads
    # the bbox image + pymupdf's raw cell guess and returns clean
    # structured rows/headers. Rejects diagrams/equations/etc. that
    # find_tables() misclassified.
    registered_table_ids: list[str] = []
    if table_candidates:
        parsed = _parse_tables_parallel(table_candidates, manifest, ctx)
        registered_table_ids = _register_tables(
            candidates=table_candidates, parsed=parsed, ctx=ctx, pdf_path=fp,
        )

    all_registered = registered_layer_ids + registered_table_ids
    return {
        "file": str(fp), "type": "pdf", "manifest": manifest,
        "registered_layer_ids": all_registered,
        "registered_figure_ids": registered_layer_ids,
        "registered_table_ids": registered_table_ids,
        "summary": f"{manifest.get('title', '?')} — "
                   f"{len(registered_layer_ids)} figure(s), "
                   f"{len(registered_table_ids)} table(s), "
                   f"{len(manifest.get('sections', []))} section(s)",
    }


def _ocr_scanned_pdf(
    fp: Path, doc: "fitz.Document", ctx: ToolContext,
) -> list[str]:
    """Render each page of a scanned PDF at `_OCR_PAGE_DPI` and ask
    Qwen-VL-Max (or whichever VLM `settings.ingest_model` points at) to
    OCR it. Returns one string per page (1-indexed: index 0 = page 1).

    Runs pages in parallel via ThreadPoolExecutor so a 40-page doc takes
    ~8 s wall time at 6 workers instead of 40 × 1 s serial. OCR failures
    on individual pages degrade to empty strings — we don't block the
    whole doc on one bad page, and the structure extractor handles
    partial text fine.
    """
    import time as _time

    n_pages = len(doc)
    log("ingest.pdf.ocr.start", file=fp.name, pages=n_pages,
        dpi=_OCR_PAGE_DPI, parallelism=_OCR_PAGE_PARALLELISM,
        model=ctx.settings.ingest_model)
    t0 = _time.monotonic()

    page_pngs: list[Path] = []
    for i in range(n_pages):
        out_path = ctx.layers_dir / f"ingest_ocr_page_{i + 1:03d}.png"
        try:
            pix = doc[i].get_pixmap(dpi=_OCR_PAGE_DPI)
            pix.save(str(out_path))
            page_pngs.append(out_path)
        except Exception as e:
            log("ingest.pdf.ocr.render_fail", page=i + 1, error=str(e))
            page_pngs.append(None)  # type: ignore[arg-type]

    def ocr_one(idx: int) -> tuple[int, str]:
        png = page_pngs[idx]
        if png is None:
            return idx, ""
        try:
            result = vlm_call_json(
                settings=ctx.settings,
                model=ctx.settings.ingest_model,
                system=_OCR_PROMPT,
                user_text=f"Page {idx + 1} of {n_pages}. OCR it.",
                images=[VlmImage.from_path(png)],
                max_tokens=4096,
                timeout_s=_OCR_PER_PAGE_TIMEOUT_S,
            )
            return idx, str(result.get("text") or "")
        except Exception as e:
            log("ingest.pdf.ocr.page_fail", page=idx + 1, error=str(e))
            return idx, ""

    results: dict[int, str] = {}
    with ThreadPoolExecutor(max_workers=_OCR_PAGE_PARALLELISM) as pool:
        futures = [pool.submit(ocr_one, i) for i in range(n_pages)]
        for fut in as_completed(futures):
            idx, text = fut.result()
            results[idx] = text

    page_texts = [results.get(i, "") for i in range(n_pages)]
    total_chars = sum(len(t) for t in page_texts)
    log("ingest.pdf.ocr.done", file=fp.name,
        wall_s=round(_time.monotonic() - t0, 1),
        total_chars=total_chars,
        pages_with_text=sum(1 for t in page_texts if t.strip()))
    return page_texts


def _extract_structure(
    page_texts: list[str],
    cover_png: Path | None,
    ctx: ToolContext,
    fp: Path,
) -> dict[str, Any]:
    import time as _time

    # Concatenate text with [PAGE N] headers; budget the total so the
    # VLM's context doesn't blow up on reference-heavy long papers.
    body_lines: list[str] = []
    used = 0
    for page_num, text in enumerate(page_texts, start=1):
        chunk = f"\n[PAGE {page_num}]\n{text.strip()}"
        if used + len(chunk) > _STRUCTURE_TOTAL_TEXT_CAP:
            body_lines.append(
                f"\n[remaining {len(page_texts) - page_num + 1} "
                f"pages omitted — cap {_STRUCTURE_TOTAL_TEXT_CAP} chars]"
            )
            break
        body_lines.append(chunk)
        used += len(chunk)
    full_text = "".join(body_lines)

    user_text = (
        "Below is the complete extracted text of the paper, followed by "
        "the cover-page image for visual grounding. Extract the "
        "structured manifest as JSON per the system prompt. Return "
        "STRICT JSON only.\n\n"
        f"{full_text}"
    )

    images = [VlmImage.from_path(cover_png)] if cover_png is not None else []

    log("ingest.pdf.structure.request",
        file=fp.name, text_chars=used, n_pages=len(page_texts),
        cover_image=cover_png is not None,
        model=ctx.settings.ingest_model,
        timeout_s=ctx.settings.ingest_http_timeout)
    t0 = _time.monotonic()
    manifest = vlm_call_json(
        settings=ctx.settings,
        model=ctx.settings.ingest_model,
        system=_INGEST_STRUCTURE_PROMPT,
        user_text=user_text,
        images=images,
        max_tokens=8192,
    )
    log("ingest.pdf.structure.response",
        file=fp.name, wall_s=round(_time.monotonic() - t0, 1))
    return manifest


def _normalize_manifest_lists(manifest: dict[str, Any]) -> None:
    """Coerce None → [] on list-shaped keys. Claude/Qwen occasionally
    emit `"figures": null` for papers they find no figures in, which
    breaks downstream len() calls."""
    for key in ("sections", "figures", "tables", "authors", "key_quotes"):
        if not isinstance(manifest.get(key), list):
            manifest[key] = []


def _match_captions_parallel(
    candidates: list[PdfFigureCandidate],
    manifest: dict[str, Any],
    ctx: ToolContext,
) -> dict[int, dict[str, Any]]:
    """Run caption matching for each candidate in a thread pool.

    Returns a dict keyed by candidate index → match result dict
    (`matched_idx`, `confidence`, `is_real_figure`, `reason`, and
    `caption_text` filled in from the manifest for convenience).
    """
    all_figs = list(manifest.get("figures", []))
    if not all_figs:
        # No captions to match — keep the raw candidates, no caption.
        return {
            i: {"matched_idx": None, "confidence": 0.0,
                "is_real_figure": True, "reason": "no captions in manifest",
                "caption_text": ""}
            for i in range(len(candidates))
        }

    results: dict[int, dict[str, Any]] = {}
    with ThreadPoolExecutor(max_workers=_CAPTION_MATCH_PARALLELISM) as ex:
        futures = {
            ex.submit(_match_one_caption, i, cand, all_figs, ctx): i
            for i, cand in enumerate(candidates)
        }
        for fut in as_completed(futures):
            i = futures[fut]
            try:
                results[i] = fut.result()
            except Exception as e:
                log("ingest.pdf.caption_match_fail",
                    cand_idx=i, error=str(e))
                results[i] = {
                    "matched_idx": None, "confidence": 0.0,
                    "is_real_figure": True,
                    "reason": f"match call failed: {e}",
                    "caption_text": "",
                }
    return results


def _match_one_caption(
    cand_idx: int,
    candidate: PdfFigureCandidate,
    all_figs: list[dict[str, Any]],
    ctx: ToolContext,
) -> dict[str, Any]:
    # Filter captions to the candidate's page ± 1, fall back to whole
    # manifest if the page window is empty (handles figures whose caption
    # ends up on the next page due to column overflow).
    near: list[tuple[int, dict[str, Any]]] = [
        (i, f) for i, f in enumerate(all_figs)
        if abs(int(f.get("page", 0)) - candidate.page) <= 1
    ]
    pool = near if near else list(enumerate(all_figs))
    # Indices in `pool` are *relative* to the pool list we show the VLM;
    # we need to remap back to the full `all_figs` index.
    local_to_global = {local_i: global_i for local_i, (global_i, _) in enumerate(pool)}

    lines = []
    for local_i, (_, fig) in enumerate(pool):
        cap = (fig.get("caption") or "").replace("\n", " ")[:240]
        lines.append(f"  [{local_i}] (p.{fig.get('page', '?')}) {cap}")
    user_text = (
        f"Candidate captions near page {candidate.page} "
        f"(strategy: {candidate.strategy}):\n"
        + "\n".join(lines)
        + "\n\nReturn the JSON described in the system prompt."
    )

    result = vlm_call_json(
        settings=ctx.settings,
        model=ctx.settings.ingest_model,
        system=_CAPTION_MATCH_PROMPT,
        user_text=user_text,
        images=[VlmImage.from_path(candidate.path)],
        max_tokens=512,
    )

    # Remap local index → global manifest index, attach caption text.
    local_idx = result.get("matched_idx")
    global_idx = None
    caption_text = ""
    if isinstance(local_idx, int) and local_idx in local_to_global:
        global_idx = local_to_global[local_idx]
        caption_text = str(all_figs[global_idx].get("caption", ""))

    return {
        "matched_idx": global_idx,
        "confidence": float(result.get("confidence", 0.0) or 0.0),
        "is_real_figure": bool(result.get("is_real_figure", True)),
        "reason": str(result.get("reason", ""))[:200],
        "caption_text": caption_text,
    }


def _register_candidates(
    *,
    candidates: list[PdfFigureCandidate],
    matches: dict[int, dict[str, Any]],
    ctx: ToolContext,
    pdf_path: Path,
) -> list[str]:
    """Apply caption matches + fake-figure filter. Register survivors
    in `ctx.state["rendered_layers"]` with the downstream contract
    schema, then rename their PNGs to `img_{layer_id}.png` so the
    layers directory stays tidy."""
    registered: list[str] = []
    next_idx = 1

    for cand_idx, cand in enumerate(candidates):
        match = matches.get(cand_idx, {})
        is_real = bool(match.get("is_real_figure", True))
        if not is_real:
            log("ingest.pdf.reject_fake",
                page=cand.page, strategy=cand.strategy,
                reason=match.get("reason", ""), path=cand.path.name)
            # Clean up the rejected PNG so we don't bloat the run dir.
            try:
                cand.path.unlink()
            except OSError:
                pass
            continue

        layer_id = f"ingest_fig_{next_idx:02d}"
        next_idx += 1
        final_path = ctx.layers_dir / f"img_{layer_id}.png"
        try:
            if cand.path.resolve() != final_path.resolve():
                shutil.move(str(cand.path), str(final_path))
        except OSError as e:
            log("ingest.pdf.rename_fail",
                layer_id=layer_id, error=str(e))
            final_path = cand.path

        ctx.state["rendered_layers"][layer_id] = {
            "layer_id": layer_id,
            "name": f"figure_{next_idx - 1}",
            "kind": "image",
            "z_index": 5,
            "bbox": None,
            "src_path": str(final_path),
            "aspect_ratio": _aspect_from_dims(cand.width_px, cand.height_px),
            "image_size": f"{cand.width_px}x{cand.height_px}",
            "sha256": sha256_file(final_path),
            "source": "ingested_pdf",
            "source_file": str(pdf_path),
            "source_page": cand.page,
            "caption": match.get("caption_text", ""),
            "extract_strategy": cand.strategy,   # for debugging
            "caption_confidence": match.get("confidence", 0.0),
        }
        registered.append(layer_id)

    log("ingest.pdf.register",
        kept=len(registered),
        dropped=len(candidates) - len(registered))
    return registered


# ───────────────────────── Table parsing ───────────────────────────────

def _parse_tables_parallel(
    candidates: list[PdfTableCandidate],
    manifest: dict[str, Any],
    ctx: ToolContext,
) -> dict[int, dict[str, Any]]:
    """Run VLM parse per table candidate in a thread pool.

    Returns a dict: candidate idx → {is_table, headers, rows, title,
    matched_idx, caption_text, reason}. Callers use it to decide
    whether to register the candidate as a `kind="table"` layer.
    """
    results: dict[int, dict[str, Any]] = {}
    with ThreadPoolExecutor(max_workers=_CAPTION_MATCH_PARALLELISM) as ex:
        futures = {
            ex.submit(_parse_one_table, i, cand, manifest, ctx): i
            for i, cand in enumerate(candidates)
        }
        for fut in as_completed(futures):
            i = futures[fut]
            try:
                results[i] = fut.result()
            except Exception as e:
                log("ingest.pdf.table_parse_fail",
                    cand_idx=i, error=str(e))
                results[i] = {
                    "is_table": False, "headers": [], "rows": [],
                    "title": "", "matched_idx": None,
                    "caption_text": "", "reason": f"parse failed: {e}",
                }
    return results


def _parse_one_table(
    cand_idx: int,
    candidate: PdfTableCandidate,
    manifest: dict[str, Any],
    ctx: ToolContext,
) -> dict[str, Any]:
    # Filter caption candidates to same page ± 1 (tables sometimes
    # caption on the following page for bottom-of-page tables). Include
    # entries with missing/unknown `page` so appendix tables the VLM
    # failed to page-number aren't silently excluded from the pool.
    all_tables = list(manifest.get("tables", []))
    near: list[tuple[int, dict[str, Any]]] = []
    for i, t in enumerate(all_tables):
        raw_page = t.get("page")
        if raw_page is None:
            near.append((i, t))
            continue
        try:
            pnum = int(raw_page)
        except (TypeError, ValueError):
            near.append((i, t))
            continue
        if abs(pnum - candidate.page) <= 1:
            near.append((i, t))
    pool = near if near else list(enumerate(all_tables))
    local_to_global = {local_i: global_i for local_i, (global_i, _) in enumerate(pool)}

    cap_lines = []
    for local_i, (_, tbl) in enumerate(pool):
        cap = (tbl.get("caption") or "").replace("\n", " ")[:240]
        cap_lines.append(f"  [{local_i}] (p.{tbl.get('page', '?')}) {cap}")
    cap_block = "\n".join(cap_lines) or "  (no caption candidates on this page)"

    # Show the VLM pymupdf's best-effort cells, truncated. Don't send
    # the full raw cells when they're huge — we trust the image more.
    raw_preview: list[str] = []
    for row in candidate.raw_cells[:12]:
        raw_preview.append(
            " | ".join(str(c)[:80].replace("\n", " ⏎ ") for c in row)
        )
    raw_block = "\n".join(raw_preview) or "  (pymupdf extracted no cells)"

    user_text = (
        f"Table candidate from page {candidate.page} "
        f"(pymupdf saw {candidate.nrows}×{candidate.ncols} cells).\n\n"
        f"pymupdf raw-cell preview (may be wrong — trust the image):\n"
        f"{raw_block}\n\n"
        f"Caption candidates near this page:\n{cap_block}\n\n"
        "Return the JSON described in the system prompt."
    )

    result = vlm_call_json(
        settings=ctx.settings,
        model=ctx.settings.ingest_model,
        system=_TABLE_PARSE_PROMPT,
        user_text=user_text,
        images=[VlmImage.from_path(candidate.image_path)],
        max_tokens=4096,
    )

    local_idx = result.get("matched_idx")
    global_idx = None
    caption_text = str(result.get("title", "")).strip()
    if isinstance(local_idx, int) and local_idx in local_to_global:
        global_idx = local_to_global[local_idx]
        caption_text = str(all_tables[global_idx].get("caption", "")) or caption_text

    rows = result.get("rows") or []
    headers = result.get("headers") or []
    col_rule = result.get("col_highlight_rule") or []
    # Light sanitation: coerce to str, drop empty trailing rows.
    rows = [[str(c) if c is not None else "" for c in row] for row in rows]
    headers = [str(c) if c is not None else "" for c in headers]
    col_rule = [str(c) if c is not None else "" for c in col_rule]
    # Normalize rule length to match header count (pad with "" / truncate).
    if headers:
        n = len(headers)
        col_rule = col_rule[:n] + [""] * max(0, n - len(col_rule))
    while rows and not any(c.strip() for c in rows[-1]):
        rows.pop()

    return {
        "is_table": bool(result.get("is_table", False)),
        "headers": headers,
        "rows": rows,
        "col_highlight_rule": col_rule,
        "title": caption_text,
        "matched_idx": global_idx,
        "caption_text": caption_text,
        "reason": str(result.get("reason", ""))[:200],
    }


def _register_tables(
    *,
    candidates: list[PdfTableCandidate],
    parsed: dict[int, dict[str, Any]],
    ctx: ToolContext,
    pdf_path: Path,
) -> list[str]:
    """Register VLM-validated tables as `kind="table"` layers in
    `ctx.state["rendered_layers"]`. Rejects go to `reject_fake` logs.

    The PSD/SVG fallback PNG (a cleanly re-drawn table image) is
    generated lazily by the renderer via `util.table_png`, so we only
    need to persist the structured rows/headers here.
    """
    registered: list[str] = []
    next_idx = 1

    for cand_idx, cand in enumerate(candidates):
        info = parsed.get(cand_idx, {})
        if not info.get("is_table", False):
            log("ingest.pdf.reject_table",
                page=cand.page, reason=info.get("reason", ""))
            try:
                cand.image_path.unlink()
            except OSError:
                pass
            continue

        rows = info.get("rows") or []
        headers = info.get("headers") or []
        col_rule = info.get("col_highlight_rule") or []
        if not rows and not headers:
            log("ingest.pdf.reject_table",
                page=cand.page, reason="empty rows+headers")
            try:
                cand.image_path.unlink()
            except OSError:
                pass
            continue

        layer_id = f"ingest_table_{next_idx:02d}"
        next_idx += 1

        # Draw a clean PIL table PNG from the VLM-parsed rows — this
        # becomes the `src_path` fallback for poster/PSD/SVG paths that
        # don't have live-table primitives. (PPTX + HTML ignore
        # src_path and use rows/headers directly.)
        final_png = ctx.layers_dir / f"img_{layer_id}.png"
        # Bundled Noto fonts — CJK-capable; render_table_png falls back
        # to system fonts when these aren't found.
        fonts_dir = ctx.settings.fonts_dir
        noto_bold = fonts_dir / "NotoSansSC-Bold.otf"
        noto_regular = noto_bold  # project only bundles a bold Noto SC
        try:
            render_table_png(
                rows=rows, headers=headers, out_path=final_png,
                width_px=cand.width_px, max_height_px=cand.height_px,
                font_path=noto_regular if noto_regular.exists() else None,
                bold_font_path=noto_bold if noto_bold.exists() else None,
                col_highlight_rule=col_rule,
            )
        except Exception as e:
            log("ingest.pdf.table_render_fail",
                layer_id=layer_id, error=str(e))
            # Fall back to the source bbox crop.
            try:
                if cand.image_path.resolve() != final_png.resolve():
                    shutil.move(str(cand.image_path), str(final_png))
            except OSError:
                final_png = cand.image_path

        # Clean up the source bbox PNG if it's distinct from final_png
        # (we successfully drew a fresh table).
        if cand.image_path.exists() and cand.image_path.resolve() != final_png.resolve():
            try:
                cand.image_path.unlink()
            except OSError:
                pass

        ctx.state["rendered_layers"][layer_id] = {
            "layer_id": layer_id,
            "name": f"table_{next_idx - 1}",
            "kind": "table",
            "z_index": 5,
            "bbox": None,
            "src_path": str(final_png),        # bbox PNG fallback
            "aspect_ratio": _aspect_from_dims(cand.width_px, cand.height_px),
            "image_size": f"{cand.width_px}x{cand.height_px}",
            "sha256": sha256_file(final_png),
            "source": "ingested_pdf",
            "source_file": str(pdf_path),
            "source_page": cand.page,
            "caption": info.get("caption_text", ""),
            # structured data:
            "rows": rows,
            "headers": headers,
            "col_highlight_rule": col_rule,
            "title": info.get("title", ""),
        }
        registered.append(layer_id)

    log("ingest.pdf.register_tables",
        kept=len(registered),
        dropped=len(candidates) - len(registered))
    return registered


# ───────────────────────── Markdown branch ─────────────────────────────

def _ingest_markdown(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    text = fp.read_text(encoding="utf-8")
    registered: list[str] = []
    skipped: list[str] = []

    for m in _MD_IMG_RE.finditer(text):
        alt_text = m.group(1)
        ref = m.group(2).strip()
        if ref.startswith(("http://", "https://", "data:")):
            skipped.append(ref[:80])
            continue
        src = Path(ref)
        if not src.is_absolute():
            src = (fp.parent / src).resolve()
        if not src.exists() or not src.is_file():
            skipped.append(str(src))
            continue
        try:
            layer_id = _register_image_file(src, ctx, name_hint=alt_text or src.stem)
        except RuntimeError:
            skipped.append(str(src))
            continue
        registered.append(layer_id)

    log("ingest.md.done", file=fp.name, chars=len(text),
        registered=len(registered), skipped=len(skipped))

    return {
        "file": str(fp), "type": "markdown",
        "raw_text": text,
        "n_chars": len(text),
        "registered_layer_ids": registered,
        "skipped_images": skipped,
        "summary": f"{fp.name} — {len(text)} chars, "
                   f"{len(registered)} image(s)"
                   + (f", {len(skipped)} skipped" if skipped else ""),
    }


# ─────────────────────────── .docx branch ─────────────────────────────

def _ingest_docx(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    """Read a Word document into the same manifest shape as PDF.

    Docx has real structural metadata (heading styles, inline images,
    captions), so unlike PDF this branch does NOT need a VLM call —
    we read the docx tree directly, which is faster, free, and more
    faithful.
    """
    from docx import Document

    doc = Document(str(fp))

    body_paras: list[tuple[str, str]] = []  # (style_name, text)
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name if para.style is not None else "") or ""
        body_paras.append((style, text))

    title = _docx_pick_title(body_paras)
    sections = _docx_build_sections(body_paras, title)

    registered_figure_ids: list[str] = []
    for rel_id, rel in doc.part.rels.items():
        if rel.reltype and "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                mime = getattr(rel.target_part, "content_type", "") or ""
                ext = _ext_for_image_mime(mime)
                layer_id = _register_image_blob(
                    blob, ext, ctx,
                    name_hint=f"{fp.stem}_{rel_id}",
                    source_file=fp,
                    source_ref=f"rel={rel_id}",
                )
                registered_figure_ids.append(layer_id)
            except (RuntimeError, OSError) as e:
                log("ingest.docx.image_skip", rel=rel_id, error=str(e))

    manifest = {
        "title": title or fp.stem,
        "authors": [],
        "venue": None,
        "abstract": sections[0]["summary"] if sections else "",
        "sections": sections,
        "figures": [],
        "tables": [],
        "key_quotes": [],
    }
    _normalize_manifest_lists(manifest)

    log("ingest.docx.done", file=fp.name,
        sections=len(sections), figures=len(registered_figure_ids))

    return {
        "file": str(fp), "type": "docx", "manifest": manifest,
        "registered_layer_ids": registered_figure_ids,
        "registered_figure_ids": registered_figure_ids,
        "registered_table_ids": [],
        "summary": f"{manifest['title']} — "
                   f"{len(registered_figure_ids)} figure(s), "
                   f"{len(sections)} section(s)",
    }


_DOCX_HEADING_PREFIXES = ("Heading", "Title")


def _docx_pick_title(paras: list[tuple[str, str]]) -> str:
    for style, text in paras:
        if style.startswith("Title"):
            return text
    for style, text in paras:
        if style.startswith("Heading 1") or style == "Heading 1":
            return text
    for _style, text in paras:
        return text
    return ""


def _docx_build_sections(
    paras: list[tuple[str, str]], title: str,
) -> list[dict[str, Any]]:
    """Group paragraphs under their nearest preceding heading. Non-
    heading paras become the section body; the first sentence (or 2-3
    sentences up to ~400 chars) becomes `summary`, and short bullet-like
    lines become `key_points`."""
    sections: list[dict[str, Any]] = []
    current_heading: str | None = None
    current_body: list[str] = []

    def flush() -> None:
        nonlocal current_heading, current_body
        if current_heading is None and not current_body:
            return
        heading = current_heading or "Body"
        body_text = "\n".join(current_body).strip()
        sections.append({
            "idx": len(sections) + 1,
            "heading": heading,
            "summary": _first_sentences(body_text, max_chars=400),
            "key_points": _pick_key_points(current_body),
        })
        current_heading = None
        current_body = []

    for style, text in paras:
        is_heading = any(style.startswith(p) for p in _DOCX_HEADING_PREFIXES)
        if is_heading and text == title:
            continue
        if is_heading:
            flush()
            current_heading = text
        else:
            current_body.append(text)
    flush()

    return sections


def _first_sentences(text: str, max_chars: int = 400) -> str:
    if not text:
        return ""
    pieces: list[str] = re.split(r"(?<=[.!?。！？])\s+", text)
    acc = ""
    for p in pieces:
        if not p:
            continue
        if len(acc) + len(p) + 1 > max_chars:
            break
        acc = (acc + " " + p).strip() if acc else p
    return acc or text[:max_chars]


def _pick_key_points(body: list[str]) -> list[str]:
    """Pull bullet-like lines (short, starts with punctuation/enumeration
    marker, OR simply short standalone paras ≤ 120 chars). Max 5."""
    out: list[str] = []
    for line in body:
        s = line.strip()
        if not s or len(s) > 160:
            continue
        if re.match(r"^[-\*•]|^\d+[\.\)]\s", s) or len(s) <= 120:
            out.append(s.lstrip("-*•").strip())
        if len(out) >= 5:
            break
    return out


def _ext_for_image_mime(mime: str) -> str:
    mime = (mime or "").lower()
    if "png" in mime:
        return ".png"
    if "jpeg" in mime or "jpg" in mime:
        return ".jpg"
    if "gif" in mime:
        return ".gif"
    if "webp" in mime:
        return ".webp"
    if "bmp" in mime:
        return ".bmp"
    if "tif" in mime:
        return ".tif"
    return ".png"


# ─────────────────────────── .pptx branch ─────────────────────────────

def _ingest_pptx(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    """Read a PowerPoint file into the same manifest shape. Each slide
    becomes one section; title placeholder → section heading, body
    placeholders → section body; picture shapes → ingest_fig_NN layers.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(fp))

    sections: list[dict[str, Any]] = []
    registered_figure_ids: list[str] = []
    deck_title: str | None = None

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title: str | None = None
        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    mime = shape.image.content_type or ""
                    ext = _ext_for_image_mime(mime)
                    layer_id = _register_image_blob(
                        blob, ext, ctx,
                        name_hint=f"{fp.stem}_slide{slide_idx:02d}",
                        source_file=fp,
                        source_ref=f"slide={slide_idx}",
                    )
                    registered_figure_ids.append(layer_id)
                except (RuntimeError, OSError) as e:
                    log("ingest.pptx.image_skip",
                        slide=slide_idx, error=str(e))
                continue
            if not shape.has_text_frame:
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            is_title = (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
            )
            if is_title and slide_title is None:
                slide_title = text
            else:
                body_parts.append(text)

        if slide_idx == 1 and slide_title:
            deck_title = slide_title

        heading = slide_title or f"Slide {slide_idx}"
        body_text = "\n".join(body_parts).strip()
        sections.append({
            "idx": slide_idx,
            "heading": heading,
            "summary": _first_sentences(body_text, max_chars=400),
            "key_points": _pick_key_points(body_parts),
        })

    manifest = {
        "title": deck_title or fp.stem,
        "authors": [],
        "venue": None,
        "abstract": sections[0]["summary"] if sections else "",
        "sections": sections,
        "figures": [],
        "tables": [],
        "key_quotes": [],
    }
    _normalize_manifest_lists(manifest)

    log("ingest.pptx.done", file=fp.name,
        slides=len(sections), figures=len(registered_figure_ids))

    return {
        "file": str(fp), "type": "pptx", "manifest": manifest,
        "registered_layer_ids": registered_figure_ids,
        "registered_figure_ids": registered_figure_ids,
        "registered_table_ids": [],
        "summary": f"{manifest['title']} — "
                   f"{len(sections)} slide(s), "
                   f"{len(registered_figure_ids)} figure(s)",
    }


def _register_image_blob(
    blob: bytes, ext: str, ctx: ToolContext, *,
    name_hint: str, source_file: Path, source_ref: str,
) -> str:
    """Register an in-memory image blob as an `ingest_fig_NN` layer.

    Unlike `_register_image_file` which expects an on-disk source and
    uses the sha-based layer_id shape (`ingest_img_<sha8>`), this helper
    allocates a sequential `ingest_fig_NN` id so .docx / .pptx images
    show up to the planner the same way PDF figures do — and the
    figure-cross-reference detector in composite picks them up too.
    """
    import hashlib
    from PIL import Image as _Image
    import io

    # Sequential id — peek at rendered_layers to find next free index.
    existing = [
        k for k in ctx.state["rendered_layers"]
        if k.startswith("ingest_fig_")
    ]
    next_idx = len(existing) + 1
    layer_id = f"ingest_fig_{next_idx:02d}"

    dest = ctx.layers_dir / f"img_{layer_id}{ext}"
    dest.write_bytes(blob)

    try:
        with _Image.open(io.BytesIO(blob)) as im:
            w, h = im.size
    except Exception as e:
        try:
            dest.unlink()
        except OSError:
            pass
        raise RuntimeError(f"blob not readable ({e})")

    sha = hashlib.sha256(blob).hexdigest()
    ctx.state["rendered_layers"][layer_id] = {
        "layer_id": layer_id,
        "name": _sanitize_name(name_hint) or layer_id,
        "kind": "image",
        "z_index": 5,
        "bbox": None,
        "src_path": str(dest),
        "aspect_ratio": _aspect_from_dims(w, h),
        "image_size": f"{w}x{h}",
        "sha256": sha,
        "source": "ingested_" + ("pptx" if source_ref.startswith("slide=") else "docx"),
        "source_file": str(source_file),
        "source_ref": source_ref,
        "caption": "",
        "extract_strategy": "embedded",
        "caption_confidence": 0.0,
    }
    return layer_id


# ─────────────────────────── Image branch ──────────────────────────────

def _ingest_image(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    layer_id = _register_image_file(fp, ctx, name_hint=fp.stem)
    rec = ctx.state["rendered_layers"][layer_id]
    w_s, h_s = rec["image_size"].split("x")
    log("ingest.image.done", file=fp.name, layer_id=layer_id,
        w=int(w_s), h=int(h_s))
    return {
        "file": str(fp), "type": "image",
        "registered_layer_ids": [layer_id],
        "width": int(w_s), "height": int(h_s),
        "summary": f"{fp.name} — {w_s}×{h_s}",
    }


# ───────────────────────────── helpers ─────────────────────────────────

def _register_image_file(src: Path, ctx: ToolContext, *, name_hint: str) -> str:
    from PIL import Image as _Image

    sha = sha256_file(src)
    ext = src.suffix.lower() if src.suffix else ".png"
    if ext == ".jpeg":
        ext = ".jpg"
    layer_id = f"ingest_img_{sha[:8]}"
    dest = ctx.layers_dir / f"img_{layer_id}{ext}"
    if not dest.exists():
        shutil.copy2(src, dest)

    try:
        with _Image.open(dest) as im:
            w, h = im.size
    except Exception as e:
        raise RuntimeError(f"image not readable: {src} ({e})")

    ctx.state["rendered_layers"][layer_id] = {
        "layer_id": layer_id,
        "name": _sanitize_name(name_hint) or layer_id,
        "kind": "image",
        "z_index": 5,
        "bbox": None,
        "src_path": str(dest),
        "aspect_ratio": _aspect_from_dims(w, h),
        "image_size": f"{w}x{h}",
        "sha256": sha,
        "source": "ingested",
        "source_file": str(src),
    }
    return layer_id


def _rank_figure_ids_for_planner(
    figure_ids: list[str], rendered: dict[str, dict[str, Any]],
) -> list[str]:
    """Order figure ids so the ones most useful to a poster planner
    come first. Ranking signal, in decreasing priority:

    1. has a non-empty caption (carries explicit author intent),
    2. larger min-dimension (more printable at poster scale),
    3. source strategy = "vector" (composite diagrams > raster sub-panels),
    4. smaller page number (main-paper figures before appendix ones).
    """
    def key(fid: str) -> tuple[int, int, int, int]:
        rec = rendered.get(fid) or {}
        cap = (rec.get("caption") or "").strip()
        size = rec.get("image_size") or "0x0"
        try:
            w_s, h_s = size.split("x")
            w = int(w_s); h = int(h_s)
            side = min(w, h)
        except Exception:
            side = 0
        strat_rank = 1 if rec.get("extract_strategy") == "vector" else 0
        has_caption_rank = 1 if cap else 0
        try:
            page = int(rec.get("source_page") or 999)
        except (TypeError, ValueError):
            page = 999
        # Sort DESC by caption/strategy/size, ASC by page → negate page.
        return (has_caption_rank, strat_rank, side, -page)

    return sorted(figure_ids, key=key, reverse=True)


def _aspect_from_dims(w: int, h: int) -> str:
    if h <= 0 or w <= 0:
        return "1:1"
    from math import gcd
    g = gcd(w, h)
    return f"{w // g}:{h // g}" if max(w // g, h // g) <= 32 else (
        "16:9" if w > h else "3:4" if h > w else "1:1"
    )


def _sanitize_name(s: str) -> str:
    s = (s or "").strip()
    return re.sub(r"[^\w\- ]", "", s)[:60]
