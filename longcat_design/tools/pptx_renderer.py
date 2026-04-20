"""PPTX renderer for deck artifacts (v1.0 #7).

One top-level LayerNode with `kind="slide"` per slide; each slide's `children`
hold `kind in {"text","image","background"}` elements positioned by bbox in
slide-canvas pixel coords. We emit native PowerPoint shapes with native
TextFrames — no Pillow rasterization of text — so the .pptx opens
type-editable in PowerPoint / Keynote / Google Slides.

Per-slide simplified PNGs are also written (for chat preview + critic).

Font embedding is intentionally NOT performed: .pptx delegates font rendering
to the consuming app's font engine (see `docs/GOTCHAS.md` for the CJK-on-
consumer-PowerPoint note). This mirrors Paper2Any's approach.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from ._contract import ToolContext


# 1 pixel at 96 DPI ≈ 9525 EMU. python-pptx accepts Emu(int).
PX_TO_EMU = 9525

DEFAULT_SLIDE_W = 1920
DEFAULT_SLIDE_H = 1080

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def write_pptx(spec: Any, pptx_path: Path, ctx: ToolContext) -> int:
    """Walk the slide tree and emit a .pptx file. Returns slide count."""
    canvas = spec.canvas or {}
    slide_w = int(canvas.get("w_px") or DEFAULT_SLIDE_W)
    slide_h = int(canvas.get("h_px") or DEFAULT_SLIDE_H)

    prs = Presentation()
    prs.slide_width = Emu(slide_w * PX_TO_EMU)
    prs.slide_height = Emu(slide_h * PX_TO_EMU)

    blank_layout = prs.slide_layouts[6]  # "Blank" layout — we position everything.

    slide_count = 0
    for node in (spec.layer_graph or []):
        if getattr(node, "kind", None) != "slide":
            continue
        slide = prs.slides.add_slide(blank_layout)
        _render_slide(slide, node, slide_w, slide_h, ctx)
        slide_count += 1

    prs.save(str(pptx_path))
    return slide_count


def _render_slide(slide: Any, slide_node: Any, slide_w: int, slide_h: int,
                  ctx: ToolContext) -> None:
    """Add shapes for each child element of a slide LayerNode."""
    children = list(getattr(slide_node, "children", None) or [])
    # Sort by z_index so higher z draws on top (pptx respects insertion order).
    children.sort(key=lambda c: int(getattr(c, "z_index", 0) or 0))

    for child in children:
        kind = getattr(child, "kind", None)
        if kind == "background":
            _add_background(slide, child, slide_w, slide_h)
        elif kind == "image":
            _add_picture(slide, child, slide_w, slide_h)
        elif kind == "text":
            _add_text_frame(slide, child, slide_w, slide_h)
        # silently skip unknown kinds; planner enforces vocab


def _bbox_to_emu(bbox: Any, slide_w: int, slide_h: int,
                 default_bbox: tuple[int, int, int, int] | None = None
                 ) -> tuple[Emu, Emu, Emu, Emu]:
    """Resolve a bbox (or fallback to slide-sized) to EMU left/top/width/height."""
    if bbox is not None:
        x = int(getattr(bbox, "x", 0) or 0)
        y = int(getattr(bbox, "y", 0) or 0)
        w = int(getattr(bbox, "w", slide_w) or slide_w)
        h = int(getattr(bbox, "h", slide_h) or slide_h)
    elif default_bbox is not None:
        x, y, w, h = default_bbox
    else:
        x, y, w, h = 0, 0, slide_w, slide_h
    # clamp to slide bounds
    x = max(0, min(x, slide_w - 1))
    y = max(0, min(y, slide_h - 1))
    w = max(1, min(w, slide_w - x))
    h = max(1, min(h, slide_h - y))
    return Emu(x * PX_TO_EMU), Emu(y * PX_TO_EMU), Emu(w * PX_TO_EMU), Emu(h * PX_TO_EMU)


def _add_background(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    """Full-slide picture background. Planner can pass bbox or leave None to
    cover the whole slide."""
    src = getattr(node, "src_path", None)
    if not src:
        return  # no background image available; PowerPoint default is white
    if not Path(src).exists():
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
        default_bbox=(0, 0, slide_w, slide_h),
    )
    slide.shapes.add_picture(src, left, top, width=width, height=height)


def _add_picture(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    src = getattr(node, "src_path", None)
    if not src or not Path(src).exists():
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
    )
    slide.shapes.add_picture(src, left, top, width=width, height=height)


def _add_text_frame(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    text = (getattr(node, "text", None) or "").strip()
    if not text:
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
    )
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    # First paragraph holds the initial run; subsequent lines become new paras.
    lines = text.splitlines() or [text]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ""  # clear any default run
        align = getattr(node, "align", None)
        if align in _ALIGN_MAP:
            para.alignment = _ALIGN_MAP[align]
        run = para.add_run()
        run.text = line
        font = run.font
        # python-pptx wants pts; we get px. pt ≈ px * 72/96 = px * 0.75.
        size_px = int(getattr(node, "font_size_px", None) or 36)
        font.size = Pt(max(6, round(size_px * 0.75)))
        family = getattr(node, "font_family", None)
        if family:
            font.name = family
        effects = getattr(node, "effects", None)
        fill_hex = None
        if effects is not None:
            fill_hex = getattr(effects, "fill", None)
        if fill_hex and isinstance(fill_hex, str) and fill_hex.startswith("#"):
            rgb = _hex_to_rgb(fill_hex)
            if rgb is not None:
                font.color.rgb = RGBColor(*rgb)


def _hex_to_rgb(hx: str) -> tuple[int, int, int] | None:
    s = hx.lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        return None
    try:
        return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    except ValueError:
        return None


def render_slide_preview_png(slide_node: Any, slide_w: int, slide_h: int,
                             out_path: Path, ctx: ToolContext) -> None:
    """Pillow-render a simplified preview of one slide.

    Not pixel-accurate with PowerPoint's renderer — it's an at-a-glance thumb
    for chat UX + for stitching into the grid preview. Shows the bg image (if
    any), image shapes (as resized thumbs), and text (approximated font).
    """
    # Downscale to a chat-friendly size while preserving aspect.
    max_w = 960
    scale = min(1.0, max_w / slide_w)
    w = max(1, int(slide_w * scale))
    h = max(1, int(slide_h * scale))

    img = Image.new("RGB", (w, h), (255, 255, 255))

    children = sorted(
        list(getattr(slide_node, "children", None) or []),
        key=lambda c: int(getattr(c, "z_index", 0) or 0),
    )

    for child in children:
        kind = getattr(child, "kind", None)
        if kind in ("background", "image"):
            _paste_image(img, child, slide_w, slide_h, scale)
        elif kind == "text":
            _draw_text(img, child, slide_w, slide_h, scale, ctx)

    img.save(out_path, format="PNG", optimize=True)


def _scaled_bbox(bbox: Any, slide_w: int, slide_h: int,
                 scale: float,
                 default_full: bool = False) -> tuple[int, int, int, int]:
    if bbox is not None:
        x = int(getattr(bbox, "x", 0) or 0)
        y = int(getattr(bbox, "y", 0) or 0)
        w = int(getattr(bbox, "w", slide_w) or slide_w)
        h = int(getattr(bbox, "h", slide_h) or slide_h)
    elif default_full:
        x, y, w, h = 0, 0, slide_w, slide_h
    else:
        x, y, w, h = 0, 0, slide_w // 2, slide_h // 4
    return (
        max(0, int(x * scale)),
        max(0, int(y * scale)),
        max(1, int(w * scale)),
        max(1, int(h * scale)),
    )


def _paste_image(canvas: Image.Image, node: Any, slide_w: int, slide_h: int,
                 scale: float) -> None:
    src = getattr(node, "src_path", None)
    if not src or not Path(src).exists():
        return
    try:
        tile = Image.open(src).convert("RGBA")
    except Exception:
        return
    default_full = getattr(node, "kind", None) == "background"
    sx, sy, sw, sh = _scaled_bbox(
        getattr(node, "bbox", None), slide_w, slide_h, scale,
        default_full=default_full,
    )
    tile = tile.resize((sw, sh), Image.LANCZOS)
    canvas.alpha_composite(tile, dest=(sx, sy)) if canvas.mode == "RGBA" else canvas.paste(tile, (sx, sy), tile)


def _draw_text(canvas: Image.Image, node: Any, slide_w: int, slide_h: int,
               scale: float, ctx: ToolContext) -> None:
    text = (getattr(node, "text", None) or "").strip()
    if not text:
        return
    sx, sy, sw, sh = _scaled_bbox(
        getattr(node, "bbox", None), slide_w, slide_h, scale,
        default_full=False,
    )

    # Pick a reasonable approximated font size from the node metadata.
    size_px = int(getattr(node, "font_size_px", None) or 36)
    approx = max(10, min(120, int(size_px * scale)))
    family = getattr(node, "font_family", None) or ctx.settings.default_text_font

    fonts = ctx.settings.fonts
    fname = fonts.get(family) or fonts[ctx.settings.default_text_font]
    try:
        font = ImageFont.truetype(str(ctx.settings.fonts_dir / fname), size=approx)
    except Exception:
        font = ImageFont.load_default()

    effects = getattr(node, "effects", None)
    fill_hex = getattr(effects, "fill", None) if effects is not None else None
    rgb = _hex_to_rgb(fill_hex) if isinstance(fill_hex, str) else None
    fill = rgb or (15, 23, 42)

    draw = ImageDraw.Draw(canvas)
    # Word-wrap to bbox: simple char-by-char wrap (works for CJK; latin uses spaces).
    lines = _wrap_for_width(text, font, sw, draw)
    line_h = approx + 6
    y = sy
    for line in lines:
        if y + line_h > sy + sh:
            break
        draw.text((sx, y), line, fill=fill, font=font)
        y += line_h


def _wrap_for_width(text: str, font: Any, max_w: int, draw: Any) -> list[str]:
    """Simple word/char wrap to fit a pixel width."""
    out: list[str] = []
    for paragraph in text.splitlines() or [text]:
        if " " in paragraph:
            # space-delimited: greedy word-wrap
            words = paragraph.split()
            line = ""
            for w in words:
                probe = (line + " " + w).strip()
                if _measure(probe, font, draw) <= max_w:
                    line = probe
                else:
                    if line:
                        out.append(line)
                    line = w
            if line:
                out.append(line)
        else:
            # CJK-style: char-by-char
            line = ""
            for ch in paragraph:
                probe = line + ch
                if _measure(probe, font, draw) <= max_w:
                    line = probe
                else:
                    if line:
                        out.append(line)
                    line = ch
            if line:
                out.append(line)
    return out


def _measure(s: str, font: Any, draw: Any) -> int:
    try:
        bbox = draw.textbbox((0, 0), s, font=font)
        return int(bbox[2] - bbox[0])
    except Exception:
        return len(s) * 10
