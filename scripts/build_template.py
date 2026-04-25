"""Deterministic generator for `assets/deck_templates/academic-editorial.pptx`.

Run with `uv run python scripts/build_template.py` to (re)generate the
template. The output .pptx is committed to the repo so production runs
don't need this script — the script exists so design-token edits are a
code review, not a binary diff.

Design contract (v2.5.2 — `academic-editorial`):
- Canvas 1920×1080 px @ 96 DPI (matches OpenDesign's `PX_TO_EMU` math
  in `open_design/tools/pptx_renderer.py`).
- Cream bg `#FAF7F0`, ink `#0F172A`, accent `#7F1D1D` oxblood,
  rule `#94A3B8`, muted `#475569`. Single accent everywhere.
- Title font `PlayfairDisplay`; body font `Inter`. Both are bundled in
  `assets/fonts/`. PowerPoint/Keynote/Google Slides will substitute if
  the consumer doesn't have them, which is acceptable for paper-deck
  hand-off.
- 6 layout slides at indices 0..5 (the renderer maps `slide.role` →
  index): cover / section_divider / content / content_with_figure /
  content_with_table / closing.

Each layout carries:
- Decorative shapes (cream bg, stripes, rules) — renderer leaves alone.
- Named text/image slots — renderer fills by `shape.name` lookup. Slot
  names match `LayerNode.template_slot` values.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt


REPO_ROOT = Path(__file__).resolve().parent.parent
OUT_PATH = REPO_ROOT / "assets" / "deck_templates" / "academic-editorial.pptx"

# ── Design tokens ─────────────────────────────────────────────────────
COLOR_BG = RGBColor(0xFA, 0xF7, 0xF0)        # cream
COLOR_INK = RGBColor(0x0F, 0x17, 0x2A)       # near-black
COLOR_ACCENT = RGBColor(0x7F, 0x1D, 0x1D)    # oxblood
COLOR_RULE = RGBColor(0x94, 0xA3, 0xB8)      # slate-400
COLOR_MUTED = RGBColor(0x47, 0x55, 0x69)     # slate-600
COLOR_SLOT_HINT = RGBColor(0xE8, 0xE3, 0xD8)  # darker cream — image_slot
COLOR_BADGE_FG = RGBColor(0xFA, 0xF7, 0xF0)  # cream text on accent

FONT_TITLE = "Playfair Display"
FONT_BODY = "Inter"

# ── Canvas math (matches PX_TO_EMU=9525 in pptx_renderer) ─────────────
PX = 9525  # 1 px @ 96 DPI = 9525 EMU
SLIDE_W = 1920
SLIDE_H = 1080


def px_emu(v: int) -> int:
    return Emu(v * PX)


# ── Shape helpers ─────────────────────────────────────────────────────


def add_cream_bg(slide):
    """Full-bleed cream rectangle. Sent to back so other shapes layer on top."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, px_emu(SLIDE_W), px_emu(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()
    bg.name = "background_fill"
    # Send to back via XML manipulation — python-pptx doesn't expose order.
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    # Insert at index 2 (after the required nvGrpSpPr + grpSpPr children).
    spTree.insert(2, bg._element)
    return bg


def add_rect(slide, name: str, x: int, y: int, w: int, h: int, color: RGBColor, *, has_border: bool = False):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, px_emu(x), px_emu(y), px_emu(w), px_emu(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    if has_border:
        rect.line.color.rgb = COLOR_RULE
        rect.line.width = Emu(PX)
    else:
        rect.line.fill.background()
    rect.name = name
    return rect


def add_text(
    slide, name: str, x: int, y: int, w: int, h: int, text: str,
    *, font: str = FONT_BODY, size_pt: int = 16, color: RGBColor = COLOR_INK,
    align: str = "left", italic: bool = False, bold: bool = False,
    uppercase: bool = False, letter_spacing_em: float | None = None,
):
    tb = slide.shapes.add_textbox(px_emu(x), px_emu(y), px_emu(w), px_emu(h))
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }[align]
    run = p.add_run()
    run.text = text.upper() if uppercase else text
    run.font.name = font
    run.font.size = Pt(size_pt)
    run.font.color.rgb = color
    if italic:
        run.font.italic = True
    if bold:
        run.font.bold = True
    if letter_spacing_em is not None:
        # python-pptx exposes spc on the rPr XML — convert em → 1/100 pt.
        rPr = run.font._rPr
        rPr.set("spc", str(int(letter_spacing_em * size_pt * 100)))
    return tb


def add_image_slot(slide, name: str, x: int, y: int, w: int, h: int):
    """Empty rectangle marking where an image goes. Renderer reads its
    bbox + removes this shape + adds the actual picture letterbox-fit."""
    return add_rect(slide, name, x, y, w, h, COLOR_SLOT_HINT, has_border=True)


# ── Layouts ───────────────────────────────────────────────────────────


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_cream_bg(slide)
    # Left accent stripe (decorative — not named, planner can't address)
    add_rect(slide, "accent_stripe_cover", 0, 0, 8, SLIDE_H, COLOR_ACCENT)
    # Title — left 50% (right is image)
    add_text(
        slide, "title", 96, 280, 880, 280,
        "Paper Title Goes Here",
        font=FONT_TITLE, size_pt=56, color=COLOR_INK, align="left",
    )
    # Authors row
    add_text(
        slide, "authors", 96, 580, 880, 60,
        "Author One · Author Two · Affiliation",
        font=FONT_BODY, size_pt=22, color=COLOR_MUTED, align="left",
    )
    # Conference badge — top-right oxblood pill with cream text
    add_rect(slide, "badge_fill", 1660, 80, 180, 40, COLOR_ACCENT)
    add_text(
        slide, "badge", 1660, 80, 180, 40, "NeurIPS 2026",
        font=FONT_BODY, size_pt=13, color=COLOR_BADGE_FG,
        align="center", bold=True, uppercase=True,
    )
    # Hero image slot — right ~50%, full-bleed
    add_image_slot(slide, "image_slot", 1000, 0, 920, SLIDE_H)
    return slide


def build_section_divider(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_cream_bg(slide)
    # Right-edge oxblood block (rhymes with cover's left stripe)
    add_rect(slide, "accent_block_divider", 1820, 0, 100, SLIDE_H, COLOR_ACCENT)
    # Section number above title
    add_text(
        slide, "section_number", 200, 320, 1520, 50,
        "01 ·", font=FONT_BODY, size_pt=24, color=COLOR_ACCENT,
        align="left", bold=True,
    )
    # Big section title
    add_text(
        slide, "title", 200, 380, 1520, 200,
        "Method →", font=FONT_TITLE, size_pt=120, color=COLOR_INK,
        align="left",
    )
    # Optional preview line
    add_text(
        slide, "subtitle", 200, 600, 1520, 60,
        "Section preview text in italic", font=FONT_BODY, size_pt=18,
        color=COLOR_MUTED, align="left", italic=True,
    )
    return slide


def _add_content_chrome(slide):
    """Shared decorations + named slots for content / content_with_figure /
    content_with_table layouts: section label + title + title-bar rule +
    footer hairline + footer slots."""
    add_cream_bg(slide)
    # Section label (top-left, oxblood, uppercase, tracked)
    add_text(
        slide, "section_label", 96, 80, 1728, 30,
        "01 · METHOD", font=FONT_BODY, size_pt=14, color=COLOR_ACCENT,
        align="left", bold=True, uppercase=True, letter_spacing_em=0.08,
    )
    # Title
    add_text(
        slide, "title", 96, 120, 1728, 80,
        "Slide Title", font=FONT_TITLE, size_pt=36, color=COLOR_INK,
        align="left",
    )
    # Title-bar rule (short oxblood line under title)
    add_rect(slide, "title_rule", 96, 220, 200, 1, COLOR_ACCENT)
    # Footer hairline above footer
    add_rect(slide, "footer_rule", 96, 1010, 1728, 1, COLOR_RULE)
    # Footer left — paper title (auto-filled by renderer with
    # ds.footer_text or first 60 chars of brief)
    add_text(
        slide, "footer", 96, 1020, 1200, 30,
        "Paper title here · authors · venue",
        font=FONT_BODY, size_pt=11, color=COLOR_MUTED,
        align="left", italic=True,
    )
    # Footer right — slide number (auto-filled by renderer with "N/total")
    add_text(
        slide, "slide_number", 1700, 1020, 124, 30, "N/N",
        font=FONT_BODY, size_pt=11, color=COLOR_INK, align="right",
    )


def build_content(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_content_chrome(slide)
    # Body — full-width
    add_text(
        slide, "body", 96, 260, 1728, 740,
        "Body content goes here. Multiple paragraphs supported via "
        "newlines. Inter 22pt, slate ink.",
        font=FONT_BODY, size_pt=22, color=COLOR_INK, align="left",
    )
    return slide


def build_content_with_figure(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_content_chrome(slide)
    # Body — left 920px
    add_text(
        slide, "body", 96, 260, 920, 740,
        "Left-column body text. Reference the figure on the right.",
        font=FONT_BODY, size_pt=22, color=COLOR_INK, align="left",
    )
    # Image slot — right 768px
    add_image_slot(slide, "image_slot", 1056, 260, 768, 740)
    return slide


def build_content_with_table(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_content_chrome(slide)
    # Body — left 800px
    add_text(
        slide, "body", 96, 260, 800, 740,
        "Left-column commentary. Reference the table on the right.",
        font=FONT_BODY, size_pt=22, color=COLOR_INK, align="left",
    )
    # Table anchor — right 904px (renderer adds native PPTX table at this bbox)
    add_image_slot(slide, "table_anchor", 920, 260, 904, 740)
    return slide


def build_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_cream_bg(slide)
    # Left accent stripe (rhymes with cover)
    add_rect(slide, "accent_stripe_closing", 0, 0, 8, SLIDE_H, COLOR_ACCENT)
    # "Thank you"
    add_text(
        slide, "title", 200, 360, 1520, 160, "Thank you",
        font=FONT_TITLE, size_pt=96, color=COLOR_INK, align="center",
    )
    # "Q&A"
    add_text(
        slide, "subtitle", 200, 540, 1520, 50, "Q&A",
        font=FONT_BODY, size_pt=28, color=COLOR_INK, align="center",
    )
    # arXiv link / contact (oxblood)
    add_text(
        slide, "links", 200, 620, 1520, 30,
        "arxiv.org/abs/2026.XXXXX · contact@example.org",
        font=FONT_BODY, size_pt=16, color=COLOR_ACCENT, align="center",
    )
    return slide


# ── Main ──────────────────────────────────────────────────────────────


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W * PX)
    prs.slide_height = Emu(SLIDE_H * PX)

    # Order MUST match ROLE_TO_LAYOUT_IDX in pptx_renderer.
    build_cover(prs)               # idx 0
    build_section_divider(prs)     # idx 1
    build_content(prs)             # idx 2
    build_content_with_figure(prs) # idx 3
    build_content_with_table(prs)  # idx 4
    build_closing(prs)             # idx 5

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    out = build()
    print(f"wrote {out} ({out.stat().st_size:,} bytes, 6 slides)")
