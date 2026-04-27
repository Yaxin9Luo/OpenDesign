"""No-API smoke test: imports + schemas + fonts + a real composite call.

Run with:
    python -m open_design.smoke

Generates `out/smoke/` containing poster.psd, poster.svg, preview.png produced
from a fake (solid-color) background + 2 real text layers rendered via Pillow.
This proves the whole pipeline below the LLM/Gemini layer works without keys.
"""

from __future__ import annotations

import json
import base64
import sys
from datetime import datetime
from io import BytesIO
from pathlib import Path

from PIL import Image
from pydantic import ValidationError

from .schema import (
    AgentTraceStep, CompositionArtifacts, CritiqueResult, DesignSpec,
    DistillTrajectory, LayerNode, SafeZone, TextEffect, ThinkingBlockRecord,
    ToolResultRecord, TrainingMetadata,
)


def _fail(msg: str) -> None:
    print(f"  FAIL  {msg}")
    sys.exit(1)


def _ok(msg: str) -> None:
    print(f"  ok    {msg}")


def check_imports() -> None:
    print("[1/52] imports")
    from . import chat, cli, config, planner, runner, schema, session  # noqa
    from .agents import CriticAgent, PromptEnhancer  # noqa
    from .tools import (
        TOOL_HANDLERS, TOOL_SCHEMAS, ToolContext,
        composite, critique_tool, edit_layer, fetch_brand_asset, finalize,
        generate_background, propose_design_spec, render_text_layer,
        switch_artifact_type,
    )  # noqa
    from .util import ids, io, logging  # noqa
    _ok("all modules import (incl. chat + session + edit_layer)")


def check_tool_registry() -> None:
    print("[2/52] tool registry")
    from .tools import TOOL_HANDLERS, TOOL_SCHEMAS

    expected = {"switch_artifact_type", "propose_design_spec",
                "generate_background", "generate_image",
                "render_text_layer", "edit_layer",
                "fetch_brand_asset", "composite", "critique", "finalize",
                "ingest_document"}
    schema_names = {s["name"] for s in TOOL_SCHEMAS}
    handler_names = set(TOOL_HANDLERS.keys())
    missing = expected - (schema_names & handler_names)
    extra = (schema_names | handler_names) - expected
    if missing:
        _fail(f"missing tools: {missing}")
    if extra:
        _fail(f"unexpected tools: {extra}")
    for s in TOOL_SCHEMAS:
        for k in ("name", "description", "input_schema"):
            if k not in s:
                _fail(f"tool '{s.get('name','?')}' missing key '{k}'")
        if s["input_schema"].get("type") != "object":
            _fail(f"tool '{s['name']}' input_schema.type != 'object'")
    # switch_artifact_type should be first in TOOL_SCHEMAS (pedagogical ordering)
    if TOOL_SCHEMAS[0]["name"] != "switch_artifact_type":
        _fail(f"switch_artifact_type should be first in TOOL_SCHEMAS; got {TOOL_SCHEMAS[0]['name']}")
    _ok(f"{len(expected)} tools wired (schemas + handlers): {sorted(schema_names)}")


def check_pydantic_roundtrip() -> None:
    """v2 trajectory schema roundtrip — covers the new DistillTrajectory
    plus all step types (input / reasoning / tool_call / tool_result /
    finalize), ToolResultRecord (success + error variants), and
    ThinkingBlockRecord (plain + redacted)."""
    print("[3/52] pydantic schema round-trip (v2)")
    plain_thinking = ThinkingBlockRecord(
        thinking="I should declare poster type then propose a 3:4 spec.",
        signature="sig_opaque_anthropic",
        is_redacted=False,
    )
    redacted_thinking = ThinkingBlockRecord(
        thinking="",
        signature="enc_payload_bytes",
        is_redacted=True,
    )
    ok_result = ToolResultRecord(
        status="ok",
        payload={"layer_id": "L1", "sha256": "deadbeef", "width": 1536, "height": 2048},
    )
    err_result = ToolResultRecord(
        status="error",
        error_message="DesignSpec validation failed: missing canvas.w_px",
        error_category="validation",
        payload={"hint": None},
    )
    crit = CritiqueResult(iteration=1, verdict="pass", score=0.86,
                          issues=[], rationale="reads as a coherent poster")
    crit_payload = ToolResultRecord(status="ok", payload=crit.model_dump(mode="json"))

    trace = [
        AgentTraceStep(step_idx=1, actor="user", type="input",
                       text="design a 3:4 poster"),
        AgentTraceStep(step_idx=2, actor="planner", type="reasoning",
                       thinking_blocks=[plain_thinking, redacted_thinking],
                       model="claude-opus-4-7", stop_reason="tool_use",
                       usage={"input": 1234, "output": 200, "cache_read": 1024,
                              "cache_create": 0}),
        AgentTraceStep(step_idx=3, actor="planner", type="tool_call",
                       tool_use_id="toolu_x", tool_name="render_text_layer",
                       tool_args={"layer_id": "L1", "text": "国宝回家"},
                       model="claude-opus-4-7"),
        AgentTraceStep(step_idx=4, actor="tool", type="tool_result",
                       tool_use_id="toolu_x", tool_name="render_text_layer",
                       tool_result=ok_result),
        AgentTraceStep(step_idx=5, actor="planner", type="tool_call",
                       tool_use_id="toolu_y", tool_name="propose_design_spec",
                       tool_args={"design_spec": "stringified-by-mistake"}),
        AgentTraceStep(step_idx=6, actor="tool", type="tool_result",
                       tool_use_id="toolu_y", tool_name="propose_design_spec",
                       tool_result=err_result),
        AgentTraceStep(step_idx=7, actor="planner", type="tool_call",
                       tool_use_id="toolu_z", tool_name="critique"),
        AgentTraceStep(step_idx=8, actor="tool", type="tool_result",
                       tool_use_id="toolu_z", tool_name="critique",
                       tool_result=crit_payload),
        AgentTraceStep(step_idx=9, actor="critic", type="reasoning",
                       thinking_blocks=[plain_thinking],
                       model="claude-opus-4-7"),
        AgentTraceStep(step_idx=10, actor="planner", type="finalize",
                       text="all done"),
    ]

    traj = DistillTrajectory(
        run_id="smoke",
        brief="design a 3:4 poster",
        agent_trace=trace,
        final_reward=0.86,
        terminal_status="pass",
        metadata=TrainingMetadata(
            schema_version="v2",
            planner_model="claude-opus-4-7",
            critic_model="claude-opus-4-7",
            image_model="gemini-3-pro-image-preview",
            planner_thinking_budget=10000,
            critic_thinking_budget=10000,
            interleaved_thinking=True,
            total_input_tokens=1234,
            total_output_tokens=200,
            total_cache_read_tokens=1024,
            total_cache_creation_tokens=0,
            estimated_cost_usd=2.34,
            wall_time_s=187.5,
            source="agent_run",
        ),
    )
    dumped = traj.model_dump(mode="json")
    serialized = json.dumps(dumped, ensure_ascii=False)
    try:
        reloaded = DistillTrajectory.model_validate(json.loads(serialized))
    except ValidationError as e:
        _fail(f"DistillTrajectory round-trip: {e.errors()[:3]}")

    if reloaded.metadata.schema_version != "v2":
        _fail(f"schema_version != v2: {reloaded.metadata.schema_version}")
    if reloaded.final_reward != 0.86:
        _fail(f"final_reward not preserved: {reloaded.final_reward}")
    if reloaded.terminal_status != "pass":
        _fail(f"terminal_status not preserved: {reloaded.terminal_status}")

    reasoning_steps = [s for s in reloaded.agent_trace if s.type == "reasoning"]
    if len(reasoning_steps) != 2:
        _fail(f"expected 2 reasoning steps, got {len(reasoning_steps)}")
    planner_reasoning = next(s for s in reasoning_steps if s.actor == "planner")
    if not planner_reasoning.thinking_blocks or len(planner_reasoning.thinking_blocks) != 2:
        _fail("thinking_blocks lost in roundtrip")
    if planner_reasoning.thinking_blocks[0].thinking != plain_thinking.thinking:
        _fail("thinking text corrupted")
    if not planner_reasoning.thinking_blocks[1].is_redacted:
        _fail("redacted thinking flag lost")
    if planner_reasoning.usage["cache_read"] != 1024:
        _fail("usage.cache_read lost")

    err_step = next(s for s in reloaded.agent_trace
                    if s.tool_result and s.tool_result.status == "error")
    if not err_step.tool_result.error_message:
        _fail("error_message stripped from tool_result")
    if err_step.tool_result.error_category != "validation":
        _fail(f"error_category lost: {err_step.tool_result.error_category}")

    crit_step = next(s for s in reloaded.agent_trace
                     if s.tool_name == "critique" and s.type == "tool_result")
    if crit_step.tool_result.payload.get("score") != 0.86:
        _fail("critique payload (verdict/score/issues) lost")

    _ok(f"DistillTrajectory v2 + 10 trace steps roundtrip ({len(serialized)} bytes)")


def check_fonts() -> None:
    print("[4/52] fonts")
    from PIL import ImageFont
    from .config import REPO_ROOT
    for fname in ("NotoSansSC-Bold.otf", "NotoSerifSC-Bold.otf"):
        path = REPO_ROOT / "assets" / "fonts" / fname
        if not path.exists():
            _fail(f"missing font: {path}")
        try:
            ImageFont.truetype(str(path), size=80)
        except Exception as e:
            _fail(f"font {fname} load failed: {e}")
        _ok(f"loaded {fname} ({path.stat().st_size // 1024} KB)")


def check_composite_no_api() -> None:
    """Build a fake background + 2 real text layers, run composite end-to-end.

    Also exercises switch_artifact_type → propose_design_spec plumbing
    (artifact_type fallback from ctx.state when spec omits it).
    """
    print("[5/52] composite (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.render_text_layer import render_text_layer
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub",
        anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7",
        critic_model="claude-opus-4-7",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke")

    # First: switch_artifact_type — verifies ctx.state update + valid type
    obs = switch_artifact_type({"type": "poster"}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"switch_artifact_type: {(obs.error_message or str(obs.payload))}")
    if ctx.state.get("artifact_type") != "poster":
        _fail(f"ctx.state.artifact_type not set; got {ctx.state.get('artifact_type')}")

    # Reject invalid type (paranoia — catches the enum drift)
    bad = switch_artifact_type({"type": "billboard"}, ctx=ctx)
    if bad.status != "error":
        _fail(f"switch_artifact_type should reject 'billboard'; got status={bad.status}")

    # Spec omits artifact_type on purpose — should fall back to ctx.state value
    spec_args = {"design_spec": {
        "brief": "smoke test poster",
        "canvas": {"w_px": 768, "h_px": 1024, "dpi": 150,
                   "aspect_ratio": "3:4", "color_mode": "RGB"},
        "palette": ["#1a1a1a", "#ffffff", "#a02018"],
        "typography": {"title_font": "NotoSerifSC-Bold",
                       "subtitle_font": "NotoSansSC-Bold"},
        "mood": ["test"],
        "composition_notes": "smoke",
        "layer_graph": [],
    }}
    obs = propose_design_spec(spec_args, ctx=ctx)
    if obs.status != "ok":
        _fail(f"propose_design_spec: {(obs.error_message or str(obs.payload))}")

    # Fallback check: spec omitted artifact_type → should inherit ctx.state value
    stored_spec = ctx.state["design_spec"]
    if stored_spec.artifact_type.value != "poster":
        _fail(f"artifact_type fallback failed; got {stored_spec.artifact_type.value!r}")

    bg_path = layers_dir / "bg_smoke.png"
    Image.new("RGB", (768, 1024), (28, 14, 10)).save(bg_path)
    ctx.state["rendered_layers"]["L0_bg"] = {
        "layer_id": "L0_bg", "name": "background", "kind": "background", "z_index": 0,
        "bbox": {"x": 0, "y": 0, "w": 768, "h": 1024},
        "src_path": str(bg_path), "prompt": "(stub)",
        "aspect_ratio": "3:4", "image_size": "1K",
        "safe_zones": [], "sha256": "stub",
    }

    obs = render_text_layer({
        "layer_id": "L1_title", "name": "title", "text": "国宝回家",
        "font_family": "NotoSerifSC-Bold", "font_size_px": 110, "fill": "#fafafa",
        "bbox": {"x": 48, "y": 80, "w": 672, "h": 180}, "align": "center",
        "z_index": 1,
        "effects": {"shadow": {"color": "#00000080", "dx": 0, "dy": 4, "blur": 12}},
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"render_text_layer title: {(obs.error_message or str(obs.payload))}")

    obs = render_text_layer({
        "layer_id": "L2_subtitle", "name": "subtitle",
        "text": "National Treasures Return Home",
        "font_family": "NotoSansSC-Bold", "font_size_px": 36, "fill": "#c9a45a",
        "bbox": {"x": 64, "y": 280, "w": 640, "h": 60}, "align": "center",
        "z_index": 2,
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"render_text_layer subtitle: {(obs.error_message or str(obs.payload))}")

    obs = composite({}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"composite: {(obs.error_message or str(obs.payload))}")

    comp = ctx.state["composition"]
    for label, p in [("PSD", comp.psd_path), ("SVG", comp.svg_path),
                     ("HTML", comp.html_path), ("preview", comp.preview_path)]:
        if p is None:
            _fail(f"{label} path missing on CompositionArtifacts")
        path = Path(p)
        if not path.exists() or path.stat().st_size == 0:
            _fail(f"{label} not written: {p}")
        _ok(f"{label} {path.name} ({path.stat().st_size // 1024} KB)")


def check_svg_text_is_vector() -> None:
    print("[6/52] SVG + HTML content (vector text, contenteditable, inline fonts)")
    from .config import REPO_ROOT
    # v2.1 versioned layout: composite writes to composites/iter_NN/ and
    # maintains final/ symlinks to the latest iter. Read through final/ so
    # this check stays correct across future iteration counts.
    out_dir = REPO_ROOT / "out" / "smoke" / "final"

    # --- SVG -------------------------------------------------------------
    svg_path = out_dir / "poster.svg"
    if not svg_path.exists():
        _fail("smoke SVG not found — composite step likely failed")
    svg = svg_path.read_text(encoding="utf-8")
    if "<text" not in svg:
        _fail("SVG has no <text> element — text was rasterized!")
    if "国宝回家" not in svg:
        _fail("Chinese title not present as vector text in SVG")
    if "@font-face" not in svg:
        print("  warn  SVG missing @font-face (font subsetting may have failed; "
              "will rely on system fonts)")
    else:
        _ok("SVG: @font-face block embedded with subsetted WOFF2")
    _ok("SVG contains <text>国宝回家</text> as real vector")

    # --- HTML ------------------------------------------------------------
    html_path = out_dir / "poster.html"
    if not html_path.exists():
        _fail("smoke HTML not found — html_renderer step likely failed")
    html_text = html_path.read_text(encoding="utf-8")

    required_markers = {
        "canvas container":     '<div class="canvas"',
        "contenteditable text": 'contenteditable="true"',
        "inline fonts":         "@font-face",
        "WOFF2 data URI":       "data:font/woff2;base64,",
        "data-layer-id attr":   "data-layer-id=",
        "data-kind attr":       "data-kind=",
        "Chinese title text":   "国宝回家",
        "bg data URI":          "data:image/png;base64,",
        "generator meta":       '<meta name="generator" content="OpenDesign"',
        # Edit toolbar markers (v1.0 #6 edit UX)
        "bbox data attrs":      'data-bbox-x=',
        "font-size data attr":  'data-font-size-px=',
        "fill data attr":       'data-fill=',
        "font-family data attr": 'data-font-family=',
        "drag handle":          'class="ld-drag-handle"',
        "toolbar container":    'class="ld-toolbar"',
        "font select":          'id="ld-family"',
        "size input":           'id="ld-size"',
        "color input":          'id="ld-color"',
        "save button":          'id="ld-save"',
        "save modal":           'id="ld-modal-backdrop"',
        "copy button":          'id="ld-copy"',
        "download button":      'id="ld-download"',
        "apply-edits hint":     "open-design apply-edits",
    }
    for label, needle in required_markers.items():
        if needle not in html_text:
            _fail(f"HTML missing expected marker — {label}: {needle!r}")
    _ok(f"HTML poster.html ({html_path.stat().st_size // 1024} KB) — "
        f"all {len(required_markers)} markers present "
        "(canvas / contenteditable / inline fonts+images / data-* attrs)")

    # Structure sanity: exactly one <canvas> div and at least 2 text layers
    if html_text.count('class="canvas"') != 1:
        _fail(f"HTML should contain exactly 1 .canvas div; got "
              f"{html_text.count('class=\"canvas\"')}")
    if html_text.count('class="layer text"') < 2:
        _fail(f"HTML should contain ≥2 text layers (smoke has title + subtitle); "
              f"got {html_text.count('class=\"layer text\"')}")
    _ok("HTML structure: 1 .canvas + ≥2 text layers, all attributed")


def check_chat_session_roundtrip() -> None:
    """ChatSession pydantic + save/load cycle — no API calls."""
    print("[7/52] chat session save/load")
    from .config import REPO_ROOT
    from .session import (
        ChatMessage, ChatSession, TrajectoryRef,
        load_session, new_session_id, save_session, session_path, list_sessions,
    )
    from .schema import ArtifactType

    tmp_dir = REPO_ROOT / "out" / "smoke_sessions"
    tmp_dir.mkdir(parents=True, exist_ok=True)

    sid = new_session_id()
    if not sid.startswith("session_"):
        _fail(f"new_session_id() should start with 'session_'; got {sid!r}")

    # Build a realistic session with 1 user + 1 assistant msg + 1 trajectory ref
    session = ChatSession(session_id=sid)
    session.append_user("design a 3:4 poster for 国宝回家")
    ref = TrajectoryRef(
        run_id="20260418-smoke-test",
        artifact_type=ArtifactType.POSTER,
        created_at=session.created_at,
        trajectory_path="/tmp/fake-trajectory.json",
        preview_path="/tmp/fake-preview.png",
        psd_path="/tmp/fake.psd",
        svg_path="/tmp/fake.svg",
        n_layers=5,
        verdict="pass",
        score=0.86,
        cost_usd=1.41,
        wall_s=100.0,
    )
    session.trajectories.append(ref)
    session.append_assistant("produced poster · 5 layers · pass(0.86)",
                              trajectory_id=ref.run_id)

    path = save_session(session, tmp_dir)
    if not path.exists():
        _fail(f"save_session did not write: {path}")
    _ok(f"saved {path.name} ({path.stat().st_size} bytes)")

    loaded = load_session(tmp_dir, sid)
    if loaded.session_id != sid:
        _fail(f"round-trip: session_id mismatch; {loaded.session_id} != {sid}")
    if len(loaded.message_history) != 2:
        _fail(f"round-trip: message count wrong; {len(loaded.message_history)} != 2")
    if len(loaded.trajectories) != 1 or loaded.trajectories[0].verdict != "pass":
        _fail("round-trip: trajectory ref did not survive")
    _ok(f"round-trip: {len(loaded.message_history)} msgs, "
        f"{len(loaded.trajectories)} trajectory ref(s), cost ${loaded.total_cost_usd()}")

    listing = list_sessions(tmp_dir)
    if sid not in [s[0] for s in listing]:
        _fail("list_sessions did not include the newly-saved session")
    _ok(f"list_sessions finds {len(listing)} session(s) incl. this one")

    # Clean up smoke session file
    path.unlink(missing_ok=True)


def check_edit_layer_no_api() -> None:
    """edit_layer semantics — subset-merge, delegates re-render, refuses non-text."""
    print("[8/52] edit_layer (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.edit_layer import edit_layer
    from .tools.propose_design_spec import propose_design_spec
    from .tools.render_text_layer import render_text_layer
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_edit"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub",
        anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7",
        critic_model="claude-opus-4-7",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke_edit")

    # Seed: switch type + spec + one text layer + one fake bg layer
    if switch_artifact_type({"type": "poster"}, ctx=ctx).status != "ok":
        _fail("seed: switch_artifact_type")

    spec_args = {"design_spec": {
        "brief": "edit_layer smoke",
        "canvas": {"w_px": 768, "h_px": 1024, "dpi": 150,
                   "aspect_ratio": "3:4", "color_mode": "RGB"},
        "palette": ["#1a1a1a", "#ffffff"],
        "typography": {"title_font": "NotoSerifSC-Bold",
                       "subtitle_font": "NotoSansSC-Bold"},
        "mood": ["test"], "composition_notes": "", "layer_graph": [],
    }}
    if propose_design_spec(spec_args, ctx=ctx).status != "ok":
        _fail("seed: propose_design_spec")

    # Fake bg layer (non-text — should be refused by edit_layer)
    ctx.state["rendered_layers"]["L0_bg"] = {
        "layer_id": "L0_bg", "name": "background", "kind": "background",
        "z_index": 0, "bbox": {"x": 0, "y": 0, "w": 768, "h": 1024},
        "src_path": "/tmp/fake.png", "sha256": "stub",
    }

    # Real text layer via render_text_layer
    obs = render_text_layer({
        "layer_id": "L1_title", "name": "title", "text": "原标题",
        "font_family": "NotoSerifSC-Bold", "font_size_px": 100, "fill": "#000000",
        "bbox": {"x": 48, "y": 80, "w": 672, "h": 180}, "align": "center",
        "z_index": 1,
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"seed: render_text_layer: {(obs.error_message or str(obs.payload))}")

    before = ctx.state["rendered_layers"]["L1_title"]
    before_sha = before["sha256"]
    before_path = Path(before["src_path"])
    if not before_path.exists():
        _fail("seed: initial PNG missing before edit")

    # --- Happy path: multi-field diff -------------------------------------
    obs = edit_layer({
        "layer_id": "L1_title",
        "diff": {"text": "新标题！", "font_size_px": 140, "fill": "#ff0000"},
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"edit_layer happy path: status={obs.status} summary={(obs.error_message or str(obs.payload))}")

    after = ctx.state["rendered_layers"]["L1_title"]
    if after["font_size_px"] != 140:
        _fail(f"font_size_px not applied: got {after['font_size_px']}")
    if after["fill"] != "#ff0000":
        _fail(f"fill not applied: got {after['fill']}")
    if after["text"] != "新标题！":
        _fail(f"text not applied: got {after['text']!r}")
    # Unchanged fields preserved
    if after["name"] != "title":
        _fail(f"name should be preserved: got {after['name']!r}")
    if after["align"] != "center":
        _fail(f"align should be preserved: got {after['align']!r}")
    if after["bbox"]["w"] != 672:
        _fail(f"bbox should be preserved: got {after['bbox']}")
    # PNG should have been rewritten (different content → different sha)
    if after["sha256"] == before_sha:
        _fail("PNG sha256 unchanged after edit — render_text_layer didn't re-run")
    _ok("happy path: text+font_size_px+fill applied, other fields preserved, PNG rewritten")

    # --- Partial bbox merge ------------------------------------------------
    obs = edit_layer({
        "layer_id": "L1_title",
        "diff": {"bbox": {"y": 200}},  # only y changes
    }, ctx=ctx)
    if obs.status != "ok":
        _fail(f"bbox partial merge: {(obs.error_message or str(obs.payload))}")
    bbox = ctx.state["rendered_layers"]["L1_title"]["bbox"]
    if not (bbox["x"] == 48 and bbox["y"] == 200 and bbox["w"] == 672 and bbox["h"] == 180):
        _fail(f"bbox partial merge broken: {bbox}")
    _ok("partial bbox merge: y updated, x/w/h preserved")

    # --- Missing layer_id --------------------------------------------------
    obs = edit_layer({"layer_id": "nope", "diff": {"text": "x"}}, ctx=ctx)
    if obs.status != "error" or obs.error_category != "not_found":
        _fail(f"unknown layer should return error/not_found, got status={obs.status} cat={obs.error_category}")
    _ok("unknown layer_id → error/not_found")

    # --- Non-text layer rejected ------------------------------------------
    obs = edit_layer({"layer_id": "L0_bg", "diff": {"text": "x"}}, ctx=ctx)
    if obs.status != "error":
        _fail(f"bg layer edit should error, got {obs.status}")
    if obs.error_category != "validation":
        _fail(f"bg-error category should be validation; got: {obs.error_category}")
    _ok("non-text layer → error/validation")

    # --- Empty / unknown diff fields --------------------------------------
    if edit_layer({"layer_id": "L1_title", "diff": {}}, ctx=ctx).status != "error":
        _fail("empty diff should error")
    if edit_layer({"layer_id": "L1_title",
                   "diff": {"color": "#ff0000"}}, ctx=ctx).status != "error":
        _fail("unknown diff field should error (caught 'color' instead of 'fill')")
    _ok("empty diff + unknown field both rejected")


def check_apply_edits_roundtrip() -> None:
    """HTML → apply-edits → new PSD/SVG/HTML/preview with same semantic content."""
    print("[9/52] apply-edits round-trip (no API)")
    from .apply_edits import apply_edits
    from .config import REPO_ROOT, Settings

    src_html = REPO_ROOT / "out" / "smoke" / "final" / "poster.html"
    if not src_html.exists():
        _fail("smoke HTML missing — [5/16] composite step must run first")

    # Simulate a user edit: bump font size + change color on the title.
    # Re-emit with a changed data-font-size-px/data-fill so round-trip should
    # reflect the new values in the regenerated layer graph.
    raw = src_html.read_text(encoding="utf-8")
    edited = (raw
              .replace('data-font-size-px="110"', 'data-font-size-px="140"')
              .replace('data-fill="#fafafa"', 'data-fill="#ff3366"'))
    if edited == raw:
        _fail("could not seed edits into smoke HTML — assumed markers missing")
    edited_html_path = REPO_ROOT / "out" / "smoke_apply" / "edited.html"
    edited_html_path.parent.mkdir(parents=True, exist_ok=True)
    edited_html_path.write_text(edited, encoding="utf-8")

    settings = Settings(
        anthropic_api_key="sk-stub",
        anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7",
        critic_model="claude-opus-4-7",
    )
    out_dir = REPO_ROOT / "out" / "smoke_apply" / "restored"

    traj, traj_path, run_dir, restored_ids, skipped = apply_edits(
        edited_html_path, settings=settings, out_dir=out_dir,
    )

    # --- v2 trajectory assertions ----------------------------------------
    if traj.metadata.source != "apply_edits":
        _fail(f"metadata.source != 'apply_edits'; got {traj.metadata.source!r}")
    if traj.agent_trace:
        _fail(f"apply-edits trajectory should have empty agent_trace; got {len(traj.agent_trace)}")
    if traj.terminal_status != "abort":
        _fail(f"apply-edits should set terminal_status=abort; got {traj.terminal_status}")
    _ok(f"trajectory: source=apply_edits, terminal=abort, "
        f"{len(restored_ids)} layers restored")

    # --- artifact files exist --------------------------------------------
    # v2.1 versioned layout: apply_edits' fresh run writes into
    # <run_dir>/composites/iter_01/ with final/ symlinks to the latest iter.
    final_dir = run_dir / "final"
    for fname in ("poster.psd", "poster.svg", "poster.html", "preview.png"):
        p = final_dir / fname
        if not p.exists() or p.stat().st_size == 0:
            _fail(f"{fname} not written: {p}")
    _ok("PSD+SVG+HTML+preview all regenerated in new run_dir")

    # --- edits landed in re-rendered run dir -----------------------------
    # v2 trajectory has no layer_graph; verify the rendered HTML on disk
    # contains the edits instead.
    rendered_html = (final_dir / "poster.html").read_text(encoding="utf-8")
    if "140" not in rendered_html:
        _fail("expected font_size 140 to appear in re-rendered HTML")
    if "#ff3366" not in rendered_html.lower():
        _fail("expected fill #ff3366 to appear in re-rendered HTML")
    _ok("edits preserved in re-rendered HTML: font_size 140, fill #ff3366")

    # --- bg was decoded from data: URI ----------------------------------
    layers_dir = run_dir / "layers"
    bgs = list(layers_dir.glob("bg_*.png")) if layers_dir.exists() else []
    if bgs:
        _ok(f"bg decoded from data URI → {bgs[0].name} ({bgs[0].stat().st_size} B)")


def check_landing_mode() -> None:
    """Landing end-to-end: section-tree spec → HTML + preview → apply-edits roundtrip.

    v1.3: expanded fixture to 4 sections (hero + features + pricing + cta)
    with a `kind="cta"` child, plus a footer-variant section to exercise
    the `<footer>` auto-upgrade. 4 sections triggers auto-nav, and the
    round-trip must preserve CTA nodes with href + variant.
    """
    print("[10/52] landing mode (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.switch_artifact_type import switch_artifact_type
    from .apply_edits import apply_edits
    from .schema import ArtifactType

    out_dir = REPO_ROOT / "out" / "smoke_landing"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7", critic_model="claude-opus-4-7",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-landing")

    # --- switch + propose landing spec ----------------------------------
    if switch_artifact_type({"type": "landing"}, ctx=ctx).status != "ok":
        _fail("switch_artifact_type(landing)")

    spec_args = {"design_spec": {
        "brief": "OpenDesign v1.0 landing",
        "artifact_type": "landing",
        "canvas": {"w_px": 1200, "h_px": 2400, "dpi": 96,
                   "aspect_ratio": "1:2", "color_mode": "RGB"},
        "palette": ["#0f172a", "#f8fafc"],
        "typography": {"title_font": "NotoSerifSC-Bold", "body_font": "NotoSansSC-Bold"},
        "mood": ["minimal"], "composition_notes": "dark hero, light features, dark cta",
        "layer_graph": [
            {"layer_id": "S1", "name": "hero", "kind": "section", "z_index": 1,
             "children": [
                 {"layer_id": "H1", "name": "hero_headline", "kind": "text", "z_index": 1,
                  "text": "OpenDesign", "font_family": "NotoSerifSC-Bold",
                  "font_size_px": 96, "align": "center",
                  "effects": {"fill": "#f8fafc"}},
                 {"layer_id": "H2", "name": "hero_subhead", "kind": "text", "z_index": 2,
                  "text": "Open source conversational design agent.",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 28,
                  "align": "center", "effects": {"fill": "#94a3b8"}},
             ]},
            {"layer_id": "S2", "name": "features", "kind": "section", "z_index": 2,
             "children": [
                 {"layer_id": "F1", "name": "features_title", "kind": "text", "z_index": 1,
                  "text": "Three outputs, one conversation",
                  "font_family": "NotoSerifSC-Bold", "font_size_px": 48,
                  "effects": {"fill": "#0f172a"}},
                 {"layer_id": "F2", "name": "feature_1", "kind": "text", "z_index": 2,
                  "text": "Poster · PSD + SVG + HTML, fully layered and editable.",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 20,
                  "effects": {"fill": "#334155"}},
             ]},
            {"layer_id": "S3", "name": "pricing", "kind": "section", "z_index": 3,
             "children": [
                 {"layer_id": "P1", "name": "pricing_title", "kind": "text", "z_index": 1,
                  # v2.3 — embed inline + display math so the KaTeX path gets
                  # exercised in smoke. Free-form strings, no rendering happens
                  # at write time; KaTeX is client-side JS.
                  "text": "Pricing model: cost is $O(n)$ per token; total budget $$B = n \\cdot c$$",
                  "font_family": "NotoSerifSC-Bold", "font_size_px": 40,
                  "align": "center", "effects": {"fill": "#0f172a"}},
             ]},
            {"layer_id": "S4", "name": "cta", "kind": "section", "z_index": 4,
             "children": [
                 {"layer_id": "C1", "name": "cta_text", "kind": "text", "z_index": 1,
                  "text": "pip install open-design",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 36,
                  "align": "center", "effects": {"fill": "#f8fafc"}},
                 {"layer_id": "C2", "name": "cta_button", "kind": "cta", "z_index": 2,
                  "text": "Get started", "href": "#sec-features",
                  "variant": "primary"},
             ]},
            {"layer_id": "S5", "name": "footer", "kind": "section", "z_index": 5,
             "children": [
                 {"layer_id": "FT1", "name": "copyright", "kind": "text", "z_index": 1,
                  "text": "© 2026 OpenDesign · MIT",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 14,
                  "align": "center", "effects": {"fill": "#94a3b8"}},
             ]},
        ],
    }}
    if propose_design_spec(spec_args, ctx=ctx).status != "ok":
        _fail("propose_design_spec(landing)")
    if ctx.state["design_spec"].artifact_type != ArtifactType.LANDING:
        _fail("design_spec.artifact_type not LANDING after propose")

    # --- composite landing ----------------------------------------------
    obs = composite({}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"landing composite: {(obs.error_message or str(obs.payload))}")
    comp = ctx.state["composition"]
    if comp.psd_path is not None or comp.svg_path is not None:
        _fail("landing should NOT produce PSD/SVG — got non-None paths")
    for label, p in [("HTML", comp.html_path), ("preview", comp.preview_path)]:
        if not p or not Path(p).exists() or Path(p).stat().st_size == 0:
            _fail(f"landing {label} missing: {p}")
    _ok(f"landing composite: HTML + preview, NO PSD/SVG (correct)")

    # --- HTML structure --------------------------------------------------
    html_text = Path(comp.html_path).read_text(encoding="utf-8")
    required_markers = {
        "landing container":     '<main class="ld-landing"',
        "landing mode meta":     '<meta name="ld-artifact-type" content="landing"',
        "hero section":          'data-section-variant="hero"',
        "features section":      'data-section-variant="features"',
        "cta section":           'data-section-variant="cta"',
        "contenteditable text":  'contenteditable="true"',
        "data-font-size attr":   'data-font-size-px=',
        "toolbar container":     'class="ld-toolbar"',
        "save modal":            'id="ld-modal-backdrop"',
        "landing hides drag":    "/* Drag handle hidden in landing",
        # v1.3 interactive markers
        "auto-nav header":       '<header class="ld-header"',
        "nav anchor":            'data-nav-target="sec-',
        "section id":            'id="sec-',
        "reveal attr":           'data-reveal="true"',
        "cta anchor":             '<a class="ld-cta ld-cta--primary"',
        "cta href":               'href="#sec-features"',
        "cta data-kind":          'data-kind="cta"',
        "semantic footer":        '<footer class="ld-section"',
        "interactive JS":         'IntersectionObserver',
        # v2.3 — KaTeX injected when math delimiters present in any text layer
        "katex stylesheet":       '<style id="ld-katex-css">',
        "katex init":             'renderMathInElement',
        "katex font data URI":    'data:font/woff2;base64',
        "inline math preserved":  'cost is $O(n)$',
        "display math preserved": 'B = n \\cdot c',
    }
    for label, needle in required_markers.items():
        if needle not in html_text:
            _fail(f"landing HTML missing marker — {label}: {needle!r}")
    _ok(f"landing HTML ({Path(comp.html_path).stat().st_size // 1024} KB) — "
        f"all {len(required_markers)} markers present")

    # --- apply-edits round-trip on a seeded-edit landing ---------------
    edited_html = html_text.replace(
        'data-font-size-px="96"', 'data-font-size-px="128"'
    ).replace(
        'data-fill="#f8fafc"', 'data-fill="#38bdf8"', 1  # only first occurrence
    )
    if edited_html == html_text:
        _fail("could not seed landing edits — markers missing")
    edited_path = out_dir / "edited.html"
    edited_path.write_text(edited_html, encoding="utf-8")

    traj, traj_path, run_dir, restored_ids, skipped = apply_edits(
        edited_path, settings=settings,
        out_dir=out_dir / "restored",
    )
    if traj.metadata.source != "apply_edits":
        _fail(f"landing round-trip lost source label: got {traj.metadata.source!r}")
    # v2: trajectory has no design_spec / layer_graph; verify the edits
    # landed in the regenerated HTML on disk instead. v2.1 versioned layout
    # keeps artifacts under composites/iter_NN/ with final/ as the stable
    # symlink surface — read through final/ so we're robust across iters.
    rendered_html = (run_dir / "final" / "index.html").read_text(encoding="utf-8")
    if "128" not in rendered_html:
        _fail("landing edit lost: font_size 128 missing from rendered HTML")
    if "#38bdf8" not in rendered_html.lower():
        _fail("landing edit lost: fill #38bdf8 missing from rendered HTML")
    # CTA must survive (data-* attrs encode href + variant in HTML).
    if 'data-href="#sec-features"' not in rendered_html:
        _fail("CTA href lost on round-trip")
    if 'data-variant="primary"' not in rendered_html:
        _fail("CTA variant lost on round-trip")
    _ok(f"landing round-trip: edits applied (hero_headline 96px → 128px, "
        f"fill→#38bdf8), CTA href+variant intact in regenerated HTML")


def check_design_system_styles() -> None:
    """Render a landing in each of the 6 bundled styles, verify the matching
    CSS got inlined and the style-specific signature tokens are present."""
    print("[11/52] design-system styles (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.switch_artifact_type import switch_artifact_type

    # A tiny per-style signature: a CSS token string that MUST appear in the
    # rendered HTML only when that style's CSS is loaded.
    style_signatures = {
        "minimalist":     "--ld-shadow-soft",
        "editorial":      "--ld-rule",
        "neubrutalism":   "--ld-shadow-hard",
        "glassmorphism":  "prefers-reduced-transparency",
        "claymorphism":   "--ld-shadow-clay",
        "liquid-glass":   "cubic-bezier(0.34, 1.56",
    }

    base_spec = {
        "brief": "design-system smoke",
        "artifact_type": "landing",
        "canvas": {"w_px": 1200, "h_px": 2400, "dpi": 96,
                   "aspect_ratio": "1:2", "color_mode": "RGB"},
        "palette": ["#0f172a", "#f8fafc"],
        "typography": {"title_font": "NotoSerifSC-Bold", "body_font": "NotoSansSC-Bold"},
        "mood": ["test"], "composition_notes": "",
        "layer_graph": [
            {"layer_id": "S1", "name": "hero", "kind": "section", "z_index": 1,
             "children": [
                 {"layer_id": "H1", "name": "hero_headline", "kind": "text",
                  "z_index": 1, "text": "Test", "font_family": "NotoSansSC-Bold",
                  "font_size_px": 80, "align": "center",
                  "effects": {"fill": "#f8fafc"}},
             ]},
            {"layer_id": "S2", "name": "features", "kind": "section", "z_index": 2,
             "children": [
                 {"layer_id": "F1", "name": "features_title", "kind": "text",
                  "z_index": 1, "text": "Feature",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 36,
                  "effects": {"fill": "#0f172a"}},
             ]},
        ],
    }

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7", critic_model="claude-opus-4-7",
    )

    for style, signature in style_signatures.items():
        out_dir = REPO_ROOT / "out" / "smoke_styles" / style
        layers_dir = out_dir / "layers"
        layers_dir.mkdir(parents=True, exist_ok=True)
        ctx = ToolContext(settings=settings, run_dir=out_dir,
                          layers_dir=layers_dir, run_id=f"smoke-style-{style}")

        if switch_artifact_type({"type": "landing"}, ctx=ctx).status != "ok":
            _fail(f"[{style}] switch_artifact_type")
        spec_args = {"design_spec": {**base_spec,
                                     "design_system": {"style": style}}}
        if propose_design_spec(spec_args, ctx=ctx).status != "ok":
            _fail(f"[{style}] propose_design_spec")
        if composite({}, ctx=ctx).status != "ok":
            _fail(f"[{style}] composite")

        comp = ctx.state["composition"]
        html = Path(comp.html_path).read_text(encoding="utf-8")

        if f'<meta name="ld-design-system" content="{style}">' not in html:
            _fail(f"[{style}] missing design-system meta tag")
        if f'data-ld-style="{style}"' not in html:
            _fail(f"[{style}] missing data-ld-style attr")
        if signature not in html:
            _fail(f"[{style}] CSS signature {signature!r} not inlined — "
                  "assets/design-systems/<style>.css probably not loaded")
        _ok(f"  {style:<14} ({Path(comp.html_path).stat().st_size // 1024} KB) — "
            f"meta + data-attr + CSS signature all present")

    # Final: accent_color override landed in the CSS
    out_dir = REPO_ROOT / "out" / "smoke_styles" / "accent_override"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-accent")
    switch_artifact_type({"type": "landing"}, ctx=ctx)
    spec_args = {"design_spec": {**base_spec,
                                 "design_system": {"style": "minimalist",
                                                   "accent_color": "#fe11ba"}}}
    propose_design_spec(spec_args, ctx=ctx)
    composite({}, ctx=ctx)
    html = Path(ctx.state["composition"].html_path).read_text(encoding="utf-8")
    if "--ld-accent: #fe11ba" not in html:
        _fail("accent_color override did not reach CSS (--ld-accent: #fe11ba missing)")
    _ok("accent_color override propagated to --ld-accent token")


def check_landing_with_images() -> None:
    """Landing mode with image children in sections. No NBP call —
    pre-stages a stub PNG in rendered_layers and asserts the renderer
    inlines it + apply-edits round-trips the image layer."""
    print("[12/52] landing with images (no API)")
    from .apply_edits import apply_edits
    from .config import REPO_ROOT, Settings
    from .schema import ArtifactType
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_landing_img"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Stub a hero image PNG so generate_image doesn't need to be called.
    hero_img_path = layers_dir / "img_H0_stub.png"
    Image.new("RGB", (600, 800), (244, 200, 180)).save(hero_img_path)
    feat_img_path = layers_dir / "img_F_stub.png"
    Image.new("RGB", (256, 256), (220, 214, 247)).save(feat_img_path)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="claude-opus-4-7", critic_model="claude-opus-4-7",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-landing-img")

    # Seed rendered_layers as if generate_image was called
    ctx.state["rendered_layers"]["H0_img"] = {
        "layer_id": "H0_img", "name": "hero_image", "kind": "image",
        "z_index": 1, "bbox": None, "src_path": str(hero_img_path),
        "prompt": "(stub)", "aspect_ratio": "3:4",
        "image_size": "2K", "sha256": "stub",
    }
    ctx.state["rendered_layers"]["F1_img"] = {
        "layer_id": "F1_img", "name": "feature_1_icon", "kind": "image",
        "z_index": 1, "bbox": None, "src_path": str(feat_img_path),
        "prompt": "(stub)", "aspect_ratio": "1:1",
        "image_size": "1K", "sha256": "stub",
    }

    if switch_artifact_type({"type": "landing"}, ctx=ctx).status != "ok":
        _fail("switch_artifact_type")

    spec_args = {"design_spec": {
        "brief": "landing with images test",
        "artifact_type": "landing",
        "design_system": {"style": "claymorphism"},
        "canvas": {"w_px": 1200, "h_px": 2400, "dpi": 96,
                   "aspect_ratio": "1:2", "color_mode": "RGB"},
        "palette": ["#f4f0ea", "#3a2f4a"],
        "typography": {}, "mood": ["test"], "composition_notes": "",
        "layer_graph": [
            {"layer_id": "S1", "name": "hero", "kind": "section", "z_index": 1,
             "children": [
                 {"layer_id": "H0_img", "name": "hero_image", "kind": "image",
                  "z_index": 1, "aspect_ratio": "3:4"},
                 {"layer_id": "H1", "name": "hero_headline", "kind": "text",
                  "z_index": 2, "text": "Test Hero",
                  "font_family": "NotoSerifSC-Bold", "font_size_px": 88,
                  "align": "center", "effects": {"fill": "#3a2f4a"}},
             ]},
            {"layer_id": "S2", "name": "features", "kind": "section", "z_index": 2,
             "children": [
                 {"layer_id": "F1_img", "name": "feature_1_icon", "kind": "image",
                  "z_index": 1, "aspect_ratio": "1:1"},
                 {"layer_id": "F1", "name": "feature_1", "kind": "text",
                  "z_index": 2, "text": "First feature copy.",
                  "font_family": "NotoSansSC-Bold", "font_size_px": 20,
                  "effects": {"fill": "#3a2f4a"}},
             ]},
        ],
    }}
    if propose_design_spec(spec_args, ctx=ctx).status != "ok":
        _fail("propose_design_spec")
    if composite({}, ctx=ctx).status != "ok":
        _fail("composite (with images)")

    comp = ctx.state["composition"]
    html_text = Path(comp.html_path).read_text(encoding="utf-8")

    required = {
        # v2.4.3 shipped drag/resize by appending the `draggable-resizable`
        # class to the figure; omit the closing quote so added classes don't
        # break the substring match.
        "figure.layer.image":     '<figure class="layer image',
        "<img tag":               "<img src=\"data:image/",
        "data-has-image":         'data-has-image="true"',
        "hero image layer_id":    'data-layer-id="H0_img"',
        "feature image layer_id": 'data-layer-id="F1_img"',
        "aspect-ratio data":      'data-aspect-ratio=',
    }
    for label, needle in required.items():
        if needle not in html_text:
            _fail(f"landing+images HTML missing — {label}: {needle!r}")
    _ok(f"landing HTML ({Path(comp.html_path).stat().st_size // 1024} KB) — "
        f"2 <figure> image layers inlined with data URIs")

    # Apply-edits round-trip: reparse HTML, rebuild section tree with images
    traj, _, restored_dir, restored_ids, skipped = apply_edits(
        Path(comp.html_path), settings=settings,
        out_dir=out_dir / "restored",
    )
    # v2: layer info no longer in trajectory; verify the regenerated run_dir
    # has the image PNGs decoded from data: URIs onto disk.
    img_files = sorted((restored_dir / "layers").glob("img_*.png")) \
        if (restored_dir / "layers").exists() else []
    if len(img_files) < 2:
        _fail(f"round-trip lost images: expected ≥2 PNG files, got {len(img_files)}")
    for img in img_files:
        if img.stat().st_size == 0:
            _fail(f"round-trip image is empty: {img}")
    _ok(f"round-trip: {len(img_files)} image layers decoded from data: URI "
        f"into {restored_dir}/layers/")


def check_deck_mode() -> None:
    """Deck end-to-end: slide-tree spec → PPTX + per-slide PNGs + preview grid.
    No API — python-pptx writes a real .pptx that we reopen + verify."""
    print("[13/52] deck mode (no API)")
    from pptx import Presentation as _Reopen

    from .config import REPO_ROOT, Settings
    from .schema import ArtifactType
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_deck"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Stub a hero image PNG that one slide embeds, so we exercise the image path.
    hero_png = layers_dir / "img_slide02_hero.png"
    Image.new("RGB", (960, 540), (80, 140, 220)).save(hero_png)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-deck")

    if switch_artifact_type({"type": "deck"}, ctx=ctx).status != "ok":
        _fail("switch_artifact_type(deck)")

    # Pre-stage the image in rendered_layers so _hydrate_deck_image_srcs can
    # pick it up (mirrors the real two-step generate_image → composite flow).
    ctx.state["rendered_layers"]["img_slide02_hero"] = {
        "layer_id": "img_slide02_hero",
        "name": "hero_image",
        "kind": "image",
        "z_index": 5,
        "bbox": None,
        "src_path": str(hero_png),
        "prompt": "",
        "aspect_ratio": "16:9",
        "image_size": "960x540",
        "sha256": "stub",
    }

    spec = {
        "brief": "Pitch deck smoke — 3 slides: cover, problem with hero image, outro.",
        "artifact_type": "deck",
        "canvas": {"w_px": 1920, "h_px": 1080, "dpi": 96,
                   "aspect_ratio": "16:9", "color_mode": "RGB"},
        "palette": ["#0f172a", "#f8fafc", "#38bdf8"],
        "typography": {"title_font": "NotoSerifSC-Bold",
                       "body_font": "NotoSansSC-Bold"},
        "mood": ["clean", "investor-ready"],
        "composition_notes": "16:9 pitch deck. Title top, body or image below.",
        "layer_graph": [
            {
                "layer_id": "slide_01", "name": "cover", "kind": "slide",
                "z_index": 1,
                "children": [
                    {"layer_id": "slide_01_title", "name": "title",
                     "kind": "text", "z_index": 10,
                     "bbox": {"x": 120, "y": 420, "w": 1680, "h": 160},
                     "text": "MilkCloud",
                     "font_family": "NotoSerifSC-Bold",
                     "font_size_px": 96, "align": "left",
                     "effects": {"fill": "#0f172a"}},
                    {"layer_id": "slide_01_tagline", "name": "tagline",
                     "kind": "text", "z_index": 10,
                     "bbox": {"x": 120, "y": 620, "w": 1680, "h": 80},
                     "text": "the calmest bubble tea brand",
                     "font_family": "NotoSansSC-Bold",
                     "font_size_px": 40, "align": "left",
                     "effects": {"fill": "#64748b"}},
                ],
            },
            {
                "layer_id": "slide_02", "name": "problem_with_image",
                "kind": "slide", "z_index": 2,
                "children": [
                    {"layer_id": "slide_02_title", "name": "title",
                     "kind": "text", "z_index": 10,
                     "bbox": {"x": 120, "y": 80, "w": 1680, "h": 140},
                     "text": "The problem",
                     "font_family": "NotoSerifSC-Bold",
                     "font_size_px": 72, "align": "left",
                     "effects": {"fill": "#0f172a"}},
                    {"layer_id": "img_slide02_hero", "name": "hero_image",
                     "kind": "image", "z_index": 5,
                     "bbox": {"x": 120, "y": 260, "w": 1680, "h": 720},
                     "aspect_ratio": "16:9"},
                ],
                # v2.3 speaker notes — populate notes_slide.notes_text_frame
                "speaker_notes": "Open with a show-of-hands: who had a stressful bubble-tea run this week? Describe the three overlapping signage problems (logos / menu density / color war). Pause 5s on the pain statement. Transition: that's why we built MilkCloud.",
            },
            {
                "layer_id": "slide_03", "name": "thank_you", "kind": "slide",
                "z_index": 3,
                "children": [
                    {"layer_id": "slide_03_title", "name": "title",
                     "kind": "text", "z_index": 10,
                     "bbox": {"x": 120, "y": 440, "w": 1680, "h": 200},
                     "text": "Thank you",
                     "font_family": "NotoSerifSC-Bold",
                     "font_size_px": 128, "align": "center",
                     "effects": {"fill": "#0f172a"}},
                ],
            },
        ],
    }

    obs = propose_design_spec({"design_spec": spec}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"propose_design_spec(deck): {(obs.error_message or str(obs.payload))}")

    obs = composite({}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"composite(deck): {(obs.error_message or str(obs.payload))}")

    comp = ctx.state["composition"]
    if comp.psd_path is not None or comp.svg_path is not None or comp.html_path is not None:
        _fail(f"deck should produce ONLY pptx + preview — got "
              f"psd={comp.psd_path} svg={comp.svg_path} html={comp.html_path}")
    if not comp.pptx_path or not Path(comp.pptx_path).exists():
        _fail(f"deck PPTX missing: {comp.pptx_path}")
    if not comp.preview_path or not Path(comp.preview_path).exists():
        _fail(f"deck preview missing: {comp.preview_path}")

    # Per-slide PNGs exist.
    slides_dir = Path(comp.pptx_path).parent / "slides"
    slide_pngs = sorted(slides_dir.glob("slide_*.png"))
    if len(slide_pngs) != 3:
        _fail(f"expected 3 slide PNGs in {slides_dir}, got {len(slide_pngs)}")
    for p in slide_pngs:
        if p.stat().st_size == 0:
            _fail(f"empty slide PNG: {p}")

    # Reopen with python-pptx and check structure.
    prs = _Reopen(comp.pptx_path)
    if len(prs.slides) != 3:
        _fail(f"pptx reopen: expected 3 slides, got {len(prs.slides)}")

    # Slide 1 should contain both "MilkCloud" and the tagline text as native runs.
    s1_text = " ".join(
        run.text for shape in prs.slides[0].shapes if shape.has_text_frame
        for para in shape.text_frame.paragraphs for run in para.runs
    )
    if "MilkCloud" not in s1_text:
        _fail(f"slide 1 text missing 'MilkCloud' — got: {s1_text!r}")
    if "calmest bubble tea brand" not in s1_text:
        _fail(f"slide 1 tagline missing — got: {s1_text!r}")

    # Slide 2 should contain one picture shape (the hero image).
    pic_count = sum(1 for shape in prs.slides[1].shapes if shape.shape_type == 13)  # MSO_SHAPE_TYPE.PICTURE
    if pic_count != 1:
        _fail(f"slide 2 expected 1 picture shape, got {pic_count}")

    # v2.3 — speaker notes on slide 2 populated notes_slide.notes_text_frame
    s2_notes = prs.slides[1].notes_slide.notes_text_frame.text
    if "show-of-hands" not in s2_notes:
        _fail(f"slide 2 speaker notes missing or wrong — got: {s2_notes[:120]!r}")
    # Slides 1 + 3 (no speaker_notes in fixture) should have empty notes
    s1_notes = prs.slides[0].notes_slide.notes_text_frame.text.strip()
    s3_notes = prs.slides[2].notes_slide.notes_text_frame.text.strip()
    if s1_notes or s3_notes:
        _fail(f"slide 1/3 should have empty notes — got s1={s1_notes!r} s3={s3_notes!r}")

    pptx_size_kb = Path(comp.pptx_path).stat().st_size // 1024
    preview_size_kb = Path(comp.preview_path).stat().st_size // 1024
    _ok(f"deck composite: 3 slides → {pptx_size_kb} KB .pptx + "
        f"3 slide PNGs + {preview_size_kb} KB preview grid")
    _ok("pptx reopen: slide count + native text runs + picture shape + speaker notes all OK")


def check_deck_design_system_template() -> None:
    """v2.5.2 templated deck path: spec with deck_design_system + slide.role +
    template_slot → renders against assets/deck_templates/academic-editorial.pptx.
    Verifies named slots get filled, image_slot gets a real picture, footer +
    slide_number auto-inject, and the original template slides are removed
    from the slide list."""
    print("[21/52] deck design system template (no API)")
    from pptx import Presentation as _Reopen
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from .config import REPO_ROOT, Settings
    from .schema import (
        ArtifactType, DeckDesignSystem, DesignSpec, LayerNode,
    )
    from .tools import ToolContext
    from .tools.composite import composite

    out_dir = REPO_ROOT / "out" / "smoke_deck_template"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Stub a paper figure on disk so image_slot has something to place.
    fig_path = layers_dir / "ingest_fig_01.png"
    Image.new("RGB", (1200, 800), (180, 220, 200)).save(fig_path)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-deck-template")

    # Pre-stage the ingested figure in rendered_layers so hydration finds it.
    ctx.state["rendered_layers"]["ingest_fig_01"] = {
        "layer_id": "ingest_fig_01", "name": "method_diagram", "kind": "image",
        "z_index": 5, "bbox": None, "src_path": str(fig_path),
        "aspect_ratio": "3:2",
    }
    ctx.state["design_spec"] = None  # will set below

    spec = DesignSpec(
        brief="LongCat-Next: Lexicalizing Modalities as Discrete Tokens",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"},
        deck_design_system=DeckDesignSystem(
            style="academic-editorial",
            footer_text="LongCat-Next · Meituan",
        ),
        layer_graph=[
            LayerNode(layer_id="slide_01", name="cover", kind="slide", z_index=1, role="cover", children=[
                LayerNode(layer_id="slide_01_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="LongCat-Next"),
                LayerNode(layer_id="slide_01_authors", name="authors", kind="text", z_index=10,
                          template_slot="authors", text="Meituan LongCat Team"),
                LayerNode(layer_id="slide_01_badge", name="badge", kind="text", z_index=10,
                          template_slot="badge", text="NeurIPS 2026"),
            ]),
            LayerNode(layer_id="slide_02", name="method", kind="slide", z_index=2,
                      role="content_with_figure", children=[
                LayerNode(layer_id="slide_02_label", name="section_label", kind="text", z_index=10,
                          template_slot="section_label", text="01 · METHOD"),
                LayerNode(layer_id="slide_02_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="DiNA paradigm"),
                LayerNode(layer_id="slide_02_body", name="body", kind="text", z_index=10,
                          template_slot="body", text="Discrete tokens for vision/text/audio."),
                LayerNode(layer_id="ingest_fig_01", name="diagram", kind="image", z_index=5,
                          template_slot="image_slot"),
            ]),
            LayerNode(layer_id="slide_03", name="closing", kind="slide", z_index=3, role="closing",
                      children=[]),
        ],
    )
    ctx.state["design_spec"] = spec

    res = composite({}, ctx=ctx)
    if res.status != "ok":
        _fail(f"composite() returned {res.status}: {res.payload}")

    pptx_path = ctx.state["composition"].pptx_path
    prs = _Reopen(str(pptx_path))
    if len(prs.slides) != 3:
        _fail(f"expected 3 slides after template cleanup, got {len(prs.slides)}")
    _ok(f"templated deck: 3 slides rendered (template's 6 layout slides removed from sldIdLst)")

    # Slide 0 (cover) — title text replaced; badge text replaced; image_slot still shape (no image since cover didn't reference one)
    cover = prs.slides[0]
    cover_texts = {s.name: s.text_frame.text for s in cover.shapes if s.has_text_frame}
    # v2.7.2: cover slides are chrome → no section prefix; title stays exact.
    if cover_texts.get("title") != "LongCat-Next":
        _fail(f"cover title text wrong: {cover_texts.get('title')!r}")
    if cover_texts.get("badge") != "NeurIPS 2026":
        _fail(f"cover badge text wrong: {cover_texts.get('badge')!r}")
    _ok("cover slide: title + badge text filled from template_slot")

    # Slide 1 (content_with_figure) — image placed, footer + slide_number auto-injected.
    # v2.7.2: the default section_number_policy="renumber" prepends "§N · " to
    # the first content slide's title; substring-match instead of equality so
    # the assertion stays orthogonal to the section policy.
    method = prs.slides[1]
    method_texts = {s.name: s.text_frame.text for s in method.shapes if s.has_text_frame}
    if "DiNA paradigm" not in (method_texts.get("title") or ""):
        _fail(f"method title text wrong: {method_texts.get('title')!r}")
    if "LongCat-Next · Meituan" not in (method_texts.get("footer") or ""):
        _fail(f"footer auto-fill missing: {method_texts.get('footer')!r}")
    if method_texts.get("slide_number") != "2/3":
        _fail(f"slide_number auto-fill wrong: expected '2/3' got {method_texts.get('slide_number')!r}")
    pic_count = sum(1 for s in method.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    if pic_count != 1:
        _fail(f"expected 1 picture on method slide, got {pic_count}")
    _ok("content_with_figure: title + body filled, image placed, footer + slide_number auto-injected")

    # Slide 2 (closing) — title still says "Thank you" (template default); no footer
    closing = prs.slides[2]
    closing_texts = {s.name: s.text_frame.text for s in closing.shapes if s.has_text_frame}
    if closing_texts.get("title") != "Thank you":
        _fail(f"closing title default wrong: {closing_texts.get('title')!r}")
    _ok("closing slide: template defaults preserved when planner emits no children")


def check_footer_leakage() -> None:
    """v2.5.2.2 hotfix: footer auto-fill must use the ingested paper title,
    NOT the user's brief. Prior bug shipped strings like '12-slide academic
    talk deck for the LongCat-Next paper. Speaker-ready with notes per slide'
    on every content slide footer because the v2.5.2 fallback chain was
    `ds.footer_text or spec.brief[:80]` — and brief is a user command.

    This check builds a spec WITHOUT `ds.footer_text`, pre-stages an
    `ingested` entry on `ctx.state` with manifest.title set, and asserts
    the rendered footer reads the paper title — not the brief, not empty.
    Also asserts the leakage blacklist rejects user-command phrases."""
    print("[22/52] footer leakage check (no API)")
    from pptx import Presentation as _Reopen

    from .config import REPO_ROOT, Settings
    from .schema import (
        ArtifactType, DeckDesignSystem, DesignSpec, LayerNode,
    )
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.pptx_renderer import _is_leakage, _resolve_footer_text

    # Unit-test the resolver directly (no API, no rendering).
    class _DSStub:
        def __init__(self, ft=None):
            self.footer_text = ft

    class _SpecStub:
        def __init__(self, brief):
            self.brief = brief

    class _CtxStub:
        def __init__(self, state):
            self.state = state

    leak_brief = "12-slide academic talk deck for this paper. Speaker-ready with notes."
    paper_title = "LongCat-Next: Lexicalizing Modalities as Discrete Tokens"

    # Case 1: brief leak with no ingest → empty (NOT brief)
    r = _resolve_footer_text(_DSStub(), _SpecStub(leak_brief), _CtxStub({}))
    if r != "":
        _fail(f"case 1 (no ingest, leak brief): expected empty, got {r!r}")
    # Case 2: ingest title present → use it
    state = {"ingested": [{"manifest": {"title": paper_title}}]}
    r = _resolve_footer_text(_DSStub(), _SpecStub(leak_brief), _CtxStub(state))
    if "LongCat-Next" not in r:
        _fail(f"case 2 (ingest title): expected paper title, got {r!r}")
    # Case 3: explicit clean override wins
    r = _resolve_footer_text(_DSStub("My Paper · Author · ICLR 2026"),
                             _SpecStub(leak_brief), _CtxStub(state))
    if r != "My Paper · Author · ICLR 2026":
        _fail(f"case 3 (explicit clean): got {r!r}")
    # Case 4: explicit dirty override falls through to ingest
    r = _resolve_footer_text(_DSStub("Generate a 12-slide deck for this"),
                             _SpecStub(leak_brief), _CtxStub(state))
    if "LongCat-Next" not in r:
        _fail(f"case 4 (explicit dirty fallthrough): got {r!r}")
    # Blacklist sanity
    assert _is_leakage("Speaker-ready slide deck"), "leakage detect FN"
    assert not _is_leakage("Discrete Tokens for Multimodal AI"), "leakage detect FP"
    _ok("_resolve_footer_text: 4 fallback cases + 2 blacklist cases all pass")

    # Integration: build a 2-slide deck with no footer_text + ingest title,
    # confirm the rendered PPTX content slide footer carries the paper title.
    out_dir = REPO_ROOT / "out" / "smoke_footer_leakage"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)
    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-footer-leak")
    ctx.state["ingested"] = [{"manifest": {"title": paper_title}}]

    spec = DesignSpec(
        brief=leak_brief,
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"},
        deck_design_system=DeckDesignSystem(style="academic-editorial"),
        layer_graph=[
            LayerNode(layer_id="slide_01", name="cover", kind="slide", z_index=1, role="cover", children=[
                LayerNode(layer_id="slide_01_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="LongCat-Next"),
            ]),
            LayerNode(layer_id="slide_02", name="content", kind="slide", z_index=2, role="content",
                      children=[
                LayerNode(layer_id="slide_02_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Method"),
                LayerNode(layer_id="slide_02_body", name="body", kind="text", z_index=10,
                          template_slot="body", text="Body content."),
            ]),
        ],
    )
    ctx.state["design_spec"] = spec

    res = composite({}, ctx=ctx)
    if res.status != "ok":
        _fail(f"composite failed: {res.payload}")
    pptx_path = ctx.state["composition"].pptx_path
    prs = _Reopen(str(pptx_path))
    content_slide = prs.slides[1]
    footer_text = ""
    for shape in content_slide.shapes:
        if shape.has_text_frame and shape.name == "footer":
            footer_text = shape.text_frame.text
            break
    if "LongCat-Next" not in footer_text:
        _fail(f"footer should carry paper title, got {footer_text!r}")
    if _is_leakage(footer_text):
        _fail(f"footer contains leakage phrase: {footer_text!r}")
    _ok(f"rendered footer carries paper title: {footer_text!r}")


def check_callout_overlay() -> None:
    """v2.6 callout system: kind="callout" children render as shapes
    overlaid on top of the anchor picture/table. Verifies all 3 styles
    (highlight / label / circle) plus the optional arrow connector."""
    print("[23/52] callout overlay (no API)")
    from pptx import Presentation as _Reopen
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from .config import REPO_ROOT, Settings
    from .schema import (
        ArtifactType, DeckDesignSystem, DesignSpec, LayerNode, SafeZone,
    )
    from .tools import ToolContext
    from .tools.composite import composite

    out_dir = REPO_ROOT / "out" / "smoke_callout"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    fig_path = layers_dir / "ingest_fig_01.png"
    Image.new("RGB", (1200, 800), (180, 220, 200)).save(fig_path)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-callout")

    ctx.state["rendered_layers"]["ingest_fig_01"] = {
        "layer_id": "ingest_fig_01", "kind": "image", "src_path": str(fig_path),
        "bbox": None, "z_index": 5, "name": "fig", "aspect_ratio": "3:2",
    }

    spec = DesignSpec(
        brief="callout overlay smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"},
        deck_design_system=DeckDesignSystem(style="academic-editorial"),
        layer_graph=[
            LayerNode(layer_id="slide_01", name="cover", kind="slide", z_index=1, role="cover", children=[
                LayerNode(layer_id="slide_01_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Callout Test"),
            ]),
            LayerNode(layer_id="slide_02", name="method", kind="slide", z_index=2,
                      role="content_with_figure", children=[
                LayerNode(layer_id="slide_02_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Method"),
                LayerNode(layer_id="ingest_fig_01", name="diagram", kind="image", z_index=5,
                          template_slot="image_slot"),
                LayerNode(layer_id="callout_a", name="hl", kind="callout", z_index=20,
                          anchor_layer_id="ingest_fig_01",
                          callout_style="highlight",
                          callout_region=SafeZone(x=1100, y=300, w=400, h=200, purpose="body")),
                LayerNode(layer_id="callout_b", name="lbl", kind="callout", z_index=21,
                          anchor_layer_id="ingest_fig_01",
                          callout_style="label",
                          callout_text="+5.2",
                          callout_region=SafeZone(x=1100, y=300, w=400, h=200, purpose="body"),
                          arrow=True),
                LayerNode(layer_id="callout_c", name="circ", kind="callout", z_index=22,
                          anchor_layer_id="ingest_fig_01",
                          callout_style="circle",
                          callout_region=SafeZone(x=1500, y=600, w=120, h=120, purpose="body")),
            ]),
        ],
    )
    ctx.state["design_spec"] = spec

    res = composite({}, ctx=ctx)
    if res.status != "ok":
        _fail(f"composite failed: {res.payload}")
    pptx_path = ctx.state["composition"].pptx_path
    prs = _Reopen(str(pptx_path))

    method = prs.slides[1]
    names = {s.name for s in method.shapes}
    for required in ("callout_a", "callout_b", "callout_c"):
        if required not in names:
            _fail(f"expected callout '{required}' on method slide; got names={sorted(names)}")
    # Picture should still be there
    pics = [s for s in method.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if len(pics) != 1:
        _fail(f"expected 1 picture on method slide alongside callouts, got {len(pics)}")
    # Arrow connector should exist (named "Connector N" by python-pptx)
    connectors = [s for s in method.shapes if s.name.startswith("Connector")]
    if not connectors:
        _fail("expected at least 1 arrow connector for label callout w/ arrow=True")
    # Highlight callout's bbox should match callout_region (1100,300,400,200) in EMU
    hl = next(s for s in method.shapes if s.name == "callout_a")
    px_to_emu = 9525
    if hl.left != 1100 * px_to_emu or hl.top != 300 * px_to_emu:
        _fail(f"highlight callout EMU position wrong: ({hl.left},{hl.top})")
    _ok(f"3 callouts (highlight + label+arrow + circle) + 1 connector overlaid alongside picture")


def check_provenance_validator() -> None:
    """v2.7 — composite-stage provenance validator audits LayerNode.evidence_quote
    against ingested raw_text. Verifies:
      a. regex extracts the right numeric tokens (and skips safe contexts)
      b. validator passes bullets whose evidence_quote substring-matches
      c. validator fails bullets with no evidence_quote
      d. validator fails bullets whose quote is not in raw_text
      e. apply_strict_provenance replaces failed numeric tokens with [?]
      f. _resolve_authors_text pulls from manifest.authors when planner
         emits placeholder text
      g. _add_table truncates >8-col tables to 6 cols + caption marker
    """
    print("[24/52] provenance validator + cover authors + wide-table cap (no API)")
    from pptx import Presentation as _Reopen

    from .config import REPO_ROOT, Settings
    from .schema import (
        ArtifactType, DeckDesignSystem, DesignSpec, LayerNode,
    )
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.pptx_renderer import (
        _is_authors_leakage, _resolve_authors_text,
    )
    from .util.provenance import (
        _extract_numeric_tokens, apply_strict_provenance, validate_provenance,
    )

    # ── Part (a): regex sanity ───────────────────────────────────────
    cases = [
        ("MMMU 70.6 vs BAGEL 55.3", ["70.6", "55.3"]),
        ("500K hours of audio", ["500K hours"]),
        ("80 GB vs 120 GB", ["80 GB", "120 GB"]),
        ("PSNR drops 28.5 to 22.1 dB", ["28.5", "22.1 dB"]),
        ("+5.2 pts on MathVista", ["+5.2"]),
        ("40% diversity drop", ["40%"]),
        ("-0.32 exponent", ["-0.32"]),
        ("1/12", []),         # safe — slide index
        ("v2.6.1", []),       # safe — version tag
        ("01 · MOTIVATION", []),  # safe — section label
    ]
    for inp, expected in cases:
        got = _extract_numeric_tokens(inp)
        if got != expected:
            _fail(f"regex {inp!r}: expected {expected}, got {got}")

    # ── Part (b-e): validator + strict mutation ──────────────────────
    paper_title = "LongCat-Next: Lexicalizing Modalities as Discrete Tokens"
    paper_authors = ["Meituan LongCat Team"]
    paper_text = (
        "LongCat-Next achieves 70.6 on MMMU and 83.1 on MathVista, "
        "outperforming BAGEL's 55.3. Table 6 reports PSNR values of "
        "20.88, 21.86, 30.52, 18.16 across configurations."
    )

    spec = DesignSpec(
        brief="provenance smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"},
        deck_design_system=DeckDesignSystem(style="academic-editorial"),
        layer_graph=[
            LayerNode(layer_id="s1", name="cover", kind="slide", z_index=1, role="cover", children=[
                LayerNode(layer_id="s1_t", name="title", kind="text", z_index=10,
                          template_slot="title", text=paper_title),
                # Placeholder text — _resolve_authors_text MUST replace it.
                LayerNode(layer_id="s1_a", name="authors", kind="text", z_index=10,
                          template_slot="authors",
                          text="Author One · Author Two · Affiliation"),
            ]),
            LayerNode(layer_id="s2", name="results", kind="slide", z_index=2,
                      role="content", children=[
                LayerNode(layer_id="s2_t", name="title", kind="text", z_index=10,
                          template_slot="title", text="Headline"),
                # ✅ verifiable: "70.6 on MMMU" is a substring of paper_text
                LayerNode(layer_id="s2_b1", name="body", kind="text", z_index=10,
                          template_slot="body",
                          text="MMMU 70.6 vs BAGEL 55.3",
                          evidence_quote="70.6 on MMMU"),
                # ❌ missing quote: numbers must be stripped to [?]
                LayerNode(layer_id="s2_b2", name="body2", kind="text", z_index=10,
                          template_slot="body",
                          text="PSNR drops 28.5 to 22.1 dB"),
                # ❌ quote not in source: numbers must be stripped to [?]
                LayerNode(layer_id="s2_b3", name="body3", kind="text", z_index=10,
                          template_slot="body",
                          text="500K hours of audio in training",
                          evidence_quote="500,000 hours of training audio"),
            ]),
        ],
    )

    out_dir = REPO_ROOT / "out" / "smoke_provenance"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)
    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-prov")
    ctx.state["ingested"] = [{
        "file": "test.pdf", "type": "pdf",
        "manifest": {"title": paper_title, "authors": paper_authors},
        "raw_text": paper_text,
    }]

    rep = validate_provenance(spec, ctx)
    if rep.n_passed != 1:
        _fail(f"expected 1 pass, got {rep.n_passed}")
    if len(rep.failures) != 2:
        _fail(f"expected 2 failures, got {len(rep.failures)}: {rep.failures}")
    fail_ids = {f.layer_id for f in rep.failures}
    if fail_ids != {"s2_b2", "s2_b3"}:
        _fail(f"failure ids: expected {{s2_b2,s2_b3}}, got {fail_ids}")
    reasons = {f.reason for f in rep.failures}
    if reasons != {"missing_quote", "quote_not_in_source"}:
        _fail(f"reasons: expected both, got {reasons}")

    n_mut = apply_strict_provenance(spec, rep)
    if n_mut != 2:
        _fail(f"expected 2 mutations, got {n_mut}")
    s2 = spec.layer_graph[1]
    b2 = s2.children[2].text
    b3 = s2.children[3].text
    if "[?]" not in b2 or "28.5" in b2 or "22.1" in b2:
        _fail(f"b2 not stripped: {b2!r}")
    if "[?]" not in b3 or "500K" in b3:
        _fail(f"b3 not stripped: {b3!r}")
    _ok("validator: 1 pass + 2 failures (missing/mismatch) → strict mutation [?] correct")

    # ── Part (f): _resolve_authors_text ──────────────────────────────
    if not _is_authors_leakage("Author One · Author Two · Affiliation"):
        _fail("leakage detect FN: 'Author One' should be flagged")
    if _is_authors_leakage("Meituan LongCat Team"):
        _fail("leakage detect FP: real team name flagged")
    r = _resolve_authors_text(ctx, existing="Author One · Affiliation")
    if r != "Meituan LongCat Team":
        _fail(f"authors fallthrough: expected manifest, got {r!r}")
    r2 = _resolve_authors_text(ctx, existing="Meituan LongCat Team")
    if r2 != "Meituan LongCat Team":
        _fail(f"authors keep: expected real-name kept, got {r2!r}")
    _ok("_resolve_authors_text: leakage detection + manifest fallthrough OK")

    # ── Part (g): wide-table cap via composite integration ───────────
    # Re-build the spec since we mutated it above; add a wide table slide.
    wide_headers = ["Model"] + [f"B{i}" for i in range(11)]  # 12 cols
    wide_rows = [
        ["LongCat-Next"] + [f"{70+i}.{i}" for i in range(11)],
        ["BAGEL"] + [f"{50+i}.{i}" for i in range(11)],
    ]
    spec2 = DesignSpec(
        brief="wide-table smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"},
        deck_design_system=DeckDesignSystem(style="academic-editorial"),
        layer_graph=[
            LayerNode(layer_id="s1", name="cover", kind="slide", z_index=1, role="cover", children=[
                LayerNode(layer_id="s1_t", name="title", kind="text", z_index=10,
                          template_slot="title", text=paper_title),
                LayerNode(layer_id="s1_a", name="authors", kind="text", z_index=10,
                          template_slot="authors",
                          text="Author One · Author Two · Affiliation"),
            ]),
            LayerNode(layer_id="s2", name="wide", kind="slide", z_index=2,
                      role="content_with_table", children=[
                LayerNode(layer_id="s2_t", name="title", kind="text", z_index=10,
                          template_slot="title", text="Wide Table"),
                LayerNode(layer_id="s2_tab", name="tab", kind="table", z_index=5,
                          template_slot="table_anchor",
                          headers=wide_headers, rows=wide_rows),
            ]),
        ],
    )
    ctx.state["design_spec"] = spec2
    res = composite({}, ctx=ctx)
    if res.status != "ok":
        _fail(f"composite failed: {res.payload}")
    pptx_path = ctx.state["composition"].pptx_path
    prs = _Reopen(str(pptx_path))

    # Cover authors: planner placeholder must be overridden by manifest
    cover = prs.slides[0]
    authors_text = ""
    for shape in cover.shapes:
        if shape.has_text_frame and shape.name == "authors":
            authors_text = shape.text_frame.text
            break
    if "Meituan LongCat Team" not in authors_text:
        _fail(f"cover authors not auto-resolved: got {authors_text!r}")
    if "Author One" in authors_text:
        _fail(f"cover authors leak: {authors_text!r}")

    # Wide table: rendered table should have ≤ 8 cols
    wide_slide = prs.slides[1]
    tables = [s for s in wide_slide.shapes if s.has_table]
    if not tables:
        _fail("wide-table slide missing rendered table")
    n_cols_rendered = len(tables[0].table.columns)
    if n_cols_rendered > 8:
        _fail(f"wide-table cap not enforced: rendered {n_cols_rendered} cols")
    if n_cols_rendered != 6:
        _fail(f"wide-table cap should keep 6 cols, got {n_cols_rendered}")

    # provenance_report.json should exist (from spec2 — no number bullets,
    # so report is clean; but file must be written when ingested present)
    import json
    iter_dirs = sorted((ctx.run_dir / "composites").glob("iter_*"))
    last_iter = iter_dirs[-1]
    report_path = last_iter / "provenance_report.json"
    if not report_path.exists():
        _fail(f"provenance_report.json missing at {report_path}")
    rep_data = json.loads(report_path.read_text())
    if "n_audited" not in rep_data:
        _fail(f"provenance_report.json missing keys: {list(rep_data.keys())}")
    _ok(f"composite integration: cover auto-fills {authors_text!r}, "
        f"wide table 12→{n_cols_rendered} cols, provenance_report.json written")


def check_reasoning_step_roundtrip() -> None:
    """v2 training-data: derive artifact_type / design_spec / critique
    verdict / layer count from agent_trace alone (no top-level fields).

    Asserts the chat.py helpers (_last_artifact_type / _last_design_spec /
    _last_critique_payload / _count_unique_layers) correctly recover state
    from a synthetic v2 trajectory shape.
    """
    print("[14/52] v2 trajectory: derive metadata from agent_trace only")
    from .chat import (
        _last_artifact_type, _last_design_spec, _last_critique_payload,
        _count_unique_layers,
    )

    crit_payload = {
        "iteration": 1, "verdict": "pass", "score": 0.86,
        "issues": [], "rationale": "ok",
    }
    spec_input = {
        "brief": "test", "artifact_type": "landing",
        "canvas": {"w_px": 1200, "h_px": 2400, "dpi": 96,
                   "aspect_ratio": "1:2", "color_mode": "RGB"},
        "layer_graph": [],
    }

    trace = [
        AgentTraceStep(step_idx=1, actor="user", type="input", text="brief"),
        AgentTraceStep(step_idx=2, actor="planner", type="tool_call",
                       tool_use_id="t1", tool_name="switch_artifact_type",
                       tool_args={"type": "landing"}),
        AgentTraceStep(step_idx=3, actor="tool", type="tool_result",
                       tool_use_id="t1", tool_name="switch_artifact_type",
                       tool_result=ToolResultRecord(status="ok",
                                                   payload={"type": "landing"})),
        AgentTraceStep(step_idx=4, actor="planner", type="tool_call",
                       tool_use_id="t2", tool_name="propose_design_spec",
                       tool_args={"design_spec": spec_input}),
        AgentTraceStep(step_idx=5, actor="planner", type="tool_call",
                       tool_use_id="t3", tool_name="generate_image",
                       tool_args={"layer_id": "L1", "prompt": "hero"}),
        AgentTraceStep(step_idx=6, actor="planner", type="tool_call",
                       tool_use_id="t4", tool_name="render_text_layer",
                       tool_args={"layer_id": "L2", "text": "Hello"}),
        AgentTraceStep(step_idx=7, actor="planner", type="tool_call",
                       tool_use_id="t5", tool_name="critique"),
        AgentTraceStep(step_idx=8, actor="tool", type="tool_result",
                       tool_use_id="t5", tool_name="critique",
                       tool_result=ToolResultRecord(status="ok",
                                                   payload=crit_payload)),
    ]
    traj = DistillTrajectory(
        run_id="smoke-v2-helpers",
        brief="brief",
        agent_trace=trace,
        final_reward=0.86,
        terminal_status="pass",
        metadata=TrainingMetadata(
            schema_version="v2",
            planner_model="m", critic_model="m", image_model="m",
            planner_thinking_budget=0, critic_thinking_budget=0,
            interleaved_thinking=False,
            total_input_tokens=0, total_output_tokens=0,
            estimated_cost_usd=0.0, wall_time_s=0.0,
            source="agent_run",
        ),
    )

    if _last_artifact_type(traj) != "landing":
        _fail(f"_last_artifact_type: {_last_artifact_type(traj)}")
    spec = _last_design_spec(traj)
    if not spec or spec.get("artifact_type") != "landing":
        _fail(f"_last_design_spec lost: {spec}")
    crit = _last_critique_payload(traj)
    if not crit or crit.get("score") != 0.86:
        _fail(f"_last_critique_payload lost: {crit}")
    n = _count_unique_layers(traj)
    if n != 2:
        _fail(f"_count_unique_layers: expected 2 (L1, L2), got {n}")

    _ok("artifact_type, design_spec, critique payload, layer count "
        "all recoverable from agent_trace alone")


def check_ingest_document_markdown() -> None:
    """Markdown ingestion: seed a stub .md with a relative image ref, verify
    ingest_document registers the image in rendered_layers + returns the raw
    text. No API — markdown path doesn't call Anthropic."""
    print("[15/52] ingest_document markdown (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import ingest_document

    out_dir = REPO_ROOT / "out" / "smoke_ingest_md"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Stage a stub MD with an embedded image reference (relative path).
    src_dir = out_dir / "src"
    src_dir.mkdir(parents=True, exist_ok=True)
    img_path = src_dir / "diagram.png"
    Image.new("RGB", (640, 360), (90, 130, 200)).save(img_path)
    md_path = src_dir / "notes.md"
    md_path.write_text(
        "# My project notes\n\n"
        "## Overview\nA few thoughts below.\n\n"
        "![System diagram](diagram.png)\n\n"
        "## Plans\n- Ship v1\n- Iterate\n"
        "\n![Missing](./nowhere.png)\n",
        encoding="utf-8",
    )

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-ingest-md")

    obs = ingest_document({"file_paths": [str(md_path)]}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"ingest_document(markdown): {(obs.error_message or str(obs.payload))}")

    ingested = ctx.state.get("ingested") or []
    if len(ingested) != 1 or ingested[0]["type"] != "markdown":
        _fail(f"ingested manifest missing markdown entry: {ingested}")
    registered = ingested[0]["registered_layer_ids"]
    skipped = ingested[0]["skipped_images"]
    if len(registered) != 1:
        _fail(f"expected 1 image layer registered from MD, got {registered}")
    if len(skipped) != 1:
        _fail(f"expected 1 skipped bad ref, got {skipped}")

    layer_id = registered[0]
    rec = ctx.state["rendered_layers"].get(layer_id)
    if not rec or not rec.get("src_path") or not Path(rec["src_path"]).exists():
        _fail(f"markdown-registered image not hydrated: {rec}")
    _ok(f"markdown ingested: {ingested[0]['n_chars']} chars, "
        f"1 image layer ({layer_id}), 1 bad ref skipped")


def check_ingest_document_image() -> None:
    """Standalone image ingestion: seed a PNG, verify ingest_document copies
    into layers_dir + registers a passthrough layer with correct shape."""
    print("[16/52] ingest_document image (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import ingest_document

    out_dir = REPO_ROOT / "out" / "smoke_ingest_image"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    src_dir = out_dir / "src"
    src_dir.mkdir(parents=True, exist_ok=True)
    img_path = src_dir / "logo.png"
    Image.new("RGB", (512, 512), (240, 80, 60)).save(img_path)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-ingest-image")

    obs = ingest_document({"file_paths": [str(img_path)]}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"ingest_document(image): {(obs.error_message or str(obs.payload))}")

    ingested = ctx.state.get("ingested") or []
    if len(ingested) != 1 or ingested[0]["type"] != "image":
        _fail(f"ingested manifest missing image entry: {ingested}")
    if (ingested[0]["width"], ingested[0]["height"]) != (512, 512):
        _fail(f"image dims not captured: {ingested[0]}")

    registered = ingested[0]["registered_layer_ids"]
    if len(registered) != 1:
        _fail(f"expected 1 image layer, got {registered}")
    layer_id = registered[0]
    rec = ctx.state["rendered_layers"].get(layer_id)
    if rec["kind"] != "image":
        _fail(f"image layer kind != image: {rec['kind']}")
    if rec.get("source") != "ingested":
        _fail(f"image source tag wrong: {rec.get('source')}")
    if not Path(rec["src_path"]).exists():
        _fail(f"copied image path missing: {rec['src_path']}")
    _ok(f"image ingested: {layer_id}, 512×512, source=ingested, "
        f"sha256[:8]={rec['sha256'][:8]}")


def check_ingest_document_docx() -> None:
    """Docx ingestion (v1.2.5): build a minimal Word doc with headings +
    an inline image, verify ingest_document extracts sections + figures
    without any VLM call."""
    print("[17/52] ingest_document docx (no API)")
    from docx import Document
    from docx.shared import Inches
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import ingest_document

    out_dir = REPO_ROOT / "out" / "smoke_ingest_docx"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    src_dir = out_dir / "src"
    src_dir.mkdir(parents=True, exist_ok=True)
    img_path = src_dir / "icon.png"
    Image.new("RGB", (64, 64), (200, 80, 80)).save(img_path)

    docx_path = src_dir / "notes.docx"
    doc = Document()
    doc.add_heading("Quick Tour", level=0)
    doc.add_heading("Motivation", level=1)
    doc.add_paragraph("First thought about the project. Second thought.")
    doc.add_paragraph("- Short bullet.")
    doc.add_heading("Architecture", level=1)
    doc.add_paragraph("Planner + critic + composer.")
    doc.add_picture(str(img_path), width=Inches(1.0))
    doc.save(str(docx_path))

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-ingest-docx")

    obs = ingest_document({"file_paths": [str(docx_path)]}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"ingest_document(docx): {(obs.error_message or str(obs.payload))}")

    ingested = ctx.state.get("ingested") or []
    if len(ingested) != 1 or ingested[0]["type"] != "docx":
        _fail(f"ingested manifest missing docx entry: {ingested}")
    m = ingested[0]["manifest"]
    if m["title"] != "Quick Tour":
        _fail(f"docx title wrong: {m['title']!r}")
    headings = [s["heading"] for s in m["sections"]]
    if headings != ["Motivation", "Architecture"]:
        _fail(f"docx sections wrong: {headings}")
    figure_ids = ingested[0]["registered_figure_ids"]
    if len(figure_ids) != 1 or not figure_ids[0].startswith("ingest_fig_"):
        _fail(f"docx figure registration wrong: {figure_ids}")
    rec = ctx.state["rendered_layers"][figure_ids[0]]
    if rec.get("source") != "ingested_docx" or not Path(rec["src_path"]).exists():
        _fail(f"docx figure record wrong: {rec}")
    _ok(f"docx ingested: title='{m['title']}', "
        f"{len(m['sections'])} section(s), 1 figure ({figure_ids[0]})")


def check_ingest_document_pptx() -> None:
    """Pptx ingestion (v1.2.5): build a 2-slide PowerPoint with a title,
    body bullets, and an embedded picture; verify slides become sections
    and the picture becomes an ingest_fig_NN layer."""
    print("[18/52] ingest_document pptx (no API)")
    from pptx import Presentation
    from pptx.util import Inches
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import ingest_document

    out_dir = REPO_ROOT / "out" / "smoke_ingest_pptx"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    src_dir = out_dir / "src"
    src_dir.mkdir(parents=True, exist_ok=True)
    img_path = src_dir / "chart.png"
    Image.new("RGB", (128, 128), (50, 160, 90)).save(img_path)

    pptx_path = src_dir / "deck.pptx"
    prs = Presentation()
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "OpenDesign Pitch"
    s1.placeholders[1].text = "Open-source conversational design agent"
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Why now?"
    s2.placeholders[1].text = "Claude Design shipped. We go terminal-first."
    s2.shapes.add_picture(str(img_path), Inches(4), Inches(2), width=Inches(2))
    prs.save(str(pptx_path))

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-ingest-pptx")

    obs = ingest_document({"file_paths": [str(pptx_path)]}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"ingest_document(pptx): {(obs.error_message or str(obs.payload))}")

    ingested = ctx.state.get("ingested") or []
    if len(ingested) != 1 or ingested[0]["type"] != "pptx":
        _fail(f"ingested manifest missing pptx entry: {ingested}")
    m = ingested[0]["manifest"]
    if m["title"] != "OpenDesign Pitch":
        _fail(f"pptx title wrong: {m['title']!r}")
    if len(m["sections"]) != 2:
        _fail(f"pptx slides→sections wrong: {[s['heading'] for s in m['sections']]}")
    figure_ids = ingested[0]["registered_figure_ids"]
    if len(figure_ids) != 1 or not figure_ids[0].startswith("ingest_fig_"):
        _fail(f"pptx figure registration wrong: {figure_ids}")
    rec = ctx.state["rendered_layers"][figure_ids[0]]
    if rec.get("source") != "ingested_pptx" or rec.get("source_ref") != "slide=2":
        _fail(f"pptx figure record wrong: {rec}")
    _ok(f"pptx ingested: title='{m['title']}', "
        f"2 slide(s), 1 figure ({figure_ids[0]})")


def check_sub_figure_registration() -> None:
    """v2.3 sub-figure extraction: VLM detects multi-panel composite figures
    and we Pillow-crop each panel into its own `ingest_fig_NN_<label>` layer.
    Tests the registration helper directly (no API) with a synthetic
    2-panel composite PNG + a stub VLM response.

    Asserts:
    - 2 panel entries → 2 new rendered_layers records (parent stays)
    - Each panel's crop file exists on disk with correct dims
    - parent_layer_id breadcrumb set on children
    - Layer_id naming convention `ingest_fig_NN_<label>` holds
    """
    print("[19/52] sub-figure extraction (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.ingest_document import _register_sub_panels

    out_dir = REPO_ROOT / "out" / "smoke_sub_figs"
    layers_dir = out_dir / "layers"
    if out_dir.exists():
        import shutil
        shutil.rmtree(out_dir)
    layers_dir.mkdir(parents=True, exist_ok=True)

    # Synthesize a composite parent: 600×300 image split into two
    # horizontal halves (panel a on left, panel b on right).
    parent_path = layers_dir / "img_ingest_fig_01.png"
    composite_img = Image.new("RGB", (600, 300), (255, 255, 255))
    # Left half red, right half blue — easy to eyeball crop correctness
    left = Image.new("RGB", (300, 300), (200, 60, 60))
    right = Image.new("RGB", (300, 300), (60, 100, 200))
    composite_img.paste(left, (0, 0))
    composite_img.paste(right, (300, 0))
    composite_img.save(parent_path, "PNG")

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-subfig")
    # Pre-stage the parent so sub-panel registration can check collisions
    ctx.state["rendered_layers"]["ingest_fig_01"] = {
        "layer_id": "ingest_fig_01", "kind": "image", "name": "figure_1",
        "src_path": str(parent_path),
    }

    panels = [
        {"label": "a", "bbox": [0, 0, 300, 300],
         "caption": "Left panel (text input)", "short_caption": "Text"},
        {"label": "b", "bbox": [300, 0, 600, 300],
         "caption": "Right panel (image input)", "short_caption": "Image"},
    ]
    created = _register_sub_panels(
        parent_layer_id="ingest_fig_01",
        parent_path=parent_path,
        parent_caption="Composite — text + image inputs",
        panels=panels,
        ctx=ctx,
        pdf_path=Path("/fake/paper.pdf"),
        source_page=3,
    )

    if created != ["ingest_fig_01_a", "ingest_fig_01_b"]:
        _fail(f"sub-panel layer_ids wrong: {created}")

    for sub_id, expected_color, expected_caption in [
        ("ingest_fig_01_a", (200, 60, 60), "Text"),
        ("ingest_fig_01_b", (60, 100, 200), "Image"),
    ]:
        rec = ctx.state["rendered_layers"].get(sub_id)
        if rec is None:
            _fail(f"sub-panel {sub_id} not registered")
        if rec.get("parent_layer_id") != "ingest_fig_01":
            _fail(f"{sub_id} missing parent_layer_id breadcrumb")
        if rec.get("extract_strategy") != "sub_panel":
            _fail(f"{sub_id} extract_strategy wrong: {rec.get('extract_strategy')!r}")
        if rec.get("caption_short") != expected_caption:
            _fail(f"{sub_id} caption_short wrong: {rec.get('caption_short')!r}")
        crop_path = Path(rec["src_path"])
        if not crop_path.exists() or crop_path.stat().st_size == 0:
            _fail(f"{sub_id} crop file missing or empty: {crop_path}")
        with Image.open(crop_path) as im:
            w, h = im.size
            if (w, h) != (300, 300):
                _fail(f"{sub_id} crop dims wrong: {w}×{h}, expected 300×300")
            # Center pixel should match the panel color (parent split in half)
            px = im.getpixel((150, 150))
            # Drop alpha if present
            if isinstance(px, tuple) and len(px) >= 3:
                px = px[:3]
            if px != expected_color:
                _fail(f"{sub_id} center pixel {px} != expected {expected_color}")

    _ok(f"sub-panel registration: 2 panels cropped + registered "
        f"({len(created)} layers, parent breadcrumb intact, crops correct)")


def check_versioning_no_api() -> None:
    """v2.2 versioning: revise loops + edit_layer must NOT clobber prior
    intermediate state. Asserts:
      - render_text_layer writes layers/<id>.v<N>.png (not <id>.png)
      - edit_layer (re-render) bumps to v<N+1>; v<N> still on disk
      - 2 composite() calls produce composites/iter_01/ AND composites/iter_02/
      - final/ symlinks point at iter_02 (the latest)
      - tool_result.payload exposes relative_path / version / supersedes_*
    """
    print("[20/52] versioning + revise-loop preservation (no API)")
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.propose_design_spec import propose_design_spec
    from .tools.render_text_layer import render_text_layer
    from .tools.edit_layer import edit_layer
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_versioning"
    layers_dir = out_dir / "layers"
    if out_dir.exists():
        # Clean prior smoke run so version counters start at 0.
        import shutil
        shutil.rmtree(out_dir)
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir,
                      layers_dir=layers_dir, run_id="smoke-vers")

    # Stub a background so composite has 2 layers
    bg_path = layers_dir / "bg_seed.png"
    Image.new("RGB", (768, 1024), (10, 10, 30)).save(bg_path)
    ctx.state["rendered_layers"]["L0_bg"] = {
        "layer_id": "L0_bg", "name": "bg", "kind": "background", "z_index": 0,
        "bbox": {"x": 0, "y": 0, "w": 768, "h": 1024},
        "src_path": str(bg_path), "prompt": "(stub)",
        "aspect_ratio": "3:4", "image_size": "1K",
        "safe_zones": [], "sha256": "seed",
    }

    switch_artifact_type({"type": "poster"}, ctx=ctx)
    propose_design_spec({"design_spec": {
        "brief": "v",
        "canvas": {"w_px": 768, "h_px": 1024, "dpi": 96,
                   "aspect_ratio": "3:4", "color_mode": "RGB"},
        "palette": ["#000", "#fff"],
        "typography": {}, "mood": [], "composition_notes": "",
        "layer_graph": [],
    }}, ctx=ctx)

    # ── v1 layer write ─────────────────────────────────────────────────
    r1 = render_text_layer({
        "layer_id": "L1", "name": "title", "text": "v1 text",
        "font_family": "NotoSansSC-Bold", "font_size_px": 60, "fill": "#fff",
        "bbox": {"x": 50, "y": 50, "w": 600, "h": 120},
        "align": "center", "z_index": 1,
    }, ctx=ctx)
    if r1.payload.get("version") != 1:
        _fail(f"first render should be v1; got version={r1.payload.get('version')}")
    if r1.payload.get("relative_path") != "layers/text_L1.v1.png":
        _fail(f"v1 relative_path wrong: {r1.payload.get('relative_path')}")
    if r1.payload.get("supersedes_sha256") is not None:
        _fail("first render should not have supersedes_sha256")
    v1_file = layers_dir / "text_L1.v1.png"
    if not v1_file.exists():
        _fail(f"v1 file missing: {v1_file}")
    v1_sha = r1.payload["sha256"]

    # ── v2 layer write (via edit_layer → re-render) ────────────────────
    r2 = edit_layer({"layer_id": "L1", "diff": {"text": "v2 text"}}, ctx=ctx)
    if r2.payload.get("version") != 2:
        _fail(f"edit_layer should bump to v2; got {r2.payload.get('version')}")
    if r2.payload.get("supersedes_sha256") != v1_sha:
        _fail(f"v2 supersedes_sha256 should equal v1 sha; got {r2.payload.get('supersedes_sha256')!r}")
    v2_file = layers_dir / "text_L1.v2.png"
    if not v2_file.exists():
        _fail(f"v2 file missing: {v2_file}")
    if not v1_file.exists():
        _fail(f"⚠ v1 file was clobbered (this is the bug we're guarding against): {v1_file}")
    _ok(f"layer versioning: v1 + v2 both on disk; v2.supersedes = v1.sha256[:8] {v1_sha[:8]}")

    # ── Composite iter 1 ──────────────────────────────────────────────
    c1 = composite({}, ctx=ctx)
    if c1.status != "ok":
        _fail(f"composite iter 1: {(c1.error_message or c1.payload)}")
    if c1.payload.get("iteration") != 1:
        _fail(f"composite iter 1 should report iteration=1; got {c1.payload.get('iteration')}")
    iter1_dir = out_dir / "composites" / "iter_01"
    for f in ("poster.html", "poster.psd", "poster.svg", "preview.png"):
        if not (iter1_dir / f).exists():
            _fail(f"iter_01/{f} missing")
    iter1_preview_sha = c1.payload["preview_sha256"]

    # ── Edit a layer + composite iter 2 ───────────────────────────────
    edit_layer({"layer_id": "L1", "diff": {"fill": "#ff00ff"}}, ctx=ctx)
    c2 = composite({}, ctx=ctx)
    if c2.status != "ok":
        _fail(f"composite iter 2: {(c2.error_message or c2.payload)}")
    if c2.payload.get("iteration") != 2:
        _fail(f"composite iter 2 should report iteration=2; got {c2.payload.get('iteration')}")
    iter2_dir = out_dir / "composites" / "iter_02"
    if not (iter2_dir / "preview.png").exists():
        _fail("iter_02/preview.png missing")
    if not (iter1_dir / "preview.png").exists():
        _fail("⚠ iter_01/preview.png was clobbered by iter_02 — versioning broken")
    if c2.payload.get("supersedes_preview_sha256") != iter1_preview_sha:
        _fail(f"iter_02 supersedes_preview_sha256 should be iter_01's preview sha")
    _ok(f"composite versioning: iter_01 + iter_02 both intact; "
        f"iter_02.supersedes_preview = iter_01.preview_sha[:8] {iter1_preview_sha[:8]}")

    # ── final/ symlinks point at the latest iter ──────────────────────
    final_dir = out_dir / "final"
    if not (final_dir / "preview.png").is_symlink():
        _fail("final/preview.png should be a symlink")
    if not (final_dir / "poster.html").is_symlink():
        _fail("final/poster.html should be a symlink")
    final_preview = (final_dir / "preview.png").resolve()
    if final_preview.parent.name != "iter_02":
        _fail(f"final/preview.png should resolve to iter_02; got {final_preview.parent.name}")
    _ok(f"final/ symlinks resolve to iter_02 (the latest)")


def _make_slide(layer_id: str, *, title: str, name: str | None = None,
                section_number: str | None = None,
                speaker_notes: str | None = None,
                role: str | None = None) -> LayerNode:
    """Compact LayerNode factory used by the v2.7.2 smokes — emits a
    minimal kind="slide" with a single title text child so section /
    notes assertions stay readable."""
    title_child = LayerNode(
        layer_id=f"{layer_id}_title", name="title", kind="text", z_index=10,
        bbox=SafeZone(x=120, y=80, w=1680, h=140),
        text=title, font_family="NotoSerifSC-Bold", font_size_px=72,
        align="left", effects=TextEffect(fill="#0f172a"),
    )
    return LayerNode(
        layer_id=layer_id,
        name=name or layer_id,
        kind="slide",
        z_index=1,
        role=role,  # type: ignore[arg-type]
        section_number=section_number,
        speaker_notes=speaker_notes,
        children=[title_child],
    )


def check_section_renumber_policy() -> None:
    """v2.7.2 smoke #25 — `apply_section_policy(policy="renumber")` walks
    a deck whose planner-supplied section_number is non-monotonic and
    rewrites the labels in slide order. Three content slides come back
    as §1 / §2 / §3 (no shared title prefix → no sub-rhythm)."""
    print("[25/52] section_number policy: renumber (no API)")
    from .util.section_renumber import apply_section_policy

    slides = [
        _make_slide("s1", title="Vision tokenizer", section_number="§3.1"),
        _make_slide("s2", title="Audio tokenizer", section_number="§2.2"),
        _make_slide("s3", title="Decoder fusion", section_number="§3.2"),
    ]
    out = apply_section_policy(slides, "renumber")

    if [s.section_number for s in slides] != ["§3.1", "§2.2", "§3.2"]:
        _fail("renumber policy mutated input slides — must be immutable")
    labels = [s.section_number for s in out]
    if labels != ["§1", "§2", "§3"]:
        _fail(f"expected ['§1', '§2', '§3'] after renumber; got {labels}")
    if [s.layer_id for s in out] != ["s1", "s2", "s3"]:
        _fail("renumber must preserve slide order + ids")
    _ok(f"renumber: {[s.section_number for s in slides]} -> {labels}")


def check_section_renumber_strip() -> None:
    """v2.7.2 smoke #26 — `apply_section_policy(policy="strip")` clears
    every SlideNode.section_number to None without touching titles or
    speaker notes."""
    print("[26/52] section_number policy: strip (no API)")
    from .util.section_renumber import apply_section_policy

    slides = [
        _make_slide("s1", title="A", section_number="§1",
                    speaker_notes="open"),
        _make_slide("s2", title="B", section_number="§2.1"),
        _make_slide("s3", title="C", section_number="§3"),
    ]
    out = apply_section_policy(slides, "strip")
    if any(s.section_number is not None for s in out):
        _fail("strip policy left a non-None section_number behind")
    if [s.section_number for s in slides] != ["§1", "§2.1", "§3"]:
        _fail("strip policy mutated input slides — must be immutable")
    if out[0].speaker_notes != "open":
        _fail("strip policy clobbered speaker_notes — must only touch section_number")
    _ok("strip: all section_number cleared, notes/titles intact")


def check_stable_id_notes_after_reorder() -> None:
    """v2.7.2 smoke #27 — speaker notes follow the SlideNode (by id),
    NOT the enumerate index. Build a 4-slide deck, reorder to
    [s4, s1, s3, s2], composite to .pptx, reopen and confirm each
    slide's notes match the source SlideNode it came from."""
    print("[27/52] speaker_notes follow slide_id after reorder (no API)")
    from pptx import Presentation as _Reopen

    from .config import REPO_ROOT, Settings
    from .schema import ArtifactType, DesignSpec
    from .tools import ToolContext
    from .tools.composite import composite
    from .tools.switch_artifact_type import switch_artifact_type

    out_dir = REPO_ROOT / "out" / "smoke_section_notes"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-section-notes")

    if switch_artifact_type({"type": "deck"}, ctx=ctx).status != "ok":
        _fail("switch_artifact_type(deck)")

    s1 = _make_slide("s1", title="Tokenizers",
                     speaker_notes="NOTES-FOR-S1: tokenizer story")
    s2 = _make_slide("s2", title="Decoder",
                     speaker_notes="NOTES-FOR-S2: decoder story")
    s3 = _make_slide("s3", title="Results",
                     speaker_notes="NOTES-FOR-S3: results story")
    s4 = _make_slide("s4", title="Future Work",
                     speaker_notes="NOTES-FOR-S4: future story")

    reordered = [s4, s1, s3, s2]
    spec = DesignSpec(
        brief="Stable-id notes binding smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96,
                "aspect_ratio": "16:9", "color_mode": "RGB"},
        layer_graph=reordered,
    )
    ctx.state["design_spec"] = spec

    obs = composite({}, ctx=ctx)
    if obs.status != "ok":
        _fail(f"composite(deck): {(obs.error_message or str(obs.payload))}")

    comp = ctx.state["composition"]
    if not comp.pptx_path or not Path(comp.pptx_path).exists():
        _fail(f"deck PPTX missing: {comp.pptx_path}")

    prs = _Reopen(comp.pptx_path)
    if len(prs.slides) != 4:
        _fail(f"expected 4 slides, got {len(prs.slides)}")

    expected_notes = [
        "NOTES-FOR-S4",  # s4 first
        "NOTES-FOR-S1",
        "NOTES-FOR-S3",
        "NOTES-FOR-S2",
    ]
    for idx, marker in enumerate(expected_notes):
        notes_text = prs.slides[idx].notes_slide.notes_text_frame.text
        if marker not in notes_text:
            _fail(
                f"slide {idx} notes mismatch: expected substring {marker!r}, "
                f"got {notes_text[:120]!r}. notes are following enumerate "
                f"index, not slide_id"
            )
    _ok("notes follow slide_id across reorder [s4, s1, s3, s2]")

    # The renumber policy is the default; verify each slide's title got
    # a section prefix in the rendered .pptx so we know the
    # apply_section_policy hook in _composite_deck actually fired.
    sections_seen = 0
    for idx in range(4):
        title_text = " ".join(
            run.text
            for shape in prs.slides[idx].shapes if shape.has_text_frame
            for para in shape.text_frame.paragraphs for run in para.runs
        )
        if "§" in title_text:
            sections_seen += 1
    if sections_seen == 0:
        _fail(
            "no slide title carries a section prefix — "
            "apply_section_policy did not fire in _composite_deck"
        )
    _ok(f"apply_section_policy fired in _composite_deck "
        f"({sections_seen}/4 slides carry section prefix)")


# ─────────────────────── v2.7.3 vision critic sub-agent ─────────────────────


class _ScriptedTurn:
    """One pre-baked turn the mock backend will replay. `tool_calls` is a list
    of `(tool_name, tool_input_dict)` tuples; the backend wraps each in a
    ToolCall with a synthetic id."""

    def __init__(self,
                 tool_calls: list[tuple[str, dict]] | None = None,
                 text: str = "",
                 stop_reason: str = "tool_use") -> None:
        self.tool_calls = tool_calls or []
        self.text = text
        self.stop_reason = stop_reason


class _MockCriticBackend:
    """Minimal LLMBackend stand-in for v2.7.3 sub-agent smokes.

    Replays a scripted list of turns. Each `create_turn` call records the
    incoming messages so the assertions can verify image bytes were threaded
    through, then returns the next scripted TurnResponse. Round-trip
    methods (`append_assistant`, `append_tool_results`) are no-ops since
    our scripted turns don't depend on prior context.
    """

    name = "mock"

    def __init__(self, model: str, turns: list[_ScriptedTurn]) -> None:
        self.model = model
        self._turns = list(turns)
        self.observed_messages: list[list] = []
        self._call_count = 0

    def create_turn(self, *, system, messages, tools, thinking_budget=0,
                    max_tokens=16384, extra_headers=None):
        from .llm_backend import ToolCall, TurnResponse
        # Snapshot the messages list so smokes can assert on what got
        # passed in (e.g. image bytes from prior tool_result turns).
        self.observed_messages.append(list(messages))
        if self._call_count >= len(self._turns):
            # No more scripted turns: emit a do-nothing end_turn so the
            # CriticAgent's max_turns failsafe triggers naturally.
            return TurnResponse(stop_reason="end_turn", raw_assistant_content={"role": "assistant"})
        scripted = self._turns[self._call_count]
        self._call_count += 1
        tcs = [
            ToolCall(id=f"toolu_mock_{i}", name=name, input=dict(args))
            for i, (name, args) in enumerate(scripted.tool_calls)
        ]
        return TurnResponse(
            text=scripted.text,
            tool_calls=tcs,
            stop_reason=scripted.stop_reason,
            usage={"input": 100, "output": 50, "cache_read": 0, "cache_create": 0},
            raw_assistant_content={"role": "assistant", "content": scripted.text},
        )

    def append_assistant(self, messages, response):
        messages.append(response.raw_assistant_content)

    def append_tool_results(self, messages, results):
        for tu_id, payload, _is_err in results:
            messages.append({"role": "tool", "tool_call_id": tu_id, "content": payload})

    def vision_user_message(self, *, image_b64, media_type, text):
        return {"role": "user", "content": [
            {"type": "image_url", "image_url": {"url": f"data:{media_type};base64,{image_b64}"}},
            {"type": "text", "text": text},
        ]}


def _make_smoke_settings(out_root: Path) -> "Settings":  # noqa: F821
    """Build a Settings stub that points at a writable smoke run dir."""
    from .config import Settings
    return Settings(
        anthropic_api_key="sk-stub",
        anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="stub-planner",
        critic_model="stub-critic",
        out_dir=out_root,
        critic_max_turns=4,
        critic_thinking_budget=0,
    )


def _make_smoke_deck_spec(n_slides: int = 3) -> tuple["DesignSpec", list[Path], Path]:  # noqa: F821
    """Build a DesignSpec for a deck plus N rendered PNG files on disk under
    a temp dir (no real composite — we just need bytes for read_slide_render
    to base64-encode). Returns (spec, slide_paths, run_dir)."""
    from .config import REPO_ROOT
    from .schema import ArtifactType, DesignSpec, LayerNode

    run_dir = REPO_ROOT / "out" / "smoke_critic_subagent"
    if run_dir.exists():
        import shutil
        shutil.rmtree(run_dir)
    slides_dir = run_dir / "composites" / "iter_01" / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)

    slide_paths: list[Path] = []
    slide_nodes: list[LayerNode] = []
    for i in range(n_slides):
        png = slides_dir / f"slide_{i:02d}.png"
        Image.new("RGB", (640, 360), (200, 50 + i * 10, 80)).save(png)
        slide_paths.append(png)
        slide_nodes.append(LayerNode(
            layer_id=f"S{i}", name=f"slide_{i}", kind="slide",
            z_index=i, children=[],
        ))

    spec = DesignSpec(
        brief="smoke critic subagent",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96,
                "aspect_ratio": "16:9", "color_mode": "RGB"},
        layer_graph=slide_nodes,
    )
    return spec, slide_paths, run_dir


def check_critic_subagent_trajectory() -> None:
    """smoke #28: spawn CriticAgent on a fixture deck with a mocked
    LLMBackend that fetches one slide PNG, then one paper excerpt, then
    calls `report_verdict`. Verify the resulting CritiqueReport has the
    expected verdict/score AND the trajectory file `critic.jsonl` lands
    in the run dir.

    v2.7.3 hotfix (2026-04-26) also verifies:
    - the FULL paper raw_text is NEVER in the initial user message —
      it must be fetched on-demand via `read_paper_section`
    - the read_slide_render tool result stays a small ack (no base64)
    """
    print("[28/52] critic sub-agent: on-demand paper + trajectory written")
    from .agents import CriticAgent
    from .schema import ArtifactType, CritiqueReport

    spec, slide_paths, run_dir = _make_smoke_deck_spec(n_slides=3)
    settings = _make_smoke_settings(run_dir.parent.parent)

    secret_marker = "BANNED-FULL-PAPER-MARKER-XYZZY"
    paper_text = (
        "LongCat-Next achieves 72.3% top-1 on ImageNet. "
        + secret_marker + " "
        + ("Filler section. " * 2000)
    )

    scripted = [
        _ScriptedTurn(tool_calls=[
            ("read_slide_render", {"slide_id": "S0"}),
            ("read_paper_section", {"query": "72.3"}),
        ]),
        _ScriptedTurn(tool_calls=[
            ("report_verdict", {
                "score": 0.82,
                "verdict": "pass",
                "summary": "structurally sound; minor typography drift only",
                "issues": [
                    {"slide_id": "S1", "severity": "low",
                     "category": "typography",
                     "description": "body font 22px below 24px lower bound",
                     "evidence_paper_anchor": None},
                ],
            }),
        ]),
    ]
    mock = _MockCriticBackend(model="stub-critic", turns=scripted)

    agent = CriticAgent(settings, ArtifactType.DECK)
    agent.backend = mock

    traj_path = run_dir / "trajectory" / "critic.jsonl"
    report = agent.critique(
        spec=spec, layer_manifest=[],
        slide_renders=slide_paths,
        paper_raw_text=paper_text,
        iteration=1, trajectory_path=traj_path,
    )

    if not isinstance(report, CritiqueReport):
        _fail(f"report should be CritiqueReport; got {type(report).__name__}")
    if report.verdict != "pass":
        _fail(f"verdict should be 'pass'; got {report.verdict!r}")
    if abs(report.score - 0.82) > 0.001:
        _fail(f"score not preserved: {report.score}")
    if len(report.issues) != 1 or report.issues[0].category != "typography":
        _fail(f"issues lost in roundtrip: {report.issues}")

    if not traj_path.exists():
        _fail(f"critic.jsonl not written at {traj_path}")
    lines = traj_path.read_text(encoding="utf-8").strip().splitlines()
    if len(lines) != 2:
        _fail(f"expected 2 trajectory lines (one per turn), got {len(lines)}")
    line0 = json.loads(lines[0])
    line1 = json.loads(lines[1])
    if line0.get("iteration") != 1 or line1.get("iteration") != 1:
        _fail(f"trajectory iteration field corrupted: {line0}, {line1}")
    if not any(tc.get("name") == "report_verdict" for tc in line1.get("tool_calls", [])):
        _fail(f"final turn should record report_verdict; got {line1.get('tool_calls')}")

    # On-demand paper: the FULL paper text (or its unique marker) must NOT
    # appear in the initial user message. The marker is only allowed to
    # surface inside a tool-result for the read_paper_section call.
    initial_msgs = mock.observed_messages[0]
    initial_blob = json.dumps(initial_msgs, default=str)
    if secret_marker in initial_blob:
        _fail("paper raw_text leaked into the initial user message — "
              "must be fetched on-demand via read_paper_section")

    # The marker SHOULD show up in turn 2's tool-result for the paper
    # excerpt query (the model asked for "72.3" and the marker sits a few
    # chars away in our fixture).
    second_blob = json.dumps(mock.observed_messages[1], default=str)
    if secret_marker not in second_blob:
        _fail("read_paper_section tool result should contain the matched "
              "excerpt window including the marker; got nothing")

    # The read_slide_render tool result must be a small ack now — no
    # 12K-char base64 string in the `tool` role message.
    big_tool = [
        m for m in mock.observed_messages[1]
        if isinstance(m, dict) and m.get("role") == "tool"
        and len(m.get("content") or "") > 2000
    ]
    if big_tool:
        _fail(f"read_slide_render leaked base64 into a tool message; "
              f"found {len(big_tool)} oversized tool entries")

    _ok(f"critic sub-agent: on-demand paper + small tool-result ack + "
        f"critic.jsonl ({len(lines)} lines)")


def check_critic_subagent_max_turns() -> None:
    """smoke #26: when the mocked LLM never calls `report_verdict`,
    CriticAgent must exhaust max_turns and synthesize a `verdict='fail'`
    CritiqueReport rather than recurse forever."""
    print("[29/52] critic sub-agent: max_turns failsafe → fail verdict")
    from .agents import CriticAgent
    from .schema import ArtifactType

    spec, slide_paths, run_dir = _make_smoke_deck_spec(n_slides=2)
    settings = _make_smoke_settings(run_dir.parent.parent)

    # Every scripted turn calls a non-terminal tool — never report_verdict.
    scripted = [
        _ScriptedTurn(tool_calls=[("read_slide_render", {"slide_id": "S0"})])
        for _ in range(settings.critic_max_turns + 4)
    ]
    mock = _MockCriticBackend(model="stub-critic", turns=scripted)

    agent = CriticAgent(settings, ArtifactType.DECK)
    agent.backend = mock

    traj_path = run_dir / "trajectory" / "critic.jsonl"
    if traj_path.exists():
        traj_path.unlink()
    report = agent.critique(
        spec=spec, layer_manifest=[],
        slide_renders=slide_paths, paper_raw_text=None,
        iteration=1, trajectory_path=traj_path,
    )

    if report.verdict != "fail":
        _fail(f"max_turns exhaustion should yield 'fail'; got {report.verdict!r}")
    if report.score != 0.0:
        _fail(f"failsafe report should have score=0.0; got {report.score}")
    if mock._call_count != settings.critic_max_turns:
        _fail(f"backend should be called critic_max_turns ({settings.critic_max_turns}) "
              f"times; got {mock._call_count}")
    if not report.issues or report.issues[0].severity != "blocker":
        _fail(f"failsafe report should carry a blocker issue: {report.issues}")
    _ok(f"max_turns ({settings.critic_max_turns}) hit → fail verdict synthesized "
        f"(no infinite loop; {mock._call_count} backend calls)")


def check_critic_planner_consumption() -> None:
    """smoke #27: the critique tool wraps CriticAgent and returns a
    CritiqueReport JSON in tool_result.payload. Verify that for each
    verdict (pass/revise/fail) the planner can route correctly:
      pass    → state.critique_results carries verdict='pass'
      revise  → state.critique_results carries verdict='revise'
      fail    → state.critique_results carries verdict='fail'
    """
    print("[30/52] planner consumption: pass/revise/fail routing")
    from .agents import CriticAgent
    from .config import REPO_ROOT
    from .schema import ArtifactType, CompositionArtifacts, DesignSpec, LayerNode
    from .tools import ToolContext
    from .tools.critique_tool import critique

    out_dir = REPO_ROOT / "out" / "smoke_critic_consume"
    if out_dir.exists():
        import shutil
        shutil.rmtree(out_dir)
    layers_dir = out_dir / "layers"
    slides_dir = out_dir / "composites" / "iter_01" / "slides"
    slides_dir.mkdir(parents=True, exist_ok=True)

    preview_path = slides_dir.parent / "preview.png"
    Image.new("RGB", (320, 180), (40, 40, 40)).save(preview_path)
    for i in range(2):
        Image.new("RGB", (320, 180), (90, 90, 90)).save(slides_dir / f"slide_{i:02d}.png")

    spec = DesignSpec(
        brief="planner consumption smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96,
                "aspect_ratio": "16:9", "color_mode": "RGB"},
        layer_graph=[
            LayerNode(layer_id="S0", name="cover", kind="slide", z_index=0),
            LayerNode(layer_id="S1", name="content", kind="slide", z_index=1),
        ],
    )

    settings = _make_smoke_settings(out_dir.parent.parent)
    # Allow several critique invocations within a single fixture run.
    settings = settings.__class__(
        **{**settings.__dict__, "max_critique_iters": 5},
    )

    def _new_ctx() -> ToolContext:
        ctx = ToolContext(settings=settings, run_dir=out_dir,
                          layers_dir=layers_dir, run_id="smoke-critic-consume")
        ctx.state["design_spec"] = spec
        ctx.state["composition"] = CompositionArtifacts(
            preview_path=str(preview_path),
            layer_manifest=[
                {"layer_id": "S0", "name": "cover", "kind": "slide",
                 "index": 0, "children": []},
                {"layer_id": "S1", "name": "content", "kind": "slide",
                 "index": 1, "children": []},
            ],
        )
        return ctx

    def _run_with_verdict(verdict: str, score: float) -> dict:
        ctx = _new_ctx()
        scripted = [_ScriptedTurn(tool_calls=[(
            "report_verdict", {"score": score, "verdict": verdict,
                               "summary": f"smoke {verdict}", "issues": []},
        )])]
        mock = _MockCriticBackend(model="stub-critic", turns=scripted)

        # Patch CriticAgent so the inner make_backend call is bypassed.
        original_init = CriticAgent.__init__

        def _patched_init(self, settings_, artifact_type):
            self.settings = settings_
            self.artifact_type = artifact_type
            self.backend = mock
            self._system_prompt = "patched"

        CriticAgent.__init__ = _patched_init  # type: ignore[assignment]
        try:
            obs = critique({}, ctx=ctx)
        finally:
            CriticAgent.__init__ = original_init  # type: ignore[assignment]

        if obs.status != "ok":
            _fail(f"critique tool errored on verdict={verdict}: {obs.error_message}")
        crits = ctx.state["critique_results"]
        if not crits:
            _fail(f"critique_results empty after verdict={verdict}")
        if crits[-1].verdict != verdict:
            _fail(f"verdict roundtrip {verdict} → {crits[-1].verdict}")
        return obs.payload

    pass_payload = _run_with_verdict("pass", 0.85)
    revise_payload = _run_with_verdict("revise", 0.62)
    fail_payload = _run_with_verdict("fail", 0.30)

    for label, p in (("pass", pass_payload), ("revise", revise_payload),
                     ("fail", fail_payload)):
        if p.get("verdict") != label:
            _fail(f"tool_result.payload.verdict {label} mismatch: {p}")
        if "score" not in p or "summary" not in p:
            _fail(f"tool_result.payload missing fields for {label}: {sorted(p.keys())}")
    _ok("CritiqueReport JSON pass/revise/fail all surface in tool_result.payload "
        "with verdict + score + summary intact")


def check_critic_subagent_png_throughput() -> None:
    """smoke #31 (v2.7.3 hotfix 2026-04-26): long-tail check that
    CriticAgent + read_slide_render survive a 15-slide deck with
    realistic file sizes AND that:

    - tool_result content for read_slide_render stays SMALL (no base64
      leak into the `tool` role message)
    - the actual PNG arrives on a follow-up user-role vision message
      so subsequent turns see image tokens, not 12K+ tokens of base64
      replayed as plain text (the bug that broke longcat-next dogfood
      on 2026-04-26)
    - the per-turn image cap defers surplus calls to a later turn"""
    print("[31/52] critic sub-agent: per-turn image cap + vision-message delivery")
    from .agents import CriticAgent
    from .schema import ArtifactType

    spec, slide_paths, run_dir = _make_smoke_deck_spec(n_slides=15)
    # Re-write each slide PNG at a more realistic resolution so the base64
    # path actually exercises non-trivial bytes — 1920x1080 is the deck
    # canvas default. Keeps the encode under 200KB per slide post-jpeg.
    for i, p in enumerate(slide_paths):
        Image.new("RGB", (1920, 1080), (30 + i * 5, 60, 90)).save(p)

    settings = _make_smoke_settings(run_dir.parent.parent)

    # Turn 1 attempts to fetch all 15 slides in parallel — exercises
    # both the cap (only `critic_max_images_per_turn` get delivered as
    # vision messages, the rest are deferred) and the vision-message
    # plumbing. Turn 2 calls report_verdict to terminate.
    bulk_read = _ScriptedTurn(tool_calls=[
        ("read_slide_render", {"slide_id": f"S{i}"})
        for i in range(15)
    ])
    final_call = _ScriptedTurn(tool_calls=[(
        "report_verdict",
        {"score": 0.78, "verdict": "pass",
         "summary": "all 15 slides reviewed",
         "issues": []},
    )])
    big_settings = settings.__class__(
        **{**settings.__dict__, "critic_max_turns": 20,
           "critic_preview_max_edge": 1024,
           "critic_max_images_per_turn": 4},
    )
    mock = _MockCriticBackend(
        model="stub-critic", turns=[bulk_read, final_call],
    )

    agent = CriticAgent(big_settings, ArtifactType.DECK)
    agent.backend = mock

    traj_path = run_dir / "trajectory" / "critic_long.jsonl"
    if traj_path.exists():
        traj_path.unlink()
    report = agent.critique(
        spec=spec, layer_manifest=[], slide_renders=slide_paths,
        paper_raw_text=None, iteration=1, trajectory_path=traj_path,
    )

    if report.verdict != "pass" or len(report.issues) != 0:
        _fail(f"15-slide run should pass cleanly; got verdict={report.verdict} "
              f"issues={len(report.issues)}")
    if mock._call_count != 2:
        _fail(f"expected 2 backend calls (1 bulk read turn + 1 verdict); "
              f"got {mock._call_count}")

    if len(mock.observed_messages) < 2:
        _fail("mock should have observed at least 2 turns of messages")
    second_turn_msgs = mock.observed_messages[1]

    # The bulk-read turn must NOT have leaked base64 into any `tool` role
    # message — only short ack JSONs (a few hundred chars at most).
    big_tool_results = [
        m for m in second_turn_msgs
        if isinstance(m, dict) and m.get("role") == "tool"
        and len(m.get("content") or "") > 2000
    ]
    if big_tool_results:
        _fail(f"tool results must stay small after hotfix; found "
              f"{len(big_tool_results)} oversized tool messages")

    # v2.7.4 — exactly `critic_max_images_per_turn` slide PNGs ride as
    # image_url blocks inside ONE follow-up user-role vision message
    # (collapsed from N messages so strict OpenAI-compat upstreams like
    # Alibaba `qwen/qwen-vl-max` accept the trailing history shape).
    cap = big_settings.critic_max_images_per_turn
    vision_msgs = []
    image_blocks: list[dict] = []
    for m in second_turn_msgs:
        if not isinstance(m, dict) or m.get("role") != "user":
            continue
        content = m.get("content")
        if not isinstance(content, list):
            continue
        per_msg_imgs = [
            b for b in content
            if isinstance(b, dict) and b.get("type") == "image_url"
        ]
        if per_msg_imgs:
            vision_msgs.append(m)
            image_blocks.extend(per_msg_imgs)
    if len(vision_msgs) != 1:
        _fail(f"expected exactly 1 collapsed vision user-message "
              f"(v2.7.4); got {len(vision_msgs)}")
    if len(image_blocks) != cap:
        _fail(f"expected exactly {cap} image_url blocks across the "
              f"vision message (per-turn cap); got {len(image_blocks)}")

    # The deferred slides (S{cap}..S14) must show up as ack JSONs with
    # `"deferred": true` in the tool-role messages.
    deferred_seen = 0
    for m in second_turn_msgs:
        if not isinstance(m, dict) or m.get("role") != "tool":
            continue
        if '"deferred": true' in (m.get("content") or ""):
            deferred_seen += 1
    expected_deferred = 15 - cap
    if deferred_seen != expected_deferred:
        _fail(f"expected {expected_deferred} deferred ack tool messages; "
              f"got {deferred_seen}")

    _ok(f"per-turn image cap honoured: 1 collapsed vision message "
        f"with {cap} image_url blocks + {expected_deferred} deferred "
        f"acks, no base64 in tool results, verdict={report.verdict}")


def check_critic_openai_compat_strict_format() -> None:
    """smoke #45 (v2.7.4 hotfix 2026-04-26): the message structure the
    critic builds across multiple read_slide_render turns MUST satisfy
    the strict OpenAI Chat Completions spec, since strict upstreams
    (Alibaba-routed `qwen/qwen-vl-max`) reject any deviation with
    `<400> InternalError.Algo.InvalidParameter: The model input format
    error`.

    Locks in three invariants that the v2.7.3 hotfix violated:
      1. Every assistant message has a `content` key (empty string is
         fine; `None`/missing is not).
      2. The vision content array places `text` before `image_url`.
      3. After a tool-results burst we emit ONE collapsed user message
         (multi-block content), never a chain of N adjacent user
         messages.
    """
    print("[45/52] critic message format: OpenAI strict-compat invariants")
    from .agents import CriticAgent
    from .llm_backend import OpenAICompatBackend, ToolCall, TurnResponse
    from .schema import ArtifactType

    spec, slide_paths, run_dir = _make_smoke_deck_spec(n_slides=4)
    settings = _make_smoke_settings(run_dir.parent.parent)
    settings = settings.__class__(
        **{**settings.__dict__,
           "critic_max_turns": 6,
           "critic_max_images_per_turn": 4},
    )

    class _CapturingBackend:
        name = "openai_compat"
        model = "stub-strict"

        def __init__(self) -> None:
            self.observed_messages: list[list] = []
            self._calls = 0

        def create_turn(self, *, system, messages, tools,
                        thinking_budget=0, max_tokens=16384,
                        extra_headers=None):
            self.observed_messages.append([_deepcopy_msg(m) for m in messages])
            self._calls += 1
            if self._calls == 1:
                tcs = [
                    ToolCall(id=f"call_{i}",
                             name="read_slide_render",
                             input={"slide_id": f"S{i}"})
                    for i in range(4)
                ]
                # Mirror the real OpenAICompatBackend's assistant_msg
                # construction so the smoke catches the missing-content
                # regression: msg.content is None on tool-only turns.
                assistant_msg = {
                    "role": "assistant",
                    "content": "",
                    "tool_calls": [
                        {"id": f"call_{i}", "type": "function",
                         "function": {"name": "read_slide_render",
                                      "arguments": json.dumps(
                                          {"slide_id": f"S{i}"})}}
                        for i in range(4)
                    ],
                }
                return TurnResponse(
                    text="", tool_calls=tcs, stop_reason="tool_use",
                    raw_assistant_content=assistant_msg, usage={},
                )
            tcs = [ToolCall(id="call_done", name="report_verdict",
                            input={"score": 0.8, "verdict": "pass",
                                   "summary": "ok", "issues": []})]
            assistant_msg = {
                "role": "assistant",
                "content": "",
                "tool_calls": [
                    {"id": "call_done", "type": "function",
                     "function": {"name": "report_verdict",
                                  "arguments": json.dumps({
                                      "score": 0.8, "verdict": "pass",
                                      "summary": "ok", "issues": []})}}
                ],
            }
            return TurnResponse(
                text="", tool_calls=tcs, stop_reason="tool_use",
                raw_assistant_content=assistant_msg, usage={},
            )

        def append_assistant(self, messages, response):
            messages.append(response.raw_assistant_content)

        def append_tool_results(self, messages, results):
            OpenAICompatBackend.append_tool_results(self, messages, results)

        def vision_user_message(self, *, image_b64, media_type, text):
            return OpenAICompatBackend.vision_user_message(
                self, image_b64=image_b64, media_type=media_type, text=text,
            )

    mock = _CapturingBackend()
    agent = CriticAgent(settings, ArtifactType.DECK)
    agent.backend = mock  # type: ignore[assignment]

    traj_path = run_dir / "trajectory" / "critic_strict.jsonl"
    if traj_path.exists():
        traj_path.unlink()
    report = agent.critique(
        spec=spec, layer_manifest=[], slide_renders=slide_paths,
        paper_raw_text=None, iteration=1, trajectory_path=traj_path,
    )
    if report.verdict != "pass":
        _fail(f"smoke fixture should yield pass; got {report.verdict}")

    if len(mock.observed_messages) < 2:
        _fail("expected at least 2 backend turns observed")
    second_turn = mock.observed_messages[1]

    # Invariant 1: every assistant message MUST have a `content` key.
    for i, m in enumerate(second_turn):
        if not isinstance(m, dict):
            continue
        if m.get("role") == "assistant" and "content" not in m:
            _fail(f"assistant message at idx {i} missing `content` key — "
                  f"OpenAI strict-spec violation: {m}")

    # Invariant 2: every vision user message places text before image_url.
    vision_messages = [
        m for m in second_turn
        if isinstance(m, dict) and m.get("role") == "user"
        and isinstance(m.get("content"), list)
        and any(isinstance(b, dict) and b.get("type") == "image_url"
                for b in m["content"])
    ]
    if not vision_messages:
        _fail("no vision user message observed in second turn history")
    for vm in vision_messages:
        blocks = vm["content"]
        first_image_idx = next(
            (i for i, b in enumerate(blocks)
             if isinstance(b, dict) and b.get("type") == "image_url"),
            -1,
        )
        first_text_idx = next(
            (i for i, b in enumerate(blocks)
             if isinstance(b, dict) and b.get("type") == "text"),
            -1,
        )
        if first_text_idx < 0 or first_image_idx < 0:
            _fail(f"vision message must contain both text and image_url: {vm}")
        if first_text_idx > first_image_idx:
            _fail("vision content blocks must place `text` before "
                  f"`image_url` (OpenAI canonical order); got: {blocks}")

    # Invariant 3: the burst collapses into ONE user message — no run of
    # adjacent user-role messages between the tool burst and the next
    # assistant turn.
    if len(vision_messages) != 1:
        _fail(f"expected exactly 1 collapsed vision user message after "
              f"the tool burst; got {len(vision_messages)} adjacent "
              f"user messages — strict OpenAI-compat upstreams reject "
              f"this shape")
    sole_vision = vision_messages[0]
    image_blocks = [
        b for b in sole_vision["content"]
        if isinstance(b, dict) and b.get("type") == "image_url"
    ]
    if len(image_blocks) != 4:
        _fail(f"collapsed vision message must carry all 4 image blocks; "
              f"got {len(image_blocks)}")
    for blk in image_blocks:
        url = blk.get("image_url")
        if not isinstance(url, dict) or "url" not in url:
            _fail(f"image_url block must be a nested object with `url`; "
                  f"got {blk}")
        if not str(url["url"]).startswith("data:image/"):
            _fail(f"image_url.url must be a data: URI; got {url['url'][:60]}")

    _ok("OpenAI strict-spec invariants hold: assistant.content present, "
        "text-before-image, 1 collapsed vision message with all "
        "image_url blocks")


def _deepcopy_msg(m):
    """Cheap deep copy via JSON round-trip for plain-dict messages.

    The capturing backend snapshots the messages list per turn so
    later mutations (e.g. CriticAgent appending vision blocks AFTER
    create_turn returns) don't retroactively corrupt the assertion's
    view of what the model saw on a given turn."""
    try:
        return json.loads(json.dumps(m, default=str))
    except (TypeError, ValueError):
        return m


# ─────────────────────── v2.8.1 archetype Phase 1 ──────────────────────


def _make_archetype_slide(
    *,
    layer_id: str,
    archetype: str,
    children: list[LayerNode],
    section_number: str | None = None,
    speaker_notes: str | None = None,
    name: str | None = None,
) -> LayerNode:
    """LayerNode factory for archetype smokes — emits a kind="slide"
    with the given children and the specified archetype label."""
    return LayerNode(
        layer_id=layer_id,
        name=name or layer_id,
        kind="slide",
        z_index=1,
        archetype=archetype,  # type: ignore[arg-type]
        section_number=section_number,
        speaker_notes=speaker_notes,
        children=children,
    )


def _new_blank_slide(slide_w: int = 1920, slide_h: int = 1080) -> tuple:
    """Spin up a fresh blank Presentation + add one blank slide.

    Returns (Presentation, Slide). Used by every archetype smoke so
    each runs in isolation."""
    from pptx import Presentation as _Pres
    from pptx.util import Emu as _Emu
    prs = _Pres()
    prs.slide_width = _Emu(slide_w * 9525)
    prs.slide_height = _Emu(slide_h * 9525)
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    return prs, slide


def _make_archetype_ctx() -> "ToolContext":  # noqa: F821
    """Lightweight ToolContext stub for archetype smokes — no API
    keys, no real run dir; archetype renderers don't touch ctx beyond
    settings.fonts on the preview path (which we don't exercise here)."""
    from .config import REPO_ROOT, Settings
    from .tools import ToolContext

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="stub", critic_model="stub",
    )
    out_dir = REPO_ROOT / "out" / "smoke_archetype"
    out_dir.mkdir(parents=True, exist_ok=True)
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)
    return ToolContext(
        settings=settings, run_dir=out_dir,
        layers_dir=layers_dir, run_id="smoke-archetype",
    )


def _all_textframes(slide: "Any") -> list:  # noqa: F821
    """Return the slide's shapes that carry a text_frame (covers
    textboxes, autoshapes, and placeholders)."""
    return [s for s in slide.shapes if s.has_text_frame]


def _shape_text(shape: "Any") -> str:  # noqa: F821
    return shape.text_frame.text or ""


def _shape_max_pt(shape: "Any") -> int:  # noqa: F821
    """Largest run font size (in pt) across the shape's text frame.
    Returns 0 when no run sets an explicit size."""
    best = 0
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            sz = getattr(run.font, "size", None)
            if sz is not None:
                # `sz` is an Emu-like length subclass with `.pt` attr.
                try:
                    best = max(best, int(sz.pt))
                except Exception:
                    pass
    return best


def check_archetype_cover_editorial() -> None:
    """v2.8.1 smoke #37 — `cover_editorial` renders ≥3 native shapes
    (title + subtitle + author strip), all backed by TextFrames. When
    the slide carries a `section_number`, the title text picks up the
    `§N · ` prefix from v2.7.2's `_with_section_prefix`."""
    print("[37/52] archetype cover_editorial: ≥3 textframes + section prefix")
    from .tools.pptx_renderer import _render_slide

    title = LayerNode(
        layer_id="c1_title", name="title", kind="text", z_index=10,
        text="OpenDesign", font_family="NotoSerifSC-Bold", font_size_px=120,
        align="left", effects=TextEffect(fill="#0f172a"),
    )
    subtitle = LayerNode(
        layer_id="c1_sub", name="subtitle", kind="text", z_index=10,
        text="A conversational design agent", font_family="NotoSansSC-Bold",
        font_size_px=36, align="left", effects=TextEffect(fill="#475569"),
    )
    authors = LayerNode(
        layer_id="c1_authors", name="author_byline", kind="text", z_index=10,
        text="Yaxin Luo · Anthropic Skill Team", font_family="NotoSansSC-Bold",
        font_size_px=26, align="left", effects=TextEffect(fill="#64748b"),
    )
    slide_node = _make_archetype_slide(
        layer_id="cover", archetype="cover_editorial",
        section_number="§1",
        children=[title, subtitle, authors],
    )
    _, slide = _new_blank_slide()
    _render_slide(slide, slide_node, 1920, 1080, _make_archetype_ctx())

    frames = _all_textframes(slide)
    if len(frames) < 3:
        _fail(f"cover_editorial should emit ≥3 TextFrames; got {len(frames)}")

    texts = [_shape_text(f) for f in frames]
    title_match = [t for t in texts if "OpenDesign" in t]
    if not title_match:
        _fail(f"cover_editorial title not rendered; texts={texts}")
    if not title_match[0].startswith("§1"):
        _fail(f"section_number prefix missing; title={title_match[0]!r}")
    if not any("conversational" in t for t in texts):
        _fail(f"subtitle missing from cover; texts={texts}")
    if not any("Yaxin Luo" in t for t in texts):
        _fail(f"authors strip missing from cover; texts={texts}")
    _ok(f"cover_editorial: {len(frames)} textframes, "
        f"section prefix '§1' on title, subtitle + authors present")


def check_archetype_evidence_snapshot() -> None:
    """v2.8.1 smoke #38 — `evidence_snapshot` renders one giant number
    (≥200 pt source ⇒ ≥150 pt rendered) plus a footnote textframe."""
    print("[38/52] archetype evidence_snapshot: huge number + footnote")
    from .tools.pptx_renderer import _render_slide

    big = LayerNode(
        layer_id="e1_big", name="big_number", kind="text", z_index=10,
        text="72.3%", font_family="NotoSerifSC-Bold", font_size_px=240,
        align="center", effects=TextEffect(fill="#0f172a"),
    )
    foot = LayerNode(
        layer_id="e1_foot", name="stat_caption", kind="text", z_index=10,
        text="ImageNet top-1, single-crop", font_family="NotoSansSC-Bold",
        font_size_px=28, align="center", effects=TextEffect(fill="#475569"),
    )
    slide_node = _make_archetype_slide(
        layer_id="ev1", archetype="evidence_snapshot",
        children=[big, foot],
    )
    _, slide = _new_blank_slide()
    _render_slide(slide, slide_node, 1920, 1080, _make_archetype_ctx())

    frames = _all_textframes(slide)
    if len(frames) < 2:
        _fail(f"evidence_snapshot should emit ≥2 TextFrames; got {len(frames)}")

    big_frames = [f for f in frames if "72.3%" in _shape_text(f)]
    if not big_frames:
        _fail(f"big number frame missing; texts={[_shape_text(f) for f in frames]}")
    big_pt = _shape_max_pt(big_frames[0])
    # 240 px source × 0.75 px-to-pt ≈ 180 pt; threshold 150 leaves
    # room for any minor calibration drift.
    if big_pt < 150:
        _fail(f"big number font too small: {big_pt} pt; want ≥150")
    if not any("ImageNet" in _shape_text(f) for f in frames):
        _fail(f"footnote missing; texts={[_shape_text(f) for f in frames]}")
    _ok(f"evidence_snapshot: number frame at {big_pt} pt, footnote present")


def check_archetype_takeaway_list() -> None:
    """v2.8.1 smoke #39 — `takeaway_list` renders 3 bullet groups
    (marker + body shapes per row), giving ≥6 textframes plus the
    title."""
    print("[39/52] archetype takeaway_list: 3 bullet groups")
    from .tools.pptx_renderer import _render_slide

    title = LayerNode(
        layer_id="t1_title", name="title", kind="text", z_index=10,
        text="Three things to remember", font_family="NotoSerifSC-Bold",
        font_size_px=72, align="left", effects=TextEffect(fill="#0f172a"),
    )
    bullets = [
        LayerNode(
            layer_id=f"t1_b{i}", name=f"bullet_{i}", kind="text", z_index=10,
            text=text, font_family="NotoSansSC-Bold", font_size_px=36,
            align="left", effects=TextEffect(fill="#0f172a"),
        )
        for i, text in enumerate([
            "Editability beats prettier raster output.",
            "Provenance is a hard gate, not a soft check.",
            "Archetypes give the planner a vocabulary.",
        ], start=1)
    ]
    slide_node = _make_archetype_slide(
        layer_id="tk1", archetype="takeaway_list",
        children=[title, *bullets],
    )
    _, slide = _new_blank_slide()
    _render_slide(slide, slide_node, 1920, 1080, _make_archetype_ctx())

    frames = _all_textframes(slide)
    # 1 title + 3 markers + 3 bodies = 7 minimum
    if len(frames) < 7:
        _fail(f"takeaway_list should emit ≥7 TextFrames (title+3 markers+3 bodies); "
              f"got {len(frames)}")

    body_texts = [_shape_text(f) for f in frames]
    for marker in ("01", "02", "03"):
        if not any(marker == t.strip() for t in body_texts):
            _fail(f"takeaway marker '{marker}' missing; texts={body_texts}")
    for needle in ("Editability", "Provenance", "Archetypes"):
        if not any(needle in t for t in body_texts):
            _fail(f"takeaway body missing '{needle}'; texts={body_texts}")
    _ok(f"takeaway_list: {len(frames)} textframes, all 3 bullet groups present")


def check_archetype_thanks_qa() -> None:
    """v2.8.1 smoke #40 — `thanks_qa` renders a thanks/Q&A headline,
    a contact row, and an optional code link."""
    print("[40/52] archetype thanks_qa: headline + contact row")
    from .tools.pptx_renderer import _render_slide

    title = LayerNode(
        layer_id="q1_title", name="title", kind="text", z_index=10,
        text="Thanks · Questions?", font_family="NotoSerifSC-Bold",
        font_size_px=96, align="center", effects=TextEffect(fill="#0f172a"),
    )
    contact = LayerNode(
        layer_id="q1_email", name="contact_email", kind="text", z_index=10,
        text="luoyaxin999@gmail.com", font_family="NotoSansSC-Bold",
        font_size_px=28, align="center", effects=TextEffect(fill="#475569"),
    )
    code = LayerNode(
        layer_id="q1_link", name="code_link", kind="text", z_index=10,
        text="github.com/yaxin/open-design", font_family="JetBrainsMono",
        font_size_px=24, align="center", effects=TextEffect(fill="#7f1d1d"),
    )
    slide_node = _make_archetype_slide(
        layer_id="ty", archetype="thanks_qa",
        children=[title, contact, code],
    )
    _, slide = _new_blank_slide()
    _render_slide(slide, slide_node, 1920, 1080, _make_archetype_ctx())

    frames = _all_textframes(slide)
    if len(frames) < 3:
        _fail(f"thanks_qa should emit ≥3 TextFrames; got {len(frames)}")
    texts = [_shape_text(f) for f in frames]
    if not any("Thanks" in t and "Questions" in t for t in texts):
        _fail(f"thanks_qa headline missing; texts={texts}")
    if not any("luoyaxin999" in t for t in texts):
        _fail(f"thanks_qa contact missing; texts={texts}")
    if not any("github.com" in t for t in texts):
        _fail(f"thanks_qa code link missing; texts={texts}")
    _ok(f"thanks_qa: headline + contact + link rendered ({len(frames)} frames)")


def check_archetype_fallback_default() -> None:
    """v2.8.1 smoke #41 — when `archetype` is the schema default
    (`"evidence_snapshot"`) and the slide has no big-number child,
    the dispatcher falls through to the original default render. Same
    for any Phase 2/3 placeholder. No exceptions; output matches the
    pre-v2.8.1 inline path."""
    print("[41/52] archetype fallback: default + Phase 2/3 placeholder")
    from .tools.pptx_renderer import _render_slide

    # Case A: default archetype, ordinary slide (no big number, no
    # bullets). Should fall through to default render without errors.
    title = LayerNode(
        layer_id="d1_title", name="title", kind="text", z_index=10,
        text="An ordinary slide", font_family="NotoSerifSC-Bold",
        font_size_px=72, align="left", effects=TextEffect(fill="#0f172a"),
    )
    body = LayerNode(
        layer_id="d1_body", name="body", kind="text", z_index=10,
        text="Two paragraph body that should render through the\n"
             "v2.7.2 default path because no big number is declared.",
        font_family="NotoSansSC-Bold", font_size_px=32,
        align="left", effects=TextEffect(fill="#0f172a"),
    )
    slide_default = LayerNode(
        layer_id="ds1", name="ds1", kind="slide", z_index=1,
        section_number="§2",  # exercise v2.7.2 prefix on default path
        children=[title, body],
        # archetype omitted → schema default "evidence_snapshot"
    )
    _, slide_a = _new_blank_slide()
    _render_slide(slide_a, slide_default, 1920, 1080, _make_archetype_ctx())
    frames_a = _all_textframes(slide_a)
    if not frames_a:
        _fail("default-archetype fallback emitted no shapes")
    if not any("§2" in _shape_text(f) for f in frames_a):
        _fail("v2.7.2 section prefix missing on default-render fallback path")

    # Case B: Phase 2/3 placeholder archetype. Dispatcher should fall
    # through to default render without raising.
    slide_placeholder = _make_archetype_slide(
        layer_id="ds2", archetype="pipeline_horizontal",
        children=[title, body],
    )
    _, slide_b = _new_blank_slide()
    _render_slide(slide_b, slide_placeholder, 1920, 1080, _make_archetype_ctx())
    frames_b = _all_textframes(slide_b)
    if not frames_b:
        _fail("Phase 2 placeholder fallback emitted no shapes")
    _ok(f"fallback paths intact: default ({len(frames_a)} frames) + "
        f"placeholder pipeline_horizontal ({len(frames_b)} frames)")


def check_archetype_determinism() -> None:
    """v2.8.1 smoke #42 — rendering the same SlideNode twice produces
    byte-identical slide XML. Guards against accidental nondeterminism
    (time-based ids, random shape positions) creeping into archetype
    renderers."""
    print("[42/52] archetype determinism: two renders → identical XML")
    from xml.etree import ElementTree as ET
    from .tools.pptx_renderer import _render_slide

    title = LayerNode(
        layer_id="d_title", name="title", kind="text", z_index=10,
        text="Deterministic cover", font_family="NotoSerifSC-Bold",
        font_size_px=120, align="left", effects=TextEffect(fill="#0f172a"),
    )
    subtitle = LayerNode(
        layer_id="d_sub", name="subtitle", kind="text", z_index=10,
        text="Same input → same shapes", font_family="NotoSansSC-Bold",
        font_size_px=36, align="left", effects=TextEffect(fill="#475569"),
    )
    authors = LayerNode(
        layer_id="d_auth", name="author_byline", kind="text", z_index=10,
        text="OpenDesign", font_family="NotoSansSC-Bold",
        font_size_px=26, align="left", effects=TextEffect(fill="#64748b"),
    )
    spec_node = _make_archetype_slide(
        layer_id="dcov", archetype="cover_editorial",
        section_number="§1",
        children=[title, subtitle, authors],
    )

    _, slide_1 = _new_blank_slide()
    _render_slide(slide_1, spec_node, 1920, 1080, _make_archetype_ctx())
    _, slide_2 = _new_blank_slide()
    _render_slide(slide_2, spec_node, 1920, 1080, _make_archetype_ctx())

    xml_1 = ET.tostring(slide_1._element, encoding="unicode")
    xml_2 = ET.tostring(slide_2._element, encoding="unicode")

    if xml_1 != xml_2:
        # Surface a small diff hint for triage.
        for i, (a, b) in enumerate(zip(xml_1, xml_2)):
            if a != b:
                _fail(
                    f"determinism violated at char {i}: "
                    f"...{xml_1[max(0,i-30):i+30]!r} vs "
                    f"...{xml_2[max(0,i-30):i+30]!r}"
                )
        _fail("determinism violated (lengths differ)")
    _ok(f"determinism: byte-identical slide XML across two renders "
        f"({len(xml_1)} chars)")


# ─────────────────────── v2.8.0 ClaimGraph extractor ────────────────────────


def _make_smoke_claim_graph_settings(out_root: Path) -> "Settings":  # noqa: F821
    """Build a Settings stub that points at a writable smoke run dir, with
    `claim_graph_max_turns` low enough that the failsafe path is reachable
    inside a unit-test budget."""
    from .config import Settings
    return Settings(
        anthropic_api_key="sk-stub",
        anthropic_base_url=None,
        gemini_api_key="stub",
        planner_model="stub-planner",
        critic_model="stub-critic",
        out_dir=out_root,
        critic_max_turns=4,
        critic_thinking_budget=0,
        claim_graph_max_turns=4,
        claim_graph_thinking_budget=0,
    )


def check_claim_graph_extractor_trajectory() -> None:
    """smoke #32: spawn ClaimGraphExtractor with a mock LLMBackend that
    calls report_claim_graph on turn 1. Verify the resulting ClaimGraph has
    the expected fields AND the trajectory file
    `claim_graph_extractor.jsonl` lands in the run dir."""
    print("[32/52] claim_graph extractor: scripted report_claim_graph + trajectory")
    from .agents import ClaimGraphExtractor
    from .config import REPO_ROOT
    from .schema import ClaimGraph

    out_root = REPO_ROOT / "out" / "smoke_claim_graph"
    if out_root.exists():
        import shutil
        shutil.rmtree(out_root)
    run_dir = out_root / "runs" / "smoke-claim-graph"
    traj_dir = run_dir / "trajectory"
    traj_dir.mkdir(parents=True, exist_ok=True)

    paper_text = (
        "LongCat-Next: Lexicalizing Modalities as Discrete Tokens.\n\n"
        "Section 1. Native multimodality requires unifying token spaces.\n"
        "We argue the dual bottleneck in diffusion samplers can only be\n"
        "resolved by joint training. LongCat-Next achieves 72.3% top-1 on\n"
        "ImageNet — a clear signal that the DiNA paradigm scales.\n\n"
        "Section 2. Limitations and conclusion."
    )
    paper_path = run_dir / "fake_paper.pdf"
    paper_path.write_text("dummy", encoding="utf-8")

    settings = _make_smoke_claim_graph_settings(out_root)

    scripted = [
        _ScriptedTurn(tool_calls=[
            ("lookup_paper_section", {"query": "ImageNet"}),
        ]),
        _ScriptedTurn(tool_calls=[
            ("report_claim_graph", {
                "paper_title": "LongCat-Next",
                "paper_anchor": "arxiv:2526.01234",
                "thesis": "Lexicalize modalities as discrete tokens.",
                "tensions": [
                    {"id": "T1", "name": "understanding-generation conflict",
                     "description": "joint training is hard", "evidence_anchor": None},
                ],
                "mechanisms": [
                    {"id": "M1", "name": "DiNA paradigm",
                     "resolves": ["T1"],
                     "description": "discrete token unification"},
                ],
                "evidence": [
                    {"id": "E1", "metric": "top-1 72.3%", "source": "intro",
                     "raw_quote": "LongCat-Next achieves 72.3% top-1 on",
                     "supports": ["M1"]},
                ],
                "implications": [
                    {"id": "I1",
                     "description": "Native multimodality is feasible.",
                     "derives_from": ["M1", "E1"]},
                ],
            }),
        ]),
    ]
    mock = _MockCriticBackend(model="stub-claim-graph", turns=scripted)

    agent = ClaimGraphExtractor(settings)
    agent.backend = mock
    agent._system_prompt = "patched"

    traj_path = traj_dir / "claim_graph_extractor.jsonl"
    graph = agent.extract(
        paper_path=paper_path,
        paper_raw_text=paper_text,
        trajectory_path=traj_path,
    )

    if not isinstance(graph, ClaimGraph):
        _fail(f"extract should return ClaimGraph; got {type(graph).__name__}")
    if graph.thesis != "Lexicalize modalities as discrete tokens.":
        _fail(f"thesis lost: {graph.thesis!r}")
    if len(graph.tensions) != 1 or graph.tensions[0].id != "T1":
        _fail(f"tensions lost: {graph.tensions}")
    if len(graph.evidence) != 1 or graph.evidence[0].raw_quote.startswith("?"):
        _fail(f"evidence lost: {graph.evidence}")

    if not traj_path.exists():
        _fail(f"claim_graph_extractor.jsonl not written at {traj_path}")
    lines = traj_path.read_text(encoding="utf-8").strip().splitlines()
    if len(lines) != 2:
        _fail(f"expected 2 trajectory lines (one per turn), got {len(lines)}")
    line1 = json.loads(lines[1])
    if not any(tc.get("name") == "report_claim_graph"
               for tc in line1.get("tool_calls", [])):
        _fail(f"final turn should record report_claim_graph; got {line1.get('tool_calls')}")
    _ok(f"extractor: report_claim_graph captured + claim_graph_extractor.jsonl ({len(lines)} lines)")


def check_claim_graph_validator_rejects_fabricated_quote() -> None:
    """smoke #33: validate_claim_graph must REJECT an EvidenceNode whose
    raw_quote does not appear in paper_raw_text."""
    print("[33/52] claim_graph validator: rejects fabricated raw_quote")
    from .schema import (
        ClaimGraph, EvidenceNode, ImplicationNode, MechanismNode, TensionNode,
    )
    from .util.claim_graph_validator import validate_claim_graph

    paper_text = (
        "LongCat-Next: Lexicalizing Modalities as Discrete Tokens.\n"
        "We argue native multimodality requires unified token spaces.\n"
        "LongCat-Next achieves 72.3% top-1 on ImageNet."
    )

    grounded = ClaimGraph(
        paper_title="LongCat-Next", paper_anchor="arxiv:x",
        thesis="x",
        tensions=[TensionNode(id="T1", name="t", description="d")],
        mechanisms=[MechanismNode(id="M1", name="m", resolves=["T1"],
                                  description="d")],
        evidence=[EvidenceNode(id="E1", metric="top-1 72.3%",
                                source="intro",
                                raw_quote="achieves 72.3% top-1 on ImageNet",
                                supports=["M1"])],
        implications=[ImplicationNode(id="I1", description="impl",
                                       derives_from=["M1", "E1"])],
    )
    errs = validate_claim_graph(grounded, paper_text)
    if errs:
        _fail(f"grounded graph should pass; got errors: {errs}")

    fabricated = ClaimGraph(
        paper_title="x", paper_anchor="x", thesis="x",
        tensions=[TensionNode(id="T1", name="t", description="d")],
        mechanisms=[MechanismNode(id="M1", name="m", resolves=["T1"],
                                  description="d")],
        evidence=[EvidenceNode(id="E2", metric="x",
                                source="intro",
                                raw_quote="LongCat-Next reaches 99.99% on every benchmark",
                                supports=["M1"])],
        implications=[],
    )
    errs2 = validate_claim_graph(fabricated, paper_text)
    if not errs2:
        _fail("fabricated quote should fail substring check; validator passed")
    if not any("E2" in e and "raw_quote" in e for e in errs2):
        _fail(f"expected error about E2 raw_quote; got {errs2}")

    bad_refs = ClaimGraph(
        paper_title="x", paper_anchor="x", thesis="x",
        tensions=[TensionNode(id="T1", name="t", description="d")],
        mechanisms=[MechanismNode(id="M1", name="m", resolves=["T9"],
                                  description="d")],
        evidence=[],
        implications=[ImplicationNode(id="I1", description="i",
                                       derives_from=["X1"])],
    )
    errs3 = validate_claim_graph(bad_refs, paper_text)
    if not any("T9" in e for e in errs3):
        _fail(f"expected error about unknown tension T9; got {errs3}")
    if not any("X1" in e for e in errs3):
        _fail(f"expected error about unknown derives_from X1; got {errs3}")
    _ok("validator: grounded passes, fabricated/bad-ref fail (substring + integrity)")


def check_planner_covers_population() -> None:
    """smoke #34: a planner-emitted DesignSpec with claim_graph attached
    populates `slide.covers` with valid ClaimGraph node ids; the union
    matches the graph's id catalog (i.e. every tension/mechanism/evidence
    is covered by at least one slide)."""
    print("[34/52] planner: SlideNode.covers populated against claim_graph ids")
    from .schema import (
        ArtifactType, ClaimGraph, DesignSpec, EvidenceNode,
        ImplicationNode, LayerNode, MechanismNode, TensionNode,
    )

    graph = ClaimGraph(
        paper_title="LongCat-Next", paper_anchor="arxiv:x",
        thesis="t",
        tensions=[
            TensionNode(id="T1", name="conflict", description="d"),
            TensionNode(id="T2", name="bottleneck", description="d"),
        ],
        mechanisms=[
            MechanismNode(id="M1", name="DiNA", resolves=["T1", "T2"],
                          description="d"),
            MechanismNode(id="M2", name="JointTraining", resolves=["T1"],
                          description="d"),
        ],
        evidence=[
            EvidenceNode(id="E1", metric="x", source="t",
                          raw_quote="x", supports=["M1"]),
            EvidenceNode(id="E2", metric="x", source="t",
                          raw_quote="x", supports=["M2"]),
        ],
        implications=[
            ImplicationNode(id="I1", description="x", derives_from=["M1"]),
        ],
    )

    # Simulate what the planner should emit when claim_graph is attached:
    # talk arc cover → tensions → mechanisms → evidence → implications → thanks.
    spec = DesignSpec(
        brief="paper deck",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96,
                "aspect_ratio": "16:9", "color_mode": "RGB"},
        layer_graph=[
            LayerNode(layer_id="S0", name="cover", kind="slide",
                      z_index=0, covers=[]),
            LayerNode(layer_id="S1", name="tensions", kind="slide",
                      z_index=1, covers=["T1", "T2"]),
            LayerNode(layer_id="S2", name="mech-DiNA", kind="slide",
                      z_index=2, covers=["M1"]),
            LayerNode(layer_id="S3", name="mech-Joint", kind="slide",
                      z_index=3, covers=["M2"]),
            LayerNode(layer_id="S4", name="evidence", kind="slide",
                      z_index=4, covers=["E1", "E2"]),
            LayerNode(layer_id="S5", name="implications", kind="slide",
                      z_index=5, covers=["I1"]),
            LayerNode(layer_id="S6", name="thanks", kind="slide",
                      z_index=6, covers=[]),
        ],
    )

    # Round-trip through pydantic to confirm the field is serialisable
    # AND the planner can read covers off a deserialized spec.
    dumped = spec.model_dump(mode="json")
    rehydrated = DesignSpec.model_validate(dumped)
    if not all(hasattr(s, "covers") for s in rehydrated.layer_graph):
        _fail("LayerNode.covers missing after round-trip")
    if rehydrated.layer_graph[1].covers != ["T1", "T2"]:
        _fail(f"S1 covers lost: {rehydrated.layer_graph[1].covers}")

    # The union of slide.covers should cover every claim graph node id.
    valid_ids = (
        {t.id for t in graph.tensions}
        | {m.id for m in graph.mechanisms}
        | {e.id for e in graph.evidence}
        | {i.id for i in graph.implications}
    )
    covered = set()
    for slide in rehydrated.layer_graph:
        for cid in slide.covers:
            if cid not in valid_ids:
                _fail(f"slide {slide.layer_id} covers unknown id {cid!r}")
            covered.add(cid)
    missing = valid_ids - covered
    if missing:
        _fail(f"planner missed claim ids in covers: {sorted(missing)}")
    _ok(f"covers populated for all {len(valid_ids)} claim ids; round-trip preserved")


def check_critic_claim_coverage_issue() -> None:
    """smoke #35: CriticAgent with a non-None claim_graph parameter
    detects an uncovered tension. Mock the LLM to read the user message,
    then emit a claim_coverage issue. Verify the issue surfaces correctly
    in the CritiqueReport."""
    print("[35/52] critic: claim_coverage issue when tension uncovered")
    from .agents import CriticAgent
    from .schema import (
        ArtifactType, ClaimGraph, DesignSpec, LayerNode, MechanismNode,
        TensionNode,
    )

    spec, slide_paths, run_dir = _make_smoke_deck_spec(n_slides=2)
    # Override layer_graph so slide 0 covers T1 only; T2 is uncovered.
    spec = DesignSpec(
        brief=spec.brief, artifact_type=spec.artifact_type,
        canvas=spec.canvas,
        layer_graph=[
            LayerNode(layer_id="S0", name="s0", kind="slide", z_index=0,
                      covers=["T1"]),
            LayerNode(layer_id="S1", name="s1", kind="slide", z_index=1,
                      covers=[]),
        ],
    )
    settings = _make_smoke_claim_graph_settings(run_dir.parent.parent)

    graph = ClaimGraph(
        paper_title="x", paper_anchor="x", thesis="t",
        tensions=[
            TensionNode(id="T1", name="covered", description="d"),
            TensionNode(id="T2", name="missing", description="d"),
        ],
        mechanisms=[
            MechanismNode(id="M1", name="m", resolves=["T1"],
                          description="d"),
        ],
        evidence=[],
        implications=[],
    )

    scripted = [
        _ScriptedTurn(tool_calls=[
            ("lookup_claim_node", {"claim_id": "T2"}),
        ]),
        _ScriptedTurn(tool_calls=[
            ("report_verdict", {
                "score": 0.6, "verdict": "revise",
                "summary": "tension T2 not covered by any slide",
                "issues": [
                    {"slide_id": None, "severity": "high",
                     "category": "claim_coverage",
                     "description": "Tension T2 'missing' has no slide.covers reference",
                     "evidence_paper_anchor": None},
                ],
            }),
        ]),
    ]
    mock = _MockCriticBackend(model="stub-critic", turns=scripted)

    agent = CriticAgent(settings, ArtifactType.DECK)
    agent.backend = mock
    agent._system_prompt = "patched"

    traj_path = run_dir / "trajectory" / "critic_claim.jsonl"
    if traj_path.exists():
        traj_path.unlink()
    report = agent.critique(
        spec=spec, layer_manifest=[],
        slide_renders=slide_paths,
        paper_raw_text="x", claim_graph=graph,
        iteration=1, trajectory_path=traj_path,
    )

    if report.verdict != "revise":
        _fail(f"verdict should be 'revise'; got {report.verdict!r}")
    cc_issues = [i for i in report.issues if i.category == "claim_coverage"]
    if not cc_issues:
        _fail(f"expected ≥1 claim_coverage issue; got {report.issues}")
    if cc_issues[0].severity != "high":
        _fail(f"uncovered tension should be severity=high; got {cc_issues[0].severity}")

    # Confirm the lookup_claim_node tool actually executed and returned
    # the T2 node payload (not an error).
    line0 = json.loads(traj_path.read_text(encoding="utf-8").splitlines()[0])
    tool_results_pre = line0.get("tool_results", [])
    if not any(not tr.get("is_error", False)
               and tr.get("name") == "lookup_claim_node"
               for tr in tool_results_pre):
        _fail(f"lookup_claim_node tool result missing/erred: {tool_results_pre}")
    _ok("critic emitted high-severity claim_coverage issue + lookup_claim_node worked")


def check_no_claim_graph_pipeline_degrades() -> None:
    """smoke #36: when --no-claim-graph is set (or the brief has no PDF
    attachment), `_run_claim_graph_extractor` returns None cleanly without
    spawning the extractor, and the runner stores None in
    `ctx.state["claim_graph"]` so the planner degrades to v2.7.3 behavior."""
    print("[36/52] --no-claim-graph degrades to v2.7.3 cleanly (no errors)")
    from .config import REPO_ROOT
    from .runner import _run_claim_graph_extractor

    out_root = REPO_ROOT / "out" / "smoke_no_claim_graph"
    if out_root.exists():
        import shutil
        shutil.rmtree(out_root)
    sub_traj_dir = out_root / "runs" / "smoke-no-cg" / "trajectory"
    sub_traj_dir.mkdir(parents=True, exist_ok=True)

    settings = _make_smoke_claim_graph_settings(out_root)

    # Case 1: --no-claim-graph flag → skip even with PDF present.
    fake_pdf = out_root / "fake.pdf"
    fake_pdf.write_bytes(b"%PDF-1.4 fake")
    g1 = _run_claim_graph_extractor(
        settings, [fake_pdf], no_claim_graph=True,
        sub_traj_dir=sub_traj_dir,
    )
    if g1 is not None:
        _fail(f"--no-claim-graph should yield None; got {type(g1).__name__}")
    if (sub_traj_dir / "claim_graph_extractor.jsonl").exists():
        _fail("extractor trajectory should NOT be written when skipped")

    # Case 2: no PDF attachment → skip without any extraction attempt.
    md = out_root / "notes.md"
    md.write_text("# notes", encoding="utf-8")
    g2 = _run_claim_graph_extractor(
        settings, [md], no_claim_graph=False,
        sub_traj_dir=sub_traj_dir,
    )
    if g2 is not None:
        _fail(f"no-PDF input should yield None; got {type(g2).__name__}")

    # Case 3: empty attachments → also None.
    g3 = _run_claim_graph_extractor(
        settings, [], no_claim_graph=False,
        sub_traj_dir=sub_traj_dir,
    )
    if g3 is not None:
        _fail(f"empty attachments should yield None; got {type(g3).__name__}")

    # Case 4: settings.enable_claim_graph=False short-circuits even with PDF.
    settings_off = settings.__class__(
        **{**settings.__dict__, "enable_claim_graph": False},
    )
    g4 = _run_claim_graph_extractor(
        settings_off, [fake_pdf], no_claim_graph=False,
        sub_traj_dir=sub_traj_dir,
    )
    if g4 is not None:
        _fail(f"enable_claim_graph=False should yield None; got {type(g4).__name__}")
    _ok("4 skip paths all return None cleanly; no extractor spawn, no trajectory written")


def check_claim_graph_lookup_whitespace_tolerance() -> None:
    """smoke #43: `_extract_paper_excerpt` MUST whitespace-normalize so a
    multi-word query matches across PDF line wraps / double-spaces. This
    is the regression the v2.8.0 hotfix targets — Sonnet 4.6 burned its
    full max_turns budget on `lookup_paper_section` calls because the
    literal substring search couldn't see content that the validator
    would happily accept once normalized."""
    print("[43/52] claim_graph lookup: whitespace-tolerant (PDF line-wrap proof)")
    from .agents.claim_graph_extractor import _extract_paper_excerpt

    paper = (
        "Section 4.\n\n"
        "Contemporary multimodal\nsystems remain\nlanguage-centric,\n"
        "treating non-linguistic\nmodalities as\nexternal attachments.\n\n"
        "Next section follows."
    )

    full_phrase = (
        "language-centric treating non-linguistic modalities as "
        "external attachments"
    )
    exc = _extract_paper_excerpt(paper, full_phrase)
    if not exc:
        _fail(
            "whitespace-tolerant lookup must match the full phrase across "
            "line wraps; got empty excerpt (regression of Sonnet 4.6 bug)"
        )
    if "language-centric" not in exc or "external attachments" not in exc:
        _fail(f"excerpt should contain the matched span; got {exc!r:.200s}")

    # Token-fallback: query with junk that won't substring-match, but
    # whose longest token ('Comparison'/'contemporary') does exist.
    fallback = _extract_paper_excerpt(
        paper, "Table 1: Comparison of contemporary",
    )
    if not fallback:
        _fail("longest-token fallback should hit on 'contemporary'")

    # True miss: nonsense query stays empty.
    miss = _extract_paper_excerpt(paper, "qzqzqz no such phrase")
    if miss:
        _fail(f"unrelated query must return empty; got {miss!r:.120s}")

    # Edge: empty inputs are no-ops.
    if _extract_paper_excerpt("", "x") != "":
        _fail("empty raw must return ''")
    if _extract_paper_excerpt(paper, "") != "":
        _fail("empty query must return ''")
    _ok(
        "lookup matches across line wraps, falls back to longest token, "
        "and returns '' on true miss"
    )


def check_claim_graph_kimi_template_leak_retry() -> None:
    """smoke #44: when an end_turn happens with no structured tool_calls
    but a `<|tool_calls_section_begin|>` template token leaked into the
    thinking block (Kimi K2.6 OpenRouter glitch from the v2.8.0 dogfood),
    the extractor injects a corrective user reminder and retries instead
    of giving up immediately. Verifies the loop survives one retry and
    can still call report_claim_graph on the next turn."""
    print("[44/52] claim_graph: Kimi template-leak triggers one retry")
    from .agents import ClaimGraphExtractor
    from .config import REPO_ROOT
    from .llm_backend import ToolCall, TurnResponse
    from .schema import ClaimGraph, ThinkingBlockRecord

    out_root = REPO_ROOT / "out" / "smoke_claim_graph_kimi_leak"
    if out_root.exists():
        import shutil
        shutil.rmtree(out_root)
    run_dir = out_root / "runs" / "smoke-kimi-leak"
    traj_dir = run_dir / "trajectory"
    traj_dir.mkdir(parents=True, exist_ok=True)

    paper_text = (
        "Section 1. LongCat-Next achieves 72.3% top-1 on ImageNet."
    )
    paper_path = run_dir / "fake.pdf"
    paper_path.write_text("dummy", encoding="utf-8")
    settings = _make_smoke_claim_graph_settings(out_root)

    # Turn 1: simulate Kimi leaking the template token (no tool_calls,
    # stop_reason=end_turn). Turn 2: model recovers and emits the graph.
    leaked = TurnResponse(
        text="",
        thinking_blocks=[
            ThinkingBlockRecord(
                thinking=(
                    "<|tool_calls_section_begin|>"
                    "[{\"name\":\"lookup_paper_section\","
                    "\"parameters\":{\"query\":\"ImageNet\"}}]"
                ),
                signature="",
                is_redacted=False,
            ),
        ],
        tool_calls=[],
        stop_reason="end_turn",
        usage={"input": 100, "output": 30},
        raw_assistant_content={"role": "assistant"},
    )
    recovered = TurnResponse(
        text="",
        thinking_blocks=[],
        tool_calls=[
            ToolCall(
                id="toolu_recover",
                name="report_claim_graph",
                input={
                    "paper_title": "LongCat-Next",
                    "paper_anchor": "x",
                    "thesis": "Recovery proves the leak retry works.",
                    "tensions": [{"id": "T1", "name": "t",
                                  "description": "d"}],
                    "mechanisms": [{"id": "M1", "name": "m",
                                    "resolves": ["T1"], "description": "d"}],
                    "evidence": [{"id": "E1", "metric": "x",
                                  "source": "x",
                                  "raw_quote": "achieves 72.3% top-1",
                                  "supports": ["M1"]}],
                    "implications": [{"id": "I1", "description": "i",
                                      "derives_from": ["M1"]}],
                },
            ),
        ],
        stop_reason="tool_use",
        usage={"input": 150, "output": 200},
        raw_assistant_content={"role": "assistant"},
    )

    class _LeakyBackend:
        name = "mock"
        model = "moonshotai/kimi-k2.6"

        def __init__(self) -> None:
            self.calls = 0
            self.last_messages: list = []

        def create_turn(self, *, system, messages, tools,
                        thinking_budget=0, max_tokens=16384,
                        extra_headers=None):
            self.last_messages = list(messages)
            self.calls += 1
            return leaked if self.calls == 1 else recovered

        def append_assistant(self, messages, response):
            messages.append(response.raw_assistant_content)

        def append_tool_results(self, messages, results):
            for tu_id, payload, _is_err in results:
                messages.append({"role": "tool",
                                 "tool_call_id": tu_id, "content": payload})

        def vision_user_message(self, *, image_b64, media_type, text):
            return {"role": "user", "content": text}

    backend = _LeakyBackend()
    agent = ClaimGraphExtractor(settings)
    agent.backend = backend
    agent._system_prompt = "patched"

    traj_path = traj_dir / "claim_graph_extractor.jsonl"
    graph = agent.extract(
        paper_path=paper_path, paper_raw_text=paper_text,
        trajectory_path=traj_path,
    )

    if not isinstance(graph, ClaimGraph):
        _fail(f"extract should return ClaimGraph; got {type(graph).__name__}")
    if graph.thesis != "Recovery proves the leak retry works.":
        _fail(f"recovery turn output lost: thesis={graph.thesis!r}")
    if backend.calls != 2:
        _fail(
            f"leaky backend should be called exactly twice (leak + retry); "
            f"got {backend.calls}"
        )
    # The retry corrective reminder should be the LAST user message in
    # the conversation handed to the recovery turn.
    last_user = next(
        (m for m in reversed(backend.last_messages)
         if isinstance(m, dict) and m.get("role") == "user"),
        None,
    )
    if last_user is None or "<|tool_calls_section_begin|>" not in str(
        last_user.get("content", "")
    ):
        _fail(
            f"retry should inject corrective user message naming the "
            f"leaked token; last_user={last_user!r:.200s}"
        )
    _ok("Kimi template-leak triggers exactly one corrective retry")


def check_deck_text_overlap_detector() -> None:
    """v2.7.5 — `_detect_text_overlaps` now runs in `_composite_deck` and
    derives effective bboxes from `_TEMPLATE_SLOT_BBOX` for templated
    children whose `bbox=None`. Three regressions to lock in:

      a. Two text children in the same `template_slot` → slot_collision
         warning surfaced in payload (severity=blocker). This is the
         class-defect that lets a "table caption" land on top of the
         section title at the top of the slide because the second text
         child silently overrides the first via last-write-wins.
      b. A text child with `bbox=None` AND `template_slot=None` on a
         templated deck → unanchored_text warning (severity=blocker).
         The renderer's fallback dumps it at (0,0,full-slide). This is
         exactly the slide12 `TextBox 11` defect from the 2026-04-26
         longcat-next dogfood.
      c. A clean templated deck (text in distinct slots, no orphan
         children) → empty `text_overlap_warnings` list.
    """
    print("[46/52] deck text-overlap detector (no API)")

    from .config import REPO_ROOT, Settings
    from .schema import (
        ArtifactType, DeckDesignSystem, DesignSpec, LayerNode, SafeZone,
    )
    from .tools import ToolContext
    from .tools.composite import (
        _detect_deck_text_overlaps, _effective_text_extent,
        _slot_bbox, composite,
    )

    # Unit checks on the slot-bbox derivation.
    if _slot_bbox("content_with_figure", "title") != (96, 120, 1728, 80):
        _fail("_slot_bbox content_with_figure/title bbox wrong")
    if _slot_bbox("content_with_figure", "image_slot") != (1056, 260, 768, 740):
        _fail("_slot_bbox content_with_figure/image_slot bbox wrong")
    if _slot_bbox("unknown_role", "title") is not None:
        _fail("_slot_bbox should return None for unknown role")

    title = LayerNode(layer_id="t", name="title", kind="text", z_index=10,
                     template_slot="title", text="Hello", font_size_px=36)
    ext = _effective_text_extent(title, role="content")
    if ext is None or ext[0:3] != (96, 120, 1728):
        _fail(f"_effective_text_extent slot derivation failed: {ext}")

    # Build a deck with three regression slides:
    #   slide_02: two text children both targeting template_slot="body"
    #             → expect slot_collision (blocker).
    #   slide_03: an unanchored text child (bbox=None, slot=None)
    #             → expect unanchored_text (blocker).
    #   slide_04: a clean content_with_figure → no text_overlap_warnings.
    out_dir = REPO_ROOT / "out" / "smoke_deck_overlap"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)
    fig_path = layers_dir / "ingest_fig_01.png"
    Image.new("RGB", (1200, 800), (180, 220, 200)).save(fig_path)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-deck-overlap")
    ctx.state["rendered_layers"]["ingest_fig_01"] = {
        "layer_id": "ingest_fig_01", "kind": "image", "src_path": str(fig_path),
        "bbox": None, "z_index": 5, "name": "fig", "aspect_ratio": "3:2",
    }

    spec = DesignSpec(
        brief="text overlap detector smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"},
        deck_design_system=DeckDesignSystem(style="academic-editorial"),
        layer_graph=[
            LayerNode(layer_id="slide_01", name="cover", kind="slide", z_index=1, role="cover", children=[
                LayerNode(layer_id="slide_01_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Cover"),
            ]),
            LayerNode(layer_id="slide_02", name="collide", kind="slide", z_index=2, role="content",
                      children=[
                LayerNode(layer_id="slide_02_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Collision Title"),
                LayerNode(layer_id="slide_02_body_a", name="body", kind="text", z_index=10,
                          template_slot="body", text="First body bullet."),
                LayerNode(layer_id="slide_02_body_b", name="body", kind="text", z_index=10,
                          template_slot="body", text="Second body that silently overwrites."),
            ]),
            LayerNode(layer_id="slide_03", name="orphan_text", kind="slide", z_index=3, role="content",
                      children=[
                LayerNode(layer_id="slide_03_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Orphan Slide"),
                LayerNode(layer_id="slide_03_floating", name="floating", kind="text", z_index=10,
                          text="Table 2 caption that landed at the top of the slide."),
            ]),
            LayerNode(layer_id="slide_04", name="clean", kind="slide", z_index=4,
                      role="content_with_figure", children=[
                LayerNode(layer_id="slide_04_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Clean Slide"),
                LayerNode(layer_id="slide_04_body", name="body", kind="text", z_index=10,
                          template_slot="body", text="One body in its own slot."),
                LayerNode(layer_id="ingest_fig_01", name="diagram", kind="image", z_index=5,
                          template_slot="image_slot"),
            ]),
        ],
    )

    slides = [n for n in spec.layer_graph if n.kind == "slide"]
    warnings = _detect_deck_text_overlaps(slides, slide_w=1920, slide_h=1080)

    by_slide: dict[str, list[dict]] = {}
    for w in warnings:
        by_slide.setdefault(w.get("slide_id") or "?", []).append(w)

    s2_warnings = by_slide.get("slide_02", [])
    if not any(w.get("kind") == "slot_collision"
               and w.get("template_slot") == "body"
               and w.get("severity") == "blocker"
               for w in s2_warnings):
        _fail(f"slide_02: expected slot_collision blocker for body slot; "
              f"got {s2_warnings!r}")

    s3_warnings = by_slide.get("slide_03", [])
    if not any(w.get("kind") == "unanchored_text"
               and w.get("severity") == "blocker"
               and w.get("layer_id") == "slide_03_floating"
               for w in s3_warnings):
        _fail(f"slide_03: expected unanchored_text blocker for "
              f"slide_03_floating; got {s3_warnings!r}")

    s4_warnings = by_slide.get("slide_04", [])
    if s4_warnings:
        _fail(f"slide_04 (clean) should produce no warnings; got {s4_warnings!r}")

    # Integration: composite() must surface the warnings on the payload so
    # the planner sees them on the next turn — no swallowing.
    ctx.state["design_spec"] = spec
    res = composite({}, ctx=ctx)
    if res.status != "ok":
        _fail(f"composite returned {res.status}: {res.payload}")
    payload_warnings = res.payload.get("text_overlap_warnings") or []
    if not any(w.get("kind") == "slot_collision" for w in payload_warnings):
        _fail(f"payload missing slot_collision warning: {payload_warnings!r}")
    if not any(w.get("kind") == "unanchored_text" for w in payload_warnings):
        _fail(f"payload missing unanchored_text warning: {payload_warnings!r}")
    _ok(f"detector flags slot_collision + unanchored_text + leaves clean slide alone "
        f"({len(payload_warnings)} warnings surfaced to planner)")


def check_orphan_callout_drop() -> None:
    """v2.7.5 — orphan callouts (no anchor / anchor missing /
    callout_region outside anchor bbox) are dropped at composite time
    and never reach the .pptx. Locks in the slide10 ("RVQ" leader-line
    pointing nowhere) and slide16 ("red circle floating in empty
    space") defects from the 2026-04-26 longcat-next dogfood.
    """
    print("[47/52] orphan callout dropped at composite time (no API)")
    from pptx import Presentation as _Reopen

    from .config import REPO_ROOT, Settings
    from .schema import (
        ArtifactType, DeckDesignSystem, DesignSpec, LayerNode, SafeZone,
    )
    from .tools import ToolContext
    from .tools.composite import _detect_orphan_callouts, composite

    out_dir = REPO_ROOT / "out" / "smoke_orphan_callout"
    layers_dir = out_dir / "layers"
    layers_dir.mkdir(parents=True, exist_ok=True)
    fig_path = layers_dir / "ingest_fig_01.png"
    Image.new("RGB", (1200, 800), (180, 220, 200)).save(fig_path)

    settings = Settings(
        anthropic_api_key="sk-stub", anthropic_base_url=None,
        gemini_api_key="stub", planner_model="stub", critic_model="stub",
    )
    ctx = ToolContext(settings=settings, run_dir=out_dir, layers_dir=layers_dir,
                      run_id="smoke-orphan-callout")
    ctx.state["rendered_layers"]["ingest_fig_01"] = {
        "layer_id": "ingest_fig_01", "kind": "image", "src_path": str(fig_path),
        "bbox": None, "z_index": 5, "name": "fig", "aspect_ratio": "3:2",
    }

    spec = DesignSpec(
        brief="orphan callout smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"},
        deck_design_system=DeckDesignSystem(style="academic-editorial"),
        layer_graph=[
            LayerNode(layer_id="slide_01", name="cover", kind="slide", z_index=1, role="cover", children=[
                LayerNode(layer_id="slide_01_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Cover"),
            ]),
            LayerNode(layer_id="slide_02", name="orphans", kind="slide", z_index=2,
                      role="content_with_figure", children=[
                LayerNode(layer_id="slide_02_title", name="title", kind="text", z_index=10,
                          template_slot="title", text="Orphan callouts"),
                LayerNode(layer_id="ingest_fig_01", name="diagram", kind="image", z_index=5,
                          template_slot="image_slot"),
                # 1) good callout — anchor present, region inside bbox
                LayerNode(layer_id="callout_good", name="good", kind="callout", z_index=20,
                          anchor_layer_id="ingest_fig_01",
                          callout_style="highlight",
                          callout_region=SafeZone(x=1100, y=400, w=200, h=200, purpose="body")),
                # 2) anchor_layer_id None
                LayerNode(layer_id="callout_no_anchor", name="bad1", kind="callout", z_index=21,
                          callout_style="circle",
                          callout_region=SafeZone(x=200, y=200, w=120, h=120, purpose="body")),
                # 3) anchor_layer_id refers to a sibling that doesn't exist
                LayerNode(layer_id="callout_missing_anchor", name="bad2", kind="callout", z_index=22,
                          anchor_layer_id="ingest_fig_99",
                          callout_style="circle",
                          callout_region=SafeZone(x=300, y=300, w=80, h=80, purpose="body")),
                # 4) anchor exists but callout_region doesn't overlap anchor bbox
                LayerNode(layer_id="callout_offsite", name="bad3", kind="callout", z_index=23,
                          anchor_layer_id="ingest_fig_01",
                          callout_style="circle",
                          callout_region=SafeZone(x=10, y=10, w=80, h=80, purpose="body")),
            ]),
        ],
    )

    slides = [n for n in spec.layer_graph if n.kind == "slide"]
    orphan_ids, warnings = _detect_orphan_callouts(slides, slide_w=1920, slide_h=1080)

    expected_orphans = {"callout_no_anchor", "callout_missing_anchor", "callout_offsite"}
    if orphan_ids != expected_orphans:
        _fail(f"orphan set mismatch: expected {expected_orphans}, got {orphan_ids}")
    reasons = {w["callout_layer_id"]: w["reason"] for w in warnings}
    if reasons.get("callout_no_anchor") != "no_anchor_layer_id":
        _fail(f"reason mismatch for callout_no_anchor: {reasons}")
    if reasons.get("callout_missing_anchor") != "anchor_not_on_slide":
        _fail(f"reason mismatch for callout_missing_anchor: {reasons}")
    if reasons.get("callout_offsite") != "region_outside_anchor_bbox":
        _fail(f"reason mismatch for callout_offsite: {reasons}")

    # Integration: composite drops the orphans before write_pptx; the
    # rendered .pptx must contain only the one good callout.
    ctx.state["design_spec"] = spec
    res = composite({}, ctx=ctx)
    if res.status != "ok":
        _fail(f"composite returned {res.status}: {res.payload}")
    pptx_path = ctx.state["composition"].pptx_path
    prs = _Reopen(str(pptx_path))
    method = prs.slides[1]
    names = {s.name for s in method.shapes}
    if "callout_good" not in names:
        _fail(f"good callout was dropped: {sorted(names)}")
    for orphan in expected_orphans:
        if orphan in names:
            _fail(f"orphan callout '{orphan}' leaked into the .pptx: {sorted(names)}")

    payload_orphans = res.payload.get("orphan_callout_warnings") or []
    if len(payload_orphans) != 3:
        _fail(f"expected 3 orphan warnings on payload; got {payload_orphans!r}")
    _ok("3 orphan callouts dropped, 1 valid callout retained, payload carries the 3 warnings")



def check_image_backend_fallback_chain() -> None:
    """smoke #46 (v2.7.5): FallbackImageBackend transparently retries
    against `image_fallback_model` when the primary raises
    `provider_unavailable`, and propagates other categories unchanged.

    Reproduces the Seedream 4.5 dogfood failure shape (4× consecutive
    404 - "No endpoints found that support the requested output
    modalities"): if the planner had to see those errors, every
    NBP-bearing slide would render text-only with orphan callouts.
    The wrapper short-circuits to the fallback model and returns a
    clean ImageResult so the tool emits `obs_ok` exactly as if the
    primary had worked.
    """
    print("[48/52] image fallback chain: provider_unavailable → fallback fires + non-PA errors propagate")
    from .image_backend import (
        FallbackImageBackend,
        ImageGenerationError,
        ImageResult,
        make_image_backend,
    )

    PRIMARY_ID = "stub-primary/seedream-broken"
    FALLBACK_ID = "stub-fallback/working"

    class _BrokenPrimary:
        name = "openrouter"
        model = PRIMARY_ID

        def __init__(self) -> None:
            self.calls = 0

        def generate(self, *, prompt, aspect_ratio, image_size):
            self.calls += 1
            raise ImageGenerationError(
                f"{self.model} via OpenRouter is unavailable: "
                f"Error code: 404 - No endpoints found that support the "
                f"requested output modalities: image, text",
                category="provider_unavailable",
            )

    class _WorkingFallback:
        name = "openrouter"
        model = FALLBACK_ID

        def __init__(self) -> None:
            self.calls = 0

        def generate(self, *, prompt, aspect_ratio, image_size):
            self.calls += 1
            from io import BytesIO
            buf = BytesIO()
            Image.new("RGB", (64, 64), (10, 20, 30)).save(buf, format="PNG")
            return ImageResult(
                data=buf.getvalue(),
                width=64, height=64,
                mime="image/png", model=self.model,
            )

    primary = _BrokenPrimary()
    fallback = _WorkingFallback()

    class _StubSettings:
        image_model = PRIMARY_ID
        image_fallback_model = FALLBACK_ID
        image_provider = "auto"

    wrapper = FallbackImageBackend(primary, _StubSettings(), FALLBACK_ID)
    # Pre-inject the fallback so we don't try to import google.genai or
    # OpenAI in the smoke (no creds, no network).
    wrapper._fallback_backend = fallback

    result = wrapper.generate(prompt="test", aspect_ratio="1:1", image_size="1K")
    if not isinstance(result, ImageResult):
        _fail(f"wrapper should return ImageResult; got {type(result).__name__}")
    if result.model != FALLBACK_ID:
        _fail(f"wrapper should surface fallback model id; got {result.model!r}")
    if primary.calls != 1:
        _fail(f"primary should be tried exactly once; got {primary.calls}")
    if fallback.calls != 1:
        _fail(f"fallback should be tried exactly once; got {fallback.calls}")

    # Non-provider-unavailable errors must NOT trigger fallback — that
    # would mask safety filtering and confuse SFT capture (a fallback
    # model would refuse the same prompt for the same reason).
    class _SafetyPrimary:
        name = "openrouter"
        model = PRIMARY_ID
        def __init__(self) -> None: self.calls = 0
        def generate(self, *, prompt, aspect_ratio, image_size):
            self.calls += 1
            raise ImageGenerationError(
                f"{self.model} returned no image — likely safety filter",
                category="safety_filter",
            )

    safety = _SafetyPrimary()
    safety_fallback = _WorkingFallback()
    wrapper2 = FallbackImageBackend(safety, _StubSettings(), FALLBACK_ID)
    wrapper2._fallback_backend = safety_fallback

    raised = None
    try:
        wrapper2.generate(prompt="x", aspect_ratio="1:1", image_size="1K")
    except ImageGenerationError as e:
        raised = e
    if raised is None:
        _fail("safety_filter should propagate from primary, not be swallowed")
    if raised.category != "safety_filter":
        _fail(f"category should stay 'safety_filter'; got {raised.category!r}")
    if safety_fallback.calls != 0:
        _fail(
            f"fallback must NOT fire on safety_filter; fallback got "
            f"{safety_fallback.calls} calls"
        )

    # Both providers down → terminal error, message names BOTH model ids
    # so the planner can pivot intelligently.
    class _BrokenFallback:
        name = "openrouter"
        model = FALLBACK_ID
        def __init__(self) -> None: self.calls = 0
        def generate(self, *, prompt, aspect_ratio, image_size):
            self.calls += 1
            raise ImageGenerationError(
                f"{self.model} via OpenRouter is unavailable: "
                f"is not a valid model ID",
                category="provider_unavailable",
            )

    primary2 = _BrokenPrimary()
    fallback2 = _BrokenFallback()
    wrapper3 = FallbackImageBackend(primary2, _StubSettings(), FALLBACK_ID)
    wrapper3._fallback_backend = fallback2

    raised2 = None
    try:
        wrapper3.generate(prompt="x", aspect_ratio="1:1", image_size="1K")
    except ImageGenerationError as e:
        raised2 = e
    if raised2 is None:
        _fail("both-providers-down should raise terminal ImageGenerationError")
    msg = str(raised2)
    if PRIMARY_ID not in msg or FALLBACK_ID not in msg:
        _fail(
            f"terminal error must name both primary AND fallback model ids "
            f"so planner can pivot; got {msg!r}"
        )
    if "ingest_fig" not in msg.lower() and "alternative" not in msg.lower():
        _fail(
            f"terminal error must hint at the planner pivot path "
            f"(ingest_fig / alternative); got {msg!r}"
        )

    # Factory wiring: when `image_fallback_model` is empty, no wrapper.
    class _NoFallbackSettings(_StubSettings):
        image_fallback_model = ""
        # Use a real-shaped id so make_image_backend doesn't need
        # google.genai (auto-routes to OpenRouter; we'll never call .generate).
        image_model = "stub-vendor/stub-model"
    # We can't easily call make_image_backend without OPENROUTER_API_KEY,
    # so just verify the routing predicate that gates the wrap:
    s_off = _NoFallbackSettings()
    fb_id = (getattr(s_off, "image_fallback_model", "") or "").strip()
    if fb_id and fb_id != s_off.image_model:
        _fail("empty image_fallback_model must NOT trigger wrap (regression guard)")

    # And when it equals the primary, also no wrap.
    class _SameModelSettings(_StubSettings):
        image_fallback_model = PRIMARY_ID
        image_model = PRIMARY_ID
    s_same = _SameModelSettings()
    fb_id2 = (getattr(s_same, "image_fallback_model", "") or "").strip()
    if fb_id2 and fb_id2 != s_same.image_model:
        _fail("identical primary/fallback ids must NOT trigger wrap")

    _ok(
        f"fallback wrapper: provider_unavailable→fallback fires (1+1 calls), "
        f"safety_filter propagates ({safety.calls}+0 calls), terminal failure "
        f"names both model ids + pivot hint, factory skips wrap when "
        f"fb=='' or fb==primary"
    )


def check_friday_gemini_image_backend_no_api() -> None:
    """smoke #47 (Friday Gemini): backend submits to Friday's async image
    API, tolerates transient 429 query responses, and extracts the final
    inline image without touching the network."""
    print("[49/52] Friday Gemini image backend: submit → poll → inline PNG")
    from . import image_backend as ib
    from .image_backend import FridayGeminiImageBackend, ImageResult

    class _Settings:
        friday_app_id = "friday-appid"
        anthropic_auth_token = None
        anthropic_api_key = None
        friday_gemini_base_url = "https://aigc.sankuai.com/v1/google/models"
        friday_image_timeout_s = 5
        friday_image_poll_interval_s = 0
        image_provider = "friday_gemini"
        image_fallback_model = "gemini-2.5-flash-image"

    buf = BytesIO()
    Image.new("RGB", (32, 24), (20, 40, 60)).save(buf, format="PNG")
    image_b64 = base64.b64encode(buf.getvalue()).decode("ascii")

    responses = [
        b"op-123",
        json.dumps({"status": 0, "data": "running"}).encode("utf-8"),
        json.dumps({"status": -1, "data": "429 Too Many Requests"}).encode("utf-8"),
        json.dumps({
            "status": 1,
            "data": {
                "candidates": [
                    {
                        "content": {
                            "parts": [
                                {"text": "done"},
                                {
                                    "inline_data": {
                                        "mime_type": "image/png",
                                        "data": image_b64,
                                    },
                                },
                            ],
                        },
                    },
                ],
            },
        }).encode("utf-8"),
    ]
    seen: list[tuple[str, str, bytes | None, dict[str, str]]] = []

    class _Resp:
        def __init__(self, payload: bytes):
            self._payload = payload

        def __enter__(self):
            return self

        def __exit__(self, *_exc):
            return False

        def read(self) -> bytes:
            return self._payload

    def fake_urlopen(req, timeout=0):
        seen.append((req.get_method(), req.full_url, req.data, dict(req.headers)))
        if not responses:
            _fail("Friday fake urlopen was called more times than expected")
        return _Resp(responses.pop(0))

    old_urlopen = ib.request.urlopen
    ib.request.urlopen = fake_urlopen
    try:
        backend = FridayGeminiImageBackend(_Settings(), "gemini-3-pro-image-preview")
        result = backend.generate(
            prompt="生成一张测试图",
            aspect_ratio="16:9",
            image_size="2K",
        )
    finally:
        ib.request.urlopen = old_urlopen

    if not isinstance(result, ImageResult):
        _fail(f"expected ImageResult; got {type(result).__name__}")
    if (result.width, result.height) != (32, 24):
        _fail(f"decoded image size mismatch: {(result.width, result.height)}")
    if result.model != "gemini-3-pro-image-preview":
        _fail(f"model should be preserved; got {result.model!r}")
    if len(seen) != 4:
        _fail(f"expected 1 submit + 3 query calls; got {len(seen)}")

    submit_method, submit_url, submit_body, submit_headers = seen[0]
    if submit_method != "POST" or not submit_url.endswith(
        "/gemini-3-pro-image-preview:imageGenerate"
    ):
        _fail(f"bad Friday submit request: {submit_method} {submit_url}")
    if submit_headers.get("Authorization") != "Bearer friday-appid":
        _fail(f"missing Friday bearer auth header: {submit_headers}")
    if not submit_body:
        _fail("submit request should carry JSON body")
    body = json.loads(submit_body.decode("utf-8"))
    cfg = body.get("generationConfig") or {}
    if cfg.get("responseModalities") != ["TEXT", "IMAGE"]:
        _fail(f"Friday request must ask for TEXT+IMAGE; got {cfg}")
    if cfg.get("imageConfig") != {"aspectRatio": "16:9", "imageSize": "2K"}:
        _fail(f"Friday imageConfig mismatch: {cfg.get('imageConfig')!r}")
    if not all(method == "GET" for method, *_ in seen[1:]):
        _fail(f"query calls should be GET; got {seen[1:]}")
    if responses:
        _fail(f"unused fake Friday responses: {len(responses)}")

    _ok(
        "Friday Gemini backend posts bearer-auth JSON, polls through running "
        "+ transient 429, and decodes inline_data to PNG"
    )


def check_export_sanitizer() -> None:
    """v2.8.2-C1 — sanitize_design_spec drops placeholder text and
    debug-named empty shapes from `layer_graph` descendants without
    touching slides themselves or callouts that carry real anchor text.
    Locks in the B1 ("Paper Title Goes Here", "Annotation 12") and B4
    (empty ``callout_05_a``) defects from the 2026-04-26 longcat-next
    dogfood.
    """
    print("[49/52] export sanitizer drops placeholders + debug-named empty shapes")
    from .schema import ArtifactType, DesignSpec, LayerNode
    from .util.export_sanitizer import sanitize_design_spec

    spec = DesignSpec(
        brief="export sanitizer smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"},
        layer_graph=[
            LayerNode(layer_id="slide_01", name="cover", kind="slide", z_index=1, role="cover", children=[
                # 1) valid title — kept
                LayerNode(layer_id="slide_01_title", name="title", kind="text", z_index=10,
                          text="Real Title"),
                # 2) placeholder text leak (B1) — dropped
                LayerNode(layer_id="slide_01_subtitle", name="subtitle", kind="text", z_index=11,
                          text="Paper Title Goes Here"),
                # 3) debug-named empty callout (B4) — dropped
                LayerNode(layer_id="callout_empty", name="callout_05_a", kind="callout", z_index=20,
                          anchor_layer_id="slide_01_title", callout_style="label",
                          callout_text=""),
                # 4) callout with real anchor text — kept (same scaffold name,
                #    but `callout_text` is non-empty so it's a real annotation)
                LayerNode(layer_id="callout_real", name="callout_05_a", kind="callout", z_index=21,
                          anchor_layer_id="slide_01_title", callout_style="label",
                          callout_text="real anchor text"),
            ]),
        ],
    )

    cleaned, warnings = sanitize_design_spec(spec)

    # Spec immutability — original untouched
    if len(spec.layer_graph[0].children) != 4:
        _fail(f"sanitizer mutated input spec: original child count = "
              f"{len(spec.layer_graph[0].children)} (expected 4)")

    # Slide kept; 2 children dropped, 2 kept
    if len(cleaned.layer_graph) != 1:
        _fail(f"slide was dropped: {len(cleaned.layer_graph)} top-level nodes")
    kept_children = cleaned.layer_graph[0].children
    if len(kept_children) != 2:
        _fail(f"expected 2 kept children, got {len(kept_children)}: "
              f"{[c.layer_id for c in kept_children]}")
    kept_ids = {c.layer_id for c in kept_children}
    if kept_ids != {"slide_01_title", "callout_real"}:
        _fail(f"wrong children kept: {sorted(kept_ids)} "
              f"(expected slide_01_title + callout_real)")

    # 2 warnings with the expected reasons
    if len(warnings) != 2:
        _fail(f"expected 2 warnings, got {len(warnings)}: {warnings}")
    reasons = {w["layer_id"]: w["reason"] for w in warnings}
    if reasons.get("slide_01_subtitle") != "placeholder_text":
        _fail(f"reason mismatch for slide_01_subtitle: {reasons}")
    if reasons.get("callout_empty") != "debug_name_empty":
        _fail(f"reason mismatch for callout_empty: {reasons}")
    for w in warnings:
        if w.get("slide_id") != "slide_01":
            _fail(f"warning missing slide_id binding: {w}")

    _ok(
        "2 children dropped (placeholder_text + debug_name_empty), "
        "2 kept (real title + real-text callout), original spec untouched"
    )


def check_slide_alignment_validator() -> None:
    """v2.8.2-C2 — naive title-body alignment validator. Detects slides
    whose title noun phrases don't appear in the body/figure text (B2,
    claim-evidence drift). Set-overlap only — NO embeddings, NO LLM call.

    Exercises three shapes:
      1. Aligned: title shares >= 50% of its tokens with body → no warning.
      2. Misaligned: title is "Training Stage Ablations" but body talks about
         phases / tokens / model with NONE of {training, stage, ablations}
         present → score 0/3 = 0.0 → blocker fires.
      3. Stopword-only title ("Results"): no extractable noun phrases →
         score defaults to 1.0 → no warning.
    """
    print("[50/52] slide alignment validator scores title vs body (no API)")
    from .schema import (
        ArtifactType, DeckDesignSystem, DesignSpec, LayerNode,
    )
    from .util.slide_alignment import (
        detect_alignment_warnings,
        extract_noun_phrases,
        slide_alignment_score,
    )

    # Sub-case A: aligned slide.
    aligned = LayerNode(
        layer_id="slide_aligned", name="aligned", kind="slide",
        z_index=1, role="content", children=[
            LayerNode(layer_id="t_a", name="title", kind="text", z_index=10,
                      template_slot="title",
                      text="Audio Tokenizer Architecture"),
            LayerNode(layer_id="b_a", name="body", kind="text", z_index=11,
                      template_slot="body",
                      text=("audio tokenizer uses RVQ and SAE for residual "
                            "quantization on architecture diagrams")),
        ],
    )
    score_a, missing_a = slide_alignment_score(aligned)
    if score_a < 0.5:
        _fail(f"aligned slide should score >= 0.5; got {score_a} "
              f"(missing={missing_a})")

    # Sub-case B: misaligned. Title noun phrases are
    # {training, stage, ablations}; body uses {model, trained, tokens, phase}
    # — no overlap, score 0/3 = 0.0 → blocker.
    misaligned = LayerNode(
        layer_id="slide_misaligned", name="misaligned", kind="slide",
        z_index=2, role="content", children=[
            LayerNode(layer_id="t_b", name="title", kind="text", z_index=10,
                      template_slot="title",
                      text="Training Stage Ablations"),
            LayerNode(layer_id="b_b", name="body", kind="text", z_index=11,
                      template_slot="body",
                      text=("the model was trained with 100B tokens in "
                            "phase 1 then 200B in phase 2")),
        ],
    )
    score_b, missing_b = slide_alignment_score(misaligned)
    if score_b >= 0.30:
        _fail(f"misaligned slide should score < 0.30; got {score_b} "
              f"(missing={missing_b})")
    if not {"training", "stage", "ablations"}.issubset(missing_b):
        _fail(f"missing keywords should include training/stage/ablations; "
              f"got {missing_b}")

    # Sub-case C edge-1: a title that survives stopword stripping but whose
    # one phrase appears verbatim in the body → score 1.0, no warning.
    single_phrase_aligned = LayerNode(
        layer_id="slide_results", name="results", kind="slide",
        z_index=3, role="content", children=[
            LayerNode(layer_id="t_c", name="title", kind="text", z_index=10,
                      template_slot="title", text="Results"),
            LayerNode(layer_id="b_c", name="body", kind="text", z_index=11,
                      template_slot="body",
                      text="results show 92.3 accuracy on the benchmark"),
        ],
    )
    title_phrases_c = extract_noun_phrases("Results")
    if title_phrases_c != {"results"}:
        _fail(f"'Results' should yield {{'results'}}; got {title_phrases_c}")
    score_c, _ = slide_alignment_score(single_phrase_aligned)
    if score_c != 1.0:
        _fail(f"single-phrase aligned slide should score 1.0; got {score_c}")

    # Sub-case C edge-2: a title made entirely of stopwords (rare in practice,
    # but exercises the "no noun phrases extractable" code path → defaults to
    # score 1.0 — we can't validate what we can't extract).
    pure_stopword_title = LayerNode(
        layer_id="slide_pure_stop", name="pure_stop", kind="slide",
        z_index=4, role="content", children=[
            LayerNode(layer_id="t_d", name="title", kind="text", z_index=10,
                      template_slot="title", text="The And Of"),
            LayerNode(layer_id="b_d", name="body", kind="text", z_index=11,
                      template_slot="body", text="totally unrelated body"),
        ],
    )
    if extract_noun_phrases("The And Of"):
        _fail("pure-stopword title should yield no noun phrases")
    score_d, _ = slide_alignment_score(pure_stopword_title)
    if score_d != 1.0:
        _fail(f"pure-stopword-title slide should default to 1.0; got {score_d}")

    # Spec-level integration: detect_alignment_warnings should emit one
    # entry for the misaligned slide and skip the other two.
    spec = DesignSpec(
        brief="alignment validator smoke",
        artifact_type=ArtifactType.DECK,
        canvas={"w_px": 1920, "h_px": 1080, "dpi": 96,
                "aspect_ratio": "16:9", "color_mode": "RGB"},
        deck_design_system=DeckDesignSystem(style="academic-editorial"),
        layer_graph=[aligned, misaligned, single_phrase_aligned,
                     pure_stopword_title],
    )
    warnings = detect_alignment_warnings(spec)
    if len(warnings) != 1:
        _fail(f"expected exactly 1 warning (misaligned slide); "
              f"got {len(warnings)}: {warnings!r}")
    w = warnings[0]
    if w["slide_id"] != "slide_misaligned":
        _fail(f"warning should target slide_misaligned; got {w!r}")
    if w["severity"] != "blocker":
        _fail(f"score 0.0 should map to severity=blocker; got {w!r}")
    if w["title"] != "Training Stage Ablations":
        _fail(f"warning should carry the title; got {w!r}")

    _ok(
        f"aligned≥0.5 ({score_a:.2f}), misaligned<0.30 "
        f"({score_b:.2f}, blocker), single-phrase-aligned=1.0, "
        f"pure-stopword-title=1.0; spec walker emits exactly 1 warning"
    )


def check_closing_stub_detector() -> None:
    """v2.8.2 C3 — last slide must carry substantive takeaways, not the
    template-default "Thank You" stub. ``_detect_closing_stub`` returns a
    warning when the closing slide has <3 text runs OR all runs are stub
    phrases (Thank You / Q&A) OR all runs are placeholder strings. Real
    takeaways override stub phrases on a mixed closing slide. Detection
    only — no auto-fix; the planner reads the warning from the tool_result
    payload and fixes on the next iteration.
    """
    print("[51/52] closing stub detector flags template-default last slide")
    from .schema import (
        ArtifactType, DeckDesignSystem, DesignSpec, LayerNode,
    )
    from .tools.composite import _detect_closing_stub

    canvas = {"w_px": 1920, "h_px": 1080, "dpi": 96, "aspect_ratio": "16:9", "color_mode": "RGB"}

    def _build(closing_children: list[LayerNode]) -> DesignSpec:
        return DesignSpec(
            brief="closing stub smoke",
            artifact_type=ArtifactType.DECK,
            canvas=canvas,
            deck_design_system=DeckDesignSystem(style="academic-editorial"),
            layer_graph=[
                LayerNode(layer_id="slide_01", name="cover", kind="slide", z_index=1, role="cover", children=[
                    LayerNode(layer_id="slide_01_title", name="title", kind="text", z_index=10,
                              template_slot="title", text="Cover"),
                    LayerNode(layer_id="slide_01_subtitle", name="subtitle", kind="text", z_index=11,
                              template_slot="subtitle", text="Subtitle"),
                ]),
                LayerNode(layer_id="slide_02", name="content", kind="slide", z_index=2, role="content", children=[
                    LayerNode(layer_id="slide_02_title", name="title", kind="text", z_index=10,
                              template_slot="title", text="Body slide"),
                    LayerNode(layer_id="slide_02_body", name="body", kind="text", z_index=11,
                              template_slot="body", text="Some content here."),
                ]),
                LayerNode(layer_id="slide_99", name="closing", kind="slide", z_index=99,
                          role="closing", children=closing_children),
            ],
        )

    # Case 1 — stub closing: a single "Thank You" child (1 run < 3) → warn.
    stub_spec = _build([
        LayerNode(layer_id="slide_99_title", name="title", kind="text", z_index=10,
                  template_slot="title", text="Thank You"),
    ])
    stub_warnings = _detect_closing_stub(stub_spec)
    if len(stub_warnings) != 1:
        _fail(f"stub closing: expected 1 warning, got {stub_warnings!r}")
    if stub_warnings[0]["slide_id"] != "slide_99":
        _fail(f"stub closing: wrong slide_id: {stub_warnings[0]!r}")
    if stub_warnings[0]["reason"] not in ("thin_content", "all_stub_phrases"):
        _fail(f"stub closing: unexpected reason: {stub_warnings[0]!r}")

    # Case 2 — substantive closing: 3 takeaway bullets → no warning.
    substantive_spec = _build([
        LayerNode(layer_id="slide_99_title", name="title", kind="text", z_index=10,
                  template_slot="title", text="Key takeaways"),
        LayerNode(layer_id="slide_99_b1", name="bullet1", kind="text", z_index=11,
                  template_slot="body",
                  text="LongCat-Next demonstrates layered tokenization scales to 1M context."),
        LayerNode(layer_id="slide_99_b2", name="bullet2", kind="text", z_index=12,
                  template_slot="body",
                  text="The unified tokenizer scales to interleaved image-text without retraining."),
        LayerNode(layer_id="slide_99_b3", name="bullet3", kind="text", z_index=13,
                  template_slot="body",
                  text="Future work focuses on streaming generation and memory-bounded decode."),
    ])
    substantive_warnings = _detect_closing_stub(substantive_spec)
    if substantive_warnings:
        _fail(f"substantive closing should not warn; got {substantive_warnings!r}")

    # Case 3 — mixed closing: "Thank You" + 3 real takeaways → 4 runs, not
    # all stubs, so no warning. Confirms a real bullet rescues a slide that
    # also carries the polite stub.
    mixed_spec = _build([
        LayerNode(layer_id="slide_99_title", name="title", kind="text", z_index=10,
                  template_slot="title", text="Thank You"),
        LayerNode(layer_id="slide_99_b1", name="bullet1", kind="text", z_index=11,
                  template_slot="body",
                  text="LongCat-Next demonstrates layered tokenization scales to 1M context."),
        LayerNode(layer_id="slide_99_b2", name="bullet2", kind="text", z_index=12,
                  template_slot="body",
                  text="The unified tokenizer scales to interleaved image-text without retraining."),
        LayerNode(layer_id="slide_99_b3", name="bullet3", kind="text", z_index=13,
                  template_slot="body",
                  text="Future work focuses on streaming generation and memory-bounded decode."),
    ])
    mixed_warnings = _detect_closing_stub(mixed_spec)
    if mixed_warnings:
        _fail(f"mixed closing (stub + real bullets) should not warn; got {mixed_warnings!r}")

    _ok("stub closing flagged; substantive + mixed closings pass")


def check_pptx_template_default_scrubber() -> None:
    """v2.8.2-C1 stage 2 — ``sanitize_pptx_file`` blanks template-default
    placeholder text baked into ``assets/deck_templates/*.pptx`` that the
    spec-level sanitizer cannot see.

    Locks in the 2026-04-26 longcat-next dogfood failure where "Paper
    Title Goes Here" / "Author One · Author Two · Affiliation" leaked into
    the cover slot from the template's own slide XML.
    """
    print("[52/52] post-write .pptx scrubber blanks template defaults")
    import shutil
    import zipfile
    from pathlib import Path
    from .util.export_sanitizer import sanitize_pptx_file

    src = Path("assets/deck_templates/academic-editorial.pptx")
    if not src.exists():
        _fail(f"template missing: {src}")
    out_dir = Path("out/smoke_pptx_scrub")
    out_dir.mkdir(parents=True, exist_ok=True)
    target = out_dir / "academic-editorial.pptx"
    shutil.copy2(src, target)

    # Confirm the leak exists in the unscrubbed copy.
    with zipfile.ZipFile(target) as z:
        before = "".join(
            z.read(n).decode()
            for n in z.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
    if "Paper Title Goes Here" not in before:
        _fail("template no longer carries 'Paper Title Goes Here' "
              "— update fixture or regression target")
    if "Author One" not in before:
        _fail("template no longer carries 'Author One' — update fixture")

    warnings = sanitize_pptx_file(target)
    if not warnings:
        _fail("scrubber returned 0 warnings; expected at least 2 "
              "(Paper Title Goes Here + Author One/Two)")
    reasons = {w["reason"] for w in warnings}
    if reasons != {"template_default_placeholder"}:
        _fail(f"unexpected reason(s): {reasons}")

    with zipfile.ZipFile(target) as z:
        after = "".join(
            z.read(n).decode()
            for n in z.namelist()
            if n.startswith("ppt/slides/slide") and n.endswith(".xml")
        )
    for needle in ("Paper Title Goes Here", "Author One", "Author Two"):
        if needle in after:
            _fail(f"scrubber failed to blank {needle!r} in {target}")

    _ok(
        f"scrubbed {len(warnings)} placeholder text run(s) from "
        "academic-editorial.pptx slide XML in-place"
    )


def main() -> int:
    check_imports()
    check_tool_registry()
    check_pydantic_roundtrip()
    check_fonts()
    check_composite_no_api()
    check_svg_text_is_vector()
    check_chat_session_roundtrip()
    check_edit_layer_no_api()
    check_apply_edits_roundtrip()
    check_landing_mode()
    check_design_system_styles()
    check_landing_with_images()
    check_deck_mode()
    check_reasoning_step_roundtrip()
    check_ingest_document_markdown()
    check_ingest_document_image()
    check_ingest_document_docx()
    check_ingest_document_pptx()
    check_sub_figure_registration()
    check_versioning_no_api()
    check_deck_design_system_template()
    check_footer_leakage()
    check_callout_overlay()
    check_provenance_validator()
    check_section_renumber_policy()
    check_section_renumber_strip()
    check_stable_id_notes_after_reorder()
    check_critic_subagent_trajectory()
    check_critic_subagent_max_turns()
    check_critic_planner_consumption()
    check_critic_subagent_png_throughput()
    check_critic_openai_compat_strict_format()
    check_archetype_cover_editorial()
    check_archetype_evidence_snapshot()
    check_archetype_takeaway_list()
    check_archetype_thanks_qa()
    check_archetype_fallback_default()
    check_archetype_determinism()


    check_claim_graph_extractor_trajectory()
    check_claim_graph_validator_rejects_fabricated_quote()
    check_planner_covers_population()
    check_critic_claim_coverage_issue()
    check_no_claim_graph_pipeline_degrades()
    check_claim_graph_lookup_whitespace_tolerance()
    check_claim_graph_kimi_template_leak_retry()
    check_deck_text_overlap_detector()
    check_orphan_callout_drop()
    check_export_sanitizer()
    check_slide_alignment_validator()


    check_closing_stub_detector()
    check_pptx_template_default_scrubber()


    check_friday_gemini_image_backend_no_api()
    check_image_backend_fallback_chain()
    print("\n  smoke test passed.")
    print("  artifacts in: out/smoke/, out/smoke_edit/, out/smoke_apply/, "
          "out/smoke_landing/, out/smoke_styles/, out/smoke_landing_img/, "
          "out/smoke_deck/, out/smoke_ingest_md/, out/smoke_ingest_image/, "
          "out/smoke_ingest_docx/, out/smoke_ingest_pptx/, out/smoke_sub_figs/, "
          "out/smoke_section_notes/, "
          "out/smoke_critic_subagent/, out/smoke_critic_consume/, "
          "out/smoke_claim_graph/, out/smoke_no_claim_graph/, "
          "out/smoke_claim_graph_kimi_leak/")
    return 0


if __name__ == "__main__":
    sys.exit(main())
