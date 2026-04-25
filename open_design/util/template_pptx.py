"""Template-based PPTX manipulation primitives (v2.5.2).

Inspired by Anthropic's pptx skill (proprietary; not vendored — its
LICENSE.txt forbids reproduction). This module reimplements only the
operations OpenDesign needs, using public python-pptx APIs only:

- `inventory_by_name(slide)` — find named shapes for slot replacement.
- `clone_template_slide(prs, src_idx)` — append a copy of a template
  slide to the end of the presentation, returning the new Slide.
- `replace_text_in_shape(shape, paragraphs)` — clear a text frame and
  rewrite paragraphs while preserving the template's run-level
  formatting (font / size / color) unless explicitly overridden.

These three primitives + python-pptx's existing `add_picture` /
`add_table` cover the deck templating use case fully. We do NOT
support OOXML-level edits, image replacement-in-place, or HTML→PPTX
conversion — out of scope for v2.5.2.
"""

from __future__ import annotations

from copy import deepcopy
from typing import Any, TypedDict

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as _Presentation
from pptx.shapes.base import BaseShape
from pptx.slide import Slide
from pptx.util import Pt


class ParagraphSpec(TypedDict, total=False):
    """One paragraph in a `replace_text_in_shape` call. Only `text` is
    required; everything else falls back to template defaults captured
    from the original first run of the shape's text frame."""
    text: str
    bullet: bool          # adds a bullet character (level=0 unless overridden)
    level: int            # bullet/indent level
    alignment: str        # "left" | "center" | "right"
    bold: bool
    italic: bool
    color: str            # hex like "FF0000" or "#FF0000"
    font_size: float      # in points
    font_name: str


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


# ── Inventory ─────────────────────────────────────────────────────────


def inventory_by_name(slide: Slide) -> dict[str, BaseShape]:
    """Return all shapes on `slide` keyed by `shape.name`.

    OpenDesign's templates assign deterministic names (`title`, `body`,
    `image_slot`, etc.) so the renderer addresses slots by name rather
    than by visual position. The Anthropic skill's positional shape-N
    indexing is fragile across template edits; named lookup is not.

    Last-write wins on duplicate names — the generator should not emit
    duplicates, but this matches python-pptx's underlying behavior.
    """
    return {s.name: s for s in slide.shapes if s.name}


# ── Cloning ───────────────────────────────────────────────────────────


def clone_template_slide(prs: _Presentation, src_idx: int) -> Slide:
    """Append a deep copy of `prs.slides[src_idx]` to the end of the
    presentation. Returns the new Slide.

    Implementation: add a new blank slide via the standard
    `prs.slides.add_slide(blank_layout)` path so all the relationship
    plumbing is correct, then copy each shape XML element into the new
    slide's spTree. Picture/media relationships are NOT copied — the
    template should not contain pictures (only shapes + text frames);
    callers add real pictures via `slide.shapes.add_picture` after.

    This matches what the Anthropic skill's `rearrange.py` does for
    the use case where the template is purely vector + text. For
    pictures-in-template, an OOXML-level rels copy would be required.
    Out of scope for v2.5.2.
    """
    if src_idx < 0 or src_idx >= len(prs.slides):
        raise IndexError(f"src_idx {src_idx} out of range [0, {len(prs.slides)})")

    src_slide = prs.slides[src_idx]

    # Pick a blank layout for the new slide (last layout is conventionally
    # blank in the default template).
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    new_slide = prs.slides.add_slide(blank_layout)

    # Remove any placeholder shapes the blank layout inherited so the new
    # slide starts clean.
    for shape in list(new_slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # Deep-copy each shape XML element from the source slide into the
    # new slide's shape tree.
    src_spTree = src_slide.shapes._spTree
    new_spTree = new_slide.shapes._spTree
    for child in list(src_spTree):
        # Skip the required group properties (nvGrpSpPr, grpSpPr) — those
        # are already on the new slide. Only copy actual shape elements.
        tag = child.tag.rsplit("}", 1)[-1]
        if tag in ("nvGrpSpPr", "grpSpPr", "extLst"):
            continue
        new_spTree.append(deepcopy(child))

    return new_slide


# ── Text replacement ──────────────────────────────────────────────────


def replace_text_in_shape(shape: BaseShape, paragraphs: list[ParagraphSpec]) -> None:
    """Replace `shape`'s text content with `paragraphs`, preserving the
    template's run-level formatting unless a paragraph spec overrides it.

    Captures the FIRST run's font (name / size / color / bold / italic /
    underline) from the existing text frame as the template default,
    then clears, then writes each new paragraph applying defaults +
    spec overrides.

    No-op (silent return) if the shape has no `text_frame` (e.g. it's a
    Picture or unsupported shape kind).
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame

    # Capture template formatting from first run of first paragraph.
    template = _capture_template_format(tf)

    # python-pptx's clear() leaves one empty paragraph behind; we'll reuse it.
    tf.clear()

    for i, p_spec in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        align = p_spec.get("alignment")
        if align and align in _ALIGN_MAP:
            p.alignment = _ALIGN_MAP[align]
        if p_spec.get("bullet"):
            level = p_spec.get("level", 0)
            p.level = level
        run = p.add_run()
        run.text = p_spec.get("text", "")
        _apply_run_format(run, template, p_spec)


def _capture_template_format(tf) -> dict[str, Any]:
    """Pull (font_name, font_size_pt, color_rgb, bold, italic) from the
    first run of the first paragraph if it exists."""
    cap: dict[str, Any] = {}
    if not tf.paragraphs:
        return cap
    first_p = tf.paragraphs[0]
    if not first_p.runs:
        return cap
    r = first_p.runs[0]
    f = r.font
    if f.name:
        cap["font_name"] = f.name
    if f.size is not None:
        cap["font_size_pt"] = f.size.pt
    if f.bold is not None:
        cap["bold"] = f.bold
    if f.italic is not None:
        cap["italic"] = f.italic
    try:
        if f.color and f.color.rgb is not None:
            cap["color_rgb"] = f.color.rgb
    except Exception:
        # color may be theme-color or unset; skip silently
        pass
    return cap


def _apply_run_format(run, template: dict[str, Any], p_spec: ParagraphSpec) -> None:
    """Apply template defaults then override with paragraph spec."""
    # Defaults from template
    if "font_name" in template:
        run.font.name = template["font_name"]
    if "font_size_pt" in template:
        run.font.size = Pt(template["font_size_pt"])
    if "bold" in template:
        run.font.bold = template["bold"]
    if "italic" in template:
        run.font.italic = template["italic"]
    if "color_rgb" in template:
        run.font.color.rgb = template["color_rgb"]

    # Overrides from spec
    if "font_name" in p_spec:
        run.font.name = p_spec["font_name"]
    if "font_size" in p_spec:
        run.font.size = Pt(p_spec["font_size"])
    if "bold" in p_spec:
        run.font.bold = p_spec["bold"]
    if "italic" in p_spec:
        run.font.italic = p_spec["italic"]
    if "color" in p_spec:
        run.font.color.rgb = _hex_to_rgb(p_spec["color"])


def _hex_to_rgb(hex_str: str) -> RGBColor:
    """Parse `#RRGGBB` or `RRGGBB` into RGBColor."""
    s = hex_str.strip().lstrip("#")
    if len(s) != 6:
        raise ValueError(f"hex color must be 6 chars after '#', got {hex_str!r}")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
