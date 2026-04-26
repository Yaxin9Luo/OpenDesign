"""Shared helpers for v2.8.1 slide archetypes.

Every archetype renderer in this package is a self-contained layout
function. They all need the same handful of primitives:

- A theme-token resolver that turns a `LayerNode.effects.fill` (or a
  spec-level palette default) into an `RGBColor`.
- An `add_textbox` helper that mirrors the inline `_add_text_frame`
  primitive in `pptx_renderer.py` but takes explicit pixel coordinates
  and font knobs instead of pulling them from a child node.
- A re-export of v2.7.2's `_with_section_prefix` so each archetype can
  decorate its title with `slide.section_number` without re-implementing
  the substring check.

Determinism rules:

- No `time.time()` / `uuid` calls; every shape position is computed from
  the input slide.
- Helpers never mutate the input `LayerNode`; `_with_section_prefix`
  already returns a new pydantic copy.
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

# Re-exported so archetype renderers depend on this module, not the
# private `pptx_renderer` namespace.
from ..pptx_renderer import (  # noqa: F401
    PX_TO_EMU,
    _hex_to_rgb,
    _is_title_child,
    _with_section_prefix,
)


_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def px_to_emu(px: int) -> int:
    """Pixel → EMU. Centralized so any future calibration tweak hits
    every archetype simultaneously."""
    return int(px) * PX_TO_EMU


def hex_to_rgbcolor(hex_str: str | None, default: tuple[int, int, int]
                    ) -> RGBColor:
    """Parse a `#rrggbb` string to RGBColor; fall back to `default` on
    any parse failure or `None` input."""
    rgb = _hex_to_rgb(hex_str) if isinstance(hex_str, str) else None
    if rgb is None:
        rgb = default
    return RGBColor(*rgb)


def add_textbox(
    slide: Any,
    *,
    text: str,
    x_px: int,
    y_px: int,
    w_px: int,
    h_px: int,
    font_family: str | None = None,
    font_size_px: int = 36,
    align: str = "left",
    fill_hex: str | None = None,
    bold: bool = False,
    italic: bool = False,
    name: str | None = None,
) -> Any:
    """Add a native textbox at pixel coordinates with the given font knobs.

    Uses python-pptx primitives only — no rasterization. Returns the
    underlying shape so callers can attach extra attributes (e.g. a
    layer-id-style `name` for inspection).

    `text` may contain newlines; each line becomes its own paragraph
    with the same alignment / font.
    """
    shape = slide.shapes.add_textbox(
        Emu(px_to_emu(x_px)), Emu(px_to_emu(y_px)),
        Emu(px_to_emu(max(1, w_px))), Emu(px_to_emu(max(1, h_px))),
    )
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    lines = (text or "").splitlines() or [text or ""]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ""
        para.alignment = _ALIGN_MAP.get(align, PP_ALIGN.LEFT)
        run = para.add_run()
        run.text = line
        font = run.font
        # python-pptx wants pts; we get px. pt ≈ px * 0.75.
        font.size = Pt(max(6, round(int(font_size_px) * 0.75)))
        if font_family:
            font.name = font_family
        font.bold = bool(bold)
        font.italic = bool(italic)
        if fill_hex:
            font.color.rgb = hex_to_rgbcolor(fill_hex, (15, 23, 42))
    return shape


def collect_text_children(slide_node: Any) -> list[Any]:
    """Return the slide's text children sorted by z_index.

    The four Phase 1 archetypes only consume kind="text" children; image
    / background / table children are passed through untouched to the
    default render via the dispatcher's fallback path. Sorting keeps
    rendering deterministic across spec re-orderings.
    """
    out: list[Any] = []
    for c in (getattr(slide_node, "children", None) or []):
        if getattr(c, "kind", None) == "text":
            out.append(c)
    out.sort(key=lambda c: (int(getattr(c, "z_index", 0) or 0),
                            getattr(c, "layer_id", "") or ""))
    return out


def find_text_child(
    slide_node: Any, *, name_hints: tuple[str, ...] = (),
) -> Any | None:
    """Return the first kind="text" child whose `name` (case-insensitive)
    contains any of `name_hints`; None when no match. Useful for picking
    out the title / subtitle / byline children on cover slides."""
    if not name_hints:
        return None
    hints = tuple(h.lower() for h in name_hints)
    for c in collect_text_children(slide_node):
        cname = (getattr(c, "name", None) or "").lower()
        if any(h in cname for h in hints):
            return c
    return None


_NON_TITLE_NAME_BLOCKLIST = ("subtitle", "byline_sub")


def first_title_child(slide_node: Any) -> Any | None:
    """Return the first text child that looks like a title, or None.

    The renderer's `_is_title_child` substring rule ("title" in name)
    is too liberal for archetype dispatch — "subtitle" matches but is
    not a title. We layer a small blocklist on top so cover_editorial
    /  thanks_qa pick up the actual headline child.
    """
    for c in collect_text_children(slide_node):
        cname = (getattr(c, "name", None) or "").lower()
        if any(b in cname for b in _NON_TITLE_NAME_BLOCKLIST):
            continue
        if _is_title_child(c):
            return c
    return None


def write_speaker_notes(slide: Any, slide_node: Any) -> None:
    """Mirror v2.3 speaker_notes binding. Archetype renderers must call
    this so the v2.7.2 stable-id contract holds across every layout."""
    notes = getattr(slide_node, "speaker_notes", None)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def title_text_with_section(slide_node: Any, fallback: str = "") -> str:
    """Resolve the slide title for an archetype renderer that builds its
    own title textbox (cover_editorial / thanks_qa). Picks the first
    title-bearing child's text, prepends `slide.section_number` via
    `_with_section_prefix`, and falls back to `fallback` (or the slide
    `name`) when no title child exists.

    Returns the final string ready to drop into `add_textbox(text=...)`.
    """
    title_child = first_title_child(slide_node)
    section_number = getattr(slide_node, "section_number", None)
    if title_child is not None:
        decorated = _with_section_prefix(title_child, section_number)
        text = (getattr(decorated, "text", None) or "").strip()
        if text:
            return text
    base = (fallback or getattr(slide_node, "name", None) or "").strip()
    if section_number and base and not base.startswith(section_number):
        return f"{section_number} · {base}"
    return base
