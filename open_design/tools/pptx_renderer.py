"""PPTX renderer for deck artifacts (v1.0 #7).

One top-level LayerNode with `kind="slide"` per slide; each slide's `children`
hold `kind in {"text","image","background"}` elements positioned by bbox in
slide-canvas pixel coords. We emit native PowerPoint shapes with native
TextFrames — no Pillow rasterization of text — so the .pptx opens
type-editable in PowerPoint / Keynote / Google Slides.

Per-slide simplified PNGs are also written (for chat preview + critic).

Font embedding is intentionally NOT performed: .pptx delegates font rendering
to the consuming app's font engine (see `docs/GOTCHAS.md` for the CJK-on-
consumer-PowerPoint note). This mirrors Paper2Any's approach.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from ._contract import ToolContext


# 1 pixel at 96 DPI ≈ 9525 EMU. python-pptx accepts Emu(int).
PX_TO_EMU = 9525

DEFAULT_SLIDE_W = 1920
DEFAULT_SLIDE_H = 1080

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


def write_pptx(spec: Any, pptx_path: Path, ctx: ToolContext) -> int:
    """Walk the slide tree and emit a .pptx file. Returns slide count.

    v2.5.2 — when `spec.deck_design_system` is set, branches to the
    template-backed renderer that opens
    `assets/deck_templates/<style>.pptx`, clones template layouts per
    `slide.role`, and fills shapes by `template_slot` name. The
    no-template path below is preserved unchanged for back-compat.
    """
    ds = getattr(spec, "deck_design_system", None)
    if ds is not None:
        return _write_pptx_templated(spec, ds, pptx_path, ctx)

    canvas = spec.canvas or {}
    slide_w = int(canvas.get("w_px") or DEFAULT_SLIDE_W)
    slide_h = int(canvas.get("h_px") or DEFAULT_SLIDE_H)

    prs = Presentation()
    prs.slide_width = Emu(slide_w * PX_TO_EMU)
    prs.slide_height = Emu(slide_h * PX_TO_EMU)

    blank_layout = prs.slide_layouts[6]  # "Blank" layout — we position everything.

    slide_count = 0
    for node in (spec.layer_graph or []):
        if getattr(node, "kind", None) != "slide":
            continue
        slide = prs.slides.add_slide(blank_layout)
        _render_slide(slide, node, slide_w, slide_h, ctx)
        slide_count += 1

    prs.save(str(pptx_path))
    return slide_count


# ── v2.5.2 templated path ─────────────────────────────────────────────


# Roles → indices into the 6 template slides built by scripts/build_template.py.
# Order MUST match build_template.py's build() invocation order.
ROLE_TO_LAYOUT_IDX = {
    "cover": 0,
    "section_divider": 1,
    "content": 2,
    "content_with_figure": 3,
    "content_with_table": 4,
    "closing": 5,
}


def _write_pptx_templated(spec: Any, ds: Any, pptx_path: Path,
                          ctx: ToolContext) -> int:
    """Render `spec` using the template at `assets/deck_templates/<style>.pptx`.

    Per spec slide:
      1. Clone the template slide whose index matches `slide.role`.
      2. Fill named text shapes from children where `child.template_slot`
         matches `shape.name`.
      3. Place images at `image_slot` shape's bbox via existing add_picture
         + aspect-fit. Falls back to `child.bbox` if no slot found.
      4. Place tables at `table_anchor` shape's bbox via existing add_table.
      5. Auto-fill `footer` shape with ds.footer_text or first 60 chars of
         spec.brief; auto-fill `slide_number` shape with "N/total".
      6. Speaker notes path unchanged.

    After all spec slides are rendered, the original 6 template slides are
    removed so only the spec's slides remain.
    """
    from ..util.template_pptx import (
        clone_template_slide,
        inventory_by_name,
        replace_text_in_shape,
    )

    canvas = spec.canvas or {}
    slide_w = int(canvas.get("w_px") or DEFAULT_SLIDE_W)
    slide_h = int(canvas.get("h_px") or DEFAULT_SLIDE_H)

    template_path = (
        ctx.settings.repo_root / "assets" / "deck_templates" / f"{ds.style}.pptx"
    )
    if not template_path.exists():
        raise FileNotFoundError(
            f"deck template not found: {template_path}. Run "
            f"`uv run python scripts/build_template.py` to (re)generate."
        )

    prs = Presentation(str(template_path))
    n_template_slides = len(prs.slides)

    # Pre-compute footer text for auto-fill on content slides.
    footer_text = _resolve_footer_text(ds, spec, ctx)

    spec_slide_nodes = [
        n for n in (spec.layer_graph or []) if getattr(n, "kind", None) == "slide"
    ]
    total = len(spec_slide_nodes)

    new_slides: list[Any] = []
    for idx, node in enumerate(spec_slide_nodes):
        role = getattr(node, "role", None) or "content"
        layout_idx = ROLE_TO_LAYOUT_IDX.get(role, ROLE_TO_LAYOUT_IDX["content"])
        slide = clone_template_slide(prs, layout_idx)
        _render_templated_slide(
            slide, node, slide_w, slide_h, ctx, role,
            footer_text=footer_text, slide_idx=idx, total=total,
        )
        new_slides.append(slide)

    # Remove the original template slides from the slide list. We only drop
    # the sldId references in `presentation.xml`; the slide parts themselves
    # stay in the package as orphans. PowerPoint / Keynote / Google Slides
    # all ignore unreferenced parts cleanly. python-pptx's _Relationships
    # doesn't expose item deletion publicly, so a full part-drop would
    # require OOXML-level surgery — out of scope for v2.5.2 (file bloat is
    # ~30 KB, negligible vs. typical 5+ MB deck output with seedream images).
    sldIdLst = prs.slides._sldIdLst
    sldIds = list(sldIdLst)
    for i in range(n_template_slides):
        sldIdLst.remove(sldIds[i])

    prs.save(str(pptx_path))
    return total


# v2.5.2.2 — phrases that indicate the user's command leaked into footer
# fallback. The original v2.5.2 implementation used `spec.brief[:80]`,
# which silently shipped strings like "12-slide academic talk deck for
# the LongCat-Next paper. Speaker-ready with notes per slide" on every
# content slide. This blacklist is checked case-insensitively as a
# defensive filter on whatever footer text resolves; matching strings
# fall through to the next priority in `_resolve_footer_text`.
_FOOTER_LEAKAGE_PHRASES = (
    "slide deck",
    "speaker-ready",
    "speaker ready",
    "with notes",
    "academic talk deck",
    "12-slide",
    "10-slide",
    "16-slide",
    "8-slide",
    "lightning talk",
    "research talk",
    "for this paper",
    "paper. speaker",
)


# v2.7 — placeholder authors strings the planner emits when it forgets
# to read manifest.authors. The cover renderer's `_resolve_authors_text`
# uses this list to detect leakage and fall through to the manifest.
# Observed in 2026-04-25 longcat-next dogfood cover slide:
# "Author One · Author Two · Affiliation".
_AUTHORS_LEAKAGE_PHRASES = (
    "author one",
    "author two",
    "author three",
    "first author",
    "second author",
    "affiliation",
    "your affiliation",
    "your name",
    "name surname",
    "anonymous author",
    "first last",
)


def _is_leakage(text: str) -> bool:
    """True if `text` contains any user-command phrase that should not
    appear in a deck footer. Case-insensitive substring match."""
    if not text:
        return True
    low = text.lower()
    return any(phrase in low for phrase in _FOOTER_LEAKAGE_PHRASES)


def _is_authors_leakage(text: str) -> bool:
    """True if `text` contains a placeholder-author phrase. Empty string
    is also leakage so the resolver falls through to the manifest."""
    if not text or not text.strip():
        return True
    low = text.lower()
    return any(phrase in low for phrase in _AUTHORS_LEAKAGE_PHRASES)


def _resolve_footer_text(ds: Any, spec: Any, ctx: ToolContext) -> str:
    """Pick the right footer text for a deck content slide.

    Precedence (first non-leakage match wins):
      1. `ds.footer_text` — explicit planner override
      2. `ctx.state["ingested"][i]["manifest"]["title"]` — paper title
         from ingest_document (truncated to 80 chars). v2.5.2's bug was
         skipping this layer entirely and falling through to brief.
      3. Empty string. We deliberately do NOT fall back to `spec.brief`
         because brief is a user command, not slide-footer content.

    Returns the resolved string (possibly empty). Always sanitized
    against the v2.5.2.2 leakage blacklist.
    """
    # 1. Explicit planner override
    explicit = (getattr(ds, "footer_text", None) or "").strip()
    if explicit and not _is_leakage(explicit):
        return explicit[:80]

    # 2. Paper title from ingest
    state = getattr(ctx, "state", None) or {}
    ingested = state.get("ingested") or []
    for entry in ingested:
        manifest = (entry or {}).get("manifest") or {}
        title = (manifest.get("title") or "").strip()
        if title and not _is_leakage(title):
            return title[:80]

    # 3. Empty rather than leak the brief
    return ""


def _resolve_authors_text(ctx: ToolContext, *, existing: str | None = None) -> str:
    """v2.7 — resolve the cover-slide authors text.

    Precedence (first non-leakage match wins):
      1. `existing` planner-supplied text — if it's a non-placeholder
         string, keep it.
      2. `ctx.state["ingested"][i]["manifest"]["authors"]` — joined with
         " · " and truncated to 80 chars.
      3. Empty string (NOT a placeholder).

    Why a separate resolver: the 2026-04-25 longcat-next dogfood emitted
    "Author One · Author Two · Affiliation" on the cover because the
    planner skipped reading manifest.authors. The renderer becomes the
    safety net: even if the planner forgets, the manifest provides the
    truth, and the leakage filter rejects placeholder strings.
    """
    if existing and not _is_authors_leakage(existing):
        return existing.strip()[:80]
    state = getattr(ctx, "state", None) or {}
    for entry in (state.get("ingested") or []):
        manifest = (entry or {}).get("manifest") or {}
        authors = manifest.get("authors") or []
        if isinstance(authors, list) and authors:
            text = " · ".join(str(a) for a in authors if a).strip()
            if text and not _is_authors_leakage(text):
                return text[:80]
    return ""


def _render_templated_slide(
    slide: Any, slide_node: Any, slide_w: int, slide_h: int,
    ctx: ToolContext, role: str, *,
    footer_text: str, slide_idx: int, total: int,
) -> None:
    """Walk the slide spec's children, route each to its named slot in
    the cloned template slide. Fall back to absolute-bbox positioning
    when no slot match exists."""
    from ..util.template_pptx import inventory_by_name, replace_text_in_shape

    inv = inventory_by_name(slide)

    # Auto-fill footer + slide_number first so they're always populated
    # even if planner didn't address them.
    needs_footer = role in ("content", "content_with_figure", "content_with_table")
    if needs_footer:
        if "footer" in inv and footer_text:
            replace_text_in_shape(inv["footer"], [{"text": footer_text}])
        if "slide_number" in inv:
            replace_text_in_shape(
                inv["slide_number"], [{"text": f"{slide_idx + 1}/{total}"}],
            )

    # Sort children by z_index for consistent layering on fallback paths.
    children = sorted(
        list(getattr(slide_node, "children", None) or []),
        key=lambda c: int(getattr(c, "z_index", 0) or 0),
    )

    # v2.6 callout — track placed anchor bboxes (in EMU) so callouts can
    # locate their target shapes after picture/table placement. Pass 1
    # places everything except callouts; pass 2 places callouts.
    placed_anchors: dict[str, tuple[int, int, int, int]] = {}

    for child in children:
        kind = getattr(child, "kind", None)
        if kind == "callout":
            continue  # deferred to pass 2 below
        slot_name = getattr(child, "template_slot", None)

        if kind == "text":
            target = inv.get(slot_name) if slot_name else None
            if target is not None and target.has_text_frame:
                replace_text_in_shape(target, _text_to_paragraphs(child))
            else:
                # No matching slot — render as floating textbox at child.bbox.
                _add_text_frame(slide, child, slide_w, slide_h)

        elif kind == "image":
            slot = inv.get(slot_name) if slot_name else inv.get("image_slot")
            if slot is not None:
                left, top, width, height = (
                    slot.left, slot.top, slot.width, slot.height,
                )
                src = getattr(child, "src_path", None)
                placed_bbox = (left, top, width, height)
                if src and Path(src).exists():
                    fit_left, fit_top, fit_w, fit_h = _aspect_fit_emu(
                        src, left, top, width, height,
                    )
                    slide.shapes.add_picture(
                        src, fit_left, fit_top, width=fit_w, height=fit_h,
                    )
                    placed_bbox = (fit_left, fit_top, fit_w, fit_h)
                # Remove the placeholder rectangle so it doesn't show under image.
                slot._element.getparent().remove(slot._element)
                # Keep inv consistent for later children that look up image_slot.
                if slot_name and slot_name in inv:
                    del inv[slot_name]
                elif "image_slot" in inv and inv["image_slot"] is slot:
                    del inv["image_slot"]
                # v2.6 — record the picture's actual EMU bbox so callouts
                # can resolve anchor references against it.
                lid = getattr(child, "layer_id", None)
                if lid:
                    placed_anchors[lid] = placed_bbox
            else:
                # No slot — fall back to absolute bbox.
                _add_picture(slide, child, slide_w, slide_h)
                lid = getattr(child, "layer_id", None)
                if lid and getattr(child, "bbox", None) is not None:
                    placed_anchors[lid] = _bbox_to_emu(
                        child.bbox, slide_w, slide_h,
                    )

        elif kind == "table":
            slot = inv.get(slot_name) if slot_name else inv.get("table_anchor")
            if slot is not None:
                # Use the anchor's bbox by stuffing pixel-space dims back onto the
                # node temporarily — _add_table uses _bbox_to_emu(node.bbox, ...).
                ax = int(slot.left / PX_TO_EMU)
                ay = int(slot.top / PX_TO_EMU)
                aw = int(slot.width / PX_TO_EMU)
                ah = int(slot.height / PX_TO_EMU)
                anchored_node = _BboxOverride(child, ax, ay, aw, ah)
                _add_table(slide, anchored_node, slide_w, slide_h)
                # Remove the anchor rectangle.
                slot._element.getparent().remove(slot._element)
                if slot_name and slot_name in inv:
                    del inv[slot_name]
                elif "table_anchor" in inv:
                    del inv["table_anchor"]
                # v2.6 — record table EMU bbox for callout anchor lookup.
                lid = getattr(child, "layer_id", None)
                if lid:
                    placed_anchors[lid] = (slot.left, slot.top, slot.width, slot.height)
            else:
                _add_table(slide, child, slide_w, slide_h)
                lid = getattr(child, "layer_id", None)
                if lid and getattr(child, "bbox", None) is not None:
                    placed_anchors[lid] = _bbox_to_emu(
                        child.bbox, slide_w, slide_h,
                    )

        elif kind == "background":
            # If the layout's background_fill is decorative cream, prefer the
            # planner-supplied background as a full-bleed picture on top.
            _add_background(slide, child, slide_w, slide_h)

    # v2.6 — pass 2: place callouts after all anchors are known.
    for child in children:
        if getattr(child, "kind", None) == "callout":
            _add_callout(slide, child, placed_anchors, slide_w, slide_h)

    # v2.7 — cover-only safety net: if the cover layout's `authors` shape
    # is empty or carries placeholder text ("Author One · Affiliation"),
    # auto-fill from `ctx.state["ingested"][...]["manifest"]["authors"]`.
    # The planner SHOULD set this via a child with template_slot="authors",
    # but the dogfood proved it forgets. Renderer is the last line of
    # defense.
    if role == "cover" and "authors" in inv:
        authors_shape = inv["authors"]
        cur = ""
        if authors_shape.has_text_frame:
            cur = authors_shape.text_frame.text or ""
        resolved = _resolve_authors_text(ctx, existing=cur)
        if resolved and resolved != cur and authors_shape.has_text_frame:
            from ..util.template_pptx import replace_text_in_shape as _r
            _r(authors_shape, [{"text": resolved}])

    # Speaker notes — same path as blank-Presentation case.
    notes = getattr(slide_node, "speaker_notes", None)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _text_to_paragraphs(node: Any) -> list[dict[str, Any]]:
    """Convert a `kind="text"` LayerNode into a list[ParagraphSpec] for
    `replace_text_in_shape`. Multi-line text becomes multiple paragraphs.
    align / fill carry over; the template's font / size / color are
    preserved unless explicitly overridden."""
    text = (getattr(node, "text", None) or "").strip()
    if not text:
        return [{"text": ""}]
    align = getattr(node, "align", None)
    effects = getattr(node, "effects", None)
    fill_hex = getattr(effects, "fill", None) if effects is not None else None
    paragraphs: list[dict[str, Any]] = []
    for line in text.splitlines() or [text]:
        p: dict[str, Any] = {"text": line}
        if align in ("left", "center", "right"):
            p["alignment"] = align
        if isinstance(fill_hex, str) and fill_hex.startswith("#") and len(fill_hex) == 7:
            p["color"] = fill_hex
        paragraphs.append(p)
    return paragraphs


class _BboxOverride:
    """Minimal proxy that overrides `bbox` on a LayerNode for one render
    call — used when a template anchor's position should win over the
    planner's bbox for a table."""
    def __init__(self, node: Any, x: int, y: int, w: int, h: int):
        self._node = node
        self._bbox = type("B", (), {"x": x, "y": y, "w": w, "h": h})()

    def __getattr__(self, name: str) -> Any:
        if name == "bbox":
            return self._bbox
        return getattr(self._node, name)


def _render_slide(slide: Any, slide_node: Any, slide_w: int, slide_h: int,
                  ctx: ToolContext) -> None:
    """Add shapes for each child element of a slide LayerNode."""
    children = list(getattr(slide_node, "children", None) or [])
    # Sort by z_index so higher z draws on top (pptx respects insertion order).
    children.sort(key=lambda c: int(getattr(c, "z_index", 0) or 0))

    for child in children:
        kind = getattr(child, "kind", None)
        if kind == "background":
            _add_background(slide, child, slide_w, slide_h)
        elif kind == "image":
            _add_picture(slide, child, slide_w, slide_h)
        elif kind == "text":
            _add_text_frame(slide, child, slide_w, slide_h)
        elif kind == "table":
            _add_table(slide, child, slide_w, slide_h)
        # silently skip unknown kinds; planner enforces vocab

    # v2.3 — populate PowerPoint's notes pane from slide.speaker_notes.
    # `notes_slide` is auto-created on first access by python-pptx; we only
    # write when the planner provided actual text, so slides without notes
    # keep an empty (but valid) notes_slide underneath.
    notes = getattr(slide_node, "speaker_notes", None)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _bbox_to_emu(bbox: Any, slide_w: int, slide_h: int,
                 default_bbox: tuple[int, int, int, int] | None = None
                 ) -> tuple[Emu, Emu, Emu, Emu]:
    """Resolve a bbox (or fallback to slide-sized) to EMU left/top/width/height."""
    if bbox is not None:
        x = int(getattr(bbox, "x", 0) or 0)
        y = int(getattr(bbox, "y", 0) or 0)
        w = int(getattr(bbox, "w", slide_w) or slide_w)
        h = int(getattr(bbox, "h", slide_h) or slide_h)
    elif default_bbox is not None:
        x, y, w, h = default_bbox
    else:
        x, y, w, h = 0, 0, slide_w, slide_h
    # clamp to slide bounds
    x = max(0, min(x, slide_w - 1))
    y = max(0, min(y, slide_h - 1))
    w = max(1, min(w, slide_w - x))
    h = max(1, min(h, slide_h - y))
    return Emu(x * PX_TO_EMU), Emu(y * PX_TO_EMU), Emu(w * PX_TO_EMU), Emu(h * PX_TO_EMU)


def _add_background(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    """Full-slide picture background. Planner can pass bbox or leave None to
    cover the whole slide."""
    src = getattr(node, "src_path", None)
    if not src:
        return  # no background image available; PowerPoint default is white
    if not Path(src).exists():
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
        default_bbox=(0, 0, slide_w, slide_h),
    )
    slide.shapes.add_picture(src, left, top, width=width, height=height)


def _add_callout(
    slide: Any, node: Any,
    placed_anchors: dict[str, tuple[int, int, int, int]],
    slide_w: int, slide_h: int,
) -> None:
    """v2.6 — overlay an annotation shape on top of a sibling picture/table.

    Resolves `node.callout_region` (in slide-pixel coordinates, top-left
    origin) against optional `node.anchor_layer_id` (looked up in
    `placed_anchors`). Renders one of three shape styles:

    - "highlight" → MSO_SHAPE.RECTANGLE with no fill, oxblood outline 2px
    - "circle"    → MSO_SHAPE.OVAL,      no fill, oxblood outline 2px
    - "label"     → text box w/ thin border + cream fill + Inter 12pt text;
                    optional thin connector from label center to region
                    center if `arrow=True`

    All three target the same EMU bbox computed from callout_region. If
    region is None and anchor exists, uses the whole anchor bbox. If
    nothing resolves, silent no-op (defensive — callout shouldn't crash
    a slide that's otherwise fine).
    """
    from pptx.dml.color import RGBColor as _RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE

    style = (getattr(node, "callout_style", None) or "highlight").lower()
    region = getattr(node, "callout_region", None)
    anchor_id = getattr(node, "anchor_layer_id", None)
    anchor_bbox = placed_anchors.get(anchor_id) if anchor_id else None

    # Compute target EMU bbox.
    if region is not None:
        # callout_region is in slide-pixel coords (top-left origin, same
        # as every other LayerNode bbox in the codebase).
        cx_emu = Emu(int(region.x) * PX_TO_EMU)
        cy_emu = Emu(int(region.y) * PX_TO_EMU)
        cw_emu = Emu(int(region.w) * PX_TO_EMU)
        ch_emu = Emu(int(region.h) * PX_TO_EMU)
    elif anchor_bbox is not None:
        cx_emu, cy_emu, cw_emu, ch_emu = anchor_bbox
    else:
        return  # nothing to render

    accent = _RGBColor(0x7F, 0x1D, 0x1D)
    cream = _RGBColor(0xFA, 0xF7, 0xF0)
    ink = _RGBColor(0x0F, 0x17, 0x2A)

    if style == "highlight":
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx_emu, cy_emu, cw_emu, ch_emu,
        )
        rect.fill.background()
        rect.line.color.rgb = accent
        rect.line.width = Emu(2 * PX_TO_EMU)
        rect.name = getattr(node, "layer_id", None) or "callout_highlight"

    elif style == "circle":
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx_emu, cy_emu, cw_emu, ch_emu,
        )
        oval.fill.background()
        oval.line.color.rgb = accent
        oval.line.width = Emu(2 * PX_TO_EMU)
        oval.name = getattr(node, "layer_id", None) or "callout_circle"

    elif style == "label":
        text = (getattr(node, "callout_text", None) or "").strip()
        # Label dims — heuristic, ~12pt Inter at ~9px per char + padding.
        label_w_px = max(80, len(text) * 11 + 24)
        label_h_px = 36
        # Try to place to the right of the region; fall back below.
        slide_right_emu = Emu(slide_w * PX_TO_EMU)
        right_emu = cx_emu + cw_emu + Emu(8 * PX_TO_EMU)
        label_w_emu = Emu(label_w_px * PX_TO_EMU)
        label_h_emu = Emu(label_h_px * PX_TO_EMU)
        if right_emu + label_w_emu <= slide_right_emu:
            lx, ly = right_emu, cy_emu
        else:
            lx = cx_emu
            ly = cy_emu + ch_emu + Emu(8 * PX_TO_EMU)
        tb = slide.shapes.add_textbox(lx, ly, label_w_emu, label_h_emu)
        tb.fill.solid()
        tb.fill.fore_color.rgb = cream
        tb.line.color.rgb = accent
        tb.line.width = Emu(1 * PX_TO_EMU)
        tb.name = getattr(node, "layer_id", None) or "callout_label"
        tf = tb.text_frame
        tf.margin_left = Emu(4 * PX_TO_EMU)
        tf.margin_right = Emu(4 * PX_TO_EMU)
        tf.margin_top = Emu(4 * PX_TO_EMU)
        tf.margin_bottom = Emu(4 * PX_TO_EMU)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = "Inter"
        run.font.size = Pt(12)
        run.font.color.rgb = ink

        # Optional arrow: thin connector from region center to label center.
        if getattr(node, "arrow", False):
            region_cx = cx_emu + cw_emu // 2
            region_cy = cy_emu + ch_emu // 2
            label_cx = lx + label_w_emu // 2
            label_cy = ly + label_h_emu // 2
            try:
                conn = slide.shapes.add_connector(
                    MSO_CONNECTOR_TYPE.STRAIGHT,
                    region_cx, region_cy, label_cx, label_cy,
                )
                conn.line.color.rgb = accent
                conn.line.width = Emu(1 * PX_TO_EMU)
            except Exception:
                # Connector not supported on this python-pptx version — skip.
                pass


def _add_picture(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    """Place an image into the slide letterbox-fit inside the planner's bbox.

    `python-pptx`'s `add_picture(left, top, width=W, height=H)` force-stretches
    the source to W×H. For paper figures pulled from `ingest_document` the
    source aspect ratio rarely matches the planner's slot bbox, so the
    stretch makes captions / axis labels / equation glyphs unreadable
    (2026-04-25 dogfood feedback). We mirror v1.2.3 poster behavior here:
    compute contain-fit dimensions from the source's real pixel aspect,
    center inside bbox, leave letterbox bands transparent so the slide
    background shows through.
    """
    src = getattr(node, "src_path", None)
    if not src or not Path(src).exists():
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
    )
    fit_left, fit_top, fit_w, fit_h = _aspect_fit_emu(src, left, top, width, height)
    slide.shapes.add_picture(src, fit_left, fit_top, width=fit_w, height=fit_h)


def _aspect_fit_emu(
    src_path: str,
    bbox_left: int,
    bbox_top: int,
    bbox_width: int,
    bbox_height: int,
) -> tuple[int, int, int, int]:
    """Letterbox-fit `src_path` into the EMU bbox; return (left, top, w, h).

    Falls back to the original bbox if source dims are unreadable so a
    stretched render is still better than a missing image.
    """
    if bbox_width <= 0 or bbox_height <= 0:
        return bbox_left, bbox_top, bbox_width, bbox_height
    try:
        with Image.open(src_path) as im:
            sw, sh = im.size
    except Exception:
        return bbox_left, bbox_top, bbox_width, bbox_height
    if sw <= 0 or sh <= 0:
        return bbox_left, bbox_top, bbox_width, bbox_height
    src_ratio = sw / sh
    bbox_ratio = bbox_width / bbox_height
    if src_ratio > bbox_ratio:
        new_w = bbox_width
        new_h = int(round(bbox_width / src_ratio))
    else:
        new_h = bbox_height
        new_w = int(round(bbox_height * src_ratio))
    new_left = bbox_left + (bbox_width - new_w) // 2
    new_top = bbox_top + (bbox_height - new_h) // 2
    return new_left, new_top, new_w, new_h


def _add_table(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    """Render a `kind="table"` layer as a native PowerPoint table.

    Expects `node.rows: list[list[str]]` and optional `node.headers:
    list[str]`. When `headers` is empty, the first row of `rows` is
    promoted to the header. When `rows` is empty we bail out silently
    (planner is expected to not reference empty tables, but defensive).

    Sizing:
    - bbox width/height set the table's outer frame.
    - Row heights are even; header row is slightly taller.
    - Column widths are proportional to the max string length seen in
      that column — rough but prevents one wide column from collapsing
      everything else.
    - Font size auto-shrinks based on row count to keep cells legible
      at deck scale (floor 10pt, ceiling 18pt).

    Optional: `node.col_highlight_rule: list[str]` — per-column "max"
    / "min" / "". When set, the winning row per column is bolded.
    """
    rows = getattr(node, "rows", None) or []
    headers = list(getattr(node, "headers", None) or [])
    col_rule = list(getattr(node, "col_highlight_rule", None) or [])
    if not rows and not headers:
        return

    # Promote first row if headers empty.
    if not headers and rows:
        headers = [str(c) for c in rows[0]]
        rows = rows[1:]

    # Normalize all rows to header width (pad / truncate).
    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    if n_cols == 0:
        return
    headers = [str(h) for h in headers] + [""] * (n_cols - len(headers))
    headers = headers[:n_cols]
    rows = [
        [str(c) for c in r] + [""] * (n_cols - len(r))
        for r in rows
    ]
    rows = [r[:n_cols] for r in rows]

    # v2.7 — wide-table safety net. Planner.md rule #7 asks the planner
    # to subset wide tables down to 4-6 cols (deck cells go illegible
    # past ~8). Two consecutive 2026-04-25 dogfoods produced 12-15 col
    # tables anyway. Cap rendering at 8 cols, keep first 6, append a
    # marker so the audience knows columns were dropped. Loud log so
    # reviews can flag the planner.
    _WIDE_CAP = 8
    _WIDE_KEEP = 6
    if n_cols > _WIDE_CAP:
        from ..util.logging import log as _wlog
        original_cols = n_cols
        _wlog("pptx.table.truncate",
              layer_id=getattr(node, "layer_id", "?"),
              original_cols=original_cols, kept_cols=_WIDE_KEEP)
        headers = headers[:_WIDE_KEEP]
        rows = [r[:_WIDE_KEEP] for r in rows]
        n_cols = _WIDE_KEEP
        # Stuff a marker into the caption so the slide carries evidence.
        marker = (f" [Truncated: showing {_WIDE_KEEP}/{original_cols} "
                  f"cols — see paper for full table]")
        try:
            cur_cap = (getattr(node, "caption", None) or "").strip()
            node.caption = (cur_cap + marker).strip() if cur_cap else marker.strip()
        except (AttributeError, TypeError):
            pass  # frozen / proxy — best-effort

    # Normalize rule list length.
    if col_rule:
        col_rule = col_rule[:n_cols] + [""] * max(0, n_cols - len(col_rule))

    # Winner map for bold highlighting. Import here to avoid cycles.
    from ..util.table_png import _compute_winner_rows
    winner_rows = _compute_winner_rows(rows, col_rule)

    n_rows = len(rows) + 1  # +1 for header
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
    )

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    # Column widths proportional to max string length in column.
    total_w = sum(col.width for col in table.columns)
    col_weights = []
    for c in range(n_cols):
        cells = [headers[c]] + [row[c] for row in rows]
        max_len = max((len(str(v)) for v in cells), default=1)
        col_weights.append(max(1, min(max_len, 30)))
    total_weight = sum(col_weights)
    for c, col in enumerate(table.columns):
        col.width = Emu(int(total_w * col_weights[c] / total_weight))

    # Font-size autoscale: smaller when the table has many rows.
    body_pt = 18 if n_rows <= 6 else 14 if n_rows <= 12 else 11
    header_pt = body_pt + 1

    _fill_table_row(table, 0, headers, font_pt=header_pt,
                    is_header=True, winner_cols=set())
    for r_idx, row in enumerate(rows, start=1):
        # Data row idx in the rows list is r_idx - 1.
        winning_cols = {c for c, win_r in winner_rows.items()
                        if win_r == r_idx - 1}
        _fill_table_row(table, r_idx, row, font_pt=body_pt,
                        is_header=False, winner_cols=winning_cols)


def _fill_table_row(table: Any, r: int, values: list[str],
                    *, font_pt: int, is_header: bool,
                    winner_cols: set[int] | None = None) -> None:
    winner_cols = winner_cols or set()
    for c, val in enumerate(values):
        cell = table.cell(r, c)
        cell.text = ""  # clear default
        tf = cell.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = val
        font = run.font
        font.size = Pt(font_pt)
        # Bold for header row OR for the winning data cell per column.
        font.bold = is_header or (c in winner_cols)
        if is_header:
            font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            # Fill header cell with a dark accent (python-pptx solid-fill API).
            from pptx.oxml.ns import qn
            tcPr = cell._tc.get_or_add_tcPr()
            for existing in tcPr.findall(qn("a:solidFill")):
                tcPr.remove(existing)
            from lxml import etree
            fill = etree.SubElement(tcPr, qn("a:solidFill"))
            etree.SubElement(fill, qn("a:srgbClr"), val="1F2A44")


def _add_text_frame(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    text = (getattr(node, "text", None) or "").strip()
    if not text:
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
    )
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    # First paragraph holds the initial run; subsequent lines become new paras.
    lines = text.splitlines() or [text]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ""  # clear any default run
        align = getattr(node, "align", None)
        if align in _ALIGN_MAP:
            para.alignment = _ALIGN_MAP[align]
        run = para.add_run()
        run.text = line
        font = run.font
        # python-pptx wants pts; we get px. pt ≈ px * 72/96 = px * 0.75.
        size_px = int(getattr(node, "font_size_px", None) or 36)
        font.size = Pt(max(6, round(size_px * 0.75)))
        family = getattr(node, "font_family", None)
        if family:
            font.name = family
        effects = getattr(node, "effects", None)
        fill_hex = None
        if effects is not None:
            fill_hex = getattr(effects, "fill", None)
        if fill_hex and isinstance(fill_hex, str) and fill_hex.startswith("#"):
            rgb = _hex_to_rgb(fill_hex)
            if rgb is not None:
                font.color.rgb = RGBColor(*rgb)


def _hex_to_rgb(hx: str) -> tuple[int, int, int] | None:
    s = hx.lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        return None
    try:
        return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    except ValueError:
        return None


def render_slide_preview_png(slide_node: Any, slide_w: int, slide_h: int,
                             out_path: Path, ctx: ToolContext) -> None:
    """Pillow-render a simplified preview of one slide.

    Not pixel-accurate with PowerPoint's renderer — it's an at-a-glance thumb
    for chat UX + for stitching into the grid preview. Shows the bg image (if
    any), image shapes (as resized thumbs), and text (approximated font).
    """
    # Downscale to a chat-friendly size while preserving aspect.
    max_w = 960
    scale = min(1.0, max_w / slide_w)
    w = max(1, int(slide_w * scale))
    h = max(1, int(slide_h * scale))

    img = Image.new("RGB", (w, h), (255, 255, 255))

    children = sorted(
        list(getattr(slide_node, "children", None) or []),
        key=lambda c: int(getattr(c, "z_index", 0) or 0),
    )

    for child in children:
        kind = getattr(child, "kind", None)
        if kind in ("background", "image", "table"):
            # Tables use their pre-rendered src_path (PIL-drawn PNG) —
            # ingest_document baked it. PPTX itself holds a live table
            # shape; this preview just needs a raster for the thumbnail.
            _paste_image(img, child, slide_w, slide_h, scale)
        elif kind == "text":
            _draw_text(img, child, slide_w, slide_h, scale, ctx)

    img.save(out_path, format="PNG", optimize=True)


def _scaled_bbox(bbox: Any, slide_w: int, slide_h: int,
                 scale: float,
                 default_full: bool = False) -> tuple[int, int, int, int]:
    if bbox is not None:
        x = int(getattr(bbox, "x", 0) or 0)
        y = int(getattr(bbox, "y", 0) or 0)
        w = int(getattr(bbox, "w", slide_w) or slide_w)
        h = int(getattr(bbox, "h", slide_h) or slide_h)
    elif default_full:
        x, y, w, h = 0, 0, slide_w, slide_h
    else:
        x, y, w, h = 0, 0, slide_w // 2, slide_h // 4
    return (
        max(0, int(x * scale)),
        max(0, int(y * scale)),
        max(1, int(w * scale)),
        max(1, int(h * scale)),
    )


def _paste_image(canvas: Image.Image, node: Any, slide_w: int, slide_h: int,
                 scale: float) -> None:
    src = getattr(node, "src_path", None)
    if not src or not Path(src).exists():
        return
    try:
        tile = Image.open(src).convert("RGBA")
    except Exception:
        return
    kind = getattr(node, "kind", None)
    default_full = kind == "background"
    sx, sy, sw, sh = _scaled_bbox(
        getattr(node, "bbox", None), slide_w, slide_h, scale,
        default_full=default_full,
    )
    if kind == "image" and tile.size != (sw, sh):
        # Letterbox-fit content figures so the preview matches the PPTX
        # render path (same v1.2.3-style aspect-preserve as poster/SVG).
        # Backgrounds keep cover-fit (force-resize) since the cover is
        # always full-bleed by design and any minor seedream aspect drift
        # is better cropped than letterboxed with white bars on a slide.
        src_w, src_h = tile.size
        if src_w > 0 and src_h > 0 and sw > 0 and sh > 0:
            src_ratio = src_w / src_h
            bbox_ratio = sw / sh
            if src_ratio > bbox_ratio:
                new_w = sw
                new_h = max(1, int(round(sw / src_ratio)))
            else:
                new_h = sh
                new_w = max(1, int(round(sh * src_ratio)))
            tile = tile.resize((new_w, new_h), Image.LANCZOS)
            sx += (sw - new_w) // 2
            sy += (sh - new_h) // 2
    else:
        tile = tile.resize((sw, sh), Image.LANCZOS)
    canvas.alpha_composite(tile, dest=(sx, sy)) if canvas.mode == "RGBA" else canvas.paste(tile, (sx, sy), tile)


def _draw_text(canvas: Image.Image, node: Any, slide_w: int, slide_h: int,
               scale: float, ctx: ToolContext) -> None:
    text = (getattr(node, "text", None) or "").strip()
    if not text:
        return
    sx, sy, sw, sh = _scaled_bbox(
        getattr(node, "bbox", None), slide_w, slide_h, scale,
        default_full=False,
    )

    # Pick a reasonable approximated font size from the node metadata.
    size_px = int(getattr(node, "font_size_px", None) or 36)
    approx = max(10, min(120, int(size_px * scale)))
    family = getattr(node, "font_family", None) or ctx.settings.default_text_font

    fonts = ctx.settings.fonts
    fname = fonts.get(family) or fonts[ctx.settings.default_text_font]
    try:
        font = ImageFont.truetype(str(ctx.settings.fonts_dir / fname), size=approx)
    except Exception:
        font = ImageFont.load_default()

    effects = getattr(node, "effects", None)
    fill_hex = getattr(effects, "fill", None) if effects is not None else None
    rgb = _hex_to_rgb(fill_hex) if isinstance(fill_hex, str) else None
    fill = rgb or (15, 23, 42)

    draw = ImageDraw.Draw(canvas)
    # Word-wrap to bbox: simple char-by-char wrap (works for CJK; latin uses spaces).
    lines = _wrap_for_width(text, font, sw, draw)
    line_h = approx + 6
    y = sy
    for line in lines:
        if y + line_h > sy + sh:
            break
        draw.text((sx, y), line, fill=fill, font=font)
        y += line_h


def _wrap_for_width(text: str, font: Any, max_w: int, draw: Any) -> list[str]:
    """Simple word/char wrap to fit a pixel width."""
    out: list[str] = []
    for paragraph in text.splitlines() or [text]:
        if " " in paragraph:
            # space-delimited: greedy word-wrap
            words = paragraph.split()
            line = ""
            for w in words:
                probe = (line + " " + w).strip()
                if _measure(probe, font, draw) <= max_w:
                    line = probe
                else:
                    if line:
                        out.append(line)
                    line = w
            if line:
                out.append(line)
        else:
            # CJK-style: char-by-char
            line = ""
            for ch in paragraph:
                probe = line + ch
                if _measure(probe, font, draw) <= max_w:
                    line = probe
                else:
                    if line:
                        out.append(line)
                    line = ch
            if line:
                out.append(line)
    return out


def _measure(s: str, font: Any, draw: Any) -> int:
    try:
        bbox = draw.textbbox((0, 0), s, font=font)
        return int(bbox[2] - bbox[0])
    except Exception:
        return len(s) * 10
